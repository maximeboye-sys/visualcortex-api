"""
layouts.py — Visual Cortex
Bibliothèque de layouts pré-testés pour le pipeline v3.
Claude choisit un layout et retourne du JSON ; ces fonctions appliquent charte + contenu.
Chaque fonction : (prs, content, palette) → slide
"""
from __future__ import annotations

import lxml.etree as etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


# ─────────────────────────────────────────────────────────────
# HELPERS INTERNES
# ─────────────────────────────────────────────────────────────

def _hex(hex_str: str) -> RGBColor:
    try:
        h = str(hex_str).lstrip('#').strip()
        if len(h) == 3:
            h = h[0]*2 + h[1]*2 + h[2]*2
        if len(h) != 6:
            return RGBColor(0x1A, 0x3A, 0x6B)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return RGBColor(0x1A, 0x3A, 0x6B)


def _blank(prs: Presentation):
    """
    Slide vierge qui hérite correctement du slide master.
    Garantit l'affichage du logo et des éléments du master
    quel que soit le template uploadé.
    """
    # Chercher un layout 'Vide' / 'Blank' / 'Empty' sans showMasterSp="0"
    target = None
    for layout in prs.slide_layouts:
        name = layout.name.lower()
        if any(x in name for x in ['vide', 'blank', 'empty']):
            xml = etree.tostring(layout._element, pretty_print=False).decode()
            if 'showMasterSp="0"' not in xml:
                target = layout
                break

    # Fallback : premier layout sans showMasterSp="0"
    if target is None:
        for layout in prs.slide_layouts:
            xml = etree.tostring(layout._element, pretty_print=False).decode()
            if 'showMasterSp="0"' not in xml:
                target = layout
                break

    if target is None:
        target = prs.slide_layouts[-1]

    slide = prs.slides.add_slide(target)

    # Supprimer uniquement les placeholders (pas les shapes héritées du master)
    sp_tree = slide.shapes._spTree
    for ph in list(slide.placeholders):
        try:
            sp_tree.remove(ph._element)
        except Exception:
            pass

    # Forcer l'affichage des shapes du master (logo, ligne déco, etc.)
    cSld = slide._element.find(qn('p:cSld'))
    if cSld is not None:
        cSld.set('showMasterSp', '1')

    W = prs.slide_width  / 914400.0
    H = prs.slide_height / 914400.0
    return slide, W, H


def _rect(slide, l, t, w, h, color: str):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(color)
    shape.line.fill.background()
    return shape


def _rrect(slide, l, t, w, h, color: str, radius: float = 0.06):
    shape = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(color)
    shape.line.fill.background()
    try:
        shape.adjustments[0] = max(0.0, min(0.5, radius))
    except Exception:
        pass
    return shape


def _circle(slide, cx, cy, r, color: str):
    shape = slide.shapes.add_shape(9, Inches(cx-r), Inches(cy-r),
                                   Inches(r*2), Inches(r*2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex(color)
    shape.line.fill.background()
    return shape


def _set_lsp(p, pct: int):
    """Interligne en pourcents (120 = 1.2×)."""
    try:
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(pct * 1000)))
    except Exception:
        pass


def _apply_text_gradient(txBox, color_start: str, color_end: str):
    """
    Applique un gradient linéaire (gauche → droite) sur tous les runs du textbox.
    color_start / color_end = hex 'RRGGBB' (sans #).
    Utilisé pour reproduire les titres en dégradé des templates corporate.
    """
    _NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    cs = color_start.lstrip('#').strip()
    ce = color_end.lstrip('#').strip()
    if len(cs) != 6 or len(ce) != 6:
        return
    grad_xml = (
        f'<a:gradFill xmlns:a="{_NS}" flip="none" rotWithShape="1">'
        f'<a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{cs}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{ce}"/></a:gs>'
        f'</a:gsLst>'
        f'<a:lin ang="5400000" scaled="0"/>'
        f'</a:gradFill>'
    )
    try:
        grad_el = etree.fromstring(grad_xml)
        for para in txBox.text_frame.paragraphs:
            for run in para.runs:
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = etree.SubElement(run._r, qn('a:rPr'))
                for sf in rPr.findall(qn('a:solidFill')):
                    rPr.remove(sf)
                rPr.insert(0, etree.fromstring(grad_xml))
    except Exception:
        pass


def _txt(slide, text: str, l, t, w, h, font: str, size: float, color: str,
         bold=False, italic=False, align='left', lsp=120):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    amap = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    p = tf.paragraphs[0]
    p.alignment = amap.get(align, PP_ALIGN.LEFT)
    _set_lsp(p, lsp)
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _hex(color)
    return tb


def _div(slide, l, t, w, color: str, h=0.032):
    """Ligne séparatrice horizontale fine."""
    _rect(slide, l, t, w, h, color)


def _footer_line(slide, text: str, W, H, font: str, color: str = 'AAAAAA'):
    if text:
        _txt(slide, text, 0.5, H - 0.42, W - 1.0, 0.35, font, 10, color, lsp=110)


# ─────────────────────────────────────────────────────────────
# LAYOUTS — 15 fonctions
# ─────────────────────────────────────────────────────────────

def cover_dark(prs: Presentation, content: dict, palette: dict):
    """Fond primary plein, titre centré large, sous-titre, footer discret."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, palette['primary'])
    # Accent: barre décorative centrée sous le titre
    bw = W * 0.22
    _rect(slide, (W - bw) / 2, H * 0.585, bw, 0.042, palette['accent'])

    _txt(slide, content.get('title', ''),
         W*0.08, H*0.20, W*0.84, H*0.32,
         f, 44, 'FFFFFF', bold=True, align='center', lsp=105)

    if content.get('subtitle'):
        _txt(slide, content['subtitle'],
             W*0.15, H*0.63, W*0.70, H*0.12,
             f, 18, 'FFFFFF', align='center', lsp=125)

    if content.get('footer'):
        _txt(slide, content['footer'],
             W*0.1, H - 0.44, W*0.8, 0.35,
             f, 10, 'AAAAAA', align='center', lsp=110)
    return slide


def cover_split(prs: Presentation, content: dict, palette: dict):
    """Barre latérale primary gauche (35 %) + fond blanc droite."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')
    sp = W * 0.35

    _rect(slide, 0,  0, sp, H, palette['primary'])
    _rect(slide, 0, H*0.78, sp, H*0.22, palette['accent'])
    _rect(slide, sp, 0, W - sp, H, 'FFFFFF')

    _txt(slide, content.get('title', ''),
         sp + 0.5, H*0.27, W - sp - 0.65, H*0.32,
         f, 38, palette['primary'], bold=True, lsp=108)
    _div(slide, sp + 0.5, H*0.62, (W - sp - 1.0)*0.55, palette['accent'], h=0.05)

    if content.get('subtitle'):
        _txt(slide, content['subtitle'],
             sp + 0.5, H*0.66, W - sp - 0.65, H*0.16,
             f, 16, palette['text'], lsp=128)

    _footer_line(slide, content.get('footer', ''), W, H, f, 'AAAAAA')
    return slide


def section(prs: Presentation, content: dict, palette: dict):
    """Séparateur de chapitre — fond primary, grand numéro déco, titre."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, palette['primary'])

    num = str(content.get('number', ''))
    if num:
        # Grand numéro en accent, légèrement décalé à droite
        _txt(slide, num, W*0.52, H*0.02, W*0.46, H*0.88,
             f, 80, palette['accent'], bold=True, align='left', lsp=100)

    _div(slide, 0.5, H*0.47, W*0.28, palette['accent'], h=0.05)
    _txt(slide, content.get('title', ''),
         0.5, H*0.52, W*0.62, H*0.32,
         f, 36, 'FFFFFF', bold=True, lsp=108)
    return slide


def kpi_grid(prs: Presentation, content: dict, palette: dict):
    """Grille 4-6 KPIs sur fond sombre."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, palette['primary'])
    _txt(slide, content.get('title', ''),
         0.5, 0.32, W - 1.0, 0.62,
         f, 26, 'FFFFFF', bold=True, lsp=108)
    _div(slide, 0.5, 0.98, W*0.13, palette['accent'], h=0.045)

    kpis = (content.get('kpis') or [])[:6]
    n = len(kpis)
    if not n:
        return slide

    cols = 3 if n > 4 else (2 if n > 2 else n)
    rows = (n + cols - 1) // cols
    cw = (W - 1.0) / cols
    ch = (H - 1.75) / rows

    for i, kpi in enumerate(kpis):
        col, row = i % cols, i // cols
        lx = 0.5 + col * cw
        ty = 1.30 + row * ch

        _txt(slide, str(kpi.get('value', '')),
             lx, ty, cw - 0.15, 0.65,
             f, 34, palette['accent'], bold=True, align='center', lsp=100)
        _txt(slide, str(kpi.get('label', '')),
             lx, ty + 0.73, cw - 0.15, 0.34,
             f, 11, 'FFFFFF', align='center', lsp=115)
        if kpi.get('sublabel'):
            _txt(slide, kpi['sublabel'],
                 lx, ty + 1.14, cw - 0.15, 0.28,
                 f, 9, 'BBBBBB', align='center', lsp=112)
        # Séparateur vertical entre colonnes
        if col < cols - 1:
            _rect(slide, lx + cw - 0.018, ty + 0.05, 0.018, ch - 0.28, palette.get('dark', '0D1F3C'))

    _footer_line(slide, content.get('footer', ''), W, H, f, 'AAAAAA')
    return slide


def kpi_row(prs: Presentation, content: dict, palette: dict):
    """Ligne 3-4 KPIs sur fond clair."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, 'FFFFFF')
    _rect(slide, 0, 0, W, 0.07, palette['primary'])

    _txt(slide, content.get('title', ''),
         0.5, 0.35, W - 1.0, 0.62,
         f, 28, palette['primary'], bold=True, lsp=108)
    _div(slide, 0.5, 1.02, W*0.13, palette['accent'], h=0.045)

    kpis = (content.get('kpis') or [])[:4]
    n = len(kpis)
    if not n:
        return slide

    cw = (W - 1.0) / n
    ty = H * 0.42

    for i, kpi in enumerate(kpis):
        lx = 0.5 + i * cw
        _rect(slide, lx + 0.12, ty - 0.16, cw - 0.40, 0.048, palette['accent'])
        _txt(slide, str(kpi.get('value', '')),
             lx, ty, cw - 0.15, 0.72,
             f, 38, palette['primary'], bold=True, align='center', lsp=100)
        _txt(slide, str(kpi.get('label', '')),
             lx, ty + 0.82, cw - 0.15, 0.34,
             f, 12, palette['text'], align='center', lsp=115)
        if kpi.get('sublabel'):
            _txt(slide, kpi['sublabel'],
                 lx, ty + 1.24, cw - 0.15, 0.28,
                 f, 10, 'AAAAAA', align='center', lsp=112)
        if i < n - 1:
            _rect(slide, lx + cw - 0.015, ty, 0.015, 1.5, palette['light'])

    _footer_line(slide, content.get('footer', ''), W, H, f)
    return slide


def timeline_h(prs: Presentation, content: dict, palette: dict):
    """Frise chronologique horizontale, 4-5 jalons."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, 'FFFFFF')
    _rect(slide, 0, 0, W, 0.07, palette['primary'])

    _txt(slide, content.get('title', ''),
         0.5, 0.35, W - 1.0, 0.62,
         f, 28, palette['primary'], bold=True, lsp=108)
    _div(slide, 0.5, 1.02, W*0.13, palette['accent'], h=0.045)

    steps = (content.get('steps') or [])[:5]
    n = len(steps)
    if not n:
        return slide

    axis_y = H * 0.52
    margin = 0.9
    avail = W - margin * 2
    step_w = avail / max(n - 1, 1)

    # Axe principal
    _rect(slide, margin, axis_y - 0.014, avail, 0.028, palette['primary'])

    for i, step in enumerate(steps):
        cx = margin if n == 1 else margin + i * step_w

        # Dot
        _circle(slide, cx, axis_y, 0.17, palette['accent'])
        _circle(slide, cx, axis_y, 0.09, 'FFFFFF')

        bw = min(step_w * 0.88, 1.75) if n > 1 else 1.75
        lx = cx - bw / 2

        # Date au-dessus
        _txt(slide, str(step.get('date', '')),
             lx, axis_y - 0.68, bw, 0.36,
             f, 12, palette['accent'], bold=True, align='center', lsp=110)

        # Titre en dessous
        _txt(slide, str(step.get('title', '')),
             lx, axis_y + 0.26, bw, 0.40,
             f, 12, palette['primary'], bold=True, align='center', lsp=115)

        if step.get('body'):
            _txt(slide, step['body'],
                 lx, axis_y + 0.74, bw, 0.70,
                 f, 10, palette['text'], align='center', lsp=118)

    _footer_line(slide, content.get('footer', ''), W, H, f)
    return slide


def two_col(prs: Presentation, content: dict, palette: dict):
    """Deux colonnes symétriques avec items à puces."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, 'FFFFFF')
    _rect(slide, 0, 0, W, 0.07, palette['primary'])

    _txt(slide, content.get('title', ''),
         0.5, 0.35, W - 1.0, 0.62,
         f, 28, palette['primary'], bold=True, lsp=108)
    _div(slide, 0.5, 1.02, W*0.13, palette['accent'], h=0.045)

    cw = (W - 1.3) / 2
    col_a = content.get('col_a') or {}
    col_b = content.get('col_b') or {}

    for j, (col, lx, bar_color) in enumerate([
        (col_a, 0.5,            palette['primary']),
        (col_b, 0.5 + cw + 0.3, palette['secondary']),
    ]):
        _rect(slide, lx, 1.22, cw, 0.040, bar_color)
        _txt(slide, col.get('title', ''),
             lx, 1.32, cw, 0.40,
             f, 14, bar_color, bold=True, lsp=115)

        items = (col.get('items') or [])[:5]
        for k, item in enumerate(items):
            iy = 1.82 + k * 0.72
            _circle(slide, lx + 0.17, iy + 0.17, 0.095, bar_color)
            _txt(slide, str(item),
                 lx + 0.40, iy, cw - 0.48, 0.60,
                 f, 12, palette['text'], lsp=122)

    _footer_line(slide, content.get('footer', ''), W, H, f)
    return slide


def quote_dark(prs: Presentation, content: dict, palette: dict):
    """Citation forte centrée sur fond sombre."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, palette['primary'])

    # Guillemet décoratif grand
    _txt(slide, '\u201c', 0.25, H*0.05, 2.0, 1.8,
         f, 130, palette['accent'], bold=True, lsp=100)

    _txt(slide, content.get('quote', ''),
         W*0.11, H*0.26, W*0.78, H*0.40,
         f, 22, 'FFFFFF', italic=True, align='center', lsp=132)

    _div(slide, W*0.38, H*0.69, W*0.24, palette['accent'], h=0.048)

    if content.get('author'):
        _txt(slide, f"— {content['author']}",
             W*0.18, H*0.73, W*0.64, 0.40,
             f, 12, palette['accent'], bold=True, align='center', lsp=115)

    _footer_line(slide, content.get('footer', ''), W, H, f, 'AAAAAA')
    return slide


def list_numbered(prs: Presentation, content: dict, palette: dict):
    """Liste numérotée avec cercles colorés, 3-5 items."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, 'FFFFFF')
    _rect(slide, 0, 0, W, 0.07, palette['primary'])

    _txt(slide, content.get('title', ''),
         0.5, 0.35, W - 1.0, 0.62,
         f, 28, palette['primary'], bold=True, lsp=108)
    _div(slide, 0.5, 1.02, W*0.13, palette['accent'], h=0.045)

    items = (content.get('items') or [])[:5]
    n = len(items)
    if not n:
        return slide

    avail_h = H - 1.88
    item_h = min(1.10, avail_h / n)

    for i, item in enumerate(items):
        iy = 1.28 + i * item_h
        cy = iy + item_h * 0.42

        _circle(slide, 0.82, cy, 0.26, palette['primary'])
        _txt(slide, str(i + 1),
             0.56, cy - 0.22, 0.52, 0.44,
             f, 13, 'FFFFFF', bold=True, align='center', lsp=100)

        title = item.get('title', '') if isinstance(item, dict) else str(item)
        _txt(slide, title, 1.28, iy + item_h*0.10, W - 1.72, 0.34,
             f, 13, palette['primary'], bold=True, lsp=115)

        if isinstance(item, dict) and item.get('body'):
            _txt(slide, item['body'],
                 1.28, iy + item_h*0.10 + 0.40, W - 1.72, item_h - 0.56,
                 f, 12, palette['text'], lsp=122)

    _footer_line(slide, content.get('footer', ''), W, H, f)
    return slide


def list_cards(prs: Presentation, content: dict, palette: dict):
    """Grille 2×2 de cartes sur fond clair."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, palette['light'])
    _rect(slide, 0, 0, W, 0.07, palette['primary'])

    _txt(slide, content.get('title', ''),
         0.5, 0.35, W - 1.0, 0.62,
         f, 28, palette['primary'], bold=True, lsp=108)
    _div(slide, 0.5, 1.02, W*0.13, palette['accent'], h=0.045)

    cards = (content.get('cards') or [])[:4]
    cw = (W - 1.3) / 2
    ch = (H - 2.05) / 2 - 0.12

    for i, card in enumerate(cards):
        col, row = i % 2, i // 2
        lx = 0.5 + col * (cw + 0.3)
        ty = 1.28 + row * (ch + 0.18)

        _rrect(slide, lx, ty, cw, ch, 'FFFFFF', radius=0.04)
        _rrect(slide, lx, ty, cw, 0.075, palette['accent'], radius=0.04)

        title = card.get('title', '') if isinstance(card, dict) else str(card)
        _txt(slide, title, lx + 0.18, ty + 0.14, cw - 0.36, 0.38,
             f, 13, palette['primary'], bold=True, lsp=115)

        if isinstance(card, dict) and card.get('body'):
            _txt(slide, card['body'],
                 lx + 0.18, ty + 0.64, cw - 0.36, ch - 0.76,
                 f, 12, palette['text'], lsp=122)

    _footer_line(slide, content.get('footer', ''), W, H, f)
    return slide


def image_split(prs: Presentation, content: dict, palette: dict):
    """Bloc coloré gauche (42 %) + titre + bullets à droite."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')
    sp = W * 0.42

    _rect(slide, 0,  0,  sp, H, palette['primary'])
    _rect(slide, 0,  H*0.72, sp, H*0.28, palette['secondary'])
    _rect(slide, sp, 0,  W - sp, H, 'FFFFFF')

    _txt(slide, content.get('title', ''),
         0.45, H*0.28, sp - 0.6, H*0.38,
         f, 28, 'FFFFFF', bold=True, lsp=108)
    _div(slide, 0.45, H*0.70, sp*0.40, palette['accent'], h=0.048)

    points = (content.get('points') or [])[:5]
    n_pts = len(points)
    step = 0.92
    total_h = max(n_pts - 1, 0) * step + 0.80
    py_start = max(0.80, (H - total_h) / 2.0)
    for i, pt in enumerate(points):
        py = py_start + i * step
        _rect(slide, sp + 0.45, py + 0.14, 0.052, 0.052, palette['accent'])
        _txt(slide, str(pt),
             sp + 0.62, py, W - sp - 0.78, 0.78,
             f, 13, palette['text'], lsp=126)

    _footer_line(slide, content.get('footer', ''), W, H, f)
    return slide


def full_text(prs: Presentation, content: dict, palette: dict):
    """Fond blanc épuré, 2-3 paragraphes aérés."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, 'FFFFFF')
    _rect(slide, 0, 0, W, 0.07, palette['primary'])
    _rect(slide, 0, 0, 0.07, H, palette['primary'])

    _txt(slide, content.get('title', ''),
         0.55, 0.35, W - 0.75, 0.62,
         f, 28, palette['primary'], bold=True, lsp=108)
    _div(slide, 0.55, 1.02, W*0.13, palette['accent'], h=0.045)

    paragraphs = (content.get('paragraphs') or [])[:3]
    n = len(paragraphs)
    if not n:
        return slide

    avail = H - 1.80
    ph = avail / n

    for i, para in enumerate(paragraphs):
        py = 1.28 + i * ph
        _txt(slide, str(para),
             0.55, py, W - 1.05, ph * 0.82,
             f, 13, palette['text'], lsp=132)
        if i < n - 1:
            _div(slide, 0.55, py + ph - 0.08, W - 1.05, palette['light'], h=0.014)

    _footer_line(slide, content.get('footer', ''), W, H, f)
    return slide


def stat_hero(prs: Presentation, content: dict, palette: dict):
    """Un seul grand chiffre central, très impactant."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, 'FFFFFF')
    _rect(slide, 0, 0,      W, 0.07, palette['primary'])
    _rect(slide, 0, H-0.07, W, 0.07, palette['primary'])

    # Cercle décoratif en arrière-plan
    _circle(slide, W*0.74, H*0.44, H*0.33, palette['light'])

    _txt(slide, str(content.get('value', '')),
         W*0.04, H*0.12, W*0.92, H*0.52,
         f, 72, palette['primary'], bold=True, align='center', lsp=100)

    _div(slide, W*0.30, H*0.67, W*0.40, palette['accent'], h=0.048)

    _txt(slide, str(content.get('label', '')),
         W*0.10, H*0.69, W*0.80, 0.52,
         f, 16, palette['text'], bold=True, align='center', lsp=115)

    if content.get('context'):
        _txt(slide, content['context'],
             W*0.15, H*0.78, W*0.70, 0.50,
             f, 11, palette['text'], align='center', lsp=122)

    _footer_line(slide, content.get('footer', ''), W, H, f)
    return slide


def closing_dark(prs: Presentation, content: dict, palette: dict):
    """Fond primary, titre court centré, sous-titre sources."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')

    _rect(slide, 0, 0, W, H, palette['primary'])
    # Bandes accent gauche et droite
    _rect(slide, 0,        0, W*0.10, H, palette['accent'])
    _rect(slide, W*0.90,   0, W*0.10, H, palette['accent'])

    _txt(slide, content.get('title', ''),
         W*0.16, H*0.28, W*0.68, H*0.32,
         f, 40, 'FFFFFF', bold=True, align='center', lsp=108)
    _div(slide, W*0.36, H*0.63, W*0.28, palette['accent'], h=0.048)

    if content.get('subtitle'):
        _txt(slide, content['subtitle'],
             W*0.18, H*0.68, W*0.64, H*0.14,
             f, 14, 'FFFFFF', align='center', lsp=122)
    return slide


def closing_split(prs: Presentation, content: dict, palette: dict):
    """Split primary/blanc, titre à gauche, sous-titre + CTA à droite."""
    slide, W, H = _blank(prs)
    f = palette.get('font', 'Calibri')
    sp = W * 0.45

    _rect(slide, 0,  0,  sp, H, palette['primary'])
    _rect(slide, 0,  H*0.76, sp, H*0.24, palette['accent'])
    _rect(slide, sp, 0,  W - sp, H, 'FFFFFF')

    _txt(slide, content.get('title', ''),
         0.5, H*0.24, sp - 0.65, H*0.40,
         f, 34, 'FFFFFF', bold=True, lsp=108)
    _div(slide, 0.5, H*0.67, sp*0.42, 'FFFFFF', h=0.038)

    if content.get('subtitle'):
        _txt(slide, content['subtitle'],
             sp + 0.5, H*0.28, W - sp - 0.70, H*0.38,
             f, 15, palette['text'], lsp=132)

    # Bouton CTA
    btn_l = sp + 0.5
    btn_t = H * 0.70
    _rrect(slide, btn_l, btn_t, 2.4, 0.52, palette['accent'], radius=0.08)
    _txt(slide, 'Merci',
         btn_l, btn_t + 0.03, 2.4, 0.46,
         f, 16, 'FFFFFF', bold=True, align='center', lsp=110)
    return slide


# ─────────────────────────────────────────────────────────────
# REGISTRE & DESCRIPTIONS
# ─────────────────────────────────────────────────────────────

LAYOUT_REGISTRY: dict = {
    'cover_dark':    cover_dark,
    'cover_split':   cover_split,
    'section':       section,
    'kpi_grid':      kpi_grid,
    'kpi_row':       kpi_row,
    'timeline_h':    timeline_h,
    'two_col':       two_col,
    'quote_dark':    quote_dark,
    'list_numbered': list_numbered,
    'list_cards':    list_cards,
    'image_split':   image_split,
    'full_text':     full_text,
    'stat_hero':     stat_hero,
    'closing_dark':  closing_dark,
    'closing_split': closing_split,
}

LAYOUT_DESCRIPTIONS: dict = {
    'cover_dark':    'Couverture fond primary — title, subtitle, footer',
    'cover_split':   'Couverture split gauche/blanc — title, subtitle, footer',
    'section':       'Séparateur chapitre fond primary — number, title',
    'kpi_grid':      'Grille 4-6 KPIs fond sombre — title, kpis[{value,label,sublabel}], footer',
    'kpi_row':       'Ligne 3-4 KPIs fond clair — title, kpis[{value,label,sublabel}], footer',
    'timeline_h':    'Frise chronologique horizontale — title, steps[{date,title,body}], footer',
    'two_col':       '2 colonnes symétriques — title, col_a{title,items[]}, col_b{title,items[]}, footer',
    'quote_dark':    'Citation fond sombre — quote, author, footer',
    'list_numbered': 'Liste numérotée cercles — title, items[{title,body}], footer',
    'list_cards':    'Grille 2×2 cartes — title, cards[{title,body}], footer',
    'image_split':   'Split coloré/texte — title, points[], footer',
    'full_text':     'Texte long fond blanc — title, paragraphs[], footer',
    'stat_hero':     'Grand chiffre héro — value, label, context, footer',
    'closing_dark':  'Clôture fond primary — title, subtitle',
    'closing_split': 'Clôture split + CTA — title, subtitle',
}

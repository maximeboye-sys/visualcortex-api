"""
test_v4_layouts.py — Test exhaustif du pipeline V4.

Sections :
  A) analyze_template_v4  — extraction thème, couleurs, font, layout_map
  B) Prompt planner       — construction du user_prompt (KeyError sur accolades)
  C) 54 layout_*_v4      — aucune exception tolérée

Usage : python3 test_v4_layouts.py
Critère de succès : exit code 0 (zéro échec dans toutes les sections).
"""
import sys
import traceback
import io
from pptx import Presentation
from pptx.util import Inches

# ── Import de main.py ─────────────────────────────────────────────────────────
import importlib.util
spec = importlib.util.spec_from_file_location("main", "main.py")
main = importlib.util.module_from_spec(spec)
spec.loader.exec_module(main)

# ══════════════════════════════════════════════════════════════════════════════
# TP DICT — fidèle à analyze_template_v4
# ══════════════════════════════════════════════════════════════════════════════
THEME = {
    'dk1':     '374649',
    'lt1':     'FFFFFF',
    'dk2':     '1F3864',
    'lt2':     'E9EFF7',
    'accent1': '009CEA',
    'accent2': 'ED0000',
    'accent3': '40A900',
    'accent4': 'F66A00',
    'accent5': '7B2FBE',
    'accent6': '00B4D8',
}

ACCENT_CYCLE = [THEME[f'accent{i}'] for i in range(1, 7)]

TP = {
    'theme':           THEME,
    'layout_map':      {
        'blank': 6, 'text': 1, 'title': 0,
        'cover': 0, 'section': 2, 'closing': 0,
        'kpi': 1, 'two_col': 1, 'content': 1,
    },
    'logo_zone':       None,
    'font':            'Calibri',
    'accent_cycle':    ACCENT_CYCLE,
    'card_bg_light':   'EEF4FB',
    'card_bg_mid':     'DCE9F5',
    'bg_type':         'plain',
    'bg_colors':       ['FFFFFF'],
    'bg_is_dark':      False,
    'bg_rich':         False,
    'rich_layout_idx': 1,
    'W':               13.33,
    'H':               7.5,
    'seed':            42,
}

# ══════════════════════════════════════════════════════════════════════════════
# CONTENU DE TEST PAR LAYOUT
# ══════════════════════════════════════════════════════════════════════════════
CONTENT = {
    'cover':             {'title': 'Stratégie 2025', 'subtitle': 'Analyse & perspectives', 'footer': 'VC'},
    'section':           {'number': '01', 'title': 'Contexte', 'footer': 'VC'},
    'fulltext':          {'title': 'Analyse', 'paragraphs': ['Para 1.', 'Para 2.'], 'footer': 'VC'},
    'closing':           {'title': 'Merci', 'subtitle': 'Questions ?', 'footer': 'VC'},
    'quote':             {'quote': 'Citation test.', 'author': 'Auteur', 'footer': 'VC'},
    'list_numbered':     {'title': 'Liste', 'items': [
        {'title': 'Point A', 'body': 'Détail A'},
        {'title': 'Point B', 'body': 'Détail B'},
        {'title': 'Point C', 'body': 'Détail C'},
    ], 'footer': 'VC'},
    'list_cards':        {'title': 'Cartes', 'cards': [
        {'title': 'Carte 1', 'body': 'Desc 1'},
        {'title': 'Carte 2', 'body': 'Desc 2'},
        {'title': 'Carte 3', 'body': 'Desc 3'},
    ], 'footer': 'VC'},
    'col3':              {'title': '3 Colonnes', 'columns': [
        {'title': 'Col A', 'items': ['A1', 'A2']},
        {'title': 'Col B', 'items': ['B1', 'B2']},
        {'title': 'Col C', 'items': ['C1', 'C2']},
    ], 'footer': 'VC'},
    'twocol':            {'title': 'Deux colonnes',
                          'col_a': {'title': 'Gauche', 'items': ['G1', 'G2']},
                          'col_b': {'title': 'Droite',  'items': ['D1', 'D2']}, 'footer': 'VC'},
    'stathero':          {'value': '98 %', 'label': 'Satisfaction client', 'context': 'Enquête 2024', 'footer': 'VC'},
    'infographic':       {'title': 'Infographie', 'items': [
        {'value': '1', 'label': 'Un'},
        {'value': '2', 'label': 'Deux'},
        {'value': '3', 'label': 'Trois'},
    ], 'footer': 'VC'},
    'timeline':          {'title': 'Chronologie', 'steps': [
        {'date': '2022', 'title': 'Étape 1', 'body': 'Desc'},
        {'date': '2023', 'title': 'Étape 2', 'body': 'Desc'},
        {'date': '2024', 'title': 'Étape 3', 'body': 'Desc'},
    ], 'footer': 'VC'},
    'processflow':       {'title': 'Processus', 'steps': [
        {'title': 'Étape 1', 'body': 'Desc'},
        {'title': 'Étape 2', 'body': 'Desc'},
        {'title': 'Étape 3', 'body': 'Desc'},
    ], 'footer': 'VC'},
    'kpi_grid':          {'title': 'KPIs', 'kpis': [
        {'value': '42',  'label': 'KPI A'},
        {'value': '99',  'label': 'KPI B'},
        {'value': '7',   'label': 'KPI C'},
        {'value': '100', 'label': 'KPI D'},
    ], 'footer': 'VC'},
    'funnel':            {'title': 'Entonnoir', 'steps': [
        {'label': 'Visiteurs', 'value': '10000'},
        {'label': 'Leads',     'value': '2000'},
        {'label': 'Clients',   'value': '400'},
    ], 'footer': 'VC'},
    'barchart':          {'title': 'Bar Chart', 'categories': ['A', 'B', 'C'],
                          'series': [{'name': 'S1', 'values': [10, 20, 15]}], 'footer': 'VC'},
    'linechart':         {'title': 'Line Chart', 'categories': ['Jan', 'Fév', 'Mar'],
                          'series': [{'name': 'S1', 'values': [10, 20, 15]}], 'footer': 'VC'},
    'piechart':          {'title': 'Pie Chart', 'slices': [
        {'label': 'A', 'value': 40},
        {'label': 'B', 'value': 35},
        {'label': 'C', 'value': 25},
    ], 'footer': 'VC'},
    'waterfall':         {'title': 'Waterfall', 'items': [
        {'label': 'Départ', 'value': 100},
        {'label': '+Ventes', 'value': 50},
        {'label': '-Coûts', 'value': -30},
    ], 'footer': 'VC'},
    'radar':             {'title': 'Radar', 'axes': ['A', 'B', 'C', 'D', 'E'],
                          'series': [{'name': 'S1', 'values': [4, 3, 5, 2, 4]}], 'footer': 'VC'},
    'pyramid':           {'title': 'Pyramide', 'levels': [
        {'label': 'Top',  'value': 10},
        {'label': 'Mid',  'value': 30},
        {'label': 'Base', 'value': 60},
    ], 'footer': 'VC'},
    'cycle':             {'title': 'Cycle', 'steps': [
        {'title': 'Phase 1', 'body': 'Desc'},
        {'title': 'Phase 2', 'body': 'Desc'},
        {'title': 'Phase 3', 'body': 'Desc'},
        {'title': 'Phase 4', 'body': 'Desc'},
    ], 'footer': 'VC'},
    'roadmap':           {'title': 'Roadmap', 'quarters': [
        {'label': 'Q1', 'items': ['T1', 'T2']},
        {'label': 'Q2', 'items': ['T3']},
        {'label': 'Q3', 'items': ['T4', 'T5']},
    ], 'footer': 'VC'},
    'stackedbar':        {'title': 'Stacked Bar', 'categories': ['A', 'B', 'C'],
                          'series': [
                              {'name': 'S1', 'values': [10, 20, 15]},
                              {'name': 'S2', 'values': [5, 10, 8]},
                          ], 'footer': 'VC'},
    'beforeafter':       {'title': 'Avant / Après',
                          'before': {'title': 'Avant', 'items': ['Pb A', 'Pb B']},
                          'after':  {'title': 'Après', 'items': ['Sol A', 'Sol B']}, 'footer': 'VC'},
    'entity':            {'title': 'Entités', 'entities': [
        {'name': 'Entité A', 'attributes': ['Attr 1', 'Attr 2']},
        {'name': 'Entité B', 'attributes': ['Attr 3', 'Attr 4']},
    ], 'footer': 'VC'},
    'conclusion':        {'title': 'Conclusion', 'points': ['Point 1', 'Point 2', 'Point 3'], 'footer': 'VC'},
    'highlight':         {'title': 'Highlight', 'message': 'Message clé', 'body': 'Détails.', 'footer': 'VC'},
    'agenda':            {'title': 'Agenda', 'items': [
        {'number': '01', 'title': 'Sujet A', 'duration': '15 min'},
        {'number': '02', 'title': 'Sujet B', 'duration': '20 min'},
        {'number': '03', 'title': 'Sujet C', 'duration': '10 min'},
    ], 'footer': 'VC'},
    'matrix':            {'title': 'Matrice', 'items': [
        {'label': 'Item A', 'x': 0.8, 'y': 0.7},
        {'label': 'Item B', 'x': 0.3, 'y': 0.6},
        {'label': 'Item C', 'x': 0.6, 'y': 0.2},
    ], 'footer': 'VC'},
    'swot':              {'title': 'SWOT',
                          'strengths': ['Force 1', 'Force 2'],
                          'weaknesses': ['Faiblesse 1'],
                          'opportunities': ['Opportunité 1'],
                          'threats': ['Menace 1'], 'footer': 'VC'},
    'proscons':          {'title': 'Pour / Contre',
                          'pros': ['Avantage 1', 'Avantage 2'],
                          'cons': ['Inconvénient 1'], 'footer': 'VC'},
    'table':             {'title': 'Tableau',
                          'headers': ['Col A', 'Col B', 'Col C'],
                          'rows': [['R1A', 'R1B', 'R1C'], ['R2A', 'R2B', 'R2C']], 'footer': 'VC'},
    'team_grid':         {'title': 'Équipe', 'members': [
        {'name': 'Alice', 'role': 'CEO', 'icon': '👩‍💼'},
        {'name': 'Bob',   'role': 'CTO', 'icon': '👨‍💻'},
        {'name': 'Claire','role': 'CFO', 'icon': '👩‍🔬'},
    ], 'footer': 'VC'},
    'stat_banner':       {'title': 'Stats', 'stats': [
        {'value': '200+', 'label': 'Clients'},
        {'value': '98%',  'label': 'Satisfaction'},
        {'value': '15',   'label': 'Pays'},
    ], 'footer': 'VC'},
    'icon_row':          {'title': 'Icônes', 'items': [
        {'icon': '★', 'title': 'Qualité',  'body': 'Top'},
        {'icon': '⚡', 'title': 'Rapidité', 'body': 'Fast'},
        {'icon': '✓', 'title': 'Fiabilité','body': 'Sure'},
    ], 'footer': 'VC'},
    'section_break':     {'title': 'Section 2', 'subtitle': 'Sous-titre', 'number': '02', 'footer': 'VC'},
    'photo_text':        {'title': 'Photo + Texte', 'body': 'Contenu.', 'points': ['Point 1', 'Point 2'], 'footer': 'VC'},
    'numbered_features': {'title': 'Fonctionnalités', 'features': [
        {'number': '01', 'title': 'Feature A', 'body': 'Desc A'},
        {'number': '02', 'title': 'Feature B', 'body': 'Desc B'},
        {'number': '03', 'title': 'Feature C', 'body': 'Desc C'},
    ], 'footer': 'VC'},
    'side_panel':        {'title': 'Side Panel', 'panel_title': 'Panel',
                          'panel_items': ['Item 1', 'Item 2', 'Item 3'],
                          'body': 'Contenu principal.', 'footer': 'VC'},
    'circle_stats':      {'title': 'Stats Cercles', 'stats': [
        {'value': '75%', 'label': 'Stat A'},
        {'value': '50%', 'label': 'Stat B'},
        {'value': '90%', 'label': 'Stat C'},
    ], 'footer': 'VC'},
    'mission_vision':    {'title': 'Mission & Vision',
                          'mission': 'Notre mission.',
                          'vision': 'Notre vision.',
                          'values': ['Valeur 1', 'Valeur 2'], 'footer': 'VC'},
    'photo_grid':        {'title': 'Grille Photos', 'items': [
        {'title': 'Photo 1', 'body': 'Légende 1'},
        {'title': 'Photo 2', 'body': 'Légende 2'},
        {'title': 'Photo 3', 'body': 'Légende 3'},
    ], 'footer': 'VC'},
    'pricing_table':     {'title': 'Tarification', 'tiers': [
        {'name': 'Basic', 'price': '9€',  'features': ['F1', 'F2'], 'highlight': False},
        {'name': 'Pro',   'price': '29€', 'features': ['F1', 'F2', 'F3'], 'highlight': True},
        {'name': 'Ent.',  'price': '99€', 'features': ['F1', 'F2', 'F3', 'F4'], 'highlight': False},
    ], 'footer': 'VC'},
    'hub_spoke':         {'title': 'Hub & Spoke',
                          'center': {'title': 'Centre', 'body': 'Nœud central'},
                          'items': [
                              {'title': 'Spoke 1', 'body': 'Desc'},
                              {'title': 'Spoke 2', 'body': 'Desc'},
                              {'title': 'Spoke 3', 'body': 'Desc'},
                              {'title': 'Spoke 4', 'body': 'Desc'},
                          ], 'footer': 'VC'},
    'competitor_matrix': {'title': 'Concurrence',
                          'competitors': ['Nous', 'Conc. A', 'Conc. B'],
                          'criteria': ['Prix', 'Qualité', 'Service'],
                          'scores': {
                              'Nous':    [4, 5, 4],
                              'Conc. A': [3, 4, 3],
                              'Conc. B': [5, 3, 4],
                          }, 'footer': 'VC'},
    'pest_analysis':     {'title': 'PEST',
                          'political':     ['P1', 'P2'],
                          'economic':      ['E1', 'E2'],
                          'social':        ['S1', 'S2'],
                          'technological': ['T1', 'T2'], 'footer': 'VC'},
    'diamond_icons':     {'title': 'Diamant', 'items': [
        {'title': 'Nord',  'body': 'Desc', 'icon': '▲'},
        {'title': 'Est',   'body': 'Desc', 'icon': '▶'},
        {'title': 'Sud',   'body': 'Desc', 'icon': '▼'},
        {'title': 'Ouest', 'body': 'Desc', 'icon': '◀'},
    ], 'footer': 'VC'},
    'market_sizing':     {'title': 'Taille Marché',
                          'tam': {'value': '10Md€', 'label': 'TAM'},
                          'sam': {'value': '2Md€',  'label': 'SAM'},
                          'som': {'value': '200M€', 'label': 'SOM'}, 'footer': 'VC'},
    'chevron_flow':      {'title': 'Chevrons', 'steps': [
        {'title': 'Étape 1', 'body': 'Desc'},
        {'title': 'Étape 2', 'body': 'Desc'},
        {'title': 'Étape 3', 'body': 'Desc'},
        {'title': 'Étape 4', 'body': 'Desc'},
    ], 'footer': 'VC'},
    'venn':              {'title': 'Venn', 'circles': [
        {'title': 'Cercle A', 'items': ['A1', 'A2']},
        {'title': 'Cercle B', 'items': ['B1', 'B2']},
        {'title': 'Cercle C', 'items': ['C1', 'C2']},
    ], 'intersection': 'Zone commune', 'footer': 'VC'},
    'icon_grid':         {'title': 'Grille Icônes', 'items': [
        {'icon': '★', 'title': 'Item 1', 'body': 'Desc'},
        {'icon': '⚡', 'title': 'Item 2', 'body': 'Desc'},
        {'icon': '✓', 'title': 'Item 3', 'body': 'Desc'},
        {'icon': '◆', 'title': 'Item 4', 'body': 'Desc'},
        {'icon': '●', 'title': 'Item 5', 'body': 'Desc'},
        {'icon': '▲', 'title': 'Item 6', 'body': 'Desc'},
    ], 'footer': 'VC'},
    'text_hero':         {'title': 'Hero', 'hero_word': 'IMPACT',
                          'subtitle': 'Sous-titre', 'body': 'Description.', 'footer': 'VC'},
    'org_chart':         {'title': 'Organigramme',
                          'root': {'name': 'CEO', 'title': 'DG'},
                          'children': [
                              {'name': 'CTO', 'title': 'Tech'},
                              {'name': 'CFO', 'title': 'Finance'},
                              {'name': 'CMO', 'title': 'Marketing'},
                          ], 'footer': 'VC'},
}

LAYOUT_FNS = [
    ('cover',             main.layout_cover_v4),
    ('section',           main.layout_section_v4),
    ('fulltext',          main.layout_fulltext_v4),
    ('closing',           main.layout_closing_v4),
    ('quote',             main.layout_quote_v4),
    ('list_numbered',     main.layout_list_numbered_v4),
    ('list_cards',        main.layout_list_cards_v4),
    ('col3',              main.layout_col3_v4),
    ('twocol',            main.layout_twocol_v4),
    ('stathero',          main.layout_stathero_v4),
    ('infographic',       main.layout_infographic_v4),
    ('timeline',          main.layout_timeline_v4),
    ('processflow',       main.layout_processflow_v4),
    ('kpi_grid',          main.layout_kpi_grid_v4),
    ('funnel',            main.layout_funnel_v4),
    ('barchart',          main.layout_barchart_v4),
    ('linechart',         main.layout_linechart_v4),
    ('piechart',          main.layout_piechart_v4),
    ('waterfall',         main.layout_waterfall_v4),
    ('radar',             main.layout_radar_v4),
    ('pyramid',           main.layout_pyramid_v4),
    ('cycle',             main.layout_cycle_v4),
    ('roadmap',           main.layout_roadmap_v4),
    ('stackedbar',        main.layout_stackedbar_v4),
    ('beforeafter',       main.layout_beforeafter_v4),
    ('entity',            main.layout_entity_v4),
    ('conclusion',        main.layout_conclusion_v4),
    ('highlight',         main.layout_highlight_v4),
    ('agenda',            main.layout_agenda_v4),
    ('matrix',            main.layout_matrix_v4),
    ('swot',              main.layout_swot_v4),
    ('proscons',          main.layout_proscons_v4),
    ('table',             main.layout_table_v4),
    ('team_grid',         main.layout_team_grid_v4),
    ('stat_banner',       main.layout_stat_banner_v4),
    ('icon_row',          main.layout_icon_row_v4),
    ('section_break',     main.layout_section_break_v4),
    ('photo_text',        main.layout_photo_text_v4),
    ('numbered_features', main.layout_numbered_features_v4),
    ('side_panel',        main.layout_side_panel_v4),
    ('circle_stats',      main.layout_circle_stats_v4),
    ('mission_vision',    main.layout_mission_vision_v4),
    ('photo_grid',        main.layout_photo_grid_v4),
    ('pricing_table',     main.layout_pricing_table_v4),
    ('hub_spoke',         main.layout_hub_spoke_v4),
    ('competitor_matrix', main.layout_competitor_matrix_v4),
    ('pest_analysis',     main.layout_pest_analysis_v4),
    ('diamond_icons',     main.layout_diamond_icons_v4),
    ('market_sizing',     main.layout_market_sizing_v4),
    ('chevron_flow',      main.layout_chevron_flow_v4),
    ('venn',              main.layout_venn_v4),
    ('icon_grid',         main.layout_icon_grid_v4),
    ('text_hero',         main.layout_text_hero_v4),
    ('org_chart',         main.layout_org_chart_v4),
]


# ══════════════════════════════════════════════════════════════════════════════
# SECTION A — analyze_template_v4
# ══════════════════════════════════════════════════════════════════════════════
def test_analyze_template(fails: list) -> None:
    """Vérifie que analyze_template_v4 retourne un tp dict valide."""
    print("\n─── Section A : analyze_template_v4 ───────────────────────────────")
    REQUIRED_KEYS = [
        'theme', 'layout_map', 'font', 'accent_cycle',
        'card_bg_light', 'card_bg_mid', 'bg_type', 'bg_is_dark',
        'bg_rich', 'W', 'H',
    ]
    REQUIRED_LAYOUT_KEYS = ['cover', 'blank', 'text', 'section', 'closing']

    # Test 1 — prs vierge (nouveau fichier)
    try:
        prs_blank = Presentation()
        tp = main.analyze_template_v4(prs_blank)
        missing = [k for k in REQUIRED_KEYS if k not in tp]
        if missing:
            raise AssertionError(f"Clés manquantes dans tp : {missing}")
        missing_lm = [k for k in REQUIRED_LAYOUT_KEYS if k not in tp['layout_map']]
        if missing_lm:
            raise AssertionError(f"layout_map manque : {missing_lm}")
        print("  ✓  analyze_template_v4 (prs vierge)")
    except Exception as e:
        print(f"  ✗  analyze_template_v4 (prs vierge) : {e}")
        fails.append(('analyze_template_v4[blank]', str(e), traceback.format_exc()))

    # Test 2 — sur test_v4_output.pptx si dispo (PPTX réel avec slides)
    import os
    if os.path.exists('test_v4_output.pptx'):
        try:
            prs_real = Presentation('test_v4_output.pptx')
            tp2 = main.analyze_template_v4(prs_real)
            missing2 = [k for k in REQUIRED_KEYS if k not in tp2]
            if missing2:
                raise AssertionError(f"Clés manquantes : {missing2}")
            print("  ✓  analyze_template_v4 (test_v4_output.pptx)")
        except Exception as e:
            print(f"  ✗  analyze_template_v4 (test_v4_output.pptx) : {e}")
            fails.append(('analyze_template_v4[real]', str(e), traceback.format_exc()))

    # Test 3 — le tp retourné ne contient PAS d'anciennes clés supprimées
    try:
        prs_blank2 = Presentation()
        tp3 = main.analyze_template_v4(prs_blank2)
        forbidden = [k for k in ('accents', 'dk1', 'lt1') if k in tp3]
        if forbidden:
            raise AssertionError(
                f"Clés obsolètes présentes dans tp (provoquent KeyError dans les layouts) : {forbidden}"
            )
        print("  ✓  analyze_template_v4 — pas de clés obsolètes (accents/dk1/lt1)")
    except Exception as e:
        print(f"  ✗  analyze_template_v4 (clés obsolètes) : {e}")
        fails.append(('analyze_template_v4[obsolete_keys]', str(e), traceback.format_exc()))


# ══════════════════════════════════════════════════════════════════════════════
# SECTION B — Construction du prompt planner (plan_presentation_v4)
# ══════════════════════════════════════════════════════════════════════════════
def test_planner_prompt(fails: list) -> None:
    """
    Vérifie que la construction du user_prompt de plan_presentation_v4
    ne lève pas d'exception (KeyError sur accolades non-échappées, etc.).
    N'appelle PAS l'API Claude.
    """
    print("\n─── Section B : Construction du prompt planner ────────────────────")

    # Extraire _V4_PLANNER_USER depuis main (chargé au module level)
    planner_user = getattr(main, '_V4_PLANNER_USER', None)
    planner_doc  = getattr(main, '_V4_DOC_INJECT', None)

    if planner_user is None:
        msg = '_V4_PLANNER_USER introuvable dans main.py'
        print(f"  ✗  {msg}")
        fails.append(('planner_prompt[missing]', msg, ''))
        return

    test_cases = [
        # (label, prompt, nb_slides, primary, accent, font, document_content)
        ('cas nominal',        'Présentation France 2024', 8, '009CEA', 'ED0000', 'Arial',   ''),
        ('prompt avec {accolades}', 'Sujet {test} & {foo}', 5, '64D700', '4632FF', 'Calibri', ''),
        ('avec document',      'Synthèse rapport Q1',     10, 'FFFFFF', '000000', 'Verdana',
         'Contenu {document} avec des {accolades} partout.'),
    ]

    for label, prompt, nb_slides, primary, accent, font, doc in test_cases:
        try:
            # Reproduit exactement la logique de plan_presentation_v4
            user_prompt = planner_user
            for _ph, _val in [
                ('{prompt}',    prompt),
                ('{nb_slides}', str(nb_slides)),
                ('{primary}',   primary),
                ('{accent}',    accent),
                ('{font}',      font),
            ]:
                user_prompt = user_prompt.replace(_ph, _val)
            user_prompt = user_prompt.replace('{{', '{').replace('}}', '}')

            if doc and planner_doc:
                _ = planner_doc.format(document_content=doc)

            # Vérifications minimales
            assert str(nb_slides) in user_prompt, "nb_slides absent du prompt"
            assert font in user_prompt,             "font absent du prompt"
            print(f"  ✓  prompt planner ({label})")
        except Exception as e:
            print(f"  ✗  prompt planner ({label}) : {e}")
            fails.append((f'planner_prompt[{label}]', str(e), traceback.format_exc()))

    # Test supplémentaire : vérifier qu'aucune accolade non-échappée ne subsiste
    # après substitution (autrement dit : pas de {variable_inconnue} dans le prompt final)
    try:
        import re
        user_prompt = planner_user
        for _ph, _val in [
            ('{prompt}',    'test'),
            ('{nb_slides}', '8'),
            ('{primary}',   '009CEA'),
            ('{accent}',    'ED0000'),
            ('{font}',      'Calibri'),
        ]:
            user_prompt = user_prompt.replace(_ph, _val)
        user_prompt = user_prompt.replace('{{', '{').replace('}}', '}')

        # Après substitution, il ne doit plus rester de { non appariées
        # (le prompt final est du texte pur, pas un format string)
        # On vérifie juste que le résultat ne fait pas planter .format() avec un dict vide
        try:
            user_prompt.format_map({})   # doit lever KeyError ou pas, mais PAS ValueError
        except KeyError:
            pass   # normal — certains {mot} peuvent subsister comme texte pour Claude
        except ValueError as ve:
            raise AssertionError(f"Accolades mal formées dans le prompt final : {ve}")
        print("  ✓  prompt planner — accolades bien formées après substitution")
    except AssertionError as e:
        print(f"  ✗  prompt planner (accolades) : {e}")
        fails.append(('planner_prompt[braces]', str(e), traceback.format_exc()))


# ══════════════════════════════════════════════════════════════════════════════
# SECTION C — 54 layout_*_v4
# ══════════════════════════════════════════════════════════════════════════════
def test_layouts(fails: list) -> None:
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    tp = dict(TP)
    seen  = set()
    total = 0

    print(f"\n─── Section C : layout_*_v4 ({len(LAYOUT_FNS)} fonctions) ────────────────")

    for name, fn in LAYOUT_FNS:
        if name in seen:
            continue
        seen.add(name)
        total += 1
        content = CONTENT.get(name, {'title': name, 'footer': 'VC'})
        try:
            fn(prs, content, tp)
            print(f"  ✓  {name}")
        except Exception as e:
            tb = traceback.format_exc()
            print(f"  ✗  {name}: {e}")
            print(f"     └─ {tb.strip().splitlines()[-1]}")
            fails.append((name, str(e), tb))

    out = 'test_v4_output.pptx'
    try:
        prs.save(out)
        print(f"\n  PPTX sauvegardé : {out} ({len(prs.slides)} slides)")
    except Exception as e:
        print(f"\n  ERREUR sauvegarde PPTX : {e}")
        fails.append(('pptx_save', str(e), traceback.format_exc()))

    return total


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def run():
    all_fails = []

    test_analyze_template(all_fails)
    test_planner_prompt(all_fails)
    n_layouts = test_layouts(all_fails)

    n_ok   = n_layouts - sum(1 for f in all_fails if f[0] not in
                             ('analyze_template_v4[blank]',
                              'analyze_template_v4[real]',
                              'analyze_template_v4[obsolete_keys]',
                              'planner_prompt[cas nominal]',
                              'planner_prompt[prompt avec {accolades}]',
                              'planner_prompt[avec document]',
                              'planner_prompt[braces]',
                              'pptx_save'))
    total  = len(all_fails)

    print(f"\n{'='*60}")
    if total == 0:
        print(f"  RÉSULTAT : TOUS LES TESTS OK — 0 ÉCHEC(S)")
    else:
        print(f"  RÉSULTAT : {total} ÉCHEC(S) — voir détails ci-dessous")
    print(f"{'='*60}")

    if all_fails:
        print(f"\n{'─'*60}")
        print("  DÉTAIL DES ÉCHECS :")
        print(f"{'─'*60}")
        for name, err, tb in all_fails:
            print(f"\n  ── {name} ──\n  {err}")
            if tb:
                print(tb)

    sys.exit(len(all_fails))


if __name__ == '__main__':
    run()

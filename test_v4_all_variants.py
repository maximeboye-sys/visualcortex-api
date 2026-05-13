"""
test_v4_all_variants.py — Comprehensive variant + stress test for all 54 layouts.

Section A: All-variants — each layout × 5 variants = 270 slides
Section B: Stress content — each layout × 1 run with rich/large content = 54 slides
Section C: Bounds check on both PPTXs

Usage: python test_v4_all_variants.py
Exit code 0 = all clean, 1 = failures found.
"""
import sys
import traceback
from pptx import Presentation

import importlib.util
spec = importlib.util.spec_from_file_location('main', 'main.py')
main = importlib.util.module_from_spec(spec)
spec.loader.exec_module(main)

from test_v4_layouts import TP, CONTENT, LAYOUT_FNS
import check_bounds_v4 as CB

# ── Stress content: rich data to stress-test every layout ──────────────────
STRESS = {
    'cover': {'title': 'Transformation numérique & stratégie de croissance 2025-2030',
              'subtitle': 'Analyse complète des opportunités, risques et perspectives de marché',
              'footer': 'VisualCortex'},
    'section': {'number': '03', 'title': 'Résultats & recommandations stratégiques', 'footer': 'VC'},
    'fulltext': {'title': 'Analyse approfondie des dynamiques de marché',
                 'paragraphs': [
                     'Le marché mondial connaît une transformation sans précédent, portée par la digitalisation accélérée des usages et la montée en puissance de l\'intelligence artificielle dans tous les secteurs économiques.',
                     'Face à ces évolutions, les entreprises doivent adapter leur modèle opérationnel, investir dans les compétences digitales et repenser leur chaîne de valeur pour rester compétitives à l\'horizon 2030.',
                 ], 'footer': 'VC'},
    'closing': {'title': 'Merci pour votre attention',
                'subtitle': 'Des questions ou remarques ? Contactez-nous à contact@visualcortex.io',
                'footer': 'VC'},
    'quote': {'quote': 'L\'innovation n\'est pas simplement une question de technologie — c\'est une transformation profonde de la culture, des processus et de la relation client.',
              'author': 'Marie Dupont, Directrice Innovation', 'footer': 'VC'},
    'list_numbered': {'title': 'Priorités stratégiques 2025', 'items': [
        {'title': 'Transformation digitale accélérée', 'body': 'Migration cloud, modernisation des systèmes legacy et automatisation des processus critiques.'},
        {'title': 'Excellence opérationnelle', 'body': 'Optimisation des flux, réduction des coûts de 20% et amélioration de la qualité de service.'},
        {'title': 'Expérience client premium', 'body': 'Personnalisation IA, omnicanalité et réduction du temps de réponse à moins de 2 heures.'},
        {'title': 'Talent & culture data-driven', 'body': 'Formation continue, recrutement de profils digitaux et création d\'une culture d\'innovation.'},
        {'title': 'Durabilité & RSE', 'body': 'Neutralité carbone en 2027, économie circulaire et reporting ESG conforme aux normes CSRD.'},
        {'title': 'Innovation produit continu', 'body': 'Lancement de 3 nouvelles fonctionnalités par trimestre avec cycles de feedback courts.'},
    ], 'footer': 'VC'},
    'list_cards': {'title': 'Solutions clés', 'cards': [
        {'title': 'Plateforme unifiée', 'body': 'Interface centralisée pour tous les flux de données en temps réel.'},
        {'title': 'Analytics avancée', 'body': 'Tableaux de bord prédictifs et alertes automatiques basés sur le ML.'},
        {'title': 'Sécurité renforcée', 'body': 'Chiffrement bout-en-bout, authentification multi-facteurs et audit continu.'},
        {'title': 'API ouverte', 'body': 'Intégration native avec 200+ outils du marché via REST et webhooks.'},
        {'title': 'Support 24/7', 'body': 'Équipe dédiée, SLA garanti à 99,9% et documentation exhaustive.'},
    ], 'footer': 'VC'},
    'col3': {'title': 'Trois axes stratégiques', 'columns': [
        {'title': 'Croissance', 'items': ['Nouveaux marchés', 'Partenariats stratégiques', 'Acquisition clients', 'Expansion géo']},
        {'title': 'Efficacité', 'items': ['Automatisation RPA', 'Lean management', 'KPIs temps réel', 'Réduction délais']},
        {'title': 'Innovation', 'items': ['R&D interne', 'Labs & startups', 'Brevets déposés', 'Veille techno']},
    ], 'footer': 'VC'},
    'twocol': {'title': 'Analyse comparative',
               'col_a': {'title': 'Avantages concurrentiels', 'items': ['Technologie propriétaire', 'Base clients fidèle', 'Marque reconnue', 'Équipe expérimentée']},
               'col_b': {'title': 'Axes d\'amélioration', 'items': ['Présence internationale', 'Portfolio produits', 'Canaux de vente', 'Support technique']},
               'footer': 'VC'},
    'stathero': {'value': '€4.2Md', 'label': 'Chiffre d\'affaires total 2024',
                 'context': 'Croissance de +23% vs 2023 — Record historique', 'footer': 'VC'},
    'infographic': {'title': 'Taux d\'adoption par segment', 'value': '87%', 'label': 'Adoption globale',
                    'context': 'Étude panel 2024 — 12 000 répondants',
                    'bars': [
                        {'label': 'Grandes entreprises (>500 salariés)', 'percent': 94},
                        {'label': 'ETI (50-500 salariés)', 'percent': 87},
                        {'label': 'PME (10-50 salariés)', 'percent': 72},
                        {'label': 'TPE (<10 salariés)', 'percent': 48},
                    ], 'footer': 'VC'},
    'timeline': {'title': 'Jalons du projet', 'steps': [
        {'date': 'Jan 2024', 'title': 'Lancement phase pilote', 'body': 'Déploiement auprès de 50 clients bêta avec suivi personnalisé.'},
        {'date': 'Mar 2024', 'title': 'Première levée de fonds', 'body': 'Série A de 8M€ avec trois fonds d\'investissement spécialisés.'},
        {'date': 'Juin 2024', 'title': 'Lancement commercial', 'body': 'Ouverture à tous les clients avec pricing segmenté par usage.'},
        {'date': 'Sep 2024', 'title': 'Expansion Europe', 'body': 'Présence dans 5 nouveaux pays : Allemagne, Espagne, Italie, NL, BE.'},
        {'date': 'Jan 2025', 'title': 'Objectif 1000 clients', 'body': 'Cap symbolique atteint avec un NPS de 72 et un churn < 3%.'},
    ], 'footer': 'VC'},
    'processflow': {'title': 'Processus de delivery', 'steps': [
        {'title': 'Découverte', 'body': 'Analyse des besoins client et cartographie des cas d\'usage prioritaires.'},
        {'title': 'Conception', 'body': 'Wireframes, architecture technique et validation fonctionnelle.'},
        {'title': 'Développement', 'body': 'Sprints agiles de 2 semaines avec revues quotidiennes.'},
        {'title': 'Validation', 'body': 'Tests unitaires, d\'intégration et UAT avec le client final.'},
        {'title': 'Déploiement', 'body': 'Mise en production progressive avec monitoring temps réel.'},
    ], 'footer': 'VC'},
    'kpi_grid': {'title': 'Tableau de bord opérationnel', 'kpis': [
        {'value': '€2.1Md', 'label': 'ARR total', 'sublabel': '+31% YoY'},
        {'value': '94%',    'label': 'Taux rétention', 'sublabel': 'Best-in-class'},
        {'value': '4 200',  'label': 'Clients actifs', 'sublabel': '+580 ce trimestre'},
        {'value': '72',     'label': 'NPS global', 'sublabel': 'Objectif 75 en Q2'},
    ], 'footer': 'VC'},
    'funnel': {'title': 'Entonnoir de conversion', 'steps': [
        {'label': 'Visiteurs uniques',  'value': '450 000'},
        {'label': 'Leads qualifiés',    'value': '28 000'},
        {'label': 'Opportunités',       'value': '6 200'},
        {'label': 'Propositions envoyées', 'value': '1 800'},
        {'label': 'Clients signés',     'value': '420'},
    ], 'footer': 'VC'},
    'barchart': {'title': 'Revenus par région (M€)', 'categories': ['France', 'Allemagne', 'Espagne', 'Italie', 'UK', 'Autres'],
                 'series': [
                     {'name': '2023', 'values': [42, 28, 18, 14, 31, 22]},
                     {'name': '2024', 'values': [58, 37, 24, 19, 42, 31]},
                 ], 'footer': 'VC'},
    'linechart': {'title': 'Évolution mensuelle du MRR (k€)', 'categories': ['Jan', 'Fév', 'Mar', 'Avr', 'Mai', 'Juin'],
                  'series': [
                      {'name': 'MRR', 'values': [120, 145, 162, 178, 195, 218]},
                      {'name': 'Objectif', 'values': [130, 150, 165, 180, 200, 220]},
                  ], 'footer': 'VC'},
    'piechart': {'title': 'Mix produits par CA', 'slices': [
        {'label': 'Licences SaaS', 'value': 52},
        {'label': 'Services pro', 'value': 24},
        {'label': 'Support & maintenance', 'value': 14},
        {'label': 'Formation', 'value': 10},
    ], 'footer': 'VC'},
    'waterfall': {'title': 'Pont de CA (M€)', 'items': [
        {'label': 'CA 2023', 'value': 155},
        {'label': 'Nouveaux clients', 'value': 48},
        {'label': 'Upsell / cross-sell', 'value': 27},
        {'label': 'Churn & downgrades', 'value': -18},
        {'label': 'FX & autres', 'value': -4},
    ], 'footer': 'VC'},
    'radar': {'title': 'Positionnement compétitif', 'axes': ['Prix', 'Qualité', 'Innovation', 'Support', 'Écosystème'],
              'series': [
                  {'name': 'Nous', 'values': [4, 5, 5, 4, 3]},
                  {'name': 'Concurrent', 'values': [3, 3, 4, 3, 5]},
              ], 'footer': 'VC'},
    'pyramid': {'title': 'Hiérarchie des besoins clients', 'levels': [
        {'label': 'Vision long terme', 'value': 5},
        {'label': 'Différenciation', 'value': 15},
        {'label': 'Performance',      'value': 30},
        {'label': 'Fiabilité',        'value': 50},
    ], 'footer': 'VC'},
    'cycle': {'title': 'Cycle d\'amélioration continue', 'steps': [
        {'title': 'Planifier', 'body': 'Définir les objectifs, KPIs et ressources nécessaires.'},
        {'title': 'Exécuter',  'body': 'Déployer le plan avec agilité et mesures de contrôle.'},
        {'title': 'Mesurer',   'body': 'Analyser les données en temps réel et identifier les écarts.'},
        {'title': 'Apprendre', 'body': 'Capitaliser sur les enseignements pour le prochain cycle.'},
    ], 'footer': 'VC'},
    'roadmap': {'title': 'Feuille de route produit', 'phases': [
        {'label': 'Q1 2025', 'milestones': ['Refonte UX', 'API v3', 'Intégration SSO']},
        {'label': 'Q2 2025', 'milestones': ['Mobile app', 'IA prédictive', 'Export avancé']},
        {'label': 'Q3 2025', 'milestones': ['Marketplace', 'Multi-tenant', 'Compliance EU']},
        {'label': 'Q4 2025', 'milestones': ['V4 GA', 'Enterprise tier', 'Partner program']},
    ], 'footer': 'VC'},
    'stackedbar': {'title': 'Répartition des coûts par poste', 'categories': ['2021', '2022', '2023', '2024'],
                   'series': [
                       {'name': 'R&D', 'values': [22, 28, 35, 42]},
                       {'name': 'Ventes', 'values': [18, 22, 26, 30]},
                       {'name': 'G&A', 'values': [12, 14, 15, 16]},
                   ], 'footer': 'VC'},
    'beforeafter': {'title': 'Avant / Après notre solution',
                    'before': {'title': 'Situation actuelle', 'items': [
                        'Processus manuels chronophages', 'Données fragmentées en silos',
                        'Reporting mensuel avec erreurs', 'Temps de réponse > 48h',
                    ]},
                    'after': {'title': 'Avec notre solution', 'items': [
                        'Automatisation complète end-to-end', 'Source unique de vérité',
                        'Dashboards temps réel sans erreur', 'Alertes instantanées <5min',
                    ]}, 'footer': 'VC'},
    'entity': {'title': 'Profils des segments clients', 'entities': [
        {'name': 'PME Tech', 'icon': '💡', 'items': ['10-50 employés', 'Budget SaaS élevé', 'Early adopters', 'ROI rapide']},
        {'name': 'ETI Industrie', 'icon': '🏭', 'items': ['50-500 employés', 'Processus complexes', 'Décision longue', 'Intégration ERP']},
        {'name': 'Grand Compte', 'icon': '🏢', 'items': ['>500 employés', 'Sécurité critique', 'On-premise option', 'SLA contractuel']},
    ], 'footer': 'VC'},
    'conclusion': {'title': 'Synthèse & prochaines étapes',
                   'cards': [
                       {'icon': '✓', 'title': 'Objectifs atteints', 'body': 'Les 3 KPIs prioritaires ont été dépassés avec une croissance de +23%.'},
                       {'icon': '✓', 'title': 'Équipe renforcée', 'body': 'Recrutement de 12 profils seniors en Data, Produit et Commercial.'},
                       {'icon': '✓', 'title': 'Base clients solide', 'body': '4 200 clients actifs avec un NPS de 72 et un churn historiquement bas.'},
                   ],
                   'sidebar_title': 'Ce qu\'il faut retenir',
                   'sidebar_quote': 'Nous avons posé des bases solides pour une croissance durable à 5 ans.',
                   'sidebar_cta': 'Lancer la phase 2', 'footer': 'VC'},
    'highlight': {'title': 'Message stratégique', 'highlight': 'Notre avantage compétitif repose sur une technologie propriétaire difficile à répliquer.',
                  'body': 'Cette position unique nous permet de capturer 35% de la valeur créée et de maintenir des marges brutes supérieures à 80% sur nos offres SaaS.',
                  'points': ['Brevets déposés', 'Équipe R&D senior', 'Data moat', 'Effets réseau'],
                  'footer': 'VC'},
    'agenda': {'title': 'Ordre du jour', 'items': [
        {'number': '01', 'title': 'Revue des résultats Q4 2024', 'duration': '20 min'},
        {'number': '02', 'title': 'Analyse des opportunités 2025', 'duration': '25 min'},
        {'number': '03', 'title': 'Plan d\'investissement R&D', 'duration': '20 min'},
        {'number': '04', 'title': 'Roadmap produit & ressources', 'duration': '20 min'},
        {'number': '05', 'title': 'Discussion & décisions', 'duration': '15 min'},
    ], 'footer': 'VC'},
    'matrix': {'title': 'Matrice attractivité / faisabilité', 'items': [
        {'label': 'Plateforme Enterprise', 'x': 0.9, 'y': 0.8},
        {'label': 'Mobile Premium',        'x': 0.7, 'y': 0.9},
        {'label': 'Analytics IA',          'x': 0.8, 'y': 0.5},
        {'label': 'Marketplace partenaires','x': 0.6, 'y': 0.7},
        {'label': 'Support automatisé',    'x': 0.4, 'y': 0.8},
    ], 'footer': 'VC'},
    'swot': {'title': 'Analyse SWOT stratégique',
             'strengths':     ['Technologie propriétaire brevétée', 'Marque reconnue B2B SaaS', 'Équipe technique expérimentée', 'NPS de 72 — meilleur secteur'],
             'weaknesses':    ['Couverture internationale limitée', 'Dépendance à 3 grands comptes', 'Cycle de vente enterprise long'],
             'opportunities': ['Réglementations IA favorables', 'Marché sous-pénétré en Europe Sud', 'Consolidation sectorielle en cours'],
             'threats':       ['Entrée de Big Tech dans le segment', 'Volatilité des budgets IT', 'Rétention des talents tech'],
             'footer': 'VC'},
    'proscons': {'title': 'Pour / Contre : expansion internationale',
                 'pros': ['Adressable market x5', 'Diversification du risque client', 'Recrutement de talents globaux', 'Valorisation boostée IPO'],
                 'cons': ['Investissement initial de 4M€', 'Complexité réglementaire multi-pays', 'Risque dilution attention management', 'Délai de 18 mois avant rentabilité'],
                 'footer': 'VC'},
    'table': {'title': 'Comparaison des offres',
              'headers': ['Fonctionnalité', 'Starter', 'Pro', 'Enterprise'],
              'rows': [
                  ['Utilisateurs', '5', '25', 'Illimité'],
                  ['Stockage', '10 Go', '100 Go', '1 To'],
                  ['API calls/mois', '10k', '100k', 'Illimité'],
                  ['Support', 'Email', 'Chat 9h-18h', 'Dédié 24/7'],
                  ['SLA', '-', '99,5%', '99,99%'],
              ], 'footer': 'VC'},
    'team_grid': {'title': 'L\'équipe dirigeante', 'members': [
        {'name': 'Sophie Marchand',  'role': 'CEO & Co-fondatrice', 'icon': '👩‍💼'},
        {'name': 'Thomas Lebrun',    'role': 'CTO & Co-fondateur',  'icon': '👨‍💻'},
        {'name': 'Camille Fontaine', 'role': 'CMO',                  'icon': '👩‍🎨'},
        {'name': 'Julien Petit',     'role': 'CFO',                  'icon': '👨‍💼'},
        {'name': 'Isabelle Morel',   'role': 'VP Sales',             'icon': '👩‍💼'},
        {'name': 'Antoine Girard',   'role': 'VP Product',           'icon': '👨‍💻'},
    ], 'footer': 'VC'},
    'stat_banner': {'title': 'Indicateurs clés de performance', 'stats': [
        {'value': '4 200', 'label': 'Clients actifs', 'sublabel': '+580 ce trimestre', 'icon': '🏢'},
        {'value': '€2.1Md', 'label': 'ARR annualisé', 'sublabel': '+31% YoY', 'icon': '📈'},
        {'value': '94%', 'label': 'Taux de rétention', 'sublabel': 'Best-in-class secteur', 'icon': '🔒'},
        {'value': '72', 'label': 'Net Promoter Score', 'sublabel': 'Objectif 75 Q2', 'icon': '⭐'},
    ], 'footer': 'VC'},
    'icon_row': {'title': 'Nos piliers de valeur', 'items': [
        {'icon': '🚀', 'title': 'Performance', 'body': 'Latence < 50ms, uptime 99,99% garanti contractuellement.'},
        {'icon': '🔒', 'title': 'Sécurité',    'body': 'ISO 27001, SOC 2 Type II, chiffrement AES-256 au repos.'},
        {'icon': '🤝', 'title': 'Support',     'body': 'Customer Success dédié, onboarding guidé, SLA 4h critique.'},
        {'icon': '🧠', 'title': 'Intelligence','body': 'IA prédictive intégrée, anomaly detection, recommandations.'},
    ], 'footer': 'VC'},
    'section_break': {'title': 'Résultats & Perspectives', 'subtitle': 'Une croissance soutenue confirmée par les données', 'number': '03', 'footer': 'VC'},
    'photo_text': {'title': 'Notre approche terrain',
                   'body': 'En 2024, notre équipe a conduit plus de 200 entretiens clients dans 8 pays pour comprendre en profondeur les frictions et attentes du marché.',
                   'points': ['200+ entretiens clients', '8 pays couverts', '32 workshops de co-design', 'Insights validés sur panel 500 utilisateurs'],
                   'footer': 'VC'},
    'numbered_features': {'title': 'Fonctionnalités phares', 'features': [
        {'number': '01', 'title': 'Tableau de bord unifié', 'body': 'Vue 360° de toutes vos métriques métier en temps réel.'},
        {'number': '02', 'title': 'Alertes intelligentes',  'body': 'Détection d\'anomalies par ML avec seuils personnalisables.'},
        {'number': '03', 'title': 'Connecteurs natifs',    'body': '200+ intégrations avec vos outils existants (CRM, ERP, BI).'},
        {'number': '04', 'title': 'Rapports automatisés',  'body': 'Génération et envoi automatique de rapports aux parties prenantes.'},
        {'number': '05', 'title': 'Collaboration temps réel','body': 'Annotations, commentaires et partage sécurisé avec votre équipe.'},
    ], 'footer': 'VC'},
    'side_panel': {'title': 'Stratégie d\'acquisition', 'panel_title': 'Canaux clés',
                   'panel_items': ['SEO & Content', 'Partenaires revendeurs', 'Outbound SDR', 'Product-led growth', 'Events & webinaires'],
                   'body': 'Notre mix d\'acquisition s\'appuie sur un moteur inbound fort (contenu, SEO, communauté) combiné à une équipe outbound ciblant les comptes stratégiques >200 salariés.',
                   'footer': 'VC'},
    'circle_stats': {'title': 'Performance par dimension', 'stats': [
        {'value': '92%', 'label': 'Satisfaction client', 'sublabel': 'CSAT trimestriel'},
        {'value': '88%', 'label': 'Adoption produit',    'sublabel': 'MAU/licences actives'},
        {'value': '76%', 'label': 'Taux d\'upsell',      'sublabel': 'Base clients N-1'},
        {'value': '94%', 'label': 'Rétention globale',   'sublabel': 'Renouvellements'},
    ], 'footer': 'VC'},
    'mission_vision': {'title': 'Mission & Vision',
                       'mission': 'Démocratiser l\'accès à la data intelligence pour toutes les entreprises, quelle que soit leur taille ou leur secteur.',
                       'vision': 'Devenir la plateforme de référence européenne pour la gestion intelligente des données opérationnelles d\'ici 2028.',
                       'values': ['Transparence', 'Innovation responsable', 'Impact client', 'Excellence technique'],
                       'footer': 'VC'},
    'photo_grid': {'title': 'Témoignages clients', 'items': [
        {'title': 'Sophie M., DSI', 'body': 'Réduction de 60% du temps de reporting et gain de 15% en productivité équipe.'},
        {'title': 'Marc T., CEO PME', 'body': 'ROI positif dès le 3e mois — impossible à croire mais chiffres à l\'appui.'},
        {'title': 'Lucie B., CFO', 'body': 'Consolidation financière multi-entités en 10 minutes au lieu de 3 jours.'},
    ], 'footer': 'VC'},
    'pricing_table': {'title': 'Nos offres tarifaires', 'tiers': [
        {'name': 'Starter', 'price': '29€/mois', 'features': ['5 utilisateurs', '10 Go', 'Email support', 'Rapports standard'], 'highlight': False},
        {'name': 'Business', 'price': '99€/mois', 'features': ['25 utilisateurs', '100 Go', 'Chat support', 'Rapports avancés', 'API access', 'SSO inclus'], 'highlight': True},
        {'name': 'Enterprise', 'price': 'Sur devis', 'features': ['Illimité', '1 To', 'Support dédié 24/7', 'SLA 99,99%', 'On-premise option', 'Audit logs', 'Custom integrations'], 'highlight': False},
    ], 'footer': 'VC'},
    'hub_spoke': {'title': 'Écosystème de valeur', 'center': {'title': 'Plateforme VC', 'body': 'Nœud central d\'orchestration'},
                  'items': [
                      {'title': 'CRM & Ventes',   'body': 'Sync Salesforce, HubSpot'},
                      {'title': 'Finance & ERP',  'body': 'Intégration SAP, Oracle'},
                      {'title': 'Marketing',      'body': 'Marketo, Pardot, GA4'},
                      {'title': 'Ops & Support',  'body': 'Zendesk, Jira, ServiceNow'},
                      {'title': 'BI & Analytics', 'body': 'Power BI, Tableau, Looker'},
                      {'title': 'Communication',  'body': 'Slack, Teams, Zoom API'},
                  ], 'footer': 'VC'},
    'competitor_matrix': {'title': 'Analyse comparative concurrentielle',
                          'competitors': ['VisualCortex', 'Concurrent A', 'Concurrent B', 'Concurrent C', 'Concurrent D'],
                          'features': [
                              {'name': 'Temps réel',          'values': [True,  True,  False, True,  False]},
                              {'name': 'IA intégrée',         'values': [True,  False, True,  False, False]},
                              {'name': 'API ouverte',         'values': [True,  True,  True,  False, True]},
                              {'name': 'On-premise',          'values': [True,  False, True,  True,  False]},
                              {'name': 'SSO / SAML',          'values': [True,  True,  True,  True,  False]},
                              {'name': 'Multi-tenant',        'values': [True,  False, False, True,  True]},
                              {'name': 'SLA 99,9%',           'values': [True,  True,  False, False, True]},
                          ], 'footer': 'VC'},
    'pest_analysis': {'title': 'Analyse PEST — Environnement sectoriel',
                      'political':     ['Réglementation IA Act UE', 'Subventions transformation digitale', 'Normes RGPD renforcées'],
                      'economic':      ['Inflation des coûts IT', 'Pression sur les marges SaaS', 'Croissance budgets cloud +18%'],
                      'social':        ['Besoin de transparence data', 'Montée des usages no-code', 'Attentes RSE des clients grands comptes'],
                      'technological': ['Généralisation du LLM/IA', 'Edge computing émergent', 'API-first architecture dominant'],
                      'footer': 'VC'},
    'diamond_icons': {'title': 'Quatre leviers de transformation', 'items': [
        {'title': 'Stratégie',      'body': 'Vision claire, objectifs SMART et allocation ressources optimisée.',           'icon': '🎯'},
        {'title': 'Technologie',    'body': 'Stack moderne, cloud-native et architecture API-first scalable.',              'icon': '⚙️'},
        {'title': 'Organisation',   'body': 'Culture agile, équipes cross-fonctionnelles et leadership distribué.',         'icon': '🤝'},
        {'title': 'Data & IA',      'body': 'Décisions pilotées par les données, automatisation et intelligence prédictive.','icon': '🧠'},
    ], 'footer': 'VC'},
    'market_sizing': {'title': 'Taille et opportunité de marché',
                      'tam': {'value': '€48Md', 'label': 'TAM — Marché adressable total', 'growth': '+18%/an'},
                      'sam': {'value': '€8Md',  'label': 'SAM — Segment accessible',      'growth': '+24%/an'},
                      'som': {'value': '€450M', 'label': 'SOM — Part de marché visée',    'growth': 'Objectif 3 ans'},
                      'footer': 'VC'},
    'chevron_flow': {'title': 'Parcours d\'implémentation', 'steps': [
        {'title': 'Audit & Cadrage',   'body': 'Analyse des besoins, cartographie des processus et définition du scope.',         'stat_value': '2 sem', 'stat_label': 'Durée'},
        {'title': 'Configuration',     'body': 'Paramétrage de la plateforme, création des connecteurs et formation des admins.', 'stat_value': '4 sem', 'stat_label': 'Durée'},
        {'title': 'Pilote',            'body': 'Déploiement auprès d\'un groupe test de 20 utilisateurs avec feedback continu.',  'stat_value': '3 sem', 'stat_label': 'Durée'},
        {'title': 'Déploiement',       'body': 'Rollout complet, migration des données historiques et go-live officiel.',         'stat_value': '2 sem', 'stat_label': 'Durée'},
        {'title': 'Optimisation',      'body': 'Suivi des KPIs d\'adoption, ajustements et montée en compétences continue.',     'stat_value': 'Continu', 'stat_label': 'Phase'},
    ], 'footer': 'VC'},
    'venn': {'title': 'Zones d\'intersection stratégique', 'circles': [
        {'label': 'Clients grands comptes', 'items': ['Contrats pluriannuels', 'SLA élevé', 'Intégration profonde']},
        {'label': 'Marchés internationaux', 'items': ['Expansion rapide', 'Partenaires locaux', 'Multi-langue']},
        {'label': 'Innovation produit',     'items': ['R&D interne', 'Co-construction', 'Bêta-testeurs']},
    ], 'intersection': {'label': 'Zone de croissance optimale', 'icon': '🎯'}, 'footer': 'VC'},
    'icon_grid': {'title': 'Capacités de la plateforme', 'items': [
        {'icon': '📊', 'title': 'Dashboards',      'body': 'Tableaux de bord personnalisables en drag & drop.'},
        {'icon': '🔔', 'title': 'Alertes',         'body': 'Notifications push, email, SMS et webhook.'},
        {'icon': '🤖', 'title': 'IA Prédictive',   'body': 'Modèles ML entraînés sur vos données propres.'},
        {'icon': '🔌', 'title': 'Intégrations',    'body': '200+ connecteurs natifs avec mise à jour auto.'},
        {'icon': '📱', 'title': 'Mobile',          'body': 'App iOS & Android avec mode hors-ligne.'},
        {'icon': '🛡️', 'title': 'Sécurité',        'body': 'Chiffrement, audit logs et contrôle d\'accès granulaire.'},
    ], 'footer': 'VC'},
    'text_hero': {'title': 'Notre vision', 'hero_word': 'IMPACT',
                  'subtitle': 'Créer un impact mesurable pour chaque client',
                  'body': 'En 2024, nos clients ont gagné en moyenne 8 heures par semaine grâce à l\'automatisation de leurs processus data.',
                  'accent_word': '+23% ROI', 'footer': 'VC'},
    'org_chart': {'title': 'Structure organisationnelle', 'root': {'name': 'Sophie Marchand', 'title': 'CEO & Co-fondatrice'},
                  'children': [
                      {'name': 'Thomas Lebrun', 'title': 'CTO', 'children': [
                          {'name': 'Équipe R&D', 'title': '12 ingénieurs'},
                          {'name': 'Équipe DevOps', 'title': '4 SRE'},
                      ]},
                      {'name': 'Camille Fontaine', 'title': 'CMO'},
                      {'name': 'Julien Petit', 'title': 'CFO'},
                      {'name': 'Isabelle Morel', 'title': 'VP Sales'},
                  ], 'footer': 'VC'},
}


# ── Test runner ───────────────────────────────────────────────────────────────
def _make_prs():
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width  = int(13.33 * 914400)
    prs.slide_height = int(7.50  * 914400)
    return prs


def run_all_variants(out_pptx='test_v4_all_variants.pptx'):
    """Generate 54 × 5 = 270 slides (one per layout-variant combo)."""
    print('\n══════════════════════════════════════════════════════════════')
    print('  Section A — All Variants (54 × 5 = 270 slides)')
    print('══════════════════════════════════════════════════════════════')

    prs = _make_prs()
    labels = []
    fails = []
    n_ok  = 0

    for layout_name, fn in LAYOUT_FNS:
        base_content = dict(CONTENT.get(layout_name, {}))
        for v in range(5):
            label = f'{layout_name}/v{v}'
            content = dict(base_content, style=v)
            try:
                fn(prs, content, TP)
                labels.append(label)
                n_ok += 1
                print(f'  ✓  {label}')
            except Exception as e:
                fails.append((label, str(e), traceback.format_exc()))
                labels.append(label)
                prs.slides.add_slide(prs.slide_layouts[6])  # blank placeholder
                print(f'  ✗  {label}  ERROR: {e}')

    prs.save(out_pptx)
    print(f'\n  {n_ok}/270 OK — saved → {out_pptx}')
    return prs, labels, fails


def run_stress(out_pptx='test_v4_stress.pptx'):
    """Generate 54 slides with rich stress content (random variant)."""
    print('\n══════════════════════════════════════════════════════════════')
    print('  Section B — Stress Content (54 layouts, rich data)')
    print('══════════════════════════════════════════════════════════════')

    prs = _make_prs()
    labels = []
    fails = []
    n_ok  = 0

    for layout_name, fn in LAYOUT_FNS:
        label = f'{layout_name}/stress'
        content = dict(STRESS.get(layout_name, CONTENT.get(layout_name, {})))
        try:
            fn(prs, content, TP)
            labels.append(label)
            n_ok += 1
            print(f'  ✓  {label}')
        except Exception as e:
            fails.append((label, str(e), traceback.format_exc()))
            labels.append(label)
            prs.slides.add_slide(prs.slide_layouts[6])
            print(f'  ✗  {label}  ERROR: {e}')

    prs.save(out_pptx)
    print(f'\n  {n_ok}/54 OK — saved → {out_pptx}')
    return prs, labels, fails


def run_bounds(pptx_path, labels):
    """Run bounds check and return violations."""
    print(f'\n══════════════════════════════════════════════════════════════')
    print(f'  Section C — Bounds Check: {pptx_path}')
    print(f'══════════════════════════════════════════════════════════════')
    return CB.check_pptx(pptx_path, slide_labels=labels)


if __name__ == '__main__':
    total_fails = []

    # A: all variants
    _, labels_v, fails_v = run_all_variants()
    total_fails.extend(fails_v)
    viols_v = run_bounds('test_v4_all_variants.pptx', labels_v)

    # B: stress content
    _, labels_s, fails_s = run_stress()
    total_fails.extend(fails_s)
    viols_s = run_bounds('test_v4_stress.pptx', labels_s)

    # Summary
    print('\n══════════════════════════════════════════════════════════════')
    print('  FINAL SUMMARY')
    print('══════════════════════════════════════════════════════════════')
    if total_fails:
        print(f'\n  EXCEPTIONS ({len(total_fails)}):')
        for label, msg, _ in total_fails:
            print(f'    ✗  {label}: {msg}')
    else:
        print('  ✓  No exceptions.')

    all_viols = viols_v + viols_s
    if all_viols:
        print(f'\n  BOUNDS VIOLATIONS ({len(all_viols)}):')
        for v in all_viols:
            print(f'    ✗  {v["label"]}: {v["detail"]}')
    else:
        print('  ✓  No bounds violations.')

    n_problems = len(total_fails) + len(all_viols)
    print(f'\n  Total problems: {n_problems}')
    print('══════════════════════════════════════════════════════════════\n')

    sys.exit(1 if n_problems > 0 else 0)

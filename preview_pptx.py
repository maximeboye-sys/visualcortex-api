"""
preview_pptx.py — Render PPTX slides to PNG using Pillow.

Usage:
    python preview_pptx.py <file.pptx> [--slides 0,1,2 | all] [--out /tmp/preview]

Renders shapes (rects, rounded-rects, text) from each slide.
Does not handle: charts, images, master shapes beyond background colour.
Good enough to catch layout issues: overflow, overlap, wrong colours, gaps.
"""

import sys
import os
import argparse
from lxml import etree

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    sys.exit("pip install Pillow")

from pptx import Presentation

# ── Namespaces ───────────────────────────────────────────────────────────────
_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

def _a(t): return f'{{{_A}}}{t}'
def _p(t): return f'{{{_P}}}{t}'

# ── Scale: 100 px/inch ───────────────────────────────────────────────────────
_PPI = 100

def _px(emu): return int(int(emu) * _PPI / 914400)

# ── Colour helpers ───────────────────────────────────────────────────────────
_SCHEME_APPROX = {
    'dk1': (55, 70, 73), 'lt1': (255, 255, 255),
    'dk2': (31, 56, 100), 'lt2': (233, 239, 247),
    'accent1': (0, 156, 234), 'accent2': (237, 0, 0),
    'accent3': (64, 169, 0), 'accent4': (246, 106, 0),
    'accent5': (123, 47, 190), 'accent6': (0, 180, 216),
    'bg1': (255, 255, 255), 'bg2': (240, 240, 240),
    'tx1': (30, 30, 30), 'tx2': (80, 80, 80),
    'phClr': (120, 120, 120),
}

def _parse_color(el):
    """Return (r,g,b) from first colour element found, or None."""
    if el is None:
        return None
    srgb = el.find(_a('srgbClr'))
    if srgb is not None:
        h = srgb.get('val', '')
        if len(h) == 6:
            try:
                return int(h[:2],16), int(h[2:4],16), int(h[4:],16)
            except ValueError:
                pass
    scl = el.find(_a('schemeClr'))
    if scl is not None:
        return _SCHEME_APPROX.get(scl.get('val',''), (150, 150, 150))
    sysClr = el.find(_a('sysClr'))
    if sysClr is not None:
        lv = sysClr.get('lastClr', 'FFFFFF')
        try:
            return int(lv[:2],16), int(lv[2:4],16), int(lv[4:],16)
        except Exception:
            pass
    return None

def _solid_fill(spPr):
    sf = spPr.find(f'.//{_a("solidFill")}')
    return _parse_color(sf)

def _grad_fill_approx(spPr):
    """Approximate gradient: return colour of last stop."""
    gf = spPr.find(f'.//{_a("gradFill")}')
    if gf is None:
        return None
    stops = gf.findall(f'.//{_a("gs")}')
    if not stops:
        return None
    clr = _parse_color(stops[-1])
    return clr

def _fill(spPr):
    c = _solid_fill(spPr)
    if c:
        return c
    c = _grad_fill_approx(spPr)
    if c:
        return c
    if spPr.find(f'.//{_a("noFill")}') is not None:
        return None
    return None

# ── Position ─────────────────────────────────────────────────────────────────
def _xfrm(spPr):
    xf = spPr.find(_a('xfrm'))
    if xf is None:
        return None
    off = xf.find(_a('off'))
    ext = xf.find(_a('ext'))
    if off is None or ext is None:
        return None
    try:
        return (_px(off.get('x',0)), _px(off.get('y',0)),
                _px(ext.get('cx',0)), _px(ext.get('cy',0)))
    except Exception:
        return None

def _is_rounded(spPr):
    pg = spPr.find(_a('prstGeom'))
    return pg is not None and pg.get('prst','') in (
        'roundRect','round1Rect','round2SameRect','round2DiagRect')

# ── Text ─────────────────────────────────────────────────────────────────────
def _text_runs(txBody):
    """Yield (text, rgb_or_None) per run across all paragraphs."""
    for para in txBody.findall(_a('p')):
        for run in para.findall(_a('r')):
            t = run.find(_a('t'))
            if t is None or not t.text:
                continue
            text = t.text.replace('​','').replace('︎','').strip()
            if not text:
                continue
            rPr = run.find(_a('rPr'))
            rgb = None
            if rPr is not None:
                sf = rPr.find(_a('solidFill'))
                rgb = _parse_color(sf)
            yield text, rgb

# ── Slide background ─────────────────────────────────────────────────────────
def _slide_bg(slide):
    for src in [slide._element,
                slide.slide_layout._element,
                slide.slide_layout.slide_master._element]:
        bg = src.find(_p('bg'))
        if bg is None:
            bg = src.find(f'.//{_p("bg")}')
        if bg is not None:
            c = _solid_fill(bg) or _grad_fill_approx(bg)
            if c:
                return c
            # bgRef → schemeClr
            ref = bg.find(f'.//{_a("schemeClr")}')
            if ref is not None:
                return _SCHEME_APPROX.get(ref.get('val','bg1'), (255,255,255))
    return (255, 255, 255)

# ── Draw rounded rect ─────────────────────────────────────────────────────────
def _rrect(draw, x, y, w, h, fill, r=8):
    r = min(r, max(w//2, 1), max(h//2, 1))
    draw.rectangle([x+r, y, x+w-r, y+h], fill=fill)
    draw.rectangle([x, y+r, x+w, y+h-r], fill=fill)
    for cx, cy in [(x,y),(x+w-2*r,y),(x,y+h-2*r),(x+w-2*r,y+h-2*r)]:
        draw.ellipse([cx, cy, cx+2*r, cy+2*r], fill=fill)

# ── Render one shape ──────────────────────────────────────────────────────────
def _render_sp(draw, sp, slide_bg_color, font):
    spPr = sp.find(_p('spPr'))
    if spPr is None:
        return
    pos = _xfrm(spPr)
    if pos is None:
        return
    x, y, w, h = pos
    if w <= 0 or h <= 0:
        return

    fill = _fill(spPr)
    rounded = _is_rounded(spPr)

    if fill:
        if rounded:
            _rrect(draw, x, y, w, h, fill, r=8)
        else:
            draw.rectangle([x, y, x+w, y+h], fill=fill)

    txBody = sp.find(_p('txBody'))
    if txBody is None:
        return

    # Default text colour: contrast against shape fill or slide bg
    ref_bg = fill if fill else slide_bg_color
    lum = (0.299*ref_bg[0] + 0.587*ref_bg[1] + 0.114*ref_bg[2]) / 255
    default_color = (20, 20, 20) if lum > 0.45 else (235, 235, 235)

    ty = y + 3
    for para in txBody.findall(_a('p')):
        tx = x + 4
        para_has_text = False
        for run in para.findall(_a('r')):
            t_el = run.find(_a('t'))
            if t_el is None or not t_el.text:
                continue
            text = t_el.text.replace('​','').replace('︎','')
            if not text.strip() and text != ' ':
                continue
            rPr = run.find(_a('rPr'))
            rgb = None
            if rPr is not None:
                sf = rPr.find(_a('solidFill'))
                rgb = _parse_color(sf)
            color = rgb if rgb else default_color
            try:
                draw.text((tx, ty), text[:60], fill=color, font=font)
            except Exception:
                pass
            tx += max(len(text) * 6, 4)
            para_has_text = True
        if para_has_text:
            ty += 14
        if ty > y + h + 4:
            break

# ── Render spTree recursively ─────────────────────────────────────────────────
def _render_tree(draw, spTree, slide_bg_color, font):
    for child in spTree:
        tag = child.tag.split('}')[-1]
        if tag == 'sp':
            _render_sp(draw, child, slide_bg_color, font)
        elif tag == 'grpSp':
            _render_tree(draw, child, slide_bg_color, font)

# ── Render one slide ──────────────────────────────────────────────────────────
def render_slide(slide, W_px=1333, H_px=750):
    bg = _slide_bg(slide)
    img = Image.new('RGB', (W_px, H_px), bg)
    draw = ImageDraw.Draw(img)

    try:
        font = ImageFont.load_default()
    except Exception:
        font = None

    # Draw slide border
    draw.rectangle([0, 0, W_px-1, H_px-1], outline=(200, 200, 200))

    cSld = slide._element.find(_p('cSld'))
    if cSld is None:
        return img
    spTree = cSld.find(_p('spTree'))
    if spTree is None:
        return img

    _render_tree(draw, spTree, bg, font)
    return img

# ── Main ──────────────────────────────────────────────────────────────────────
def render_pptx(pptx_path, slide_indices=None, out_dir='/tmp/pptx_preview'):
    os.makedirs(out_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    n = len(prs.slides)

    if slide_indices is None:
        slide_indices = list(range(n))

    paths = []
    for i in slide_indices:
        if i >= n:
            continue
        slide = prs.slides[i]
        img = render_slide(slide)
        path = os.path.join(out_dir, f'slide_{i:02d}.png')
        img.save(path)
        paths.append(path)
        print(f'  slide {i:02d} → {path}')

    print(f'\n{len(paths)} slides rendered → {out_dir}/')
    return paths

if __name__ == '__main__':
    ap = argparse.ArgumentParser()
    ap.add_argument('pptx')
    ap.add_argument('--slides', default='all')
    ap.add_argument('--out', default='/tmp/pptx_preview')
    args = ap.parse_args()

    prs_tmp = Presentation(args.pptx)
    n = len(prs_tmp.slides)

    if args.slides == 'all':
        idxs = list(range(n))
    else:
        idxs = [int(i) for i in args.slides.split(',') if i.strip().isdigit()]

    render_pptx(args.pptx, idxs, args.out)

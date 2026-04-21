"""
Visual Cortex — PPTX Generator API v13 (Niveau 1 + Niveau 2)
═════════════════════════════════════════════════════════════════════════════
[Niveau 1 — inchangé depuis v12]
Architecture 3 phases : Compréhension → Planification → Génération/Hydratation

[Niveau 2 — nouveau dans v13]
Architecture 4 phases :
  Phase 1 — Analyse brand (identique Niveau 1)
  Phase 2 — Planification narrative (identique Niveau 1)
  Phase 3 — Génération de code python-pptx par Claude
  Phase 4 — Exécution sécurisée (sandbox exec) + assemblage PPTX

Nouveautés v13 :
  - CLIENT_PROFILES : 5 profils visuels (finance, industrial, institutional, startup, creative)
  - Bibliothèque de helpers h2_* : shapes, textes, KPIs, dividers, numéros décoratifs
  - Sandbox exec() à namespace restreint (aucun import, os, open exposé)
  - Timeout threading 30s + fallback automatique Niveau 1 si exécution échoue
  - Routes GET /profiles et POST /generate-v2

Modèle : claude-sonnet-4-6 (configurable via CLAUDE_MODEL)
"""

import os, io, json, time, copy, re, logging, threading, zipfile, base64
from collections import defaultdict
from typing import Optional

import anthropic
from fastapi import FastAPI, File, Form, UploadFile, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
import uvicorn
from layouts import LAYOUT_REGISTRY, LAYOUT_DESCRIPTIONS

# ─────────────────────────────────────────────
# LOGGING & APP
# ─────────────────────────────────────────────
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
log = logging.getLogger("visual-cortex")

app = FastAPI(title="Visual Cortex API", version="13.0.0")
app.add_middleware(
    CORSMiddleware, allow_origins=["*"], allow_credentials=False,
    allow_methods=["*"], allow_headers=["*"],
)

@app.exception_handler(Exception)
async def global_error(request: Request, exc: Exception):
    log.error(f"Exception: {exc}", exc_info=True)
    return JSONResponse(
        status_code=500,
        content={"detail": {"message": f"Erreur serveur : {str(exc)}"}},
        headers={"Access-Control-Allow-Origin": "*"},
    )

# ─────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────
ANTHROPIC_API_KEY = os.environ.get("ANTHROPIC_API_KEY", "")
PRO_SECRET_TOKEN  = os.environ.get("PRO_SECRET_TOKEN", "change-me")
FREE_QUOTA_PER_IP = int(os.environ.get("FREE_QUOTA_PER_IP", "3"))
CLAUDE_MODEL      = os.environ.get("CLAUDE_MODEL", "claude-sonnet-4-6")

_usage: dict = defaultdict(list)
DAY_SEC = 86400

def _ip(r: Request) -> str:
    fwd = r.headers.get("x-forwarded-for")
    return fwd.split(",")[0].strip() if fwd else r.client.host

def _is_pro(auth: Optional[str]) -> bool:
    return bool(auth and auth.replace("Bearer ", "").strip() == PRO_SECRET_TOKEN)

def _quota(ip: str) -> tuple:
    now = time.time()
    _usage[ip] = [t for t in _usage[ip] if now - t < DAY_SEC]
    if len(_usage[ip]) >= FREE_QUOTA_PER_IP:
        raise HTTPException(429, {"message": "Quota gratuit épuisé. Passez en Pro."})
    _usage[ip].append(now)
    return len(_usage[ip]), FREE_QUOTA_PER_IP


# ══════════════════════════════════════════════════════════════
# UTILITAIRES FONDAMENTAUX
# ══════════════════════════════════════════════════════════════

def _emu(v: int) -> float:
    return v / 914400.0

def _clean_json(raw: str) -> str:
    """
    Extrait un bloc JSON propre depuis la réponse de Claude.
    Gère : markdown fences, préambules textuels, JSON tronqué.
    """
    s = raw.strip()

    # Cas 1 : bloc ```json ... ``` ou ``` ... ```
    if "```" in s:
        parts = s.split("```")
        for part in parts[1::2]:          # parties entre backticks
            candidate = part.strip()
            if candidate.startswith("json"):
                candidate = candidate[4:].strip()
            if candidate.startswith("{") or candidate.startswith("["):
                s = candidate
                break

    # Cas 2 : extraire depuis le premier { ou [ jusqu'au dernier } ou ]
    start_brace  = s.find("{")
    start_bracket = s.find("[")
    if start_brace == -1 and start_bracket == -1:
        return s  # on laisse planter json.loads avec un message clair

    if start_bracket == -1 or (start_brace != -1 and start_brace < start_bracket):
        start = start_brace
        end   = s.rfind("}")
    else:
        start = start_bracket
        end   = s.rfind("]")

    if start != -1 and end != -1 and end > start:
        s = s[start:end + 1]

    return s.strip()


def _parse_json_robust(raw: str, context: str = "") -> dict:
    """
    Parse JSON avec fallback sur réparation basique.
    Si le JSON est tronqué (max_tokens atteint), tente de le compléter.
    """
    cleaned = _clean_json(raw)

    # Tentative 1 : parse direct
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError as e:
        log.warning(f"JSON parse error [{context}]: {e} — tentative de réparation")

    # Tentative 2 : supprimer trailing comma avant } ou ]
    fixed = re.sub(r",\s*([}\]])", r"\1", cleaned)
    try:
        return json.loads(fixed)
    except json.JSONDecodeError:
        pass

    # Tentative 3 : JSON tronqué — essayer de fermer les accolades manquantes
    attempt = fixed
    open_braces   = attempt.count("{") - attempt.count("}")
    open_brackets = attempt.count("[") - attempt.count("]")
    # Supprimer la dernière entrée potentiellement incomplète
    last_comma = attempt.rfind(",")
    last_brace  = attempt.rfind("}")
    if last_comma > last_brace:
        attempt = attempt[:last_comma]
    attempt += "]" * open_brackets + "}" * open_braces
    try:
        result = json.loads(attempt)
        log.warning(f"JSON réparé (tronqué) [{context}] : {open_braces} accolades + {open_brackets} crochets ajoutés")
        return result
    except json.JSONDecodeError as e:
        raise ValueError(f"JSON irrécupérable [{context}] : {e}\nDébut : {cleaned[:200]}") from e

def _safe_name(s: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", s[:40].lower()).strip("-")

FOOTER_PLACEHOLDERS = [
    "date - pied de page de votre présentation",
    "pied de page de votre présentation",
    "titre de la présentation",
    "date de la présentation",
    "[date]",
    "[footer]",
    "click to edit",
]

def _is_footer_placeholder(text: str) -> bool:
    return any(p in text.lower() for p in FOOTER_PLACEHOLDERS)


# ══════════════════════════════════════════════════════════════
# TRAVERSÉE RÉCURSIVE DES SHAPES (FIX CRITIQUE GROUP SHAPES)
# ══════════════════════════════════════════════════════════════

def iter_all_shapes(shapes):
    """
    Parcourt TOUS les shapes, y compris ceux imbriqués dans des GROUP shapes.
    FIX CRITIQUE : le bug "contenu fantôme" venait de l'absence de cette traversée.
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_all_shapes(shape.shapes)
        else:
            yield shape


# ══════════════════════════════════════════════════════════════
# PHASE 1 — COMPRÉHENSION PROFONDE DU TEMPLATE
# ══════════════════════════════════════════════════════════════

SLIDE_TYPE_DESC = {
    "cover":      "Couverture — titre fort + sous-titre clair",
    "section":    "Ouverture de chapitre — numéro + titre de section",
    "two_col":    "Deux colonnes — 2 blocs symétriques côte à côte",
    "kpi":        "Chiffres clés — métriques visuelles impactantes",
    "quote":      "Citation forte — accroche ou message clé mis en valeur",
    "timeline":   "Chronologie — étapes ou jalons temporels",
    "list":       "Liste structurée — items avec labels courts",
    "image_text": "Image + texte — visuel fort côte à côte",
    "full_text":  "Contenu développé — argumentaire structuré",
    "closing":    "Conclusion — remerciement ou call-to-action",
    "complex":    "Slide graphique — diagramme ou organigramme",
    "unknown":    "Layout générique",
}

DENSITY_RULES = {
    "cover": {
        "max_title_words":    7,
        "max_subtitle_words": 12,
        "style":              "Titre : accroche forte, verbe d'action ou tension. Sous-titre : contexte ou angle en une phrase.",
        "never":              "Jamais de bullet points. Jamais plus d'un sous-titre.",
        "example_title":      "TotalEnergies sous pression",
        "example_subtitle":   "Anatomie des campagnes ONG depuis 2020",
    },
    "section": {
        "max_title_words":    6,
        "max_label_words":    2,
        "style":              "Numéro de section (01, 02…) + titre court et percutant. Pas de body.",
        "never":              "Jamais de body text sur une slide de section. Jamais plus de 2 zones.",
        "example_title":      "Liens financiers & fiscaux",
        "example_label":      "01",
    },
    "two_col": {
        "max_title_words":     8,
        "max_col_title_words": 4,
        "max_col_body_words":  30,
        "max_items_per_col":   4,
        "style":               "Titre général + 2 colonnes symétriques. Chaque colonne : label court + 2-4 items courts.",
        "never":               "Jamais de paragraphes dans les colonnes. Jamais d'asymétrie.",
        "example":             "Colonne gauche : 'LIENS FINANCIERS' + 4 items de 10 mots max. Colonne droite : 'LIENS RÉGLEMENTAIRES' + 4 items.",
    },
    "kpi": {
        "max_kpi_count":      6,
        "max_value_words":    2,
        "max_label_words":    5,
        "max_sublabel_words": 12,
        "style":              "Chiffre ou métrique très visible + label court + sous-label contextuel.",
        "never":              "Jamais plus de 6 KPIs. Jamais de phrases complètes pour les valeurs.",
        "example":            "'600 M€' / 'contribution exceptionnelle' / 'versée en 2022, loi superprofits'",
    },
    "quote": {
        "max_quote_words":    20,
        "max_author_words":   6,
        "style":              "Citation courte et percutante, ton affirmatif. Idéalement entre guillemets.",
        "never":              "Jamais de bullet points. Jamais plus d'une citation. Jamais > 20 mots.",
        "example":            "'La transition énergétique ne se fera pas sans les majors.' — Analyse 2025",
    },
    "timeline": {
        "max_steps":            6,
        "max_step_title_words": 4,
        "max_step_body_words":  12,
        "style":                "4 à 6 jalons chronologiques. Chaque étape : date + titre court + phrase optionnelle.",
        "never":                "Jamais plus de 6 étapes. Jamais de paragraphes. Jamais sans repère temporel.",
        "example":              "1924 / 'Création CFP' / 'Fondation par décret d'État'",
    },
    "list": {
        "max_items":            5,
        "max_item_title_words": 4,
        "max_item_body_words":  20,
        "style":                "3 à 5 items. Chaque item : titre en gras court + phrase de développement concise.",
        "never":                "Jamais plus de 5 items. Jamais sans titre par item.",
        "example":              "'Lobbying institutionnel' / '~2,3 M€ déclarés/an — contacts réguliers Élysée, Bercy.'",
    },
    "image_text": {
        "max_title_words":  8,
        "max_body_words":   40,
        "max_body_items":   3,
        "style":            "Titre + corps structuré en 2-3 points courts. L'image fait le travail visuel.",
        "never":            "Jamais de mur de texte. Jamais plus de 3 bullet points.",
    },
    "full_text": {
        "max_title_words":  8,
        "max_body_words":   60,
        "max_paragraphs":   3,
        "style":            "Titre + 2-3 paragraphes courts et aérés. Chaque paragraphe = une idée.",
        "never":            "Jamais de corps > 60 mots. Jamais plus de 3 paragraphes.",
    },
    "closing": {
        "max_title_words":    5,
        "max_subtitle_words": 15,
        "style":              "Message mémorable ou 'Merci !' + sous-titre : sources, contact ou CTA.",
        "never":              "Jamais de bullet points. Jamais de corps long. Simple, élégant.",
        "example_title":      "Merci !",
        "example_subtitle":   "Sources : Rapport Annuel 2023 · HATVP · Loi de vigilance 2017",
    },
    "complex": {
        "max_label_words": 4,
        "max_body_words":  15,
        "style":           "Labels ultra-courts pour les éléments du diagramme.",
        "never":           "Jamais de phrases complètes dans un diagramme.",
    },
    "unknown": {
        "max_title_words": 8,
        "max_body_words":  35,
        "style":           "Titre + corps aéré. Respecter la structure du template.",
        "never":           "Jamais de surcharge.",
    },
}

WORD_LIMITS = {
    "title":       8,
    "subtitle":   15,
    "section_num": 2,
    "label":       5,
    "kpi_value":   3,
    "kpi_label":   5,
    "body":       40,
    "footer":     10,
    "page_number": 1,
    "quote":      20,
    "list_item":  18,
    "placeholder": 8,
    "text":       30,
}


def _classify_slide(slide, idx: int, total: int, w: int, h: int) -> str:
    shapes      = list(iter_all_shapes(slide.shapes))
    raw_shapes  = list(slide.shapes)
    img_shapes  = [s for s in shapes if s.shape_type in (13, 11)]
    group_shapes = [s for s in raw_shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]

    texts = []
    for s in shapes:
        if not getattr(s, "has_text_frame", False):
            continue
        for para in s.text_frame.paragraphs:
            t = "".join(r.text for r in para.runs).strip()
            if len(t) > 3:
                texts.append({"text": t, "shape": s})

    n           = len(texts)
    total_chars = sum(len(t["text"]) for t in texts)

    if idx == 0:
        return "cover"
    if idx == total - 1:
        return "closing"
    if len(group_shapes) >= 5:
        return "complex"
    if img_shapes and n >= 2:
        return "image_text"
    if n <= 3 and total_chars < 100:
        return "section"
    if n >= 4:
        short = [t for t in texts if len(t["text"]) < 25]
        if len(short) >= 3:
            try:
                lefts = sorted([t["shape"].left for t in texts])
                spread = _emu(lefts[-1]) - _emu(lefts[0])
                if spread > 4.0:
                    return "kpi"
            except Exception:
                pass
    if n >= 3 and total_chars < 500:
        try:
            tops = sorted([t["shape"].top for t in texts])
            if _emu(tops[-1]) - _emu(tops[0]) > 2.5:
                return "timeline"
        except Exception:
            pass
    if n >= 4:
        try:
            half = w / 2
            lc = sum(1 for t in texts if t["shape"].left < half)
            rc = sum(1 for t in texts if t["shape"].left >= half)
            if lc >= 2 and rc >= 2:
                return "two_col"
        except Exception:
            pass
    if n <= 3 and any(len(t["text"]) > 60 for t in texts):
        return "quote"
    if n >= 3:
        lengths = [len(t["text"]) for t in texts]
        avg = sum(lengths) / len(lengths)
        var = sum((l - avg) ** 2 for l in lengths) / len(lengths)
        if var < 800 and avg < 150:
            return "list"
    if total_chars > 200:
        return "full_text"
    return "unknown"


def _shape_role(shape, w: int, h: int) -> str:
    """Détermine le rôle via placeholder natif PPTX, puis géométrie."""
    try:
        if shape.is_placeholder:
            ph = shape.placeholder_format
            ph_map = {
                0: "title", 1: "body", 2: "subtitle",
                3: "date", 4: "footer", 5: "page_number",
                13: "title", 15: "subtitle",
            }
            return ph_map.get(ph.idx, "placeholder")
    except Exception:
        pass

    if not getattr(shape, "has_text_frame", False):
        return "decoration"

    try:
        tr = shape.top    / h
        wr = shape.width  / w
        hr = shape.height / h

        if tr > 0.87 and wr > 0.15:
            return "footer"
        if tr > 0.87:
            return "page_number"
        if tr < 0.28 and wr > 0.45:
            return "title"
        if 0.20 < tr < 0.50 and wr > 0.35 and hr < 0.15:
            return "subtitle"
        if 0.25 < tr < 0.75 and wr < 0.40 and hr < 0.12:
            return "label"
        if 0.25 < tr < 0.87:
            return "body"
    except Exception:
        pass

    return "text"


def extract_brand(prs: Presentation) -> dict:
    """
    Extraction complète de la charte depuis le template.
    Source prioritaire : ppt/theme/theme1.xml lu depuis le ZIP PPTX.
    Fallback : scan regex de tout le XML pour les couleurs et polices.
    """
    import re as _re
    import io
    import zipfile as _zipfile
    import lxml.etree as etree
    from collections import Counter

    color_counter = Counter()
    fonts_seen: list = []
    fonts_set: set = set()

    def _add_font(f: str):
        if f and not f.startswith('+') and f.strip() and f not in fonts_set:
            fonts_set.add(f)
            fonts_seen.append(f)

    def _is_neutral(h: str) -> bool:
        try:
            r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
            if r > 235 and g > 235 and b > 235: return True
            if r < 25  and g < 25  and b < 25:  return True
            if abs(r-g) < 25 and abs(g-b) < 25: return True
            return False
        except Exception:
            return True

    # ── Étape 1 : lire ppt/theme/theme1.xml depuis le ZIP ────────────────────
    # Un PPTX est un ZIP. Le thème officiel s'y trouve toujours.
    theme_colors: dict = {}
    try:
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        with _zipfile.ZipFile(buf) as zf:
            theme_files = sorted([n for n in zf.namelist()
                                   if _re.search(r'ppt/theme/theme\d*\.xml$', n, _re.I)])
            for tf in theme_files[:1]:
                xml = zf.read(tf).decode('utf-8', errors='ignore')
                for slot in ['dk1', 'lt1', 'dk2', 'lt2',
                             'accent1', 'accent2', 'accent3',
                             'accent4', 'accent5', 'accent6']:
                    m = _re.search(
                        rf'<a:{slot}[^>]*>\s*<a:srgbClr val="([0-9A-Fa-f]{{6}})"',
                        xml)
                    if m:
                        theme_colors[slot] = m.group(1).upper()
                        continue
                    m = _re.search(
                        rf'<a:{slot}[^>]*>\s*<a:sysClr[^>]*lastClr="([0-9A-Fa-f]{{6}})"',
                        xml)
                    if m:
                        theme_colors[slot] = m.group(1).upper()
                # Polices du thème
                for m in _re.findall(r'typeface="([^"]+)"', xml):
                    _add_font(m)
    except Exception as e:
        log.warning(f'[brand] theme ZIP extraction failed: {e}')

    log.info(f'[brand] theme_colors={theme_colors}')

    # ── Étape 2 : scan regex (slides + layouts + masters) ────────────────────
    sources = list(prs.slides) + list(prs.slide_layouts) + list(prs.slide_masters)
    for source in sources:
        try:
            xml = etree.tostring(source._element, pretty_print=False).decode()
        except Exception:
            continue
        for m in _re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', xml):
            color_counter[m.upper()] += 1
        for m in _re.findall(r'typeface="([^"]+)"', xml):
            _add_font(m)

    chromatic  = [c for c, _ in color_counter.most_common() if not _is_neutral(c)]
    all_colors = [c for c, _ in color_counter.most_common()]

    w, h = prs.slide_width, prs.slide_height
    log.info(f'[brand] chromatic={chromatic[:4]} fonts={fonts_seen[:3]}')
    return {
        'fonts':           fonts_seen[:5],
        'colors':          chromatic[:8],
        'all_colors':      all_colors[:20],
        'theme_colors':    theme_colors,
        'slide_count':     len(prs.slides),
        'layouts':         [l.name for l in prs.slide_layouts],
        'slide_width_in':  round(_emu(w), 2),
        'slide_height_in': round(_emu(h), 2),
        'aspect_ratio':    (
            '16:9' if abs(w / h - 16 / 9) < 0.05 else
            '4:3'  if abs(w / h - 4 / 3)  < 0.05 else 'custom'
        ),
    }


def build_layout_library(prs: Presentation) -> list:
    library = []
    total = len(prs.slides)
    w, h  = prs.slide_width, prs.slide_height

    for idx, slide in enumerate(prs.slides):
        slide_type  = _classify_slide(slide, idx, total, w, h)
        root_groups = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP)
        has_images  = any(s.shape_type in (13, 11) for s in iter_all_shapes(slide.shapes))

        zones = []
        seen  = set()

        for shape in iter_all_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            role = _shape_role(shape, w, h)
            if role == "decoration":
                continue

            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs).strip()
                if len(text) < 2 or text in seen:
                    continue
                seen.add(text)

                is_placeholder_footer = _is_footer_placeholder(text)
                zones.append({
                    "original_text":         text,
                    "role":                  "footer" if is_placeholder_footer else role,
                    "word_count":            len(text.split()),
                    "word_limit":            WORD_LIMITS.get(
                        "footer" if is_placeholder_footer else role, 30
                    ),
                    "is_placeholder_footer": is_placeholder_footer,
                    "char_count":            len(text),
                })

        if zones:
            library.append({
                "slide_index":  idx,
                "slide_type":   slide_type,
                "description":  SLIDE_TYPE_DESC.get(slide_type, ""),
                "position":     (
                    "cover"   if idx == 0 else
                    "closing" if idx == total - 1 else
                    f"{idx+1}/{total}"
                ),
                "root_groups":  root_groups,
                "has_images":   has_images,
                "zones":        zones,
                "total_words":  sum(z["word_count"] for z in zones),
                "visual_score": (
                    has_images * 3 +
                    (1 if slide_type in ("kpi", "timeline", "two_col", "image_text", "quote") else 0) * 2 +
                    (1 if root_groups > 0 else 0) * 1
                ),
            })

    return library


def select_template_slides(library: list, nb_slides: int) -> list:
    if not library:
        return []

    cover   = [s for s in library if s["slide_type"] == "cover"]
    closing = [s for s in library if s["slide_type"] == "closing"]
    middle  = [s for s in library if s["slide_type"] not in ("cover", "closing")]

    middle_sorted = sorted(
        middle,
        key=lambda s: (
            0 if s["slide_type"] == "complex" else 1,
            s["visual_score"],
            -s["total_words"],
        ),
        reverse=True,
    )

    n_cover   = min(len(cover), 1)
    n_closing = min(len(closing), 1)
    n_middle  = nb_slides - n_cover - n_closing

    if n_middle < 0:
        n_middle = 0

    selected_middle = []
    last_type = None
    pool = middle_sorted.copy()

    while len(selected_middle) < n_middle and pool:
        for i, s in enumerate(pool):
            if s["slide_type"] != last_type or i == len(pool) - 1:
                selected_middle.append(s)
                last_type = s["slide_type"]
                pool.pop(i)
                break

    if len(selected_middle) < n_middle and middle_sorted:
        cycle = middle_sorted * 10
        for s in cycle:
            if len(selected_middle) >= n_middle:
                break
            dup = {**s, "duplicate": True}
            selected_middle.append(dup)

    result = (
        cover[:n_cover] +
        sorted(selected_middle, key=lambda s: s["slide_index"]) +
        closing[:n_closing]
    )

    result = result[:nb_slides]

    # Garantie : si closing existe dans le template, elle est TOUJOURS en dernière position
    if closing and (not result or result[-1]["slide_type"] != "closing"):
        if len(result) < nb_slides:
            result.append(closing[0])
        else:
            result[-1] = closing[0]

    return result


# ══════════════════════════════════════════════════════════════
# PHASE 2 — PLANIFICATION NARRATIVE
# ══════════════════════════════════════════════════════════════

PLANNER_SYSTEM = """Tu es Visual Cortex Planner, architecte narratif de présentations B2B professionnelles.

MISSION : concevoir une structure narrative percutante et cohérente, slide par slide.

PRINCIPES NARRATIFS :
- Cover : accroche forte, crée une tension ou une promesse. Titre ≤ 7 mots.
- Slides de contenu : progression logique et variée —
  ne jamais enchaîner deux fois le même type de slide.
  Séquences éprouvées : contexte → enjeux → solution → preuves → bénéfices → ROI → next steps
- Closing : message simple et mémorable. CTA clair ou remerciement élégant.

AFFECTATION DES TYPES AUX CONTENUS :
- données chiffrées, métriques, résultats     → kpi
- étapes, jalons, chronologie, processus       → timeline
- comparaison, dualité, pour/contre, avant/après → two_col
- message fort, conviction, prise de position  → quote
- liste d'acteurs, d'arguments, de features    → list
- concept visuel, cas d'usage, illustration    → image_text
- argumentaire développé, analyse, contexte    → full_text
- séparation de partie                         → section

DENSITÉ PAR TYPE (à indiquer dans visual_hint) :
- kpi       : 4 à 6 chiffres max, valeur courte + label + sous-label
- timeline  : 4 à 6 jalons max, date + titre court + phrase optionnelle
- two_col   : 2 colonnes symétriques, 3-4 items de ≤ 10 mots chacun
- quote     : 1 citation de ≤ 20 mots, source courte
- list      : 3 à 5 items structurés (titre gras + corps ≤ 20 mots)
- section   : numéro (01) + titre ≤ 6 mots, rien d'autre
- cover     : titre ≤ 7 mots + sous-titre ≤ 12 mots
- closing   : titre ≤ 5 mots + sous-titre/sources ≤ 15 mots

RÈGLE ABSOLUE : le tableau "slides" doit contenir EXACTEMENT {nb_slides} entrées.

Réponds UNIQUEMENT en JSON valide, sans markdown."""

PLANNER_USER = """SUJET : {prompt}
NB SLIDES SOUHAITÉ : {nb_slides}

SLIDES DISPONIBLES DANS LE TEMPLATE :
{selection_json}

CHARTE : Polices {fonts} | Couleurs {colors} | Format {ratio}

Génère le plan. Le tableau "slides" doit contenir EXACTEMENT {nb_slides} entrées.

FORMAT :
{{
  "presentation_title": "Titre accrocheur ≤ 7 mots",
  "narrative_arc": "Logique narrative en 1 phrase",
  "footer_text": "Entreprise · Contexte · Année  (≤ 8 mots)",
  "slides": [
    {{
      "plan_index": 0,
      "template_slide_index": 0,
      "slide_type": "cover",
      "narrative_angle": "Ce que cette slide accomplit dans l'histoire (1 phrase)",
      "key_message": "Le message principal ≤ 10 mots",
      "visual_hint": "Contrainte de densité et de style pour cette slide"
    }}
  ]
}}\n"""


def plan_presentation(prompt: str, nb_slides: int, selection: list, brand: dict) -> dict:
    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

    sel_light = [
        {
            "template_slide_index": s["slide_index"],
            "slide_type":           s["slide_type"],
            "description":          s["description"],
            "zone_roles":           [z["role"] for z in s["zones"]],
            "has_images":           s["has_images"],
            "visual_score":         s["visual_score"],
        }
        for s in selection
    ]

    system = PLANNER_SYSTEM.replace("{nb_slides}", str(nb_slides))

    user = PLANNER_USER.format(
        prompt         = prompt,
        nb_slides      = nb_slides,
        selection_json = json.dumps(sel_light, ensure_ascii=False, indent=2),
        fonts          = ", ".join(brand.get("fonts", [])) or "Standard",
        colors         = ", ".join(f"#{c}" for c in brand.get("colors", [])) or "non détectées",
        ratio          = brand.get("aspect_ratio", "16:9"),
    )

    # max_tokens adaptatif : ~120 tokens/slide suffisent pour le plan JSON
    planner_tokens = max(2000, nb_slides * 180)

    for attempt in range(3):
        msg = client.messages.create(
            model=CLAUDE_MODEL, max_tokens=planner_tokens,
            system=system,
            messages=[{"role": "user", "content": user}],
        )
        try:
            plan = _parse_json_robust(msg.content[0].text.strip(), context="plan")
            log.info(f"Plan: {len(plan.get('slides', []))} slides — {plan.get('narrative_arc', '')[:80]}")
            return plan
        except (ValueError, KeyError) as e:
            log.warning(f"plan_presentation attempt {attempt+1}/3 échoué : {e}")
            if attempt == 2:
                raise
    raise RuntimeError("plan_presentation : 3 tentatives échouées")


# ══════════════════════════════════════════════════════════════
# PHASE 3 — GÉNÉRATION DU CONTENU (visuel-first, Niveau 1)
# ══════════════════════════════════════════════════════════════

CORTEX_SYSTEM = """Tu es Visual Cortex, expert en présentations B2B professionnelles et visuelles.
Philosophie : une slide = une idée. Le texte est une accroche, pas un rapport.

═══════════════════════════════════════════════════
RÈGLES UNIVERSELLES (s'appliquent à toutes les slides)
═══════════════════════════════════════════════════
1. LIMITES PAR RÔLE — ne jamais dépasser :
   title       → ≤ 8 mots   subtitle    → ≤ 12 mots  label       → ≤ 5 mots
   kpi_value   → ≤ 3 mots   kpi_label   → ≤ 5 mots   body        → ≤ 40 mots
   list_item   → ≤ 18 mots  quote       → ≤ 20 mots  footer      → ≤ 8 mots
   page_number → NE PAS MODIFIER

2. COHÉRENCE : footer identique sur toutes les slides de contenu.
3. B2B : vocabulaire du secteur, ton direct, orienté valeur, zéro formule creuse.
4. PROGRESSION : chaque slide fait avancer l'histoire selon son narrative_angle.
5. ZÉRO invention de données, chiffres ou noms non fournis dans le prompt.
6. FOOTERS PLACEHOLDER : zones "is_placeholder_footer: true" → remplacer par le footer_text.
7. ZONES VIDES INTENTIONNELLES : si une zone n'a pas de contenu pertinent pour le type de slide,
   retourne "" (chaîne vide) pour la vider. Ne jamais laisser un texte de template générique.
8. SLIDE CLOSING : toujours présente en dernière position. Titre court, mémorable. Jamais de bullet points.

═══════════════════════════════════════════════════
RÈGLES PAR TYPE DE SLIDE (density rules)
═══════════════════════════════════════════════════

[cover] Titre ≤ 7 mots — accroche forte. Sous-titre ≤ 12 mots.
[section] Numéro (01, 02…) + titre ≤ 6 mots. RIEN D'AUTRE.
[kpi] 4 à 6 KPIs MAX. Valeur courte + label + sous-label contextuel.
[timeline] 4 à 6 jalons MAX. Repère temporel + titre ≤ 4 mots + phrase optionnelle.
[two_col] 2 colonnes SYMÉTRIQUES — max 4 items/colonne ≤ 18 mots.
[quote] 1 SEULE citation ≤ 20 mots. Source optionnelle ≤ 6 mots.
[list] 3 à 5 items MAX. Titre ≤ 5 mots + corps ≤ 20 mots. Structure parallèle.
[image_text] Titre ≤ 8 mots + corps 2-3 points ≤ 40 mots total.
[full_text] Titre ≤ 8 mots + 2-3 §§ courts = total body ≤ 60 mots.
[closing] Titre ≤ 5 mots. Sous-titre ≤ 15 mots. Simple, élégant.
[complex] Labels ≤ 4 mots. Titres ≤ 5 mots. Jamais de phrases complètes.

Réponds UNIQUEMENT en JSON valide, sans commentaire ni markdown."""

CORTEX_USER = """PRÉSENTATION : {title}
ARC NARRATIF : {arc}
FOOTER : "{footer}"
SUJET COMPLET : {prompt}
CHARTE : Polices {fonts} | Couleurs {colors}

═══════════════════════
SLIDES À GÉNÉRER — {n} slides
═══════════════════════
{slides_json}

FORMAT DE SORTIE (clés = template_slide_index) :
{{
  "0": {{"Texte original exact": "Nouveau texte généré"}},
  ...
}}"""


def generate_content(prompt: str, plan: dict, selection: list, brand: dict) -> dict:
    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    sel_by_idx  = {s["slide_index"]: s for s in selection}
    footer_text = plan.get("footer_text", "")

    slides_payload = []
    for sp in plan.get("slides", []):
        tidx       = sp.get("template_slide_index", 0)
        tmpl       = sel_by_idx.get(tidx, {})
        slide_type = sp.get("slide_type", "unknown")
        density    = DENSITY_RULES.get(slide_type, DENSITY_RULES["unknown"])

        slides_payload.append({
            "template_slide_index": tidx,
            "slide_type":           slide_type,
            "narrative_angle":      sp.get("narrative_angle"),
            "key_message":          sp.get("key_message"),
            "visual_hint":          sp.get("visual_hint", ""),
            "density_rules": {
                "style":  density.get("style", ""),
                "never":  density.get("never", ""),
                "limits": {
                    k: v for k, v in density.items()
                    if k not in ("style", "never", "example", "example_title",
                                 "example_subtitle", "example_label")
                },
            },
            "zones": [
                {
                    "original_text":         z["original_text"],
                    "role":                  z["role"],
                    "word_limit":            z["word_limit"],
                    "char_count":            z.get("char_count", 0),
                    "is_placeholder_footer": z.get("is_placeholder_footer", False),
                }
                for z in tmpl.get("zones", [])
            ],
        })

    user = CORTEX_USER.format(
        title       = plan.get("presentation_title", prompt[:60]),
        arc         = plan.get("narrative_arc", ""),
        footer      = footer_text,
        prompt      = prompt,
        fonts       = ", ".join(brand.get("fonts", [])) or "Standard",
        colors      = ", ".join(f"#{c}" for c in brand.get("colors", [])) or "non détectées",
        n           = len(slides_payload),
        slides_json = json.dumps(slides_payload, ensure_ascii=False, indent=2),
    )

    # max_tokens adaptatif : ~220 tokens/slide (JSON + textes courts)
    content_tokens = max(3000, len(slides_payload) * 320)

    msg = client.messages.create(
        model=CLAUDE_MODEL, max_tokens=content_tokens,
        system=CORTEX_SYSTEM,
        messages=[{"role": "user", "content": user}],
    )

    try:
        mapping = _parse_json_robust(msg.content[0].text.strip(), context="content")
    except ValueError as e:
        log.error(f"generate_content JSON irrécupérable : {e}")
        raise
    mapping = _validate_and_trim(mapping, slides_payload)

    log.info(f"Contenu généré et validé : {len(mapping)} slides.")
    return mapping


def _validate_and_trim(mapping: dict, slides_payload: list) -> dict:
    zone_limits: dict = {}
    for sp in slides_payload:
        tidx = str(sp["template_slide_index"])
        for z in sp.get("zones", []):
            key = (tidx, z["original_text"])
            zone_limits[key] = {
                "word_limit": z["word_limit"],
                "role":       z["role"],
            }

    validated = {}
    for slide_key, replacements in mapping.items():
        validated[slide_key] = {}
        for orig, new_text in replacements.items():
            if not new_text:
                validated[slide_key][orig] = new_text
                continue

            zone_info = zone_limits.get((str(slide_key), orig), {})
            role      = zone_info.get("role", "text")
            limit     = zone_info.get("word_limit", WORD_LIMITS.get(role, 40))

            if role == "page_number":
                validated[slide_key][orig] = orig
                continue

            words = new_text.split()
            if len(words) > limit:
                trimmed = " ".join(words[:limit])
                for punct in [".", ",", ";", ":", "—", "–"]:
                    last = trimmed.rfind(punct)
                    if last > len(trimmed) * 0.6:
                        trimmed = trimmed[:last + 1].strip()
                        break
                log.debug(f"Trim slide {slide_key} role={role}: {len(words)}→{len(trimmed.split())} mots")
                validated[slide_key][orig] = trimmed
            else:
                validated[slide_key][orig] = new_text

    return validated


# ══════════════════════════════════════════════════════════════
# HYDRATATION NIVEAU 1 — Injection avec traversée récursive
# ══════════════════════════════════════════════════════════════

def _replace_text_in_para(para, replacements: dict):
    para_text = "".join(r.text for r in para.runs).strip()
    if not para_text:
        return

    new_text = replacements.get(para_text)
    if new_text is None:
        normalized = para_text.replace("\u2019", "'").replace("\u2018", "'")
        for k, v in replacements.items():
            k_norm = k.replace("\u2019", "'").replace("\u2018", "'")
            if k_norm == normalized:
                new_text = v
                break

    if not new_text:
        return

    rpr_xml = None
    if para.runs:
        rpr_el = para.runs[0]._r.find(qn("a:rPr"))
        if rpr_el is not None:
            rpr_xml = copy.deepcopy(rpr_el)

    para.text = new_text

    if rpr_xml is not None:
        for run in para.runs:
            ex = run._r.find(qn("a:rPr"))
            if ex is not None:
                run._r.remove(ex)
            run._r.insert(0, copy.deepcopy(rpr_xml))


def _hydrate_slide(slide, replacements: dict):
    for shape in iter_all_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        for para in shape.text_frame.paragraphs:
            _replace_text_in_para(para, replacements)


def hydrate_presentation(
    pptx_bytes: bytes,
    mapping: dict,
    plan_slides: list,
    nb_slides: int,
) -> bytes:
    prs = Presentation(io.BytesIO(pptx_bytes))
    template_indices = [s.get("template_slide_index", 0) for s in plan_slides]
    _reorder_and_hydrate(prs, template_indices, mapping, nb_slides)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()


def _reorder_and_hydrate(prs: Presentation, template_indices: list, mapping: dict, nb_slides: int):
    total_tmpl = len(prs.slides)

    for slide_key, replacements in mapping.items():
        try:
            idx = int(str(slide_key).replace("slide_", ""))
            if idx < total_tmpl:
                _hydrate_slide(prs.slides[idx], replacements)
        except Exception as e:
            log.warning(f"Hydratation slide {slide_key}: {e}")

    xml_slides  = prs.slides._sldIdLst
    all_sld_ids = list(xml_slides)

    new_order = []
    for tidx in template_indices:
        if 0 <= tidx < len(all_sld_ids):
            new_order.append(all_sld_ids[tidx])

    seen_ids   = set()
    final_order = []
    for sld_el in new_order:
        sld_id = sld_el.get("id")
        if sld_id not in seen_ids:
            seen_ids.add(sld_id)
            final_order.append(sld_el)
        else:
            dup = _duplicate_slide_element(prs, sld_el)
            if dup is not None:
                final_order.append(dup)

    final_order = final_order[:nb_slides]

    for sld in list(xml_slides):
        xml_slides.remove(sld)
    for sld in final_order:
        xml_slides.append(sld)

    _cleanup_orphan_slides(prs, final_order)


def _duplicate_slide_element(prs: Presentation, src_sld_el):
    try:
        src_rId  = src_sld_el.get(qn("r:id"))
        src_part = prs.part.related_parts.get(src_rId)
        if src_part is None:
            return None

        blank_layout = prs.slide_layouts[-1]
        new_slide    = prs.slides.add_slide(blank_layout)

        import lxml.etree as etree
        src_xml = copy.deepcopy(src_part._element)
        new_slide._element.getparent().replace(new_slide._element, src_xml)

        return list(prs.slides._sldIdLst)[-1]
    except Exception as e:
        log.warning(f"Duplication slide: {e}")
        return None


def _cleanup_orphan_slides(prs: Presentation, kept_sld_els: list):
    # drop_rel est instable selon les versions python-pptx — la sldIdLst
    # controle les slides affichees, les rels orphelins sont ignores par PowerPoint.
    pass


# ══════════════════════════════════════════════════════════════
# PIPELINE NIVEAU 1
# ══════════════════════════════════════════════════════════════

# ══════════════════════════════════════════════════════════════
# RÉSOLUTION NB_SLIDES (Essentiel / Complet / Approfondi)
# ══════════════════════════════════════════════════════════════

NB_SLIDES_MAP = {
    # Valeurs fixes — le planner s'adapte
    "essentiel":  6,
    "complet":    9,
    "approfondi": 14,
}

def _resolve_nb_slides(value) -> int:
    """
    Accepte un int, un string numérique, ou un label
    ("Essentiel" / "Complet" / "Approfondi").
    """
    if isinstance(value, str):
        v = value.strip().lower()
        if v in NB_SLIDES_MAP:
            return NB_SLIDES_MAP[v]
        try:
            return max(2, min(int(v), 30))
        except ValueError:
            return 8
    try:
        return max(2, min(int(value), 30))
    except (ValueError, TypeError):
        return 8


def run_pipeline(pptx_bytes: bytes, prompt: str, nb_slides: int) -> tuple:
    if not ANTHROPIC_API_KEY:
        raise ValueError("Clé API Claude manquante.")

    prs = Presentation(io.BytesIO(pptx_bytes))
    nb_slides = max(2, min(nb_slides, 30))

    log.info("Phase 1 : analyse du template...")
    brand     = extract_brand(prs)
    library   = build_layout_library(prs)
    selection = select_template_slides(library, nb_slides)
    log.info(f"Template : {len(library)} slides → {len(selection)} sélectionnées pour {nb_slides} demandées")

    log.info("Phase 2 : planification narrative...")
    plan = plan_presentation(prompt, nb_slides, selection, brand)

    plan_slides = plan.get("slides", [])
    while len(plan_slides) < nb_slides:
        fallback = selection[min(len(plan_slides), len(selection) - 2)]
        plan_slides.append({
            "plan_index":           len(plan_slides),
            "template_slide_index": fallback["slide_index"],
            "slide_type":           fallback["slide_type"],
            "narrative_angle":      "Développement complémentaire",
            "key_message":          "Argument additionnel",
            "visual_hint":          "",
        })
    plan["slides"] = plan_slides[:nb_slides]

    log.info("Phase 3 : génération du contenu...")
    mapping = generate_content(prompt, plan, selection, brand)

    log.info("Hydratation PPTX...")
    final_bytes = hydrate_presentation(pptx_bytes, mapping, plan["slides"], nb_slides)

    return final_bytes, plan, brand


# ══════════════════════════════════════════════════════════════
# ROUTES NIVEAU 1
# ══════════════════════════════════════════════════════════════

@app.get("/")
def root():
    return {"status": "ok", "version": "14.0.0", "model": CLAUDE_MODEL,
            "pipeline": "V4", "levels": ["L1: /generate", "L2: /generate-v2"]}


@app.get("/health-v4")
def health_v4():
    """Diagnostic : vérifie que V4 est actif et lister les layouts disponibles."""
    layouts = [
        'cover','section','closing','full_text','list_numbered','list_cards',
        'two_col','kpi_grid','stat_hero','bar_chart','line_chart','pie_chart',
        'stacked_bar','waterfall','radar','timeline','process_flow','funnel',
        'matrix_2x2','swot','pros_cons','table','pyramid','cycle','roadmap',
        'quote','before_after','highlight_box','agenda',
    ]
    return {
        "pipeline":      "V4",
        "version":       "14.0.0",
        "model":         CLAUDE_MODEL,
        "api_key_set":   bool(ANTHROPIC_API_KEY),
        "n_layouts":     len(layouts),
        "layouts":       layouts,
        "n_variants":    2,
        "variant_mode":  "deterministic_hash",
    }


@app.post("/analyze-template")
async def analyze_template(file: UploadFile = File(...)):
    pptx_bytes = await file.read()
    prs   = Presentation(io.BytesIO(pptx_bytes))
    brand = extract_brand(prs)
    lib   = build_layout_library(prs)

    type_counts: dict = defaultdict(int)
    for s in lib:
        type_counts[s["slide_type"]] += 1

    fonts_display = ", ".join(brand["fonts"]) if brand["fonts"] else "Standard"

    return {
        "success":     True,
        "message":     (
            f"Charte détectée : {fonts_display} • "
            f"{len(brand['colors'])} couleurs • "
            f"{brand['slide_count']} slides • "
            f"{brand['aspect_ratio']}"
        ),
        "brand":       brand,
        "slide_types": dict(type_counts),
    }


@app.post("/generate-preview")
async def generate_preview(
    request:       Request,
    template:      UploadFile = File(...),
    prompt:        str        = Form(...),
    nb_slides:     str        = Form(default="complet"),
    authorization: str        = Form(default=None),
):
    pro = _is_pro(authorization)
    quota_info = (
        {"plan": "pro"} if pro
        else {"used": _quota(_ip(request))[0], "total": FREE_QUOTA_PER_IP, "plan": "free"}
    )

    import asyncio as _aio
    n          = _resolve_nb_slides(nb_slides)
    pptx_bytes = await template.read()
    prs        = Presentation(io.BytesIO(pptx_bytes))
    brand      = extract_brand(prs)
    palette    = _h2_extract_palette(brand)
    plan       = await _aio.to_thread(plan_presentation_v3, prompt, n, palette)

    return {
        "success":            True,
        "presentation_title": plan.get("presentation_title", prompt[:60]),
        "footer_text":        plan.get("footer_text", ""),
        "slides": [
            {
                "index":  i,
                "layout": s.get("layout"),
                "title":  s.get("content", {}).get("title", ""),
            }
            for i, s in enumerate(plan.get("slides", []))
        ],
        "brand": brand,
        "quota": quota_info,
    }


@app.post("/generate")
async def generate_presentation(
    request:       Request,
    template:      UploadFile = File(...),
    prompt:        str        = Form(...),
    nb_slides:     str        = Form(default="complet"),
    document:      UploadFile = File(default=None),
    authorization: str        = Form(default=None),
):
    if not _is_pro(authorization):
        _quota(_ip(request))

    n          = _resolve_nb_slides(nb_slides)
    pptx_bytes = await template.read()

    # Extraction du document source optionnel
    doc_content = ''
    if document is not None:
        try:
            doc_bytes = await document.read()
            _doc_mb   = len(doc_bytes) / (1024 * 1024)
            if _doc_mb > 20:
                raise HTTPException(400, f'Document trop volumineux ({_doc_mb:.1f} MB, max 20 MB)')
            doc_content = extract_document_content(doc_bytes, document.filename or 'doc')
        except HTTPException:
            raise
        except Exception as e:
            log.warning(f'[/generate] document extraction failed: {e}')

    import asyncio as _asyncio
    # V4 → V3 → L1
    try:
        final_bytes, plan, _, _pal = await run_pipeline_v4(pptx_bytes, prompt, n, doc_content)
        log.info('[/generate] Pipeline V4 OK')
    except Exception as e:
        log.warning(f'[/generate] V4 échoué ({e}) → fallback V3')
        try:
            final_bytes, plan, _, _pal = await _asyncio.to_thread(
                run_pipeline_v3, pptx_bytes, prompt, n
            )
            log.info('[/generate] Pipeline V3 OK (fallback)')
        except Exception as e2:
            log.warning(f'[/generate] V3 échoué ({e2}) → fallback L1')
            final_bytes, plan, _ = await _asyncio.to_thread(
                run_pipeline, pptx_bytes, prompt, n
            )

    filename = f"visualcortex-{_safe_name(prompt)}.pptx"
    return StreamingResponse(
        io.BytesIO(final_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename={filename}",
            "Content-Length":      str(len(final_bytes)),
        },
    )


# ─────────────────────────────────────────────────────────────
# GÉNÉRATION AVEC PROGRESSION SSE
# Affiche la progression en temps réel côté frontend.
# Protocole : POST multipart, réponse text/event-stream
# Chaque event : { step, label, progress (0-100) }
# Event final "done" : + file_b64 (base64), filename
# ─────────────────────────────────────────────────────────────

_PROGRESS_LABELS = {
    "start":     ("Analyse du template…",           5),
    "planned":   ("Structure narrative créée…",     35),
    "generated": ("Contenu généré…",                70),
    "hydrating": ("Assemblage de la présentation…", 90),
    "done":      ("Prêt !",                         100),
    "error":     ("Erreur…",                         0),
}

def _sse(event: str, data: dict) -> str:
    label, pct = _PROGRESS_LABELS.get(event, (event, 50))
    payload = json.dumps({"step": event, "label": label, "progress": pct, **data},
                         ensure_ascii=False)
    return f"data: {payload}\n\n"


@app.post("/generate-stream")
async def generate_stream(
    request:       Request,
    template:      UploadFile = File(...),
    prompt:        str        = Form(default=''),
    nb_slides:     int        = Form(default=8),
    document:      UploadFile = File(default=None),
    authorization: str        = Form(default=None),
):
    """
    Génération avec progression SSE (V4).
    Event final "done" contient le fichier PPTX encodé en base64.
    prompt peut être vide si un document source est fourni.
    """
    if not _is_pro(authorization):
        _quota(_ip(request))

    pptx_bytes = await template.read()

    # Extraction document optionnel
    doc_content = ''
    if document is not None:
        try:
            doc_bytes = await document.read()
            _doc_mb   = len(doc_bytes) / (1024 * 1024)
            if _doc_mb > 20:
                raise ValueError(f'Document trop volumineux ({_doc_mb:.1f} MB, max 20 MB)')
            doc_content = extract_document_content(doc_bytes, document.filename or 'doc')
        except Exception as e:
            log.warning(f'[/generate-stream] document extraction failed: {e}')

    # Prompt par défaut si document fourni sans texte libre
    prompt_val = prompt.strip()
    if not prompt_val:
        prompt_val = (
            'Crée une présentation structurée synthétisant ce document.'
            if doc_content else 'Présentation'
        )

    async def _stream():
        import asyncio as _aio, base64 as _b64
        try:
            yield _sse("start", {"nb_slides": nb_slides})

            # Phase 1 : analyze template (rapide, sync via thread)
            prs_tmp = Presentation(io.BytesIO(pptx_bytes))
            tp_tmp  = await _aio.to_thread(analyze_template_v4, prs_tmp)

            # Phase 2 : planning V4 async (non-bloquant, une seule fois)
            plan = await plan_presentation_v4(prompt_val, nb_slides, tp_tmp, doc_content)
            yield _sse("planned", {"title": plan.get("presentation_title", "")})

            # Phase 3+4 : création slides + export (plan pré-fourni → pas de double-appel)
            yield _sse("generated", {})
            yield _sse("hydrating", {})

            try:
                final_bytes, _plan, _brand, _pal = await run_pipeline_v4(
                    pptx_bytes, prompt_val, nb_slides, doc_content, plan=plan
                )
                log.info("[/generate-stream] Pipeline V4 OK")
            except Exception as e_v4:
                log.warning(f"[/generate-stream] V4 échoué ({e_v4}) → fallback V3")
                final_bytes, _plan, _brand, _pal = await _aio.to_thread(
                    run_pipeline_v3, pptx_bytes, prompt_val, nb_slides
                )

            b64      = _b64.b64encode(final_bytes).decode()
            filename = f"visualcortex-{_safe_name(prompt_val)}.pptx"
            yield _sse("done", {"file_b64": b64, "filename": filename})

        except Exception as e:
            log.error(f"[stream] Erreur : {e}", exc_info=True)
            yield _sse("error", {"message": str(e)})

    return StreamingResponse(
        _stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control":               "no-cache",
            "X-Accel-Buffering":           "no",
            "Access-Control-Allow-Origin": "*",
        },
    )




# ══════════════════════════════════════════════════════════════════════════════
#
#  ███╗   ██╗██╗██╗   ██╗███████╗ █████╗ ██╗   ██╗    ██████╗
#  ████╗  ██║██║██║   ██║██╔════╝██╔══██╗██║   ██║    ╚════██╗
#  ██╔██╗ ██║██║██║   ██║█████╗  ███████║██║   ██║     █████╔╝
#  ██║╚██╗██║██║╚██╗ ██╔╝██╔══╝  ██╔══██║██║   ██║    ██╔═══╝
#  ██║ ╚████║██║ ╚████╔╝ ███████╗██║  ██║╚██████╔╝    ███████╗
#  ╚═╝  ╚═══╝╚═╝  ╚═══╝  ╚══════╝╚═╝  ╚═╝ ╚═════╝    ╚══════╝
#
#  GÉNÉRATION CRÉATIVE — Claude voit le template, crée de A à Z
#
#  Architecture v14 :
#  Phase 1 — Analyse brand + extraction visuels (thumbnail + palette image)
#  Phase 2 — Planification narrative (commun L1)
#  Phase 3 — Claude VOIT le template → génère code python-pptx par slide
#  Phase 4 — Exécution sandbox + suppression slides originales + export
#
# ══════════════════════════════════════════════════════════════════════════════


# ─────────────────────────────────────────────────────────────
# PROFILS CLIENT
# ─────────────────────────────────────────────────────────────

CLIENT_PROFILES = {
    "finance": {
        "label":       "Finance / Conseil",
        "description": "Premium, minimaliste, données en avant",
        "style_guide": (
            "Fond blanc ou très clair. Grandes marges. Hiérarchie typographique nette. "
            "Accent doré ou bleu marine. KPIs proéminents. Zéro décoration superflue. "
            "Lignes fines. Espacement aéré. Confiance et autorité."
        ),
        "layout_prefs": ["kpi", "two_col", "full_text", "quote", "timeline"],
        "bg_dark":      False,
    },
    "industrial": {
        "label":       "Industriel / Technique",
        "description": "Clarté, données, schémas, sobriété",
        "style_guide": (
            "Fond blanc ou gris très clair. Sans-serif sobre. Données et faits en premier. "
            "Couleurs : bleu, gris, orange accent. Timelines et schémas favorisés. "
            "Structure claire, pas d'effets visuels."
        ),
        "layout_prefs": ["timeline", "kpi", "two_col", "list"],
        "bg_dark":      False,
    },
    "institutional": {
        "label":       "Institutionnel / Public",
        "description": "Formalismes respectés, sobre, structuré",
        "style_guide": (
            "Fond blanc. Typographie classique. Structure formelle et lisible. "
            "Éviter les effets visuels marqués. Hiérarchie claire. Sobriété maximale. "
            "Bleu institutionnel, blanc, accent discret."
        ),
        "layout_prefs": ["full_text", "list", "two_col", "timeline"],
        "bg_dark":      False,
    },
    "startup": {
        "label":       "Startup / Tech",
        "description": "Moderne, aéré, accents colorés",
        "style_guide": (
            "Fond très clair ou sombre. Sans-serif bold. Grands espaces blancs. "
            "1-2 couleurs vives en accent. Peu de texte, impact fort. "
            "Chiffres oversize, layouts asymétriques."
        ),
        "layout_prefs": ["cover", "quote", "kpi", "image_text", "closing"],
        "bg_dark":      True,
    },
    "creative": {
        "label":       "Créatif / Agence",
        "description": "Audacieux, typographie expressive",
        "style_guide": (
            "Fond sombre ou couleur franche. Typographie oversize. Compositions audacieuses. "
            "Couleurs inattendues. Géométrie forte. "
            "Peu de texte — chaque mot compte. Impact visuel prime."
        ),
        "layout_prefs": ["cover", "quote", "image_text", "section", "closing"],
        "bg_dark":      True,
    },
}


# ─────────────────────────────────────────────────────────────
# EXTRACTION VISUELS DU TEMPLATE (pour vision Claude)
# ─────────────────────────────────────────────────────────────

def _extract_template_thumbnail(pptx_bytes: bytes) -> dict | None:
    """
    Extrait la miniature embarquée dans le PPTX (docProps/thumbnail.*).
    Présente dans ~90% des fichiers créés par PowerPoint / Google Slides.
    Retourne {'media_type': ..., 'data': base64_string} ou None.
    """
    try:
        with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
            for name in z.namelist():
                lower = name.lower()
                if 'thumbnail' in lower and lower.endswith(('.jpeg', '.jpg', '.png')):
                    data = z.read(name)
                    mt = 'image/png' if lower.endswith('.png') else 'image/jpeg'
                    return {
                        'media_type': mt,
                        'data':       base64.standard_b64encode(data).decode(),
                    }
    except Exception as e:
        log.warning(f"[V2] Extraction thumbnail : {e}")
    return None


def _make_palette_swatch(palette: dict) -> dict | None:
    """
    Génère une image de palette couleurs (600×100px) avec Pillow.
    Permet à Claude de voir les vraies couleurs de la charte.
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
        W, H   = 700, 100
        img    = Image.new('RGB', (W, H), (255, 255, 255))
        draw   = ImageDraw.Draw(img)
        roles  = ['primary', 'secondary', 'accent', 'light', 'text']
        labels = ['Primary', 'Secondary', 'Accent', 'Light', 'Text']
        bw     = W // len(roles)
        for i, (role, label) in enumerate(zip(roles, labels)):
            hex_val = palette.get(role, 'CCCCCC').lstrip('#')
            try:
                r, g, b = int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
            except Exception:
                r, g, b = 128, 128, 128
            draw.rectangle([i * bw, 0, (i + 1) * bw - 2, 72], fill=(r, g, b))
            # Label blanc ou noir selon luminosité
            lum = 0.299 * r + 0.587 * g + 0.114 * b
            fg  = (255, 255, 255) if lum < 140 else (30, 30, 30)
            draw.text((i * bw + 6, 52), f"#{hex_val}", fill=fg)
            draw.text((i * bw + 6, 76), label, fill=(60, 60, 60))
        buf = io.BytesIO()
        img.save(buf, format='JPEG', quality=90)
        return {
            'media_type': 'image/jpeg',
            'data':       base64.standard_b64encode(buf.getvalue()).decode(),
        }
    except Exception as e:
        log.warning(f"[V2] Palette swatch : {e}")
        return None


# ─────────────────────────────────────────────────────────────
# HELPERS PALETTE
# ─────────────────────────────────────────────────────────────

def _lighten(hex_str: str, factor: float) -> str:
    """Éclaircit une couleur hex d'un facteur 0–1 vers le blanc."""
    h = hex_str.lstrip('#')
    try:
        r, g, b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
        return (f"{int(r + (255-r)*factor):02X}"
                f"{int(g + (255-g)*factor):02X}"
                f"{int(b + (255-b)*factor):02X}")
    except Exception:
        return 'AAAAAA'


def _complementary(hex_str: str) -> str:
    """Génère une couleur d'accent complémentaire sobre (décalage 150°)."""
    import colorsys
    h = hex_str.lstrip('#')
    try:
        r, g, b = int(h[0:2],16)/255, int(h[2:4],16)/255, int(h[4:6],16)/255
        hue, sat, val = colorsys.rgb_to_hsv(r, g, b)
        hue = (hue + 150/360) % 1.0
        sat = min(sat, 0.7)
        r2, g2, b2 = colorsys.hsv_to_rgb(hue, sat, min(val * 1.1, 1.0))
        return f"{int(r2*255):02X}{int(g2*255):02X}{int(b2*255):02X}"
    except Exception:
        return 'F0A500'


def _find_darkest_neutral(colors: list) -> str | None:
    """Trouve le gris le plus foncé utilisable comme primary dans un template monochrome."""
    candidates = []
    for c in colors:
        h = c.lstrip('#')
        if len(h) != 6:
            continue
        try:
            r, g, b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
            lum = 0.299*r + 0.587*g + 0.114*b
            if abs(r-g) < 20 and abs(g-b) < 20 and 30 < lum < 180:
                candidates.append((lum, h))
        except Exception:
            pass
    if candidates:
        return min(candidates, key=lambda x: x[0])[1]
    return None


# ─────────────────────────────────────────────────────────────
# EXTRACTION DE PALETTE
# ─────────────────────────────────────────────────────────────

def _h2_extract_palette(brand: dict) -> dict:
    """
    Construit la palette depuis les couleurs réelles du template.
    - CAS 1 : template coloré (≥2 chromatiques) → ses couleurs directement
    - CAS 2 : template semi-coloré (1 chromatique) → complète avec dérivés
    - CAS 3 : template monochrome → palette élégante depuis neutres
    Source de vérité : scan fréquence d'abord, theme_colors en complément.
    """
    def _lum(h: str) -> float:
        try:
            r, g, b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
            return 0.299*r + 0.587*g + 0.114*b
        except Exception:
            return 128.0

    fonts      = brand.get('fonts', [])
    chromatic  = [c.lstrip('#').strip() for c in brand.get('colors', [])
                  if len(c.lstrip('#').strip()) == 6]
    all_colors = [c.lstrip('#').strip() for c in brand.get('all_colors', [])
                  if len(c.lstrip('#').strip()) == 6]
    tc = {k: v.lstrip('#').strip() for k, v in brand.get('theme_colors', {}).items()
          if v and len(v.lstrip('#').strip()) == 6}

    # ── Supplément depuis theme si scan freq trop creux ──────────────────────
    # (templates qui utilisent des références de thème plutôt qu'RGB explicite)
    if len(chromatic) < 2 and tc:
        _seen: set = set(chromatic)
        for _slot in ['accent1', 'dk2', 'accent2']:  # slots principaux seulement
            _c = tc.get(_slot, '')
            if not _c or _c in _seen:
                continue
            _r2, _g2, _b2 = int(_c[0:2],16), int(_c[2:4],16), int(_c[4:6],16)
            _lm = 0.299*_r2 + 0.587*_g2 + 0.114*_b2
            _mx = max(_r2, _g2, _b2)
            _sat = (_mx - min(_r2,_g2,_b2)) / _mx if _mx > 0 else 0
            if _sat > 0.2 and 20 < _lm < 230:
                chromatic.append(_c)
                _seen.add(_c)

    if not all_colors and tc:
        all_colors = [v for v in tc.values() if v]

    palette: dict = {'font': fonts[0] if fonts else 'Calibri', 'theme': tc}

    # ── CAS 1 : template coloré (≥2 chromatiques) ────────────────────────────
    if len(chromatic) >= 2:
        palette['primary']   = chromatic[0]
        palette['secondary'] = chromatic[1]
        palette['accent']    = chromatic[2] if len(chromatic) >= 3 else chromatic[0]

    # ── CAS 2 : 1 couleur chromatique → dérivés ───────────────────────────────
    elif len(chromatic) == 1:
        base = chromatic[0]
        palette['primary']   = base
        palette['secondary'] = _lighten(base, 0.25)
        palette['accent']    = _complementary(base)

    # ── CAS 3 : template monochrome ───────────────────────────────────────────
    else:
        dn = _find_darkest_neutral(all_colors)
        if dn:
            palette['primary']   = dn
            palette['secondary'] = _lighten(dn, 0.35)
            palette['accent']    = _lighten(dn, 0.65)
        else:
            palette['primary']   = '1A1A1A'
            palette['secondary'] = '444444'
            palette['accent']    = '888888'

    # ── Couleur texte : dk1 du thème (le plus fiable), sinon scan ────────────
    if 'dk1' in tc and _lum(tc['dk1']) < 150:
        palette['text'] = tc['dk1']
    else:
        palette['text'] = next((c for c in all_colors if 15 < _lum(c) < 120), '1A1A1A')

    # ── Couleur sombre (fonds dark slides) : dk1 si vraiment sombre ──────────
    if 'dk1' in tc and _lum(tc['dk1']) < 100:
        palette['dark'] = tc['dk1']
    else:
        dark_cands = [(c, _lum(c)) for c in all_colors if 10 < _lum(c) < 100]
        palette['dark'] = min(dark_cands, key=lambda x: x[1])[0] if dark_cands else ''

    # ── Dérivés depuis primary si dark/light non définis ─────────────────────
    p = palette['primary']
    try:
        r, g, b = int(p[0:2],16), int(p[2:4],16), int(p[4:6],16)
        if not palette.get('dark'):
            palette['dark']  = f"{max(0,int(r*0.25)):02X}{max(0,int(g*0.25)):02X}{max(0,int(b*0.25)):02X}"
        palette['light'] = f"{int(r*0.08+255*0.92):02X}{int(g*0.08+255*0.92):02X}{int(b*0.08+255*0.92):02X}"
    except Exception:
        if not palette.get('dark'): palette['dark'] = '1A1A1A'
        palette['light'] = 'F0F4FA'

    # ── Couleur de fond des slides claires : lt1 du thème ────────────────────
    raw_lt1 = tc.get('lt1', 'FFFFFF') if tc else 'FFFFFF'
    # Exclure les valeurs invalides ou trop sombres (lt1 doit être clair)
    try:
        h = raw_lt1.lstrip('#')
        if len(h) == 6 and _lum(h) > 180:
            palette['bg'] = h.upper()
        else:
            palette['bg'] = 'FFFFFF'
    except Exception:
        palette['bg'] = 'FFFFFF'

    log.info(f"[palette] primary=#{palette['primary']} secondary=#{palette['secondary']} "
             f"accent=#{palette['accent']} light=#{palette['light']} "
             f"dark=#{palette['dark']} text=#{palette['text']} "
             f"bg=#{palette['bg']} chromatic={chromatic[:4]} font={palette['font']}")
    return palette


# ─────────────────────────────────────────────────────────────
# HELPERS H2_* — Bibliothèque de génération de shapes
# Chaque helper est exposé dans le sandbox exec() du Niveau 2.
# RÈGLE : aucun des helpers ne prend prs en argument.
#         h2_blank_slide() s'appelle SANS argument (prs est dans la closure).
# ─────────────────────────────────────────────────────────────

def _h2_parse_hex(hex_str: str) -> RGBColor:
    """Parse 'RRGGBB' ou '#RRGGBB' → RGBColor. Fallback bleu corporate si invalide."""
    try:
        h = str(hex_str).lstrip('#').strip()
        if len(h) == 3:
            h = h[0]*2 + h[1]*2 + h[2]*2
        if len(h) != 6:
            return RGBColor(0x1A, 0x3A, 0x6B)
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return RGBColor(0x1A, 0x3A, 0x6B)


def _h2_blank_slide(prs: Presentation):
    """
    Ajoute une slide entièrement vierge (sans placeholders résiduels).
    Cherche un layout 'Blank' par nom, sinon prend le layout index 6 ou le dernier.
    ⚠ Dans le code généré : appeler h2_blank_slide() SANS argument.
    Retourne (slide, W, H) en pouces.
    """
    target = None
    for layout in prs.slide_layouts:
        if 'blank' in layout.name.lower():
            target = layout
            break
    if target is None:
        idx = min(6, len(prs.slide_layouts) - 1)
        target = prs.slide_layouts[idx]

    slide   = prs.slides.add_slide(target)
    sp_tree = slide.shapes._spTree
    for ph in list(slide.placeholders):
        try:
            sp_tree.remove(ph._element)
        except Exception:
            pass

    # Garantir l'héritage du master (logo, éléments décoratifs)
    import lxml.etree as _etree
    cSld = slide._element.find(qn('p:cSld'))
    if cSld is not None:
        cSld.set('showMasterSp', '1')

    W = prs.slide_width  / 914400.0
    H = prs.slide_height / 914400.0
    return slide, W, H


def _h2_rect(slide, left: float, top: float, width: float, height: float,
             color: str, alpha: int = 0):
    """
    Rectangle coloré plein (en pouces). color = hex 'RRGGBB'.
    alpha ignoré (python-pptx ne supporte pas la transparence sur les formes).
    """
    width  = max(0.02, width)
    height = max(0.02, height)
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _h2_parse_hex(color)
    shape.line.fill.background()
    return shape


def _h2_rounded_rect(slide, left: float, top: float, width: float, height: float,
                     color: str, radius: float = 0.08):
    """
    Rectangle arrondi (pill). radius = proportion d'arrondi (0–0.5).
    """
    width  = max(0.05, width)
    height = max(0.05, height)
    from pptx.util import Emu as _Emu
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _h2_parse_hex(color)
    shape.line.fill.background()
    try:
        shape.adjustments[0] = max(0.0, min(0.5, radius))
    except Exception:
        pass
    return shape


def _h2_circle(slide, cx: float, cy: float, r: float, color: str):
    """
    Cercle centré en (cx, cy) de rayon r (en pouces).
    """
    shape = slide.shapes.add_shape(
        9,  # MSO_SHAPE.OVAL
        Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _h2_parse_hex(color)
    shape.line.fill.background()
    return shape


def _h2_text(slide, text: str,
             left: float, top: float, width: float, height: float,
             font: str, size_pt: float, color: str,
             bold: bool = False, italic: bool = False, align: str = 'left',
             line_spacing: float = 1.0):
    """
    Textbox stylée (dimensions en pouces).
    align : 'left' | 'center' | 'right'
    line_spacing : multiplicateur d'interligne (1.0 = normal, 1.2 = aéré)
    """
    # Protéger contre les dimensions invalides
    width  = max(0.1, width)
    height = max(0.1, height)
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf           = txBox.text_frame
    tf.word_wrap = True
    # Clip text at box boundary — prevents overflow into adjacent elements
    try:
        import lxml.etree as _et2
        bodyPr = tf._txBody.get_or_add_bodyPr()
        for tag in ('a:spAutoFit', 'a:normAutofit', 'a:noAutofit'):
            el = bodyPr.find(qn(tag))
            if el is not None:
                bodyPr.remove(el)
        _et2.SubElement(bodyPr, qn('a:normAutofit'))
    except Exception:
        pass

    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}

    p           = tf.paragraphs[0]
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)

    # Interligne
    try:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn as _qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        import lxml.etree as _etree
        lnSpc = _etree.SubElement(pPr, _qn('a:lnSpc'))
        spcPct = _etree.SubElement(lnSpc, _qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 100000)))
    except Exception:
        pass

    run             = p.add_run()
    run.text        = str(text)
    run.font.name   = str(font)
    run.font.size   = Pt(size_pt)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = _h2_parse_hex(color)
    return txBox


def _h2_multiline_text(slide, lines: list,
                       left: float, top: float, width: float, height: float,
                       font: str, size_pt: float, color: str,
                       bold: bool = False, align: str = 'left',
                       line_spacing: float = 1.15):
    """
    Textbox avec plusieurs paragraphes (une entrée de `lines` par paragraphe).
    Chaque entrée peut être un str ou un dict {'text', 'bold', 'size', 'color'}.
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf           = txBox.text_frame
    tf.word_wrap = True
    align_map    = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}

    for i, line in enumerate(lines):
        if isinstance(line, dict):
            txt    = line.get('text', '')
            b      = line.get('bold', bold)
            sz     = line.get('size', size_pt)
            clr    = line.get('color', color)
        else:
            txt, b, sz, clr = str(line), bold, size_pt, color

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)

        run             = p.add_run()
        run.text        = txt
        run.font.name   = font
        run.font.size   = Pt(sz)
        run.font.bold   = b
        run.font.color.rgb = _h2_parse_hex(clr)
    return txBox


def _h2_kpi(slide,
            left: float, top: float, width: float,
            value: str, label: str, sublabel: str,
            palette: dict, dark: bool = True):
    """
    Bloc KPI : grande valeur + label + sous-label.
    dark=True  → textes blancs (pour fonds sombres)
    dark=False → textes couleur charte (pour fonds clairs)
    width ≈ 2.0–3.0 pouces.
    """
    font    = palette.get('font', 'Calibri')
    v_clr   = 'FFFFFF' if dark else palette.get('primary', '1A3A6B')
    l_clr   = palette.get('accent', 'F0A500')
    sl_clr  = 'CCCCCC' if dark else '888888'

    _h2_text(slide, str(value),  left, top,        width, 0.75, font, 34, v_clr, bold=True, align='center')
    _h2_rect(slide, left + width*0.18, top + 0.76,  width*0.64, 0.035, l_clr)
    _h2_text(slide, str(label),  left, top + 0.82,  width, 0.38, font, 12, l_clr,  bold=False, align='center')
    _h2_text(slide, str(sublabel), left, top + 1.22, width, 0.52, font, 10, sl_clr, bold=False, align='center')


def _h2_divider(slide, left: float, top: float, width: float,
                color: str, thickness: float = 0.035):
    """Ligne horizontale fine (en pouces)."""
    _h2_rect(slide, left, top, width, thickness, color)


def _h2_number(slide, text: str, left: float, top: float,
               size_in: float, color: str, font: str, opacity_hint: str = ''):
    """
    Grand chiffre / lettre décoratif (numéro de section, oversize stat).
    size_in = hauteur approximative de la boîte en pouces.
    """
    font_pt = max(28, int(size_in * 70))
    _h2_text(slide, str(text), left, top, size_in * 2.0, size_in + 0.2,
             font, font_pt, color, bold=True, align='center')


def _h2_icon_circle(slide, cx: float, cy: float, r: float,
                    label: str, font: str, color_bg: str, color_fg: str,
                    font_size: float = 11):
    """
    Cercle coloré avec un label centré dedans (substitut d'icône).
    """
    _h2_circle(slide, cx, cy, r, color_bg)
    _h2_text(slide, label, cx - r, cy - r * 0.55, r * 2, r * 1.1,
             font, font_size, color_fg, bold=True, align='center')


def _h2_card(slide, left: float, top: float, width: float, height: float,
             bg_color: str, title: str, body: str,
             font: str, title_color: str, body_color: str,
             title_size: float = 13, body_size: float = 11,
             rounded: bool = True):
    """
    Carte avec fond coloré + titre en gras + corps.
    Utile pour grilles, colonnes, items de liste visuels.
    """
    if rounded:
        _h2_rounded_rect(slide, left, top, width, height, bg_color, radius=0.06)
    else:
        _h2_rect(slide, left, top, width, height, bg_color)
    pad = 0.18
    _h2_text(slide, title, left + pad, top + pad, width - pad*2, 0.42,
             font, title_size, title_color, bold=True)
    if body:
        _h2_text(slide, body, left + pad, top + pad + 0.44, width - pad*2,
                 height - pad*2 - 0.44, font, body_size, body_color,
                 line_spacing=1.2)


def _h2_progress_bar(slide, left: float, top: float, width: float,
                     value_pct: float, color_fill: str, color_bg: str = 'EEEEEE',
                     height: float = 0.12):
    """
    Barre de progression horizontale. value_pct entre 0 et 100.
    """
    _h2_rect(slide, left, top, width, height, color_bg)
    fill_w = max(0.0, min(width, width * value_pct / 100.0))
    if fill_w > 0:
        _h2_rect(slide, left, top, fill_w, height, color_fill)


def _h2_tag(slide, text: str, left: float, top: float,
            font: str, size_pt: float, bg_color: str, fg_color: str):
    """
    Pill / tag arrondi avec texte centré. Hauteur auto ≈ size_pt * 0.022 po.
    """
    h = size_pt * 0.028 + 0.08
    w = len(text) * size_pt * 0.014 + 0.3
    _h2_rounded_rect(slide, left, top, w, h, bg_color, radius=0.5)
    _h2_text(slide, text, left, top, w, h, font, size_pt, fg_color,
             bold=True, align='center')


# ─────────────────────────────────────────────────────────────
# SÉCURITÉ — Validation du code généré
# ─────────────────────────────────────────────────────────────

_FORBIDDEN_CODE_PATTERNS = [
    'import ', '__import__', 'eval(', 'exec(',
    'open(', 'os.', 'sys.', 'subprocess', 'socket',
    'urllib', 'requests', '__builtins__', '__globals__',
    '__locals__', '__class__', 'compile(', 'globals(',
    'locals(', 'vars(', 'dir(',
]

def _validate_code_safety(code: str) -> tuple:
    for pattern in _FORBIDDEN_CODE_PATTERNS:
        if pattern in code:
            return False, f"Pattern interdit: '{pattern}'"
    return True, ''


# ─────────────────────────────────────────────────────────────
# NAMESPACE SANDBOX
# ─────────────────────────────────────────────────────────────

def _build_safe_namespace(prs: Presentation, palette: dict) -> dict:
    """
    Namespace restreint exposé au code généré.
    ⚠ h2_blank_slide est une lambda SANS argument (prs dans la closure).
    """
    return {
        'brand':  palette,
        # Helpers — h2_blank_slide() s'appelle SANS argument
        'h2_blank_slide':    lambda: _h2_blank_slide(prs),
        'h2_rect':           _h2_rect,
        'h2_rounded_rect':   _h2_rounded_rect,
        'h2_circle':         _h2_circle,
        'h2_text':           _h2_text,
        'h2_multiline_text': _h2_multiline_text,
        'h2_kpi':            _h2_kpi,
        'h2_divider':        _h2_divider,
        'h2_number':         _h2_number,
        'h2_icon_circle':    _h2_icon_circle,
        'h2_card':           _h2_card,
        'h2_progress_bar':   _h2_progress_bar,
        'h2_tag':            _h2_tag,
        # Builtins sécurisés
        '__builtins__': {
            'range': range, 'len': len, 'int': int, 'float': float,
            'str': str, 'bool': bool, 'list': list, 'dict': dict,
            'tuple': tuple, 'enumerate': enumerate, 'zip': zip,
            'round': round, 'abs': abs, 'min': min, 'max': max,
            'sum': sum, 'sorted': sorted, 'reversed': reversed,
            'any': any, 'all': all, 'print': lambda *a, **k: None,
            'True': True, 'False': False, 'None': None,
        },
    }


# ─────────────────────────────────────────────────────────────
# PROMPTS NIVEAU 2
# ─────────────────────────────────────────────────────────────

_V2_SYSTEM = """\
Tu es Visual Cortex Level 2 — Créateur de présentations PowerPoint professionnelles.

Tu reçois :
1. Une image de miniature du template de l'entreprise (pour voir la charte visuelle réelle)
2. Une image de la palette de couleurs extraite
3. Le plan narratif à implémenter

TA MISSION : générer du code python-pptx pour chaque slide. Le résultat doit être :
→ Professionnel, impactant, visuellement varié
→ Fidèle à la charte de l'entreprise (ses vraies couleurs, ses proportions)
→ Digne d'une agence de communication

══════════════════════════════════════════════
FONCTIONS DISPONIBLES (et UNIQUEMENT celles-ci)
══════════════════════════════════════════════

slide, W, H = h2_blank_slide()
  → TOUJOURS la première ligne de chaque slide. AUCUN ARGUMENT.
  → W ≈ 10.0, H ≈ 5.63 pour 16:9. Dimensions en pouces.

h2_rect(slide, left, top, width, height, color)
  → Rectangle plein. color = brand["primary"] ou "FFFFFF".

h2_rounded_rect(slide, left, top, width, height, color, radius=0.08)
  → Rectangle arrondi. radius entre 0 et 0.5. Parfait pour cartes, tags.

h2_circle(slide, cx, cy, r, color)
  → Cercle centré (cx, cy) de rayon r. Tous en pouces.

h2_text(slide, text, left, top, width, height, font, size_pt, color,
        bold=False, italic=False, align="left", line_spacing=1.0)
  → Textbox. align: "left"|"center"|"right".
  → Utiliser brand["font"] pour TOUS les textes.

h2_multiline_text(slide, lines, left, top, width, height, font, size_pt, color,
                  bold=False, align="left", line_spacing=1.15)
  → Multi-paragraphes. lines = liste de str ou dict {"text","bold","size","color"}.

h2_kpi(slide, left, top, width, value, label, sublabel, brand, dark=True)
  → Bloc KPI. dark=True pour fond sombre, dark=False pour fond clair.
  → width ≈ 2.0–3.0 po.

h2_divider(slide, left, top, width, color, thickness=0.035)
  → Ligne horizontale fine. Utiliser avec parcimonie.

h2_number(slide, text, left, top, size_in, color, font)
  → Grand chiffre/lettre décoratif. size_in ≈ 1.5–2.5 po.

h2_icon_circle(slide, cx, cy, r, label, font, color_bg, color_fg, font_size=11)
  → Cercle coloré avec texte centré (substitut d'icône).

h2_card(slide, left, top, width, height, bg_color, title, body,
        font, title_color, body_color, title_size=13, body_size=11, rounded=True)
  → Carte avec fond + titre + corps. Parfait pour grilles.

h2_progress_bar(slide, left, top, width, value_pct, color_fill, color_bg="EEEEEE", height=0.12)
  → Barre de progression horizontale (0–100%).

h2_tag(slide, text, left, top, font, size_pt, bg_color, fg_color)
  → Pill / tag arrondi avec texte. Largeur automatique.

══════════════════════════════════════════════
ACCÈS À LA CHARTE (via brand dict)
══════════════════════════════════════════════
brand["primary"]   → couleur principale de la charte
brand["secondary"] → couleur secondaire
brand["accent"]    → couleur d'accent (titres, highlights)
brand["light"]     → version très claire du primary (fonds doux)
brand["text"]      → couleur de texte foncé
brand["font"]      → police principale à utiliser PARTOUT

Couleurs supplémentaires utiles (hardcodées) :
"FFFFFF" = blanc   "000000" = noir   "888888" = gris moyen
"F5F5F5" = gris très clair   "333333" = presque noir

══════════════════════════════════════════════
RÈGLES DE DENSITÉ (inviolables)
══════════════════════════════════════════════
cover    → titre ≤ 7 mots  + sous-titre ≤ 12 mots
section  → numéro décoratif + titre ≤ 6 mots UNIQUEMENT
kpi      → 4 à 6 KPIs MAX. Valeur courte + label court + sous-label contextuel
timeline → 4 à 6 jalons. Date + titre ≤ 4 mots + phrase optionnelle ≤ 12 mots
two_col  → 2 colonnes symétriques. Max 4 items/colonne ≤ 18 mots chacun
quote    → 1 citation ≤ 20 mots. Source ≤ 8 mots
list     → 3 à 5 items. Titre ≤ 5 mots + corps ≤ 20 mots
full_text→ titre ≤ 8 mots + 2-3 paragraphes ≤ 60 mots total
closing  → titre ≤ 5 mots + sous-titre ≤ 15 mots

══════════════════════════════════════════════
RÈGLES VISUELLES OBLIGATOIRES
══════════════════════════════════════════════
1. VARIÉTÉ : chaque slide DOIT avoir un look distinct des autres.
   Alterner fonds sombres (primary) et clairs (FFFFFF ou light).
   Ne jamais répéter la même composition deux fois.

2. STRUCTURE SANDWICH : cover et closing sur fond sombre/coloré.
   Slides de contenu sur fond clair. Sections sur fond primary.

3. RESPIRATION : marges min 0.4 po sur tous les bords.
   Espacement entre blocs ≥ 0.25 po.

4. HIÉRARCHIE : titre 28-44pt bold. Corps 12-16pt. Footer 10pt.
   Contraste fort entre titre et corps.

5. FOOTER : présent sur toutes les slides de contenu (pas cover ni closing).
   Toujours en bas, texte 10pt, couleur discrète.

6. COULEURS : utiliser UNIQUEMENT les couleurs brand[] + blanc + gris.
   Jamais inventer une couleur non présente dans la charte.

══════════════════════════════════════════════
PATTERNS PAR TYPE (à adapter à la charte vue)
══════════════════════════════════════════════

[cover — fond primary] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, brand["primary"])
h2_rect(slide, 0, H*0.75, W, H*0.25, brand["secondary"])
h2_text(slide, "Titre fort ≤ 7 mots", 0.7, H/2-1.2, W-1.4, 1.5,
        brand["font"], 44, "FFFFFF", bold=True)
h2_text(slide, "Sous-titre contextuel ≤ 12 mots", 0.7, H/2+0.45, W-1.4, 0.8,
        brand["font"], 20, "FFFFFF")
h2_text(slide, "Footer · Contexte · 2025", 0.7, H-0.52, W-1.4, 0.4,
        brand["font"], 10, brand["secondary"])

[cover — barre latérale sur fond clair] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, "FFFFFF")
h2_rect(slide, 0, 0, 0.55, H, brand["primary"])
h2_rect(slide, 0.55, H*0.45, W-0.55, 0.055, brand["accent"])
h2_text(slide, "Titre fort ≤ 7 mots", 0.9, H*0.45-0.9, W-1.4, 1.3,
        brand["font"], 40, brand["primary"], bold=True)
h2_text(slide, "Sous-titre ≤ 12 mots", 0.9, H*0.45+0.55, W-1.4, 0.75,
        brand["font"], 18, brand["text"])
h2_text(slide, "Footer · Contexte · 2025", 0.9, H-0.52, W-1.4, 0.4,
        brand["font"], 10, brand["text"])

[section] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, brand["primary"])
h2_number(slide, "01", 0.7, H/2-1.5, 2.0, brand["secondary"], brand["font"])
h2_text(slide, "Titre de section ≤ 6 mots", 0.7, H/2+0.6, W-1.2, 1.0,
        brand["font"], 34, "FFFFFF", bold=True)
h2_rect(slide, 0.7, H/2+0.52, 3.5, 0.055, brand["accent"])

[kpi — grille 3×2 sur fond sombre] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, brand["primary"])
h2_text(slide, "Titre KPIs ≤ 7 mots", 0.5, 0.22, W-1.0, 0.75,
        brand["font"], 30, "FFFFFF", bold=True)
h2_rect(slide, 0.5, 1.08, W-1.0, 0.04, brand["accent"])
kpi_data = [("600 M€","contribution","versée en 2022"),
            ("~5.6%","du capital","État actionnaire"),
            ("28 Md€","CA mondial","exercice 2023"),
            ("95k","collaborateurs","dont 25k en France"),
            ("80+","pays","présence opérationnelle"),
            ("1er rang","en France","par capitalisation")]
for i, (v, l, s) in enumerate(kpi_data):
    col = i % 3
    row = i // 3
    x = 0.4 + col * ((W-0.8)/3)
    y = 1.35 + row * 1.95
    h2_kpi(slide, x, y, (W-0.8)/3 - 0.15, v, l, s, brand, dark=True)
h2_text(slide, "Footer · Contexte · 2025", 0.5, H-0.46, W-1.0, 0.38,
        brand["font"], 10, brand["secondary"])

[kpi — ligne de 4 sur fond clair] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, "FFFFFF")
h2_rect(slide, 0, 0, 0.08, H, brand["primary"])
h2_text(slide, "Titre KPIs", 0.5, 0.22, W-1.0, 0.72,
        brand["font"], 28, brand["primary"], bold=True)
h2_rect(slide, 0.5, 1.05, W-1.0, 0.04, brand["accent"])
h2_rect(slide, 0.5, 1.28, W-1.0, 2.35, brand["light"])
kpi_data = [("600 M€","contribution","versée en 2022"),
            ("~5.6%","du capital","État actionnaire"),
            ("28 Md€","CA mondial","2023"),
            ("95k","collaborateurs","80 pays")]
kw = (W-1.0) / len(kpi_data)
for i, (v, l, s) in enumerate(kpi_data):
    h2_kpi(slide, 0.5 + i*kw, 1.45, kw-0.1, v, l, s, brand, dark=False)
h2_text(slide, "Footer · Contexte · 2025", 0.5, H-0.46, W-1.0, 0.38,
        brand["font"], 10, brand["primary"])

[timeline — axe horizontal sur fond clair] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, "FFFFFF")
h2_rect(slide, 0, 0, 0.08, H, brand["primary"])
h2_text(slide, "Chronologie", 0.5, 0.22, W-1.0, 0.72,
        brand["font"], 28, brand["primary"], bold=True)
h2_rect(slide, 0.5, 1.05, W-1.0, 0.04, brand["accent"])
h2_rect(slide, 0.5, 2.45, W-1.0, 0.065, brand["primary"])
steps = [("1924","Création","Fondation par l'État"),
         ("1985","Privatisation","Entrée en bourse"),
         ("2003","Fusion","Nouveau groupe mondial"),
         ("2024","Centenaire","Repositionnement stratégique")]
n = len(steps)
for i, (date, titre, detail) in enumerate(steps):
    x = 0.5 + i * ((W-1.0)/(n-1))
    h2_circle(slide, x, 2.48, 0.16, brand["primary"])
    h2_text(slide, date,  x-0.6, 1.68, 1.2, 0.52,
            brand["font"], 13, brand["primary"], bold=True, align="center")
    h2_text(slide, titre, x-0.7, 2.82, 1.4, 0.44,
            brand["font"], 12, brand["text"], bold=True, align="center")
    h2_text(slide, detail, x-0.8, 3.32, 1.6, 0.56,
            brand["font"], 10, "888888", align="center")
h2_text(slide, "Footer · Contexte · 2025", 0.5, H-0.46, W-1.0, 0.38,
        brand["font"], 10, brand["primary"])

[two_col — fond blanc, colonnes avec header coloré] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, "FFFFFF")
h2_rect(slide, 0, 0, 0.08, H, brand["primary"])
h2_text(slide, "Comparaison / Dualité", 0.5, 0.22, W-1.0, 0.72,
        brand["font"], 28, brand["primary"], bold=True)
h2_rect(slide, 0.5, 1.05, W-1.0, 0.04, brand["accent"])
col_w = (W-1.2) / 2
for ci, (label, items) in enumerate([
    ("COLONNE A", ["Aspect 1 concis", "Aspect 2 concis", "Aspect 3", "Aspect 4"]),
    ("COLONNE B", ["Aspect 1 bis", "Aspect 2 bis", "Aspect 3 bis", "Aspect 4 bis"])
]):
    x = 0.5 + ci * (col_w + 0.2)
    clr = brand["primary"] if ci == 0 else brand["secondary"]
    h2_rect(slide, x, 1.25, col_w, 0.5, clr)
    h2_text(slide, label, x+0.12, 1.30, col_w-0.24, 0.4,
            brand["font"], 13, "FFFFFF", bold=True)
    for j, item in enumerate(items):
        h2_text(slide, "→  " + item, x+0.12, 1.92+j*0.65, col_w-0.24, 0.58,
                brand["font"], 12, brand["text"])
h2_text(slide, "Footer · Contexte · 2025", 0.5, H-0.46, W-1.0, 0.38,
        brand["font"], 10, brand["primary"])

[quote — fond sombre, barre latérale accent] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, brand["primary"])
h2_rect(slide, 0.52, H/2-1.15, 0.1, 2.3, brand["accent"])
h2_text(slide, "\u00ab\u202fCitation forte et mémorable, ≤ 20 mots absolument.\u202f\u00bb",
        0.88, H/2-0.98, W-1.6, 1.85,
        brand["font"], 26, "FFFFFF", bold=True, line_spacing=1.25)
h2_text(slide, "\u2014 Auteur ou source, 2025",
        0.88, H/2+1.0, W-1.6, 0.58,
        brand["font"], 14, brand["accent"])
h2_text(slide, "Footer · Contexte · 2025", 0.5, H-0.46, W-1.0, 0.38,
        brand["font"], 10, brand["secondary"])

[list — cartes numérotées sur fond clair] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, "FFFFFF")
h2_rect(slide, 0, 0, 0.08, H, brand["primary"])
h2_text(slide, "Titre de la liste", 0.5, 0.22, W-1.0, 0.72,
        brand["font"], 28, brand["primary"], bold=True)
h2_rect(slide, 0.5, 1.05, W-1.0, 0.04, brand["accent"])
items = [("Titre 1","Corps concis, une idée forte, ≤ 20 mots."),
         ("Titre 2","Corps concis, une idée distincte."),
         ("Titre 3","Corps concis, conclusion ou conséquence."),
         ("Titre 4","Corps concis, argument complémentaire.")]
for i, (titre, corps) in enumerate(items):
    y = 1.28 + i * 0.88
    h2_circle(slide, 0.82, y+0.26, 0.22, brand["primary"])
    h2_text(slide, str(i+1), 0.60, y+0.06, 0.44, 0.4,
            brand["font"], 16, "FFFFFF", bold=True, align="center")
    h2_text(slide, titre, 1.22, y,        W-2.1, 0.38, brand["font"], 13, brand["primary"], bold=True)
    h2_text(slide, corps, 1.22, y+0.40,   W-2.1, 0.44, brand["font"], 11, brand["text"])
h2_text(slide, "Footer · Contexte · 2025", 0.5, H-0.46, W-1.0, 0.38,
        brand["font"], 10, brand["primary"])

[list — grille de cartes 2×2] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, brand["light"])
h2_text(slide, "Titre liste cartes", 0.5, 0.22, W-1.0, 0.72,
        brand["font"], 28, brand["primary"], bold=True)
h2_rect(slide, 0.5, 1.05, W-1.0, 0.04, brand["accent"])
cards = [("Titre 1","Contenu concis ≤ 20 mots.", brand["primary"], "FFFFFF"),
         ("Titre 2","Contenu concis ≤ 20 mots.", "FFFFFF", brand["primary"]),
         ("Titre 3","Contenu concis ≤ 20 mots.", brand["secondary"], "FFFFFF"),
         ("Titre 4","Contenu concis ≤ 20 mots.", "FFFFFF", brand["secondary"])]
cw = (W-1.4) / 2
ch = 1.55
for i, (titre, corps, bg, fg) in enumerate(cards):
    col = i % 2
    row = i // 2
    x = 0.5 + col*(cw+0.2)
    y = 1.28 + row*(ch+0.18)
    h2_card(slide, x, y, cw, ch, bg, titre, corps,
            brand["font"], fg, fg, title_size=13, body_size=11)
h2_text(slide, "Footer · Contexte · 2025", 0.5, H-0.46, W-1.0, 0.38,
        brand["font"], 10, brand["primary"])

[image_text — split vertical, texte à droite] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W/2, H, brand["primary"])
h2_text(slide, "Titre ≤ 8 mots", W/2+0.4, 0.35, W/2-0.75, 0.9,
        brand["font"], 26, brand["primary"], bold=True)
h2_rect(slide, W/2+0.4, 1.35, W/2-0.75, 0.04, brand["accent"])
points = ["Point clé 1 — concis et fort.", "Point clé 2 — idée distincte.", "Point clé 3 — conclusion."]
for i, pt in enumerate(points):
    h2_circle(slide, W/2+0.62, 1.78+i*0.9, 0.14, brand["accent"])
    h2_text(slide, "  " + pt, W/2+0.85, 1.62+i*0.9, W/2-1.2, 0.75,
            brand["font"], 12, brand["text"])
h2_text(slide, "Footer · Contexte · 2025", W/2+0.4, H-0.46, W/2-0.75, 0.38,
        brand["font"], 10, brand["primary"])

[full_text — fond blanc, layout aéré] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, "FFFFFF")
h2_rect(slide, 0, 0, 0.08, H, brand["primary"])
h2_text(slide, "Titre développement ≤ 8 mots", 0.5, 0.22, W-1.0, 0.72,
        brand["font"], 28, brand["primary"], bold=True)
h2_rect(slide, 0.5, 1.05, W-1.0, 0.04, brand["accent"])
h2_rounded_rect(slide, 0.5, 1.25, W-1.0, 0.9, brand["light"])
h2_text(slide, "Premier paragraphe : une idée principale, 2 phrases courtes. Ton direct, factuel et orienté valeur.",
        0.68, 1.38, W-1.36, 0.68, brand["font"], 12, brand["text"], line_spacing=1.3)
h2_text(slide, "Deuxième paragraphe : idée distincte et complémentaire. Pas redondant avec le premier.",
        0.5, 2.32, W-1.0, 0.68, brand["font"], 12, brand["text"], line_spacing=1.3)
h2_text(slide, "Troisième paragraphe : conclusion ou conséquence pratique. Ce que ça implique concrètement.",
        0.5, 3.12, W-1.0, 0.68, brand["font"], 12, brand["text"], line_spacing=1.3)
h2_text(slide, "Footer · Contexte · 2025", 0.5, H-0.46, W-1.0, 0.38,
        brand["font"], 10, brand["primary"])

[closing — fond sombre centré] :
slide, W, H = h2_blank_slide()
h2_rect(slide, 0, 0, W, H, brand["primary"])
h2_rect(slide, 0, H*0.76, W, H*0.24, brand["secondary"])
h2_text(slide, "Merci !", W/2-3.8, H/2-1.05, 7.6, 1.5,
        brand["font"], 58, "FFFFFF", bold=True, align="center")
h2_text(slide, "Sources · Contact · 2025", W/2-3.8, H/2+0.58, 7.6, 0.7,
        brand["font"], 16, "FFFFFF", align="center")
h2_text(slide, "Footer · Contexte · 2025", 0.7, H-0.46, W-1.4, 0.38,
        brand["font"], 11, "FFFFFF")
"""

_V2_USER_TEMPLATE = """\
PRÉSENTATION : {title}
ARC NARRATIF : {arc}
FOOTER UNIFIÉ (identique sur toutes les slides de contenu) : "{footer}"
SUJET COMPLET : {prompt}

CHARTE EXTRAITE DU TEMPLATE :
  brand["primary"]   = "#{primary}"
  brand["secondary"] = "#{secondary}"
  brand["accent"]    = "#{accent}"
  brand["light"]     = "#{light}"
  brand["text"]      = "#{text_color}"
  brand["font"]      = "{font}"
  Dimensions slide   : W = {W:.2f} po  ×  H = {H:.2f} po

PROFIL CLIENT : {profile_label}
Guide de style : {style_guide}

SLIDES À GÉNÉRER — {n} slides :
{slides_json}

══════════════════════════════════════════════
INSTRUCTIONS DE GÉNÉRATION
══════════════════════════════════════════════
1. Pour chaque slide, écris le code python-pptx COMPLET qui crée la slide.
2. COMMENCE chaque slide par : slide, W, H = h2_blank_slide()  ← SANS argument
3. Adapte le contenu réel au narrative_angle et key_message de chaque slide.
4. RESPECTE les density rules : titre ≤ 8 mots, body ≤ 40 mots, etc.
5. VARIE les compositions : alterne fonds sombres et clairs, change les layouts.
6. Structure SANDWICH : cover et closing sur fond sombre, sections sur fond primary,
   slides de contenu sur fond clair (FFFFFF ou light).
7. Utilise le footer_text sur toutes les slides SAUF cover et closing.
8. Les couleurs brand[] sont visibles dans les images jointes — respecte-les.
9. Utilise des boucles Python pour les grilles/listes répétitives.

FORMAT DE SORTIE : JSON valide uniquement, sans markdown, sans commentaires.
Clés = plan_index en string ("0", "1", ...).
{{
  "0": "slide, W, H = h2_blank_slide()\\nh2_rect(...)\\n...",
  "1": "slide, W, H = h2_blank_slide()\\n...",
  ...
}}
Chaque valeur = code Python complet pour UNE slide (string avec \\n comme séparateurs).
"""


# ─────────────────────────────────────────────────────────────
# GÉNÉRATION DE CODE — Appel Claude avec vision
# ─────────────────────────────────────────────────────────────

def generate_codes_v2(
    prompt:    str,
    plan:      dict,
    palette:   dict,
    brand:     dict,
    profile:   str,
    template_thumbnail: dict | None = None,
    palette_swatch:     dict | None = None,
) -> dict:
    """
    Phase 3 du pipeline Niveau 2.
    Claude reçoit les images du template ET de la palette pour voir la vraie charte.
    Retourne { "plan_index": "code_python_string" }.
    """
    client       = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)
    profile_data = CLIENT_PROFILES.get(profile, CLIENT_PROFILES['institutional'])

    slides_payload = [
        {
            'plan_index':      sp.get('plan_index', i),
            'slide_type':      sp.get('slide_type', 'unknown'),
            'narrative_angle': sp.get('narrative_angle', ''),
            'key_message':     sp.get('key_message', ''),
            'visual_hint':     sp.get('visual_hint', ''),
        }
        for i, sp in enumerate(plan.get('slides', []))
    ]

    user_text = _V2_USER_TEMPLATE.format(
        title       = plan.get('presentation_title', prompt[:60]),
        arc         = plan.get('narrative_arc', ''),
        footer      = plan.get('footer_text', ''),
        prompt      = prompt,
        primary     = palette.get('primary', '1A3A6B'),
        secondary   = palette.get('secondary', '2E6DA4'),
        accent      = palette.get('accent', 'F0A500'),
        light       = palette.get('light', 'EEF3FA'),
        text_color  = palette.get('text', '1A1A2E'),
        font        = palette.get('font', 'Calibri'),
        W           = brand.get('slide_width_in', 10.0),
        H           = brand.get('slide_height_in', 5.63),
        profile_label = profile_data['label'],
        style_guide   = profile_data['style_guide'],
        n           = len(slides_payload),
        slides_json = json.dumps(slides_payload, ensure_ascii=False, indent=2),
    )

    # Construction du message multimodal (images + texte)
    content_parts = []

    if template_thumbnail:
        content_parts.append({
            'type': 'text',
            'text': 'Voici la miniature du template de l\'entreprise (charte graphique réelle) :',
        })
        content_parts.append({
            'type':   'image',
            'source': {
                'type':       'base64',
                'media_type': template_thumbnail['media_type'],
                'data':       template_thumbnail['data'],
            },
        })

    if palette_swatch:
        content_parts.append({
            'type': 'text',
            'text': 'Palette de couleurs extraite du template :',
        })
        content_parts.append({
            'type':   'image',
            'source': {
                'type':       'base64',
                'media_type': palette_swatch['media_type'],
                'data':       palette_swatch['data'],
            },
        })

    content_parts.append({'type': 'text', 'text': user_text})

    code_tokens = max(4000, len(slides_payload) * 600)

    for attempt in range(3):
        if attempt > 0:
            log.info(f'[V2] Retry génération code ({attempt+1}/3)...')

        msg = client.messages.create(
            model      = CLAUDE_MODEL,
            max_tokens = code_tokens,
            system     = _V2_SYSTEM,
            messages   = [{'role': 'user', 'content': content_parts}],
        )

        try:
            code_map = _parse_json_robust(msg.content[0].text.strip(), context='codes_v2')
            log.info(f'[V2] Code généré pour {len(code_map)} slides (attempt {attempt+1}).')
            return code_map
        except (ValueError, KeyError) as e:
            log.warning(f'[V2] generate_codes_v2 attempt {attempt+1}/3 échoué : {e}')
            if attempt == 2:
                log.warning('[V2] 3 tentatives échouées → dict vide (fallback L1 sera déclenché)')
                return {}

    return {}


# ─────────────────────────────────────────────────────────────
# EXÉCUTION SÉCURISÉE
# ─────────────────────────────────────────────────────────────

def _execute_slide_code_v2(code: str, prs: Presentation, palette: dict) -> bool:
    """
    Exécute un bloc de code python-pptx dans le sandbox.
    Timeout 30s via threading. Retourne True si succès.
    """
    ok, reason = _validate_code_safety(code)
    if not ok:
        log.warning(f'[V2] Code rejeté (sécurité) : {reason}')
        return False

    result   = {'success': False, 'error': None}
    safe_ns  = _build_safe_namespace(prs, palette)

    def _run():
        try:
            exec(code, safe_ns)  # noqa: S102
            result['success'] = True
        except Exception as e:
            result['error'] = str(e)

    t = threading.Thread(target=_run, daemon=True)
    t.start()
    t.join(timeout=30)

    if t.is_alive():
        log.warning('[V2] exec() timeout 30s')
        return False
    if result['error']:
        log.warning(f"[V2] exec() error: {result['error']}")
        return False
    return True


def _inject_fallback_slide(prs: Presentation, slide_plan: dict, palette: dict):
    """Slide de fallback minimaliste si le code L2 échoue."""
    try:
        slide, W, H = _h2_blank_slide(prs)
        _h2_rect(slide, 0, 0, W, H, 'FFFFFF')
        _h2_rect(slide, 0, 0, 0.08, H, palette.get('primary', '1A3A6B'))
        key_msg = slide_plan.get('key_message', 'Contenu')
        _h2_text(slide, key_msg, 0.5, H/2-0.4, W-1.0, 0.8,
                 palette.get('font', 'Calibri'), 24,
                 palette.get('primary', '1A3A6B'), bold=True)
    except Exception as e:
        log.error(f'[V2] Fallback slide failed: {e}')


def _execute_all_codes_v2(
    code_map:    dict,
    plan_slides: list,
    prs:         Presentation,
    palette:     dict,
) -> int:
    """
    Exécute les codes dans l'ordre du plan narratif.
    Retourne le nombre de slides générées avec succès.
    """
    success = 0
    for sp in plan_slides:
        plan_idx = str(sp.get('plan_index', 0))
        code     = code_map.get(plan_idx, '')
        if not code:
            log.warning(f'[V2] Pas de code pour plan_index={plan_idx}')
            _inject_fallback_slide(prs, sp, palette)
            continue
        ok = _execute_slide_code_v2(code, prs, palette)
        if ok:
            success += 1
        else:
            _inject_fallback_slide(prs, sp, palette)
    return success


def _remove_original_slides_v2(prs: Presentation, n_original: int):
    """
    Supprime les n_original premières slides du template.
    Ne conserve que les slides générées par Level 2 (ajoutées après les originales).
    """
    xml_slides  = prs.slides._sldIdLst
    all_sld_ids = list(xml_slides)
    new_sld_ids = all_sld_ids[n_original:]

    if not new_sld_ids:
        log.warning('[V2] Aucune slide L2 générée — conservation des slides originales.')
        return

    log.info(f'[V2] Suppression des {n_original} slides originales, '
             f'conservation de {len(new_sld_ids)} slides L2.')

    # Reconstruire sldIdLst avec uniquement les slides L2
    for sld in list(xml_slides):
        xml_slides.remove(sld)
    for sld in new_sld_ids:
        xml_slides.append(sld)


# ─────────────────────────────────────────────────────────────
# PIPELINE NIVEAU 2
# ─────────────────────────────────────────────────────────────

def run_pipeline_v2(
    pptx_bytes: bytes,
    prompt:     str,
    nb_slides:  int,
    profile:    str = 'institutional',
) -> tuple:
    """
    Pipeline Level 2 — 4 phases.
    """
    if not ANTHROPIC_API_KEY:
        raise ValueError('Clé API Claude manquante.')

    nb_slides = max(2, min(nb_slides, 30))
    profile   = profile if profile in CLIENT_PROFILES else 'institutional'

    prs        = Presentation(io.BytesIO(pptx_bytes))
    n_original = len(prs.slides)

    # ── Phase 1 : Analyse brand + extraction visuels ──────────
    log.info(f'[V2] Phase 1 : analyse brand (profil={profile}, {nb_slides} slides)...')
    brand      = extract_brand(prs)
    palette    = _h2_extract_palette(brand)
    library    = build_layout_library(prs)
    selection  = select_template_slides(library, nb_slides)

    thumbnail  = _extract_template_thumbnail(pptx_bytes)
    swatch     = _make_palette_swatch(palette)
    log.info(f'[V2] Palette: primary=#{palette["primary"]} font={palette["font"]} '
             f'thumbnail={"oui" if thumbnail else "non"} swatch={"oui" if swatch else "non"}')

    # ── Phase 2 : Planification narrative ─────────────────────
    log.info('[V2] Phase 2 : planification narrative...')
    plan = plan_presentation(prompt, nb_slides, selection, brand)

    plan_slides = plan.get('slides', [])
    while len(plan_slides) < nb_slides:
        fb = selection[min(len(plan_slides), len(selection) - 1)]
        plan_slides.append({
            'plan_index':           len(plan_slides),
            'template_slide_index': fb['slide_index'],
            'slide_type':           fb.get('slide_type', 'full_text'),
            'narrative_angle':      'Développement complémentaire',
            'key_message':          'Argument additionnel',
            'visual_hint':          '',
        })
    plan['slides'] = plan_slides[:nb_slides]

    # ── Phase 3 : Génération code python-pptx (avec vision) ───
    log.info('[V2] Phase 3 : génération code python-pptx (Claude avec vision)...')
    code_map = generate_codes_v2(
        prompt, plan, palette, brand, profile,
        template_thumbnail = thumbnail,
        palette_swatch     = swatch,
    )

    # ── Phase 4 : Exécution + assemblage ──────────────────────
    log.info('[V2] Phase 4 : exécution des codes...')
    success_count = _execute_all_codes_v2(code_map, plan['slides'], prs, palette)
    log.info(f'[V2] {success_count}/{nb_slides} slides générées avec succès.')

    # Fallback L1 si taux < 50%
    success_rate = success_count / max(nb_slides, 1)
    if success_rate < 0.5:
        log.warning(f'[V2] Taux d\'échec élevé ({success_rate:.0%}) → fallback Level 1')
        final_bytes_l1, plan_l1, brand_l1 = run_pipeline(pptx_bytes, prompt, nb_slides)
        return final_bytes_l1, plan_l1, brand_l1, palette

    _remove_original_slides_v2(prs, n_original)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read(), plan, brand, palette


# ══════════════════════════════════════════════════════════════
# PIPELINE V3 — Layouts pré-testés
# ══════════════════════════════════════════════════════════════

_V3_PLANNER_SYSTEM = """\
Tu es un consultant senior spécialisé en communication visuelle et stratégie.
Tu crées des présentations PowerPoint de niveau agence — niveau Gamma.app ou McKinsey.

PRINCIPES ÉDITORIAUX (non négociables) :
1. Une seule idée forte par slide. Jamais deux messages.
2. Les titres sont des assertions, pas des thèmes. Pas "Contexte marché" mais "Le marché a doublé en 3 ans".
3. Les chiffres sont réels, précis et sourcés. Pas "beaucoup" mais "47 %".
4. Le contenu est celui d'un expert du domaine, pas un résumé Wikipedia.
5. Zéro phrase passive, zéro jargon vide. Chaque mot compte.

FORMAT : retourne UNIQUEMENT du JSON valide, sans markdown ni commentaire.\
"""

_V3_PLANNER_USER = """\
Crée une présentation sur : {prompt}
Nombre de slides : {nb_slides}

Charte graphique :
- Couleur primaire : #{primary}
- Couleur accent   : #{accent}
- Police           : {font}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
LAYOUTS DISPONIBLES + STRUCTURE JSON EXACTE
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

cover_dark / cover_split  →  {{"title":"...", "subtitle":"...", "footer":"..."}}
section                   →  {{"number":"01", "title":"..."}}
kpi_grid                  →  {{"title":"...", "kpis":[{{"value":"47 %","label":"Part de marché","sublabel":"2023, source XYZ"}}], "footer":"..."}}
kpi_row                   →  {{"title":"...", "kpis":[{{"value":"3,2 Md€","label":"CA 2023","sublabel":"vs 2,1 Md€ en 2022"}}], "footer":"..."}}
timeline_h                →  {{"title":"...", "steps":[{{"date":"2021","title":"Lancement","body":"Déploiement initial dans 3 pays"}}], "footer":"..."}}
two_col                   →  {{"title":"...", "section_label":"ANALYSE", "col_a":{{"title":"POUR","subtitle":"sous-titre optionnel","items":["Argument 1","Argument 2"]}}, "col_b":{{"title":"CONTRE","subtitle":"sous-titre optionnel","items":["Limite 1","Limite 2"]}}, "footer":"..."}}
col3                      →  {{"title":"...", "section_label":"STRATÉGIE", "subtitle":"Trois axes prioritaires", "columns":[{{"icon":"⚡","label":"CAPEX","title":"INVESTISSEMENTS","subtitle":"Déficit potentiel","items":["Point 1","Point 2"],"stat_value":"-28%","stat_label":"BAISSE MONDIALE"}}], "footer":"..."}}
conclusion                →  {{"title":"...", "section_label":"SYNTHÈSE", "subtitle":"...", "cards":[{{"icon":"🌐","title":"AXE 1","body":"..."}}], "sidebar_title":"Vision", "sidebar_quote":"Citation...", "sidebar_cta":"QUESTIONS & ÉCHANGES", "footer":"..."}}
entity                    →  {{"title":"...", "section_label":"COMPARAISON", "subtitle":"Analyse comparative", "entities":[{{"icon":"🇺🇸","name":"États-Unis","badge":"LEADER","items":["Point 1","Point 2"],"stat_value":"34%","stat_label":"PART MONDIALE"}}], "footer":"..."}}
infographic               →  {{"title":"...", "section_label":"CHIFFRES CLÉS", "subtitle":"...", "value":"47 %", "label":"Taux de croissance", "context":"Contexte additionnel en 20 mots.", "bars":[{{"label":"Segment A","percent":47}},{{"label":"Segment B","percent":31}},{{"label":"Segment C","percent":22}}], "footer":"..."}}
quote                     →  {{"section_label":"LE GRAND ENTRETIEN", "category":"LE GRAND ENTRETIEN", "quote":"Citation percutante ≤ 20 mots", "author":"Prénom NOM, Titre — 2024", "source":"Source — Date", "footer":"..."}}
list_numbered             →  {{"title":"...", "section_label":"ANALYSE", "subtitle":"...", "items":[{{"title":"Levier 1","body":"Explication concise en 15 mots max."}}], "footer":"..."}}
list_cards                →  {{"title":"...", "section_label":"ANALYSE", "subtitle":"...", "cards":[{{"icon":"📊","label":"KPI","title":"Axe 1","subtitle":"Sous-axe","body":"Description en 20 mots max.","stat_value":"47%","stat_label":"PART DE MARCHÉ"}}], "footer":"..."}}
image_split               →  {{"title":"...", "points":["Point 1 en 15 mots","Point 2"], "footer":"..."}}
full_text                 →  {{"title":"...", "paragraphs":["Paragraphe 1 (≤ 40 mots)","Paragraphe 2"], "footer":"..."}}
stat_hero                 →  {{"value":"€ 2,3 Md", "label":"Montant des pertes évitées", "context":"Estimation sur 5 ans — rapport Ernst & Young 2023", "footer":"..."}}
closing_dark / closing_split → {{"title":"Merci", "subtitle":"Sources : ... · Contact : ..."}}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
STRUCTURE DE RÉPONSE ATTENDUE
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
{{
  "presentation_title": "Titre accrocheur ≤ 8 mots",
  "footer_text": "Nom entreprise · Sujet · Année",
  "slides": [
    {{"layout": "cover_dark", "content": {{"title": "...", "subtitle": "...", "footer": "..."}}}},
    {{"layout": "kpi_grid",   "content": {{"title": "...", "kpis": [...], "footer": "..."}}}},
    ...
    {{"layout": "closing_dark", "content": {{"title": "Merci", "subtitle": "..."}}}}
  ]
}}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
RÈGLES IMPÉRATIVES
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
- Slide 1 : TOUJOURS cover_dark ou cover_split
- Dernière slide : TOUJOURS closing_dark ou closing_split
- Jamais deux layouts identiques consécutifs
- Alterner fonds sombres (cover/section/quote/kpi_grid) et clairs (list/two_col/full_text/stat_hero)
- kpi_grid : 4 à 6 KPIs avec valeurs chiffrées réelles
- kpi_row : 3 à 4 KPIs maximum
- timeline_h : 4 à 5 jalons avec dates réelles
- list_* : 3 à 5 items maximum
- Titres : ≤ 8 mots, assertifs, jamais nominaux
- Bodies : ≤ 25 mots, concis, factuels
- footer_text : identique sur toutes les slides sauf cover et closing
- Contenu expert : chiffres précis, sources citées, angle analytique
\
"""


def plan_presentation_v3(prompt: str, nb_slides: int, palette: dict) -> dict:
    """
    Planifie une présentation via layouts pré-testés.
    Claude retourne uniquement du JSON : [{layout, content}, ...].
    """
    client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

    layouts_block = '\n'.join(
        f'- {name}: {desc}'
        for name, desc in LAYOUT_DESCRIPTIONS.items()
    )

    user = _V3_PLANNER_USER.format(
        prompt        = prompt,
        nb_slides     = nb_slides,
        primary       = palette.get('primary', '1A3A6B'),
        accent        = palette.get('accent',  'F0A500'),
        font          = palette.get('font',    'Calibri'),
        layouts_block = layouts_block,
    )

    max_tokens = max(4000, nb_slides * 380)

    for attempt in range(3):
        msg = client.messages.create(
            model      = CLAUDE_MODEL,
            max_tokens = max_tokens,
            system     = _V3_PLANNER_SYSTEM,
            messages   = [{'role': 'user', 'content': user}],
        )
        try:
            plan = _parse_json_robust(msg.content[0].text.strip(), context='plan_v3')
            slides = plan.get('slides', [])
            log.info(f'[V3] Plan: {len(slides)} slides, title="{plan.get("presentation_title","")[:60]}"')
            return plan
        except (ValueError, KeyError) as e:
            log.warning(f'plan_presentation_v3 attempt {attempt+1}/3 failed: {e}')
            if attempt == 2:
                raise
    raise RuntimeError('plan_presentation_v3 : 3 tentatives échouées')


def run_pipeline_v3(pptx_bytes: bytes, prompt: str, nb_slides: int) -> tuple:
    """
    Pipeline V3 — layouts pré-testés, zéro génération de code.

    Phase 1 : extraction charte + palette
    Phase 2 : planification narrative par Claude (JSON layouts)
    Phase 3 : application des fonctions de layout pré-testées
    Phase 4 : suppression des slides originales du template
    """
    if not ANTHROPIC_API_KEY:
        raise ValueError('Clé API Claude manquante.')

    nb_slides = max(2, min(nb_slides, 30))

    prs        = Presentation(io.BytesIO(pptx_bytes))
    n_original = len(prs.slides)

    # ── Phase 1 ───────────────────────────────────────────────
    log.info(f'[V3] Phase 1 : extraction charte ({nb_slides} slides)...')
    brand   = extract_brand(prs)
    palette = _h2_extract_palette(brand)
    log.info(f'[V3] Palette : primary=#{palette["primary"]} accent=#{palette["accent"]} font={palette["font"]}')

    # ── Phase 2 ───────────────────────────────────────────────
    log.info('[V3] Phase 2 : planification narrative...')
    plan = plan_presentation_v3(prompt, nb_slides, palette)
    log.info(f'[V3] Plan reçu : {json.dumps(plan, ensure_ascii=False)[:1000]}')

    slides_plan = plan.get('slides', [])

    # Compléter si Claude en a généré moins que demandé
    fallback_layouts = ['list_cards', 'col3', 'two_col', 'kpi_grid']
    while len(slides_plan) < nb_slides:
        fb_name = fallback_layouts[len(slides_plan) % len(fallback_layouts)]
        slides_plan.append({
            'layout':  fb_name,
            'content': {
                'title':      'Développement complémentaire',
                'paragraphs': ['Contenu additionnel à personnaliser.'],
                'footer':     plan.get('footer_text', ''),
            },
        })
    slides_plan = slides_plan[:nb_slides]

    # Garantie : la dernière slide est TOUJOURS une closing
    closing_layouts = {'closing_dark', 'closing_split'}
    if not slides_plan or slides_plan[-1].get('layout') not in closing_layouts:
        closing_slide = {
            'layout': 'closing_dark',
            'content': {
                'title':    'Merci',
                'subtitle': plan.get('footer_text', ''),
            },
        }
        if len(slides_plan) >= nb_slides and slides_plan:
            slides_plan[-1] = closing_slide   # remplace la dernière si plan complet
        else:
            slides_plan.append(closing_slide)  # ajoute sinon

    log.info(f'[V3] {len(slides_plan)} slides à générer : {[s.get("layout") for s in slides_plan]}')

    # ── Phase 3 ───────────────────────────────────────────────
    log.info('[V3] Phase 3 : application des layouts...')
    success = 0
    slides_before = len(prs.slides)
    for sp in slides_plan:
        layout_name = sp.get('layout', 'full_text')
        content     = sp.get('content', {})

        # Injecter footer_text global si absent du content
        if not content.get('footer') and plan.get('footer_text'):
            content['footer'] = plan['footer_text']

        layout_fn = LAYOUT_REGISTRY.get(layout_name) or LAYOUT_REGISTRY['full_text']
        log.info(f'[V3] → layout "{layout_name}" (total slides avant: {len(prs.slides)})')
        try:
            layout_fn(prs, content, palette)
            success += 1
            log.info(f'[V3] ✓ "{layout_name}" OK — total slides: {len(prs.slides)}')
        except Exception as e:
            log.error(f'[V3] ✗ "{layout_name}" ÉCHOUÉ : {repr(e)}', exc_info=True)
            try:
                LAYOUT_REGISTRY['full_text'](prs, {
                    'title':      content.get('title', ''),
                    'paragraphs': [],
                    'footer':     content.get('footer', ''),
                }, palette)
                success += 1
                log.info(f'[V3] ✓ fallback full_text OK — total slides: {len(prs.slides)}')
            except Exception as e2:
                log.error(f'[V3] ✗ fallback full_text AUSSI ÉCHOUÉ : {repr(e2)}', exc_info=True)

    log.info(f'[V3] Phase 3 terminée : {success}/{nb_slides} OK — '
             f'{len(prs.slides) - slides_before} slides ajoutées au total')

    # ── Phase 4 ───────────────────────────────────────────────
    slides_added = len(prs.slides) - n_original
    log.info(f'[V3] Phase 4 : {len(prs.slides)} slides totales '
             f'({n_original} originales + {slides_added} nouvelles)')

    if slides_added == 0:
        raise RuntimeError(
            f'[V3] Aucune slide créée (success={success}) — fallback L1 déclenché'
        )

    xml_slides = prs.slides._sldIdLst
    to_remove  = list(prs.slides._sldIdLst)[:n_original]
    for sld_id in to_remove:
        xml_slides.remove(sld_id)
    log.info(f'[V3] {n_original} slides originales supprimées — '
             f'{len(prs.slides)} slides finales')

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read(), plan, brand, palette


# ══════════════════════════════════════════════════════════════
# PIPELINE V4 — TEMPLATE-NATIVE GENERATION
# ══════════════════════════════════════════════════════════════

def _fill_placeholder_preserving_style(ph, new_text: str) -> None:
    """
    Remplace le texte d'un placeholder en préservant le style XML du premier run
    (police, taille, couleur schemeClr, gras, italique, effets).
    Gère le texte multi-paragraphes (séparés par \\n).
    """
    if not new_text:
        return

    import lxml.etree as _etree
    tf    = ph.text_frame
    txBody = tf._txBody

    # Capturer le style du premier run existant
    rPr_xml = None
    pPr_xml = None
    if tf.paragraphs:
        first_p = tf.paragraphs[0]
        p_elem  = first_p._p
        pPr     = p_elem.find(qn('a:pPr'))
        if pPr is not None:
            pPr_xml = copy.deepcopy(pPr)
        if first_p.runs:
            rPr = first_p.runs[0]._r.find(qn('a:rPr'))
            if rPr is not None:
                rPr_xml = copy.deepcopy(rPr)

    # Vider le txBody de tous ses paragraphes
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    # Recréer les paragraphes avec le style capturé
    paragraphs = new_text.split('\n')
    for para_text in paragraphs:
        p_elem = _etree.SubElement(txBody, qn('a:p'))
        if pPr_xml is not None:
            p_elem.insert(0, copy.deepcopy(pPr_xml))
        r_elem = _etree.SubElement(p_elem, qn('a:r'))
        if rPr_xml is not None:
            r_elem.insert(0, copy.deepcopy(rPr_xml))
        t_elem = _etree.SubElement(r_elem, qn('a:t'))
        t_elem.text = para_text


# ── V4 Foundation Helpers ───────────────────────────────────────────────────


def analyze_template_v4(prs: Presentation) -> dict:
    """
    Analyse complète du template en un seul appel.
    Retourne un template profile `tp` utilisé par toutes les fonctions V4.
    Couvre : couleurs thème, layout_map, logo_zone, font, accent_cycle, W, H.
    """
    import re as _re2, zipfile as _zf2, io as _io2

    # ── Couleurs thème depuis ZIP ────────────────────────────────────────────
    theme: dict = {}
    try:
        buf = _io2.BytesIO()
        prs.save(buf)
        buf.seek(0)
        with _zf2.ZipFile(buf) as zf:
            theme_files = sorted([n for n in zf.namelist()
                                   if _re2.search(r'ppt/theme/theme\d*\.xml$', n, _re2.I)])
            if theme_files:
                xml = zf.read(theme_files[0]).decode('utf-8', errors='ignore')
                for slot in ['dk1', 'lt1', 'dk2', 'lt2',
                             'accent1', 'accent2', 'accent3',
                             'accent4', 'accent5', 'accent6']:
                    m = _re2.search(
                        rf'<a:{slot}[^>]*>\s*<a:srgbClr val="([0-9A-Fa-f]{{6}})"', xml)
                    if m:
                        theme[slot] = m.group(1).upper()
                        continue
                    m = _re2.search(
                        rf'<a:{slot}[^>]*>\s*<a:sysClr[^>]*lastClr="([0-9A-Fa-f]{{6}})"', xml)
                    if m:
                        theme[slot] = m.group(1).upper()
    except Exception as e:
        log.warning(f'[V4] analyze_template_v4 theme: {e}')

    # ── Layout map ──────────────────────────────────────────────────────────
    layout_map: dict = {}
    # Track which layouts have picture placeholders (ph type=pic or shape_type=13)
    def _has_picture(layout):
        try:
            for sh in layout.shapes:
                if sh.shape_type == 13:
                    return True
            for ph in layout.placeholders:
                if ph.placeholder_format.type and str(ph.placeholder_format.type) in ('PP_MEDIA', 'PP_OBJECT', 'PICTURE', '18'):
                    return True
        except Exception:
            pass
        return False

    for idx, layout in enumerate(prs.slide_layouts):
        n = layout.name.lower()
        ph_idxs = {ph.placeholder_format.idx for ph in layout.placeholders}
        if any(k in n for k in ['couverture', 'cover', 'titre de pré', 'garde', 'home', 'accueil', 'front']):
            # Prefer cover layout that has an actual picture (logo/bg image)
            if 'cover' not in layout_map or _has_picture(layout):
                layout_map['cover'] = idx
        elif any(k in n for k in ['ouverture', 'section', 'chapter', 'séparateur', 'separator', 'intertitre', 'transition']):
            layout_map.setdefault('section', idx)
        elif any(k in n for k in ['chiffres', 'kpi', 'metrics', 'stats']):
            layout_map.setdefault('kpi', idx)
        elif any(k in n for k in ['merci', 'closing', 'thank', 'fin', 'end']):
            # Prefer closing layout that has picture (matches cover design)
            if 'closing' not in layout_map or _has_picture(layout):
                layout_map['closing'] = idx
        elif any(k in n for k in ['vide', 'blank', 'vierge', 'empty']):
            layout_map.setdefault('blank', idx)
        elif any(k in n for k in ['two content', '2 contenus', 'comparison', 'deux', 'two col', '2 col', 'dual', 'compare']):
            layout_map.setdefault('two_col', idx)
        elif any(k in n for k in ['texte', 'text', 'contenu', 'content']) or (0 in ph_idxs and 1 in ph_idxs):
            layout_map.setdefault('text', idx)

    n_layouts = len(prs.slide_layouts)
    layout_map.setdefault('cover',   0)
    layout_map.setdefault('blank',   n_layouts - 1)
    layout_map.setdefault('text',    min(1, n_layouts - 1))
    layout_map.setdefault('section', layout_map['cover'])
    layout_map.setdefault('kpi',     layout_map['text'])
    layout_map.setdefault('closing', layout_map['cover'])
    layout_map.setdefault('two_col', layout_map['text'])
    layout_map.setdefault('content', layout_map['text'])  # backward compat

    # ── Logo zone (premier shape picture du master) ──────────────────────────
    logo_zone = None
    try:
        for shape in prs.slide_masters[0].shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                logo_zone = {
                    'left':   shape.left   / 914400.0,
                    'top':    shape.top    / 914400.0,
                    'right':  (shape.left + shape.width)  / 914400.0,
                    'bottom': (shape.top  + shape.height) / 914400.0,
                }
                break
    except Exception:
        pass

    # ── Font (première police non-'+' du master) ─────────────────────────────
    font = 'Calibri'
    try:
        import lxml.etree as _etree2
        master_xml = _etree2.tostring(prs.slide_masters[0]._element).decode('utf-8', errors='ignore')
        for m in _re2.findall(r'typeface="([^"]+)"', master_xml):
            if m and not m.startswith('+') and len(m) < 64:
                font = m
                break
    except Exception:
        pass

    W = prs.slide_width  / 914400.0
    H = prs.slide_height / 914400.0

    accent1 = theme.get('accent1', '009CEA')
    accent2 = theme.get('accent2', 'ED0000')
    accent3 = theme.get('accent3', accent1)   # fallback → accent1, not invented green
    accent4 = theme.get('accent4', accent2)   # fallback → accent2, not invented orange

    # Build accent_cycle from colors ACTUALLY in the template (no invented fallbacks)
    # Priority: accent1, accent2, then accent3/accent4 only if template defines them
    _cycle_raw = []
    for k in ('accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6'):
        if k in theme:
            _cycle_raw.append(theme[k])
    if not _cycle_raw:
        _cycle_raw = [accent1, accent2]
    # Ensure at least 4 entries (repeat from start so cycling works)
    while len(_cycle_raw) < 4:
        _cycle_raw.append(_cycle_raw[len(_cycle_raw) % max(len(_cycle_raw), 1)])
    accent_cycle = _cycle_raw

    # Light neutral for card backgrounds — derived from template light color if possible
    lt1 = theme.get('lt1', 'FFFFFF')
    # If lt1 is near-white (luminance > 95%), use a subtle off-white; else use lt1 itself
    try:
        r, g, b = int(lt1[0:2], 16), int(lt1[2:4], 16), int(lt1[4:6], 16)
        lum = (0.299 * r + 0.587 * g + 0.114 * b) / 255
        card_bg_light = 'F8F8F8' if lum > 0.94 else lt1
        card_bg_mid   = 'F0F0F0' if lum > 0.94 else lt1
    except Exception:
        card_bg_light, card_bg_mid = 'F8F8F8', 'F0F0F0'

    tp = {
        'theme':          theme,
        'layout_map':     layout_map,
        'logo_zone':      logo_zone,
        'font':           font,
        'accent_cycle':   accent_cycle,
        'card_bg_light':  card_bg_light,   # nearly white card background
        'card_bg_mid':    card_bg_mid,     # slightly darker alternating background
        'W':              W,
        'H':              H,
    }
    log.info(f'[V4] tp: theme={list(theme.items())[:4]}, lmap={layout_map}, font={font}, cycle={accent_cycle[:3]}, card_bg={card_bg_light}')
    return tp


def _fill_preserving_style(ph, new_text: str) -> None:
    """
    Alias de _fill_placeholder_preserving_style.
    Remplace le texte d'un placeholder en préservant police/taille/couleur du template.
    """
    _fill_placeholder_preserving_style(ph, new_text)


def _add_template_header_and_footer(slide, title: str, footer_text: str, tp: dict,
                                     content: dict | None = None) -> None:
    """
    Header/footer pattern V4 :
      - section_label : | LABEL  (9 pt, dk1, barre accent1) — si content['section_label']
      - titre         : accent1, 28 pt bold
      - subtitle      : 888888, 11 pt — si content['subtitle']
      - ligne         : accent1, épaisseur 0.04"
      - footer        : ligne #DDDDDD + texte #AAAAAA 9 pt
    Ne dépasse pas x=10.8" pour préserver la zone logo.
    """
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    accent1 = theme.get('accent1', '009CEA')
    dk1     = theme.get('dk1', '374649')
    H       = tp.get('H', 7.5)

    c             = content or {}
    section_label = c.get('section_label', '')
    subtitle      = c.get('subtitle', '')

    # Section label : barre verticale accent + texte dk1
    if section_label:
        _h2_rect(slide, left=0.60, top=0.16, width=0.055, height=0.24, color=accent1)
        _h2_text(slide, section_label.upper(),
                 left=0.72, top=0.16, width=7.0, height=0.24,
                 font=font, size_pt=9, color=dk1, bold=True, align='left')
        title_y = 0.42
    else:
        title_y = 0.20

    # Budget total header zone: 1.52" (laisse 0.03" avant CT=1.55)
    # Sans section_label : title_y=0.20, budget titre+sep = 1.32"
    # Avec section_label  : title_y=0.42, budget titre+sep = 1.10"
    # Calcul adaptatif de la hauteur titre
    if subtitle:
        title_h = 0.64   # on garde de la place pour le sous-titre
        sub_h   = 0.26
    else:
        title_h = 0.80   # jusqu'à 2 lignes 28 pt
        sub_h   = 0.0

    _h2_text(slide, title,
             left=0.6, top=title_y, width=10.2, height=title_h,
             font=font, size_pt=28, color=accent1,
             bold=True, align='left')

    sep_y = title_y + title_h + 0.05

    # Sous-titre facultatif
    if subtitle:
        _h2_text(slide, subtitle,
                 left=0.6, top=sep_y, width=10.2, height=sub_h,
                 font=font, size_pt=11, color='888888',
                 bold=False, align='left')
        sep_y += sub_h + 0.04

    # Separator ne doit pas dépasser 1.50" (CT=1.55 → 0.05" de marge)
    sep_y = min(sep_y, 1.50)

    # Ligne de séparation
    _h2_rect(slide, left=0.6, top=sep_y, width=10.4, height=0.04, color=accent1)

    # Footer
    _h2_rect(slide, left=0.0, top=H - 0.4, width=13.33, height=0.003, color='DDDDDD')
    if footer_text:
        _h2_text(slide, footer_text,
                 left=0.6, top=H - 0.38, width=10.4, height=0.32,
                 font=font, size_pt=9, color='AAAAAA',
                 bold=False, align='left')


# ── V4 Native Layout Functions ──────────────────────────────────────────────


def _add_slide_native(prs: Presentation, layout_idx: int):
    """
    Crée une slide sur le layout template à l'index donné.
    Garantit showMasterSp='1' (logo visible).
    Retourne (slide, ph_map) où ph_map = {idx: placeholder}.
    """
    layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
    slide  = prs.slides.add_slide(layout)
    cSld   = slide._element.find(qn('p:cSld'))
    if cSld is not None:
        cSld.set('showMasterSp', '1')
    ph_map = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    return slide, ph_map


def _ph_fill(ph_map: dict, idx: int, text: str) -> bool:
    """Remplit ph[idx] si présent. Retourne True si rempli."""
    if idx in ph_map and text:
        _fill_preserving_style(ph_map[idx], str(text))
        return True
    return False


def layout_cover_v4(prs: Presentation, content: dict, tp: dict):
    """
    Slide de couverture sur le layout natif 'cover'.
    ph[0] = titre principal
    ph[1] = sous-titre / tagline
    Footer : textbox custom en bas (ph[14] évité car il peut chevaucher le subtitle).
    """
    idx   = tp['layout_map']['cover']
    slide, ph_map = _add_slide_native(prs, idx)

    _ph_fill(ph_map, 0, content.get('title', ''))
    subtitle = content.get('subtitle', '')
    if not _ph_fill(ph_map, 1, subtitle):
        _ph_fill(ph_map, 2, subtitle)

    # Footer en bas via textbox custom (ph[14] peut être mal positionné sur le layout cover)
    footer = content.get('footer', '')
    if footer:
        font = tp.get('font', 'Calibri')
        H    = tp.get('H', 7.5)
        _h2_rect(slide, left=0.0, top=H - 0.4, width=13.33, height=0.003, color='DDDDDD')
        _h2_text(slide, footer,
                 left=0.6, top=H - 0.38, width=10.4, height=0.32,
                 font=font, size_pt=9, color='AAAAAA',
                 bold=False, align='left')

    return slide


def layout_section_v4(prs: Presentation, content: dict, tp: dict):
    """
    Slide séparateur de section sur le layout natif 'section'.
    ph[13] = numéro de section ("01", "02"…)
    ph[0]  = titre de section
    ph[1]  = sous-titre / description (si présent)
    """
    idx   = tp['layout_map']['section']
    slide, ph_map = _add_slide_native(prs, idx)

    number = str(content.get('number', ''))
    if not number:
        # Auto-numérotation par position dans les slides existantes
        number = f'{len(prs.slides):02d}'

    _ph_fill(ph_map, 13, number)
    _ph_fill(ph_map, 0,  content.get('title', ''))
    _ph_fill(ph_map, 1,  content.get('subtitle', ''))

    # Fallbacks si indices non standards
    if 13 not in ph_map:
        _ph_fill(ph_map, 12, number) or _ph_fill(ph_map, 11, number)

    return slide


def layout_fulltext_v4(prs: Presentation, content: dict, tp: dict):
    """
    Slide de texte riche sur le layout natif 'text'.
    ph[0]  = titre
    ph[13] = corps (multi-paragraphes joints avec \\n)
    Fallback : ph[1] si ph[13] absent.
    """
    idx   = tp['layout_map']['text']
    slide, ph_map = _add_slide_native(prs, idx)

    # Construire le corps depuis toutes les sources possibles
    # Limiter à 2 paragraphes de ≤ 30 mots chacun — jamais de mur de texte
    body = content.get('body', '')
    if not body:
        paras = content.get('paragraphs', [])
        if paras:
            body = '\n'.join(_trunc(str(p), 30) for p in paras[:2])
    if not body:
        body = content.get('subtitle', '')
    if body:
        # Tronquer chaque paragraphe existant
        lines = body.split('\n')
        body = '\n'.join(_trunc(ln, 30) for ln in lines[:2])

    _ph_fill(ph_map, 0, content.get('title', ''))

    # Essayer ph[13] en priorité (corps principal Cortex_1), sinon ph[1]
    if not _ph_fill(ph_map, 13, body):
        _ph_fill(ph_map, 1, body)

    _ph_fill(ph_map, 14, content.get('footer', ''))

    return slide


def layout_closing_v4(prs: Presentation, content: dict, tp: dict):
    """
    Slide de clôture sur le layout natif 'closing'.
    ph[0] = titre ("Merci", "Questions ?")
    Textbox additionnel pour subtitle/sources si fourni.
    """
    idx   = tp['layout_map']['closing']
    slide, ph_map = _add_slide_native(prs, idx)

    _ph_fill(ph_map, 0, content.get('title', 'Merci'))

    # Sous-titre : toujours en textbox custom (ph[1]/ph[2] chevauche ph[0] sur ce layout)
    subtitle = content.get('subtitle', '') or content.get('body', '')
    font    = tp.get('font', 'Calibri')
    H       = tp.get('H', 7.5)
    theme   = tp.get('theme', {})
    accent1 = theme.get('accent1', '009CEA')
    if subtitle:
        _h2_text(slide, subtitle,
                 left=1.5, top=4.2, width=10.0, height=1.5,
                 font=font, size_pt=18, color=accent1,
                 bold=False, align='center')

    # Footer custom en bas (ph[14] peut être mal positionné sur ce layout)
    footer = content.get('footer', '')
    if footer:
        _h2_rect(slide, left=0.0, top=H - 0.4, width=13.33, height=0.003, color='DDDDDD')
        _h2_text(slide, footer, left=0.6, top=H - 0.38, width=10.4, height=0.32,
                 font=font, size_pt=9, color='AAAAAA', bold=False, align='left')

    return slide


# ── V4 Hybrid Layout Functions (Blank + shapes) ─────────────────────────────

def _trunc(text: str, max_words: int = 12) -> str:
    """Truncate to max_words words, adding ellipsis if cut. Prevents text overflow."""
    if not text:
        return text
    words = str(text).split()
    if len(words) <= max_words:
        return text
    return ' '.join(words[:max_words]) + '…'


def _cbg(tp: dict, idx: int = 0) -> str:
    """Return a card/row background color from the template palette (neutral, no blue tint)."""
    if idx % 2 == 0:
        return tp.get('card_bg_light', 'F8F8F8')
    return tp.get('card_bg_mid', 'F0F0F0')


def _darken(hex_color: str, factor: float = 0.75) -> str:
    """Darken a hex color by factor (0.0=black, 1.0=unchanged). Used for nested dark boxes."""
    try:
        r = max(0, int(int(hex_color[0:2], 16) * factor))
        g = max(0, int(int(hex_color[2:4], 16) * factor))
        b = max(0, int(int(hex_color[4:6], 16) * factor))
        return f'{r:02X}{g:02X}{b:02X}'
    except Exception:
        return hex_color


def _layout_has_own_bg(layout) -> bool:
    """True if the layout defines its own background (overrides master gradient)."""
    try:
        cSld = layout._element.find(qn('p:cSld'))
        if cSld is None:
            return False
        bg = cSld.find(qn('p:bg'))
        return bg is not None
    except Exception:
        return False


def _blank_v4(prs: Presentation, tp: dict):
    """
    Crée une slide vide :
    - Préfère un layout qui hérite du fond maître (gradient/image préservés)
    - Supprime tous les placeholders résiduels
    - Garantit showMasterSp='1' (logo et arrière-plan maître visibles)
    Retourne slide.
    """
    # Prefer a layout that inherits the master background (no own bg override)
    # Fallback order: text → blank → last layout
    lmap     = tp['layout_map']
    n_layouts = len(prs.slide_layouts)
    preferred = [lmap.get('text'), lmap.get('blank'), n_layouts - 1]
    chosen_idx = lmap.get('blank', n_layouts - 1)  # default
    for idx in preferred:
        if idx is None:
            continue
        idx = min(idx, n_layouts - 1)
        if not _layout_has_own_bg(prs.slide_layouts[idx]):
            chosen_idx = idx
            break

    layout = prs.slide_layouts[min(chosen_idx, n_layouts - 1)]
    slide  = prs.slides.add_slide(layout)

    # Remove residual placeholders
    sp_tree = slide.shapes._spTree
    for ph in list(slide.placeholders):
        try:
            sp_tree.remove(ph._element)
        except Exception:
            pass

    # Ensure master shapes (logo, gradient bg) are visible
    cSld = slide._element.find(qn('p:cSld'))
    if cSld is not None:
        cSld.set('showMasterSp', '1')
        # Remove any slide-level background override so master gradient shows through
        bg = cSld.find(qn('p:bg'))
        if bg is not None:
            try:
                cSld.remove(bg)
            except Exception:
                pass

    return slide


def _v4_variant(content: dict, n: int = 3, seed: int = 0) -> int:
    """Sélecteur de variante visuelle — 3 sources de diversité :
    1. Champ 'style' fourni par le planner (prioritaire, 0..n-1)
    2. Seed de présentation + hash du titre (varie entre présentations ET entre slides)
    Résultat : entier 0..n-1. Déterministe par (seed, title, layout, n)."""
    explicit = content.get('style')
    if explicit is not None:
        try:
            return int(explicit) % n
        except (ValueError, TypeError):
            pass
    key = str(content.get('title', '')) + '|' + str(content.get('layout', ''))
    return (seed + abs(hash(key))) % n


class _LY:
    """Constantes de mise en page V4 — garantissent l'harmonie entre tous les layouts."""
    # Zone de contenu standard (canvas 13.33" × 7.5")
    CT    = 1.55   # content top
    CB    = 6.95   # content bottom  (H - 0.55)
    CL    = 0.60   # content left
    CR    = 12.73  # content right
    CW    = 12.13  # content width   (CR - CL)
    # Colonnes 2-col
    COL_W  = 5.80  # largeur colonne
    COL_GAP = 0.50 # gap inter-colonnes
    # Typographie
    T_TITLE  = 14  # titre de carte/item
    T_BODY   = 11  # corps standard
    T_SMALL  = 10  # texte secondaire
    T_LABEL  = 12  # label
    T_HEADER = 13  # header section/colonne
    T_KPI    = 38  # valeur KPI
    T_HERO   = 64  # valeur hero
    # Géométrie
    PAD      = 0.18  # padding interne
    R_SM     = 0.04  # border-radius petit
    R_CIRC   = 0.25  # rayon cercle numéroté
    HEAD_H   = 0.50  # hauteur header
    ITEM_H   = 0.48  # hauteur ligne d'item
    BORDER_W = 0.06  # épaisseur bordure colorée
    BAR_H    = 0.044 # hauteur barre accent
    # Espacements
    GAP_XS  = 0.06
    GAP_SM  = 0.10
    GAP_MD  = 0.15
    GAP_LG  = 0.30


def layout_quote_v4(prs: Presentation, content: dict, tp: dict):
    """
    Citation — 3 variantes (toutes sur fond template).
    v0 : centré épuré — cercles décoratifs + badge pill + grande citation bold + auteur.
    v1 : panneau accent1 gauche + guillemet géant + citation à droite (auteur + séparateur).
    v2 : encadré accent1 arrondi centré + citation blanche + auteur DDDDDD.
    content: {title, quote, author, source?, section_label?}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accent2 = theme.get('accent2', 'ED0000')
    W       = tp.get('W', 13.33)
    H       = tp.get('H', 7.5)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    quote  = content.get('quote', '')
    author = content.get('author', '')
    source = content.get('source', '')
    v      = _v4_variant(content, 3, tp.get('seed', 0))

    n_lines = max(1, len(quote) // 55 + 1)
    quote_h = min(2.8, 0.50 * n_lines)

    if v == 0:
        # Cercles décoratifs en fond — couleurs très légères dérivées du template
        circ1 = _cbg(tp, 0)   # very light neutral (card_bg_light)
        circ2 = _cbg(tp, 1)   # slightly darker neutral (card_bg_mid)
        _h2_circle(slide, cx=1.8, cy=H - 1.2, r=1.9, color=circ1)
        _h2_circle(slide, cx=W - 1.5, cy=2.0, r=1.4, color=circ2)

        # Badge pill catégorie (title utilisé comme label de rubrique)
        category = content.get('category', content.get('title', ''))
        if category:
            bw = min(3.6, len(category) * 0.13 + 0.8)
            bx = W / 2 - bw / 2
            _h2_rounded_rect(slide, left=bx, top=_LY.CT + 0.15,
                              width=bw, height=0.34, color='EEEEEE', radius=0.17)
            _h2_text(slide, category.upper(),
                     left=bx, top=_LY.CT + 0.19,
                     width=bw, height=0.26,
                     font=font, size_pt=9, color=dk1, bold=True, align='center')

        # Grande citation centrée
        q_y = _LY.CT + 0.70
        _h2_text(slide, f'\u00ab\u00a0{quote}\u00a0\u00bb',
                 left=_LY.CL + 0.8, top=q_y,
                 width=_LY.CW - 1.6, height=quote_h + 0.3,
                 font=font, size_pt=24, color=dk1,
                 bold=True, align='center', line_spacing=1.35)

        # Séparateur court accent2
        sep_x = W / 2 - 0.35
        sep_y = q_y + quote_h + 0.35
        _h2_rect(slide, left=sep_x, top=sep_y, width=0.70, height=0.055, color=accent2)

        # Auteur
        if author:
            parts = author.split(',', 1)
            _h2_text(slide, parts[0].strip(),
                     left=_LY.CL, top=sep_y + 0.14,
                     width=_LY.CW, height=0.42,
                     font=font, size_pt=15, color=dk1, bold=True, align='center')
            if len(parts) > 1:
                _h2_text(slide, parts[1].strip(),
                         left=_LY.CL, top=sep_y + 0.57,
                         width=_LY.CW, height=0.30,
                         font=font, size_pt=10, color='777777', bold=False, align='center')

        # Source
        if source:
            _h2_text(slide, source,
                     left=_LY.CL, top=H - 0.75,
                     width=_LY.CW, height=0.28,
                     font=font, size_pt=9, color='AAAAAA', bold=False, align='center')

    elif v == 1:
        # Panneau accent gauche + citation droite
        panel_w = _LY.CW * 0.27
        panel_h = _LY.CB - _LY.CT
        _h2_rect(slide, left=_LY.CL, top=_LY.CT, width=panel_w, height=panel_h, color=accent1)
        _h2_text(slide, '\u201c',
                 left=_LY.CL, top=_LY.CT + 0.05,
                 width=panel_w, height=1.8,
                 font=font, size_pt=90, color='FFFFFF', bold=True, align='center')
        x_q = _LY.CL + panel_w + 0.45
        q_w = _LY.CR - x_q
        q_y = max(_LY.CT + 0.20, (_LY.CT + _LY.CB) / 2 - quote_h / 2 - 0.25)
        _h2_text(slide, quote,
                 left=x_q, top=q_y, width=q_w, height=quote_h + 0.4,
                 font=font, size_pt=20, color=dk1,
                 bold=False, italic=True, align='left', line_spacing=1.3)
        if author:
            _h2_rect(slide, left=x_q, top=q_y + quote_h + 0.45,
                     width=1.4, height=0.05, color=accent2)
            _h2_text(slide, author,
                     left=x_q, top=q_y + quote_h + 0.58,
                     width=q_w, height=0.42,
                     font=font, size_pt=_LY.T_HEADER, color='555555',
                     bold=False, align='left')
        if source:
            _h2_text(slide, source,
                     left=x_q, top=_LY.CB - 0.30,
                     width=q_w, height=0.26,
                     font=font, size_pt=9, color='AAAAAA', bold=False, align='left')
    else:
        # Encadré accent1 centré + citation blanche
        box_h = quote_h + 1.4
        box_y = max(_LY.CT + 0.15, (H - box_h) / 2)
        _h2_text(slide, '\u201c',
                 left=_LY.CL + 0.3, top=box_y - 0.42,
                 width=0.9, height=0.72,
                 font=font, size_pt=52, color=accent2,
                 bold=True, align='left')
        _h2_rounded_rect(slide, left=_LY.CL + 0.15, top=box_y,
                          width=_LY.CW - 0.3, height=box_h,
                          color=accent1, radius=_LY.R_SM)
        _h2_text(slide, quote,
                 left=_LY.CL + 0.55, top=box_y + 0.30,
                 width=_LY.CW - 1.1, height=quote_h + 0.4,
                 font=font, size_pt=20, color='FFFFFF',
                 bold=False, italic=True, align='center', line_spacing=1.35)
        if author:
            _h2_text(slide, f'\u2014\u00a0{author}',
                     left=_LY.CL + 0.55, top=box_y + box_h - 0.50,
                     width=_LY.CW - 1.1, height=0.40,
                     font=font, size_pt=11, color='DDDDDD',
                     bold=False, align='center')

    return slide


def layout_list_numbered_v4(prs: Presentation, content: dict, tp: dict):
    """
    Liste numérotée — 5 variantes visuelles déterministes.
    v0 : cercles numérotés accent_cycle + titre bold + body.
    v1 : badges rectangulaires + lignes alternées palette + bordure gauche.
    v2 : grille 2 colonnes — grand numéro coloré + trait + titre + body.
    v3 : timeline verticale — ligne centrale, cercles + contenu à droite.
    v4 : cartes horizontales pleines — grand numéro à gauche + contenu fond EEEEEE.
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accents = tp.get('accent_cycle', [
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        theme.get('accent1', '009CEA'),
    ])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    items = content.get('items', content.get('points', []))
    if not items:
        return slide

    n = min(len(items), 4)   # max 4 items — slides visuelles, pas rapports
    v = _v4_variant(content, 5, tp.get('seed', 0))

    def _item_txt(item):
        if isinstance(item, dict):
            title = _trunc(item.get('title', str(item)), 8)
            body  = _trunc(item.get('body', ''), 12)
            return title, body
        return _trunc(str(item), 12), ''

    if v == 2:
        # v2 : grille 2 colonnes — grand numéro coloré en-tête de chaque item
        n_col  = 2
        n_rows = (n + 1) // 2
        col_w  = (_LY.CW - _LY.COL_GAP) / 2
        row_h  = (_LY.CB - _LY.CT) / max(n_rows, 1)
        for i in range(n):
            color     = accents[i % len(accents)]
            title_txt, body_txt = _item_txt(items[i])
            cx = _LY.CL + (i % n_col) * (col_w + _LY.COL_GAP)
            cy = _LY.CT + (i // n_col) * row_h
            _h2_text(slide, str(i + 1),
                     left=cx, top=cy + 0.05, width=0.65, height=0.62,
                     font=font, size_pt=36, color=color, bold=True, align='left')
            _h2_rect(slide, left=cx, top=cy + 0.68, width=0.5, height=0.03, color=color)
            _h2_text(slide, title_txt,
                     left=cx + 0.02, top=cy + 0.78,
                     width=col_w - 0.05, height=0.42,
                     font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
            if body_txt:
                _h2_text(slide, body_txt,
                         left=cx + 0.02, top=cy + 1.24,
                         width=col_w - 0.05, height=row_h - 1.34,
                         font=font, size_pt=_LY.T_BODY, color='555555',
                         bold=False, align='left', line_spacing=1.1)
        return slide

    if v == 3:
        # v3 : timeline verticale — ligne accent gauche + nœuds + contenu
        n  = min(n, 5)
        step = (_LY.CB - _LY.CT) / max(n, 1)
        lx = _LY.CL + 0.55        # x de la ligne
        _h2_rect(slide, left=lx - 0.015, top=_LY.CT, width=0.03,
                 height=_LY.CB - _LY.CT, color='DDDDDD')
        for i in range(n):
            color     = accents[i % len(accents)]
            title_txt, body_txt = _item_txt(items[i])
            cy = _LY.CT + i * step + step / 2
            # Cercle sur la ligne
            _h2_circle(slide, cx=lx, cy=cy, r=0.22, color=color)
            _h2_text(slide, str(i + 1),
                     left=lx - 0.22, top=cy - 0.22,
                     width=0.44, height=0.44,
                     font=font, size_pt=12, color='FFFFFF', bold=True, align='center')
            # Titre + body à droite
            x_t = lx + 0.34
            w_t = _LY.CR - x_t
            _h2_text(slide, title_txt,
                     left=x_t, top=cy - step / 2 + 0.10,
                     width=w_t, height=0.38,
                     font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
            if body_txt:
                _h2_text(slide, body_txt,
                         left=x_t, top=cy - step / 2 + 0.50,
                         width=w_t, height=step - 0.60,
                         font=font, size_pt=_LY.T_BODY, color='555555',
                         bold=False, align='left', line_spacing=1.1)
            # Séparateur léger entre items
            if i < n - 1:
                _h2_rect(slide, left=x_t, top=_LY.CT + (i + 1) * step - 0.04,
                         width=w_t, height=0.02, color='EEEEEE')
        return slide

    if v == 4:
        # v4 : cartes horizontales — fond EEEEEE, numéro grand gauche, contenu droite
        n  = min(n, 5)
        gap   = _LY.GAP_SM
        card_h = (_LY.CB - _LY.CT - gap * (n - 1)) / max(n, 1)
        num_w  = 0.72
        for i in range(n):
            color     = accents[i % len(accents)]
            title_txt, body_txt = _item_txt(items[i])
            cy = _LY.CT + i * (card_h + gap)
            bg = 'EEEEEE' if i % 2 == 0 else 'F0F0F0'
            _h2_rounded_rect(slide, left=_LY.CL, top=cy,
                              width=_LY.CW, height=card_h,
                              color=bg, radius=_LY.R_SM)
            _h2_rect(slide, left=_LY.CL, top=cy, width=0.055, height=card_h, color=color)
            # Grand numéro
            _h2_text(slide, str(i + 1),
                     left=_LY.CL + 0.12, top=cy + (card_h - 0.50) / 2,
                     width=num_w, height=0.50,
                     font=font, size_pt=26, color=color, bold=True, align='center')
            # Séparateur vertical léger
            _h2_rect(slide, left=_LY.CL + num_w + 0.14, top=cy + 0.10,
                     width=0.02, height=card_h - 0.20, color='CCCCCC')
            x_t = _LY.CL + num_w + 0.26
            w_t = _LY.CR - x_t
            _h2_text(slide, title_txt,
                     left=x_t, top=cy + 0.10,
                     width=w_t, height=0.38,
                     font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
            if body_txt:
                _h2_text(slide, body_txt,
                         left=x_t, top=cy + 0.50,
                         width=w_t, height=card_h - 0.60,
                         font=font, size_pt=_LY.T_BODY, color='555555',
                         bold=False, align='left', line_spacing=1.1)
        return slide

    # v0 et v1 : boucle commune
    n    = min(n, 5)
    step = (_LY.CB - _LY.CT) / max(n, 1)

    for i in range(n):
        color     = accents[i % len(accents)]
        title_txt, body_txt = _item_txt(items[i])

        if v == 0:
            # v0 : cercles numérotés
            x_circ = _LY.CL + 0.4
            cy     = _LY.CT + i * step + _LY.R_CIRC
            x_text = x_circ + _LY.R_CIRC + 0.28
            _h2_circle(slide, cx=x_circ, cy=cy, r=_LY.R_CIRC, color=color)
            _h2_text(slide, str(i + 1),
                     left=x_circ - _LY.R_CIRC, top=cy - _LY.R_CIRC - 0.05,
                     width=_LY.R_CIRC * 2, height=_LY.R_CIRC * 2 + 0.1,
                     font=font, size_pt=13, color='FFFFFF', bold=True, align='center')
            _h2_text(slide, title_txt,
                     left=x_text, top=cy - _LY.R_CIRC,
                     width=_LY.CR - x_text, height=0.38,
                     font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
            if body_txt:
                _h2_text(slide, body_txt,
                         left=x_text, top=cy - _LY.R_CIRC + 0.40,
                         width=_LY.CR - x_text, height=step - _LY.R_CIRC * 2 - 0.4,
                         font=font, size_pt=_LY.T_BODY, color='555555',
                         bold=False, align='left', line_spacing=1.1)
        else:
            # v1 : badge rectangle + fond alternée + bordure gauche
            row_y   = _LY.CT + i * step
            row_h   = step - _LY.GAP_XS
            bg_row  = 'F0F0F0' if i % 2 == 0 else 'F4F4F4'
            badge_h, badge_w = 0.38, 0.46
            badge_y = row_y + (row_h - badge_h) / 2
            x_text  = _LY.CL + badge_w + _LY.PAD * 2
            _h2_rect(slide, left=_LY.CL, top=row_y, width=_LY.CW, height=row_h, color=bg_row)
            _h2_rect(slide, left=_LY.CL, top=row_y, width=_LY.BORDER_W, height=row_h, color=color)
            _h2_rounded_rect(slide, left=_LY.CL + _LY.GAP_SM, top=badge_y,
                              width=badge_w, height=badge_h, color=color, radius=_LY.R_SM)
            _h2_text(slide, str(i + 1),
                     left=_LY.CL + _LY.GAP_SM, top=badge_y,
                     width=badge_w, height=badge_h,
                     font=font, size_pt=13, color='FFFFFF', bold=True, align='center')
            _h2_text(slide, title_txt,
                     left=x_text, top=badge_y,
                     width=_LY.CR - x_text - 0.2, height=badge_h,
                     font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
            if body_txt:
                _h2_text(slide, body_txt,
                         left=x_text, top=badge_y + badge_h + 0.04,
                         width=_LY.CR - x_text - 0.2,
                         height=max(0.28, row_h - badge_h - badge_y + row_y - 0.08),
                         font=font, size_pt=_LY.T_BODY, color='555555',
                         bold=False, align='left', line_spacing=1.1)

    return slide


def _card_data(card) -> tuple:
    """Extract (icon, label, title, subtitle, body, items, stat_value, stat_label) from a card dict."""
    if not isinstance(card, dict):
        return '', '', _trunc(str(card), 10), '', '', [], '', ''
    stat = card.get('stat', {})
    if isinstance(stat, dict):
        sv, sl = stat.get('value', ''), stat.get('label', '')
    else:
        sv, sl = str(stat), ''
    sv = card.get('stat_value', sv)
    sl = card.get('stat_label', sl)
    items = [_trunc(str(it), 10) for it in card.get('items', [])[:4]]
    body  = _trunc(card.get('body', ''), 15)  # ≤ 15 mots par card body
    return (
        card.get('icon', ''),
        card.get('label', ''),
        _trunc(card.get('title', ''), 10),
        _trunc(card.get('subtitle', ''), 10),
        body,
        items,
        sv,
        sl,
    )


def layout_list_cards_v4(prs: Presentation, content: dict, tp: dict):
    """
    Cartes — 5 variantes visuelles déterministes.
    v0 : grille 2×2 — barre accent top + icon + title + body/items + stat bas.
    v1 : colonnes verticales pleine hauteur — bande colorée + icon + title + body/items + stat.
    v2 : bandeaux horizontaux — rayure accent gauche + icon + title + body.
    v3 : spotlight icons — grand icon centré + fond EBF9F3 + title + body, 2-3 cartes larges.
    v4 : steps flow — connecteur horizontal + numéro circulaire + title + body.
    card schema: {icon?, label?, title, subtitle?, body?, items?:[str], stat_value?, stat_label?}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accents = tp.get('accent_cycle', [
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        theme.get('accent1', '009CEA'),
        theme.get('accent2', 'ED0000'),
    ])
    _TPL = ['F0F0F0', 'EEEEEE', 'F4F4F4']

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    cards = content.get('cards', content.get('items', []))
    if not cards:
        return slide

    n = min(len(cards), 4)
    v = _v4_variant(content, 5, tp.get('seed', 0))

    def _has_stat(card):
        d = _card_data(card)
        return bool(d[6])  # stat_value

    any_stat = any(_has_stat(c) for c in cards[:n])

    if v == 0:
        # Variante 0 : grille 2×2
        card_h = (_LY.CB - _LY.CT - _LY.GAP_LG) / 2
        x_cols = [_LY.CL, _LY.CL + _LY.COL_W + _LY.GAP_LG]
        y_rows = [_LY.CT, _LY.CT + card_h + _LY.GAP_LG]
        for i in range(n):
            cx, cy = x_cols[i % 2], y_rows[i // 2]
            color  = accents[i % len(accents)]
            icon, label, title_txt, subtitle, body_txt, items, sv, sl = _card_data(cards[i])
            _h2_rounded_rect(slide, left=cx, top=cy,
                              width=_LY.COL_W, height=card_h,
                              color='F8F8F8', radius=_LY.R_SM)
            _h2_rect(slide, left=cx, top=cy, width=_LY.COL_W, height=0.07, color=color)
            # Bordure gauche colorée
            _h2_rect(slide, left=cx, top=cy + 0.07, width=0.055, height=card_h - 0.07, color=color)
            y_cur = cy + 0.16
            if icon:
                _h2_text(slide, icon,
                         left=cx + 0.18, top=y_cur,
                         width=0.44, height=0.38,
                         font=font, size_pt=16, color=dk1, bold=False, align='left')
                y_cur += 0.40
            _h2_text(slide, title_txt,
                     left=cx + 0.18, top=y_cur,
                     width=_LY.COL_W - 0.26, height=0.40,
                     font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
            y_cur += 0.42
            body_h = card_h - (y_cur - cy) - (0.60 if sv else 0.10) - 0.08
            body_src = body_txt or ('\n'.join(f'• {it}' for it in items[:5]) if items else '')
            if body_src:
                _h2_text(slide, body_src,
                         left=cx + 0.18, top=y_cur,
                         width=_LY.COL_W - 0.26, height=max(0.30, body_h),
                         font=font, size_pt=_LY.T_SMALL, color=dk1,
                         bold=False, align='left', line_spacing=1.2)
            if sv:
                _h2_rect(slide, left=cx + 0.18, top=cy + card_h - 0.62,
                         width=_LY.COL_W - 0.26, height=0.025, color='DDDDDD')
                _h2_text(slide, sv,
                         left=cx + 0.18, top=cy + card_h - 0.58,
                         width=_LY.COL_W - 0.26, height=0.36,
                         font=font, size_pt=20, color=color, bold=True, align='left')
                if sl:
                    _h2_text(slide, sl,
                             left=cx + 0.18, top=cy + card_h - 0.22,
                             width=_LY.COL_W - 0.26, height=0.20,
                             font=font, size_pt=8, color='888888', bold=True, align='left')
    elif v == 1:
        # Variante 1 : colonnes verticales (2 ou 3 cartes pleine hauteur)
        card_h = _LY.CB - _LY.CT
        card_w = (_LY.CW - _LY.GAP_LG * (n - 1)) / n
        for i in range(n):
            cx    = _LY.CL + i * (card_w + _LY.GAP_LG)
            color = accents[i % len(accents)]
            icon, label, title_txt, subtitle, body_txt, items, sv, sl = _card_data(cards[i])
            _h2_rounded_rect(slide, left=cx, top=_LY.CT,
                              width=card_w, height=card_h,
                              color='F8F8F8', radius=_LY.R_SM)
            # Bande colorée haute
            _h2_rect(slide, left=cx, top=_LY.CT, width=card_w, height=0.07, color=color)
            y_cur = _LY.CT + 0.16
            if icon:
                _h2_text(slide, icon,
                         left=cx + _LY.PAD, top=y_cur,
                         width=0.44, height=0.38,
                         font=font, size_pt=18, color=dk1, bold=False, align='left')
                if label:
                    _h2_text(slide, label,
                             left=cx + _LY.PAD + 0.46, top=y_cur + 0.08,
                             width=card_w - _LY.PAD - 0.50, height=0.24,
                             font=font, size_pt=8, color='999999', bold=True, align='right')
                y_cur += 0.42
            _h2_text(slide, title_txt,
                     left=cx + _LY.PAD, top=y_cur,
                     width=card_w - _LY.PAD * 2, height=0.42,
                     font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
            y_cur += 0.44
            if subtitle:
                _h2_text(slide, subtitle,
                         left=cx + _LY.PAD, top=y_cur,
                         width=card_w - _LY.PAD * 2, height=0.26,
                         font=font, size_pt=9, color='777777', bold=False, align='left')
                _h2_rect(slide, left=cx + _LY.PAD, top=y_cur + 0.28,
                         width=card_w - _LY.PAD * 2, height=0.025, color='E8E8E8')
                y_cur += 0.34
            body_h = card_h - (y_cur - _LY.CT) - (0.72 if sv else 0.10) - 0.08
            body_src = body_txt or ('\n'.join(f'• {it}' for it in items[:6]) if items else '')
            if body_src:
                _h2_text(slide, body_src,
                         left=cx + _LY.PAD, top=y_cur,
                         width=card_w - _LY.PAD * 2, height=max(0.30, body_h),
                         font=font, size_pt=_LY.T_SMALL, color=dk1,
                         bold=False, align='left', line_spacing=1.2)
            if sv:
                _h2_rect(slide, left=cx + _LY.PAD, top=_LY.CB - 0.68,
                         width=card_w - _LY.PAD * 2, height=0.025, color='DDDDDD')
                _h2_text(slide, sv,
                         left=cx + _LY.PAD, top=_LY.CB - 0.65,
                         width=card_w - _LY.PAD * 2, height=0.38,
                         font=font, size_pt=22, color=color, bold=True, align='left')
                if sl:
                    _h2_text(slide, sl,
                             left=cx + _LY.PAD, top=_LY.CB - 0.26,
                             width=card_w - _LY.PAD * 2, height=0.20,
                             font=font, size_pt=8, color='888888', bold=True, align='left')
    else:
        # Variante 2 : bandeaux horizontaux — rayure accent à gauche + fond palette alternée
        stripe_w = 0.38
        gap      = _LY.GAP_SM
        card_h   = (_LY.CB - _LY.CT - gap * (n - 1)) / max(n, 1)
        for i in range(n):
            cy     = _LY.CT + i * (card_h + gap)
            color  = accents[i % len(accents)]
            bg     = _TPL[i % len(_TPL)]
            icon, label, title_txt, subtitle, body_txt, items, sv, sl = _card_data(cards[i])
            # Fond principal
            _h2_rounded_rect(slide, left=_LY.CL, top=cy,
                              width=_LY.CW, height=card_h,
                              color=bg, radius=_LY.R_SM)
            # Rayure accent gauche
            _h2_rect(slide, left=_LY.CL, top=cy, width=stripe_w, height=card_h, color=color)
            tx = _LY.CL + stripe_w + _LY.PAD
            tw = _LY.CW - stripe_w - _LY.PAD * 2
            y_title = cy + 0.10
            if icon:
                _h2_text(slide, icon,
                         left=tx, top=y_title,
                         width=0.38, height=0.38,
                         font=font, size_pt=15, color=dk1, bold=False, align='left')
                tx += 0.40
                tw -= 0.40
            _h2_text(slide, title_txt,
                     left=tx, top=y_title,
                     width=tw, height=0.38,
                     font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
            body_src = body_txt or ('\n'.join(f'• {it}' for it in items[:4]) if items else '')
            if body_src:
                _h2_text(slide, body_src,
                         left=tx, top=cy + 0.52,
                         width=tw, height=card_h - 0.62,
                         font=font, size_pt=_LY.T_SMALL, color=dk1,
                         bold=False, align='left', line_spacing=1.2)
        return slide

    if v == 3:
        # Variante 3 : spotlight icons — grand icon centré + fond palette + title + body
        n3 = min(n, 3)
        card_w = (_LY.CW - _LY.GAP_LG * (n3 - 1)) / n3
        card_h = _LY.CB - _LY.CT
        bgs    = ['F0F0F0', 'EEEEEE', 'F4F4F4']
        for i in range(n3):
            cx    = _LY.CL + i * (card_w + _LY.GAP_LG)
            color = accents[i % len(accents)]
            bg    = bgs[i % len(bgs)]
            icon, label, title_txt, subtitle, body_txt, items, sv, sl = _card_data(cards[i])
            _h2_rounded_rect(slide, left=cx, top=_LY.CT,
                              width=card_w, height=card_h, color=bg, radius=_LY.R_SM)
            _h2_rect(slide, left=cx, top=_LY.CT, width=card_w, height=0.07, color=color)
            if label:
                _h2_rounded_rect(slide, left=cx + card_w - 1.30, top=_LY.CT + 0.12,
                                  width=1.10, height=0.24, color=color, radius=0.06)
                _h2_text(slide, label, left=cx + card_w - 1.28, top=_LY.CT + 0.13,
                         width=1.06, height=0.22,
                         font=font, size_pt=7, color='FFFFFF', bold=True, align='center')
            y_icon = _LY.CT + 0.22
            if icon:
                _h2_text(slide, icon, left=cx, top=y_icon, width=card_w, height=0.60,
                         font=font, size_pt=28, color=color, bold=False, align='center')
                y_icon += 0.64
            else:
                y_icon += 0.10
            _h2_rect(slide, left=cx + _LY.PAD, top=y_icon,
                     width=card_w - _LY.PAD * 2, height=0.04, color=color)
            y_icon += 0.10
            _h2_text(slide, title_txt, left=cx + _LY.PAD, top=y_icon,
                     width=card_w - _LY.PAD * 2, height=0.40,
                     font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='center')
            y_icon += 0.44
            body_src = body_txt or ('\n'.join(f'• {it}' for it in items[:5]) if items else '')
            if body_src:
                body_h = card_h - (y_icon - _LY.CT) - (0.60 if sv else 0.10) - 0.08
                _h2_text(slide, body_src, left=cx + _LY.PAD, top=y_icon,
                         width=card_w - _LY.PAD * 2, height=max(0.30, body_h),
                         font=font, size_pt=_LY.T_SMALL, color=dk1,
                         bold=False, align='left', line_spacing=1.2)
            if sv:
                _h2_rect(slide, left=cx + _LY.PAD, top=_LY.CB - 0.62,
                         width=card_w - _LY.PAD * 2, height=0.025, color='CCCCCC')
                _h2_text(slide, sv, left=cx + _LY.PAD, top=_LY.CB - 0.60,
                         width=card_w - _LY.PAD * 2, height=0.36,
                         font=font, size_pt=22, color=color, bold=True, align='center')
                if sl:
                    _h2_text(slide, sl, left=cx + _LY.PAD, top=_LY.CB - 0.24,
                             width=card_w - _LY.PAD * 2, height=0.18,
                             font=font, size_pt=7, color='999999', bold=True, align='center')
        return slide

    # Variante 4 : steps flow — ligne connecteur horizontal + numéro circulaire + title + body
    n4     = min(n, 4)
    card_w = (_LY.CW - _LY.GAP_LG * (n4 - 1)) / n4
    conn_y = _LY.CT + 0.22
    circle_r = 0.28
    card_top = conn_y + circle_r + 0.14
    card_h   = _LY.CB - card_top
    _h2_rect(slide, left=_LY.CL, top=conn_y + circle_r - 0.02,
             width=_LY.CW, height=0.04, color='E0E0E0')
    for i in range(n4):
        cx    = _LY.CL + i * (card_w + _LY.GAP_LG)
        color = accents[i % len(accents)]
        icon, label, title_txt, subtitle, body_txt, items, sv, sl = _card_data(cards[i])
        num_cx = cx + card_w / 2 - circle_r
        _h2_rounded_rect(slide, left=num_cx, top=conn_y,
                          width=circle_r * 2, height=circle_r * 2,
                          color=color, radius=circle_r)
        _h2_text(slide, str(i + 1), left=num_cx, top=conn_y + 0.02,
                 width=circle_r * 2, height=circle_r * 2,
                 font=font, size_pt=14, color='FFFFFF', bold=True, align='center')
        _h2_rounded_rect(slide, left=cx, top=card_top, width=card_w, height=card_h,
                          color='F8F8F8', radius=_LY.R_SM)
        _h2_rect(slide, left=cx, top=card_top, width=card_w, height=0.055, color=color)
        y_cur = card_top + 0.12
        if icon:
            _h2_text(slide, icon, left=cx + _LY.PAD, top=y_cur,
                     width=0.36, height=0.34,
                     font=font, size_pt=14, color=color, bold=False, align='left')
            y_cur += 0.36
        _h2_text(slide, title_txt, left=cx + _LY.PAD, top=y_cur,
                 width=card_w - _LY.PAD * 2, height=0.38,
                 font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
        y_cur += 0.40
        body_src = body_txt or ('\n'.join(f'• {it}' for it in items[:5]) if items else '')
        if body_src:
            body_h = card_h - (y_cur - card_top) - 0.08
            _h2_text(slide, body_src, left=cx + _LY.PAD, top=y_cur,
                     width=card_w - _LY.PAD * 2, height=max(0.30, body_h),
                     font=font, size_pt=_LY.T_SMALL, color=dk1,
                     bold=False, align='left', line_spacing=1.2)

    return slide


def layout_col3_v4(prs: Presentation, content: dict, tp: dict):
    """
    3 colonnes enrichies — barre accent top + icon + label + title + subtitle + bullets + stat bas.
    Idéal pour comparaisons thématiques, piliers stratégiques, axes d'analyse.
    content: {title, section_label?, subtitle?, columns:[{icon?,label?,title,subtitle?,items:[str],stat_value?,stat_label?}]}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accents = tp.get('accent_cycle', [
        theme.get('accent1', '009CEA'),
        theme.get('accent2', 'ED0000'),
        theme.get('accent3', '40A900'),
    ])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    columns = content.get('columns', [])
    if not columns:
        return slide

    n      = min(len(columns), 3)
    card_w = (_LY.CW - _LY.GAP_LG * (n - 1)) / n
    card_h = _LY.CB - _LY.CT
    bar_h  = 0.07

    for i, col in enumerate(columns[:n]):
        cx    = _LY.CL + i * (card_w + _LY.GAP_LG)
        color = accents[i % len(accents)]
        icon     = col.get('icon', '') if isinstance(col, dict) else ''
        label    = col.get('label', '') if isinstance(col, dict) else ''
        title    = col.get('title', '') if isinstance(col, dict) else str(col)
        subtitle = col.get('subtitle', '') if isinstance(col, dict) else ''
        items    = col.get('items', []) if isinstance(col, dict) else []
        body     = col.get('body', '') if isinstance(col, dict) else ''
        stat     = col.get('stat', {}) if isinstance(col, dict) else {}
        if isinstance(stat, dict):
            sv, sl = stat.get('value', ''), stat.get('label', '')
        else:
            sv, sl = str(stat), ''
        sv = col.get('stat_value', sv) if isinstance(col, dict) else sv
        sl = col.get('stat_label', sl) if isinstance(col, dict) else sl

        # Fond carte
        _h2_rounded_rect(slide, left=cx, top=_LY.CT, width=card_w, height=card_h,
                          color='F8F8F8', radius=_LY.R_SM)
        # Barre accent top
        _h2_rect(slide, left=cx, top=_LY.CT, width=card_w, height=bar_h, color=color)

        y = _LY.CT + bar_h + 0.14

        # Icône (gauche) + label catégorie (droite, petite)
        if icon:
            _h2_text(slide, icon,
                     left=cx + _LY.PAD, top=y,
                     width=0.44, height=0.38,
                     font=font, size_pt=18, color=dk1, bold=False, align='left')
        if label:
            _h2_text(slide, label,
                     left=cx + _LY.PAD + 0.46, top=y + 0.08,
                     width=card_w - _LY.PAD - 0.52, height=0.24,
                     font=font, size_pt=8, color='999999', bold=True, align='right')
        if icon or label:
            y += 0.44

        # Titre
        _h2_text(slide, title,
                 left=cx + _LY.PAD, top=y,
                 width=card_w - _LY.PAD * 2, height=0.42,
                 font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
        y += 0.44

        # Sous-titre + séparateur
        if subtitle:
            _h2_text(slide, subtitle,
                     left=cx + _LY.PAD, top=y,
                     width=card_w - _LY.PAD * 2, height=0.26,
                     font=font, size_pt=9, color='777777', bold=False, align='left')
            _h2_rect(slide, left=cx + _LY.PAD, top=y + 0.28,
                     width=card_w - _LY.PAD * 2, height=0.025, color='E0E0E0')
            y += 0.34

        # Corps / bullets
        body_src = body or ('\n'.join(f'• {it}' for it in items[:6]) if items else '')
        stat_reserved = 0.78 if sv else 0.14
        body_h = max(0.10, _LY.CB - y - stat_reserved - 0.06)
        if body_src:
            _h2_text(slide, body_src,
                     left=cx + _LY.PAD, top=y,
                     width=card_w - _LY.PAD * 2, height=body_h,
                     font=font, size_pt=_LY.T_SMALL, color=dk1,
                     bold=False, align='left', line_spacing=1.2)

        # Stat bas
        if sv:
            _h2_rect(slide, left=cx + _LY.PAD, top=_LY.CB - 0.74,
                     width=card_w - _LY.PAD * 2, height=0.025, color='DDDDDD')
            _h2_text(slide, sv,
                     left=cx + _LY.PAD, top=_LY.CB - 0.70,
                     width=card_w - _LY.PAD * 2, height=0.38,
                     font=font, size_pt=24, color=color, bold=True, align='left')
            if sl:
                _h2_text(slide, sl,
                         left=cx + _LY.PAD, top=_LY.CB - 0.30,
                         width=card_w - _LY.PAD * 2, height=0.22,
                         font=font, size_pt=8, color='888888', bold=True, align='left')

    return slide


def layout_twocol_v4(prs: Presentation, content: dict, tp: dict):
    """
    Deux colonnes — 3 variantes visuelles déterministes.
    v0 : headers rectangulaires colorés + items sur fonds alternés.
    v1 : fond palette (E8EEF5 / EBF9F3) + puce ronde colorée + items aérés.
    v2 : tiret accent gauche + titre coloré, pas de fond de colonne.
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accent3 = theme.get('accent3', '40A900')

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    col_a = content.get('col_a', {})
    col_b = content.get('col_b', {})
    if not col_a and not col_b:
        return slide

    v       = _v4_variant(content, 5, tp.get('seed', 0))
    x_a     = _LY.CL
    x_b     = _LY.CL + _LY.COL_W + _LY.COL_GAP
    y_head  = _LY.CT
    cols    = [(x_a, col_a, accent1), (x_b, col_b, accent3)]

    # Détecter si une colonne a un sous-titre pour ajuster l'espace header
    any_subtitle = any(
        bool((c.get('subtitle', '') if isinstance(c, dict) else ''))
        for _, c, _ in cols
    )
    head_h  = _LY.HEAD_H + (0.32 if any_subtitle else 0)
    y_items = y_head + head_h + _LY.GAP_SM

    def _col_items(col):
        raw = col.get('items', []) if isinstance(col, dict) else []
        return [_trunc(str(it), 10) for it in raw[:4]]  # max 4 items, ≤ 10 mots

    def _col_label(col):
        return col.get('title', '') if isinstance(col, dict) else ''

    def _col_sub(col):
        return col.get('subtitle', '') if isinstance(col, dict) else ''

    if v == 0:
        # Variante 0 : header rectangulaire coloré + fonds alternés
        for x_col, col, color in cols:
            label, subtitle, items = _col_label(col), _col_sub(col), _col_items(col)
            _h2_rect(slide, left=x_col, top=y_head,
                     width=_LY.COL_W, height=head_h, color=color)
            _h2_text(slide, label,
                     left=x_col + _LY.PAD, top=y_head + 0.07,
                     width=_LY.COL_W - _LY.PAD * 2, height=_LY.HEAD_H - 0.1,
                     font=font, size_pt=_LY.T_HEADER, color='FFFFFF',
                     bold=True, align='left')
            if subtitle:
                _h2_text(slide, subtitle,
                         left=x_col + _LY.PAD, top=y_head + _LY.HEAD_H - 0.02,
                         width=_LY.COL_W - _LY.PAD * 2, height=0.30,
                         font=font, size_pt=_LY.T_SMALL, color='FFFFFF',
                         bold=False, italic=True, align='left')
            for j, item in enumerate(items[:8]):
                item_txt = item.get('title', str(item)) if isinstance(item, dict) else str(item)
                bg_color = 'EEEEEE' if j % 2 == 0 else 'F8F8F8'
                iy = y_items + j * _LY.ITEM_H
                _h2_rect(slide, left=x_col, top=iy,
                         width=_LY.COL_W, height=_LY.ITEM_H - 0.04, color=bg_color)
                _h2_rect(slide, left=x_col, top=iy,
                         width=_LY.BORDER_W, height=_LY.ITEM_H - 0.04, color=color)
                _h2_text(slide, item_txt,
                         left=x_col + _LY.PAD, top=iy + 0.06,
                         width=_LY.COL_W - _LY.PAD - 0.04, height=_LY.ITEM_H - 0.1,
                         font=font, size_pt=_LY.T_BODY, color=dk1, bold=False, align='left')
    elif v == 1:
        # Variante 1 : fond palette pleine hauteur + puce ronde colorée + items aérés
        bg_cols = ['F0F0F0', 'F4F4F4']
        for ci, (x_col, col, color) in enumerate(cols):
            label, subtitle, items = _col_label(col), _col_sub(col), _col_items(col)
            col_h = _LY.CB - y_head
            _h2_rounded_rect(slide, left=x_col, top=y_head,
                              width=_LY.COL_W, height=col_h,
                              color=bg_cols[ci], radius=_LY.R_SM)
            _h2_text(slide, label,
                     left=x_col + _LY.PAD, top=y_head + 0.07,
                     width=_LY.COL_W - _LY.PAD * 2, height=_LY.HEAD_H - 0.1,
                     font=font, size_pt=_LY.T_HEADER, color=color,
                     bold=True, align='left')
            if subtitle:
                _h2_text(slide, subtitle,
                         left=x_col + _LY.PAD, top=y_head + _LY.HEAD_H - 0.02,
                         width=_LY.COL_W - _LY.PAD * 2, height=0.30,
                         font=font, size_pt=_LY.T_SMALL, color=color,
                         bold=False, italic=True, align='left')
            _h2_rect(slide, left=x_col + _LY.PAD, top=y_head + head_h - 0.05,
                     width=_LY.COL_W - _LY.PAD * 2, height=0.03, color=color)
            item_h2 = (_LY.CB - y_items) / max(len(items[:8]), 1)
            for j, item in enumerate(items[:8]):
                item_txt = item.get('title', str(item)) if isinstance(item, dict) else str(item)
                iy = y_items + j * item_h2
                _h2_circle(slide, cx=x_col + _LY.PAD + 0.10,
                            cy=iy + item_h2 / 2, r=0.07, color=color)
                _h2_text(slide, item_txt,
                         left=x_col + _LY.PAD + 0.28, top=iy + (item_h2 - 0.36) / 2,
                         width=_LY.COL_W - _LY.PAD - 0.35, height=0.36,
                         font=font, size_pt=_LY.T_BODY, color=dk1,
                         bold=False, align='left', line_spacing=1.1)
    elif v == 2:
        # Variante 2 : tiret accent gauche + titre coloré, fond blanc aéré
        for x_col, col, color in cols:
            label, subtitle, items = _col_label(col), _col_sub(col), _col_items(col)
            _h2_rect(slide, left=x_col, top=y_head,
                     width=0.05, height=_LY.CB - y_head, color=color)
            _h2_text(slide, label,
                     left=x_col + 0.18, top=y_head + 0.07,
                     width=_LY.COL_W - 0.22, height=_LY.HEAD_H - 0.1,
                     font=font, size_pt=_LY.T_HEADER, color=color,
                     bold=True, align='left')
            if subtitle:
                _h2_text(slide, subtitle,
                         left=x_col + 0.18, top=y_head + _LY.HEAD_H - 0.02,
                         width=_LY.COL_W - 0.22, height=0.30,
                         font=font, size_pt=_LY.T_SMALL, color=dk1,
                         bold=False, italic=True, align='left')
            item_h2 = (_LY.CB - y_items) / max(len(items[:8]), 1)
            for j, item in enumerate(items[:8]):
                item_txt = item.get('title', str(item)) if isinstance(item, dict) else str(item)
                iy = y_items + j * item_h2
                _h2_rect(slide, left=x_col + 0.18, top=iy + item_h2 / 2 - 0.02,
                         width=0.18, height=0.04, color=color)
                _h2_text(slide, item_txt,
                         left=x_col + 0.42, top=iy + (item_h2 - 0.36) / 2,
                         width=_LY.COL_W - 0.45, height=0.36,
                         font=font, size_pt=_LY.T_BODY, color=dk1,
                         bold=False, align='left', line_spacing=1.1)
    elif v == 3:
        # Variante 3 : cartes aérées sur fond blanc — grand label accent + séparateur + items numérotés
        for x_col, col, color in cols:
            label, subtitle, items = _col_label(col), _col_sub(col), _col_items(col)
            # Grande étiquette de colonne
            _h2_text(slide, label,
                     left=x_col + _LY.PAD, top=y_head + 0.06,
                     width=_LY.COL_W - _LY.PAD * 2, height=_LY.HEAD_H - 0.1,
                     font=font, size_pt=_LY.T_HEADER + 2, color=color,
                     bold=True, align='left')
            if subtitle:
                _h2_text(slide, subtitle,
                         left=x_col + _LY.PAD, top=y_head + _LY.HEAD_H - 0.02,
                         width=_LY.COL_W - _LY.PAD * 2, height=0.30,
                         font=font, size_pt=_LY.T_SMALL, color='888888',
                         bold=False, italic=True, align='left')
            _h2_rect(slide, left=x_col + _LY.PAD, top=y_head + head_h,
                     width=1.8, height=0.05, color=color)
            # Items numérotés dans des cartes
            n_it   = min(len(items), 6)
            gap    = _LY.GAP_SM
            it_h   = (_LY.CB - y_items - gap * (n_it - 1)) / max(n_it, 1)
            for j, item in enumerate(items[:n_it]):
                item_txt = item.get('title', str(item)) if isinstance(item, dict) else str(item)
                iy = y_items + j * (it_h + gap)
                _h2_rounded_rect(slide, left=x_col, top=iy,
                                  width=_LY.COL_W, height=it_h,
                                  color='F8F8F8', radius=_LY.R_SM)
                _h2_text(slide, str(j + 1),
                         left=x_col + 0.10, top=iy + (it_h - 0.34) / 2,
                         width=0.28, height=0.34,
                         font=font, size_pt=11, color=color, bold=True, align='center')
                _h2_text(slide, item_txt,
                         left=x_col + 0.44, top=iy + (it_h - 0.34) / 2,
                         width=_LY.COL_W - 0.50, height=0.34,
                         font=font, size_pt=_LY.T_BODY, color=dk1,
                         bold=False, align='left', line_spacing=1.1)
    else:
        # Variante 4 : fond EEEEEE colonne gauche / E8EEF5 colonne droite + items avec icône ✓/→
        icons_map = [('→', accent1), ('✓', accent3)]
        for ci, (x_col, col, color) in enumerate(cols):
            label, subtitle, items = _col_label(col), _col_sub(col), _col_items(col)
            bg = 'EEEEEE' if ci == 0 else 'F0F0F0'
            icon_chr, icon_col = icons_map[ci]
            _h2_rounded_rect(slide, left=x_col, top=y_head,
                              width=_LY.COL_W, height=_LY.CB - y_head,
                              color=bg, radius=_LY.R_SM)
            _h2_rect(slide, left=x_col, top=y_head,
                     width=_LY.COL_W, height=0.06, color=color)
            _h2_text(slide, label,
                     left=x_col + _LY.PAD, top=y_head + 0.10,
                     width=_LY.COL_W - _LY.PAD * 2, height=_LY.HEAD_H - 0.14,
                     font=font, size_pt=_LY.T_HEADER, color=dk1,
                     bold=True, align='left')
            if subtitle:
                _h2_text(slide, subtitle,
                         left=x_col + _LY.PAD, top=y_head + _LY.HEAD_H - 0.02,
                         width=_LY.COL_W - _LY.PAD * 2, height=0.28,
                         font=font, size_pt=_LY.T_SMALL, color='777777',
                         bold=False, italic=True, align='left')
            item_h2 = (_LY.CB - y_items) / max(len(items[:8]), 1)
            for j, item in enumerate(items[:8]):
                item_txt = item.get('title', str(item)) if isinstance(item, dict) else str(item)
                iy = y_items + j * item_h2
                _h2_text(slide, icon_chr,
                         left=x_col + 0.12, top=iy + (item_h2 - 0.34) / 2,
                         width=0.26, height=0.34,
                         font=font, size_pt=12, color=icon_col, bold=True, align='center')
                _h2_text(slide, item_txt,
                         left=x_col + 0.44, top=iy + (item_h2 - 0.34) / 2,
                         width=_LY.COL_W - 0.50, height=0.34,
                         font=font, size_pt=_LY.T_BODY, color=dk1,
                         bold=False, align='left', line_spacing=1.1)

    return slide


def layout_stathero_v4(prs: Presentation, content: dict, tp: dict):
    """
    Grande statistique centrée — 5 variantes visuelles déterministes.
    v0 : fond blanc, valeur accent1 centrée, label dk1.
    v1 : bande verticale accent gauche (1.8") + valeur centrée à droite.
    v2 : cadre centré E8EEF5 (arrondi) contenant valeur + label.
    v3 : split écran — panneau dk1 gauche (valeur blanche) + contexte droite.
    v4 : multi-stats — 2-3 valeurs côte à côte si 'values' fourni.
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accent2 = theme.get('accent2', 'ED0000')
    W       = tp.get('W', 13.33)
    v       = _v4_variant(content, 5, tp.get('seed', 0))

    value   = str(content.get('value', ''))
    label   = content.get('label', '')
    context = content.get('context', '')
    y_center = (_LY.CT + _LY.CB) / 2
    val_h    = 1.1
    val_y    = y_center - val_h / 2 - (0.4 if label else 0)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    if v == 0:
        # Variante 0 : fond blanc, valeur accent1 centrée
        _h2_text(slide, value,
                 left=0.5, top=val_y, width=W - 1.0, height=val_h,
                 font=font, size_pt=_LY.T_HERO, color=accent1, bold=True, align='center')
        if label:
            line_w = min(3.0, len(value) * 0.55 + 0.5)
            _h2_rect(slide, left=(W - line_w) / 2, top=val_y + val_h,
                     width=line_w, height=_LY.BAR_H, color=accent2)
            _h2_text(slide, label,
                     left=0.5, top=val_y + val_h + 0.1, width=W - 1.0, height=0.55,
                     font=font, size_pt=18, color=dk1, bold=False, align='center')
        if context:
            _h2_text(slide, context,
                     left=1.5, top=val_y + val_h + (0.75 if label else 0.15),
                     width=W - 3.0, height=0.6,
                     font=font, size_pt=_LY.T_LABEL, color='666666',
                     bold=False, align='center', line_spacing=1.2)
    elif v == 1:
        # Variante 1 : bande verticale accent1 à gauche + valeur à droite
        band_w = 1.80
        _h2_rect(slide, left=0, top=0, width=band_w, height=7.5, color=accent1)
        rx = band_w + 0.5
        rw = W - band_w - 0.8
        _h2_text(slide, value,
                 left=rx, top=val_y, width=rw, height=val_h,
                 font=font, size_pt=_LY.T_HERO, color=accent1, bold=True, align='center')
        if label:
            line_w = min(3.0, len(value) * 0.55 + 0.5)
            _h2_rect(slide, left=rx + (rw - line_w) / 2, top=val_y + val_h,
                     width=line_w, height=_LY.BAR_H, color=accent2)
            _h2_text(slide, label,
                     left=rx, top=val_y + val_h + 0.1, width=rw, height=0.55,
                     font=font, size_pt=18, color=dk1, bold=False, align='center')
        if context:
            _h2_text(slide, context,
                     left=rx, top=val_y + val_h + (0.75 if label else 0.15),
                     width=rw, height=0.6,
                     font=font, size_pt=_LY.T_LABEL, color='666666',
                     bold=False, align='center', line_spacing=1.2)
    elif v == 2:
        # Variante 2 : cadre centré E8EEF5 contenant valeur + label
        frame_w = 7.0
        frame_h = 2.8
        fx = (W - frame_w) / 2
        fy = y_center - frame_h / 2
        _h2_rounded_rect(slide, left=fx, top=fy,
                         width=frame_w, height=frame_h,
                         color='F0F0F0', radius=_LY.R_MD)
        _h2_rect(slide, left=fx, top=fy, width=frame_w, height=0.08, color=accent1)
        _h2_text(slide, value,
                 left=fx, top=fy + 0.18, width=frame_w, height=val_h,
                 font=font, size_pt=_LY.T_HERO, color=accent1, bold=True, align='center')
        if label:
            _h2_text(slide, label,
                     left=fx, top=fy + 0.18 + val_h + 0.08, width=frame_w, height=0.55,
                     font=font, size_pt=18, color=dk1, bold=False, align='center')
        if context:
            ctx_y = fy + frame_h + 0.18
            _h2_text(slide, context,
                     left=1.5, top=ctx_y, width=W - 3.0, height=0.6,
                     font=font, size_pt=_LY.T_LABEL, color='666666',
                     bold=False, align='center', line_spacing=1.2)

    elif v == 3:
        # Variante 3 : split screen — panneau dk1 gauche + contexte droite
        split_x = W * 0.45
        _h2_rect(slide, left=0, top=0, width=split_x, height=7.5, color=dk1)
        _h2_text(slide, value,
                 left=0.3, top=val_y, width=split_x - 0.4, height=val_h,
                 font=font, size_pt=_LY.T_HERO, color=accent1, bold=True, align='center')
        if label:
            _h2_rect(slide, left=(split_x - 2.0) / 2, top=val_y + val_h,
                     width=2.0, height=_LY.BAR_H, color=accent2)
            _h2_text(slide, label,
                     left=0.3, top=val_y + val_h + 0.12, width=split_x - 0.4, height=0.50,
                     font=font, size_pt=16, color='DDDDDD', bold=False, align='center')
        rx = split_x + 0.5
        rw = W - rx - 0.4
        if context:
            _h2_text(slide, context,
                     left=rx, top=val_y, width=rw, height=1.6,
                     font=font, size_pt=13, color=dk1,
                     bold=False, align='left', line_spacing=1.4)
        points = content.get('points', [])
        pt_y = val_y + (1.80 if context else 0.20)
        for j, pt in enumerate(points[:5]):
            _h2_rect(slide, left=rx, top=pt_y + j * 0.44 + 0.10,
                     width=0.06, height=0.06, color=accent1)
            _h2_text(slide, str(pt),
                     left=rx + 0.16, top=pt_y + j * 0.44,
                     width=rw - 0.16, height=0.40,
                     font=font, size_pt=11, color=dk1, bold=False, align='left')

    else:
        # Variante 4 : multi-stats — jusqu'à 3 valeurs côte à côte
        values_list = content.get('values', [])
        if not values_list:
            values_list = [{'value': value, 'label': label, 'context': context}]
        nv      = min(len(values_list), 3)
        stat_w  = (_LY.CW - _LY.GAP_LG * (nv - 1)) / nv
        bgs     = ['F0F0F0', 'EEEEEE', 'F4F4F4']
        accs    = [accent1, accent2, theme.get('accent3', '40A900')]
        for i, vs in enumerate(values_list[:nv]):
            sx    = _LY.CL + i * (stat_w + _LY.GAP_LG)
            color = accs[i % len(accs)]
            bg    = bgs[i % len(bgs)]
            sv    = str(vs.get('value', vs) if isinstance(vs, dict) else vs)
            sl    = vs.get('label', '') if isinstance(vs, dict) else ''
            sc    = vs.get('context', '') if isinstance(vs, dict) else ''
            _h2_rounded_rect(slide, left=sx, top=_LY.CT,
                              width=stat_w, height=_LY.CB - _LY.CT,
                              color=bg, radius=_LY.R_MD)
            _h2_rect(slide, left=sx, top=_LY.CT, width=stat_w, height=0.08, color=color)
            mid_y = (_LY.CT + _LY.CB) / 2
            _h2_text(slide, sv,
                     left=sx + 0.2, top=mid_y - 0.70, width=stat_w - 0.4, height=1.10,
                     font=font, size_pt=52, color=color, bold=True, align='center')
            if sl:
                _h2_rect(slide, left=sx + (stat_w - 1.6) / 2, top=mid_y + 0.44,
                         width=1.6, height=_LY.BAR_H, color=color)
                _h2_text(slide, sl,
                         left=sx + 0.2, top=mid_y + 0.56, width=stat_w - 0.4, height=0.42,
                         font=font, size_pt=13, color=dk1, bold=False, align='center')
            if sc:
                _h2_text(slide, sc,
                         left=sx + 0.2, top=mid_y + (1.06 if sl else 0.56),
                         width=stat_w - 0.4, height=0.60,
                         font=font, size_pt=9, color='888888',
                         bold=False, align='center', line_spacing=1.2)

    return slide


def layout_infographic_v4(prs: Presentation, content: dict, tp: dict):
    """
    Infographie hybride — grande statistique + barres segmentées + contexte.
    v0 : valeur hero gauche + barres horizontales droite (label + bar + %).
    v1 : valeur hero centré haut + grille de barres verticales en bas.
    v2 : donut visuel simplifié (rectangles empilés) + légende + contexte.
    content: {title, section_label?, subtitle?, value, label?, context?,
              bars:[{label, percent, color?}], footer}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accent2 = theme.get('accent2', 'ED0000')
    accents = tp.get('accent_cycle', [
        theme.get('accent1', '009CEA'),
        theme.get('accent2', 'ED0000'),
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
    ])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    value   = str(content.get('value', ''))
    label   = content.get('label', '')
    context = content.get('context', '')
    bars    = content.get('bars', [])
    v       = _v4_variant(content, 3, tp.get('seed', 0))

    if not value and not bars:
        return slide

    if v == 0:
        # Hero valeur gauche + barres horizontales droite
        hero_w  = _LY.CW * 0.38
        bars_x  = _LY.CL + hero_w + _LY.GAP_LG
        bars_w  = _LY.CR - bars_x
        y_mid   = (_LY.CT + _LY.CB) / 2

        _h2_rounded_rect(slide, left=_LY.CL, top=_LY.CT,
                          width=hero_w, height=_LY.CB - _LY.CT,
                          color='F0F0F0', radius=_LY.R_MD)
        _h2_rect(slide, left=_LY.CL, top=_LY.CT, width=hero_w, height=0.07, color=accent1)
        _h2_text(slide, value,
                 left=_LY.CL, top=y_mid - 0.60, width=hero_w, height=1.10,
                 font=font, size_pt=52, color=accent1, bold=True, align='center')
        if label:
            _h2_rect(slide, left=_LY.CL + (hero_w - 1.6) / 2, top=y_mid + 0.56,
                     width=1.6, height=_LY.BAR_H, color=accent2)
            _h2_text(slide, label,
                     left=_LY.CL, top=y_mid + 0.68, width=hero_w, height=0.44,
                     font=font, size_pt=13, color=dk1, bold=False, align='center')
        if context:
            _h2_text(slide, context,
                     left=_LY.CL + 0.15, top=y_mid + (1.20 if label else 0.62),
                     width=hero_w - 0.30, height=0.72,
                     font=font, size_pt=9, color='666666',
                     bold=False, align='center', line_spacing=1.3)

        n_bars  = min(len(bars), 6)
        bar_gap = (_LY.CB - _LY.CT - 0.10) / max(n_bars, 1)
        bar_h   = min(0.36, bar_gap - 0.16)
        max_w   = bars_w - 0.10
        for i, b in enumerate(bars[:n_bars]):
            bl    = b.get('label', '') if isinstance(b, dict) else str(b)
            pct   = float(b.get('percent', b.get('pct', 50)) if isinstance(b, dict) else 50) / 100
            pct   = max(0.0, min(1.0, pct))
            color = b.get('color', accents[i % len(accents)]) if isinstance(b, dict) else accents[i % len(accents)]
            by    = _LY.CT + 0.10 + i * bar_gap
            _h2_text(slide, bl, left=bars_x, top=by,
                     width=bars_w - 0.60, height=bar_h,
                     font=font, size_pt=10, color=dk1, bold=False, align='left')
            _h2_rect(slide, left=bars_x, top=by + bar_h + 0.04,
                     width=max_w, height=0.14, color='EEEEEE')
            _h2_rect(slide, left=bars_x, top=by + bar_h + 0.04,
                     width=max(0.08, max_w * pct), height=0.14, color=color)
            _h2_text(slide, f'{int(pct * 100)} %',
                     left=bars_x + max_w + 0.04, top=by + bar_h + 0.02,
                     width=0.52, height=0.20,
                     font=font, size_pt=9, color=color, bold=True, align='left')
        return slide

    if v == 1:
        # Valeur hero centré haut + barres verticales en bas
        hero_h  = (_LY.CB - _LY.CT) * 0.38
        bars_y  = _LY.CT + hero_h + _LY.GAP_MD
        bars_h  = _LY.CB - bars_y

        _h2_text(slide, value,
                 left=_LY.CL, top=_LY.CT + 0.10, width=_LY.CW, height=hero_h - 0.30,
                 font=font, size_pt=60, color=accent1, bold=True, align='center')
        if label:
            _h2_rect(slide, left=(_LY.CL + _LY.CR - 2.0) / 2, top=_LY.CT + hero_h - 0.22,
                     width=2.0, height=_LY.BAR_H, color=accent2)
            _h2_text(slide, label,
                     left=_LY.CL, top=_LY.CT + hero_h - 0.10, width=_LY.CW, height=0.36,
                     font=font, size_pt=14, color=dk1, bold=False, align='center')

        n_bars  = min(len(bars), 6)
        bar_w   = (_LY.CW - _LY.GAP_SM * (n_bars - 1)) / max(n_bars, 1)
        max_bar = bars_h - 0.36
        for i, b in enumerate(bars[:n_bars]):
            bl    = b.get('label', '') if isinstance(b, dict) else str(b)
            pct   = float(b.get('percent', b.get('pct', 50)) if isinstance(b, dict) else 50) / 100
            pct   = max(0.0, min(1.0, pct))
            color = b.get('color', accents[i % len(accents)]) if isinstance(b, dict) else accents[i % len(accents)]
            bx    = _LY.CL + i * (bar_w + _LY.GAP_SM)
            bh    = max(0.08, max_bar * pct)
            by    = _LY.CB - 0.30 - bh
            _h2_rect(slide, left=bx, top=bars_y, width=bar_w, height=max_bar, color='EEEEEE')
            _h2_rect(slide, left=bx, top=by, width=bar_w, height=bh, color=color)
            _h2_text(slide, f'{int(pct * 100)}%', left=bx, top=by - 0.26,
                     width=bar_w, height=0.24,
                     font=font, size_pt=9, color=color, bold=True, align='center')
            _h2_text(slide, bl, left=bx, top=_LY.CB - 0.28,
                     width=bar_w, height=0.26,
                     font=font, size_pt=8, color=dk1, bold=False, align='center')
        return slide

    # Variante 2 : blocs empilés (donut simplifié) + légende + contexte
    n_bars   = min(len(bars), 5)
    total    = sum(float(b.get('percent', b.get('pct', 20)) if isinstance(b, dict) else 20) for b in bars[:n_bars])
    total    = total or 100.0
    stack_x  = _LY.CL
    stack_w  = 1.60
    stack_h  = _LY.CB - _LY.CT
    leg_x    = stack_x + stack_w + _LY.GAP_LG
    leg_w    = _LY.CW * 0.44
    ctx_x    = leg_x + leg_w + _LY.GAP_LG
    ctx_w    = _LY.CR - ctx_x

    # Barre empilée verticale
    _h2_rect(slide, left=stack_x, top=_LY.CT, width=stack_w, height=stack_h, color='F0F0F0')
    cursor = _LY.CT
    for i, b in enumerate(bars[:n_bars]):
        pct   = float(b.get('percent', b.get('pct', 20)) if isinstance(b, dict) else 20)
        color = b.get('color', accents[i % len(accents)]) if isinstance(b, dict) else accents[i % len(accents)]
        seg_h = stack_h * (pct / total)
        _h2_rect(slide, left=stack_x, top=cursor, width=stack_w, height=seg_h, color=color)
        cursor += seg_h

    # Légende
    leg_item_h = (_LY.CB - _LY.CT) / max(n_bars, 1)
    for i, b in enumerate(bars[:n_bars]):
        bl    = b.get('label', '') if isinstance(b, dict) else str(b)
        pct   = float(b.get('percent', b.get('pct', 20)) if isinstance(b, dict) else 20)
        color = b.get('color', accents[i % len(accents)]) if isinstance(b, dict) else accents[i % len(accents)]
        ly    = _LY.CT + i * leg_item_h + (leg_item_h - 0.40) / 2
        _h2_rounded_rect(slide, left=leg_x, top=ly + 0.04,
                          width=0.22, height=0.22, color=color, radius=0.04)
        _h2_text(slide, bl, left=leg_x + 0.28, top=ly,
                 width=leg_w - 0.36, height=0.30,
                 font=font, size_pt=10, color=dk1, bold=False, align='left')
        _h2_text(slide, f'{int(pct)}%', left=leg_x + leg_w - 0.50, top=ly,
                 width=0.48, height=0.30,
                 font=font, size_pt=11, color=color, bold=True, align='right')

    # Valeur hero + contexte droite
    if value:
        _h2_text(slide, value,
                 left=ctx_x, top=(_LY.CT + _LY.CB) / 2 - 0.60, width=ctx_w, height=1.10,
                 font=font, size_pt=48, color=accent1, bold=True, align='center')
    if label:
        _h2_text(slide, label,
                 left=ctx_x, top=(_LY.CT + _LY.CB) / 2 + 0.56, width=ctx_w, height=0.40,
                 font=font, size_pt=12, color=dk1, bold=False, align='center')
    if context:
        _h2_text(slide, context,
                 left=ctx_x, top=(_LY.CT + _LY.CB) / 2 + 1.02, width=ctx_w, height=0.80,
                 font=font, size_pt=9, color='666666',
                 bold=False, align='center', line_spacing=1.3)

    return slide


# ── V4 Schémas (timeline, process, KPI grid, funnel) ────────────────────────

def layout_timeline_v4(prs: Presentation, content: dict, tp: dict):
    """
    Frise chronologique — 2 variantes visuelles déterministes.
    v0 : frise horizontale, jalons alternés dessus/dessous.
    v1 : frise verticale, ligne accent1 à gauche, blocs à droite.
    """
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accents = tp.get('accent_cycle', [
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        accent1,
    ])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    steps = content.get('steps', [])
    n = min(len(steps), 6)
    if n == 0:
        return slide

    v = _v4_variant(content, 2, tp.get('seed', 0))

    if v == 0:
        # Variante 0 : frise horizontale jalons alternés
        axis_y   = (_LY.CT + _LY.CB) / 2
        ax_left  = _LY.CL + 0.1
        ax_right = _LY.CR - 0.1
        ax_w     = ax_right - ax_left
        _h2_rect(slide, left=ax_left, top=axis_y - 0.022,
                 width=ax_w, height=0.044, color=accent1)

        slot_w  = ax_w / n
        r_circ  = _LY.R_CIRC - 0.06
        gap     = _LY.GAP_MD
        blk_w   = slot_w - 0.2
        above_h = axis_y - r_circ - gap - _LY.CT
        below_h = _LY.CB - (axis_y + r_circ + gap)

        for i, step in enumerate(steps[:n]):
            sx    = ax_left + (i + 0.5) * slot_w
            color = accents[i % len(accents)]
            date  = step.get('date', '') if isinstance(step, dict) else ''
            title = step.get('title', '') if isinstance(step, dict) else str(step)
            body  = step.get('body', '') if isinstance(step, dict) else ''

            _h2_circle(slide, cx=sx, cy=axis_y, r=r_circ, color=color)
            above = (i % 2 == 0)
            if above:
                bx, by, bh = sx - blk_w / 2, _LY.CT + 0.05, above_h
                _h2_rect(slide, left=sx - 0.012, top=by + bh,
                         width=0.024, height=axis_y - r_circ - (by + bh), color=color)
            else:
                bx, by, bh = sx - blk_w / 2, axis_y + r_circ + gap, below_h
                _h2_rect(slide, left=sx - 0.012, top=axis_y + r_circ,
                         width=0.024, height=gap, color=color)

            _h2_rounded_rect(slide, left=bx, top=by,
                              width=blk_w, height=bh, color='F0F0F0', radius=0.04)
            y_cur = by + 0.1
            if date:
                pill_h = 0.28
                _h2_rounded_rect(slide, left=bx + 0.08, top=y_cur,
                                  width=blk_w - 0.16, height=pill_h, color=color, radius=0.5)
                _h2_text(slide, date,
                         left=bx + 0.08, top=y_cur, width=blk_w - 0.16, height=pill_h,
                         font=font, size_pt=9, color='FFFFFF', bold=True, align='center')
                y_cur += pill_h + 0.1
            _h2_text(slide, title,
                     left=bx + 0.08, top=y_cur, width=blk_w - 0.16,
                     height=max(0.28, bh - (y_cur - by) - 0.08),
                     font=font, size_pt=11, color=dk1,
                     bold=True, align='center', line_spacing=1.1)
            if body and bh - (y_cur - by) > 0.55:
                _h2_text(slide, body,
                         left=bx + 0.08, top=y_cur + 0.38,
                         width=blk_w - 0.16,
                         height=max(0.28, bh - (y_cur - by) - 0.46),
                         font=font, size_pt=9, color='555555',
                         align='center', line_spacing=1.1)
    else:
        # Variante 1 : frise verticale — ligne accent1 à gauche, blocs à droite
        axis_x  = _LY.CL + 0.50
        blk_x   = axis_x + 0.55
        blk_w   = _LY.CR - blk_x - 0.05
        step_h  = (_LY.CB - _LY.CT) / n
        r_circ  = _LY.R_CIRC - 0.05
        # Ligne verticale pleine hauteur
        _h2_rect(slide, left=axis_x - 0.020, top=_LY.CT,
                 width=0.040, height=_LY.CB - _LY.CT, color=accent1)

        for i, step in enumerate(steps[:n]):
            sy    = _LY.CT + (i + 0.5) * step_h
            color = accents[i % len(accents)]
            date  = step.get('date', '') if isinstance(step, dict) else ''
            title = step.get('title', '') if isinstance(step, dict) else str(step)
            body  = step.get('body', '') if isinstance(step, dict) else ''

            # Cercle sur la ligne
            _h2_circle(slide, cx=axis_x, cy=sy, r=r_circ, color=color)

            # Ligne de connexion horizontale
            _h2_rect(slide, left=axis_x + r_circ, top=sy - 0.018,
                     width=blk_x - (axis_x + r_circ) - 0.04, height=0.036, color=color)

            # Bloc texte
            bh  = step_h - _LY.GAP_SM
            by  = sy - bh / 2
            _h2_rounded_rect(slide, left=blk_x, top=by,
                              width=blk_w, height=bh, color='F0F0F0', radius=0.04)
            y_cur = by + 0.10
            if date:
                pill_h = 0.25
                _h2_rounded_rect(slide, left=blk_x + 0.08, top=y_cur,
                                  width=min(blk_w * 0.55, 2.5), height=pill_h,
                                  color=color, radius=0.5)
                _h2_text(slide, date,
                         left=blk_x + 0.08, top=y_cur,
                         width=min(blk_w * 0.55, 2.5), height=pill_h,
                         font=font, size_pt=8, color='FFFFFF', bold=True, align='center')
                y_cur += pill_h + 0.06
            _h2_text(slide, title,
                     left=blk_x + 0.10, top=y_cur, width=blk_w - 0.18,
                     height=min(0.44, bh - (y_cur - by) - 0.06),
                     font=font, size_pt=12, color=dk1,
                     bold=True, align='left', line_spacing=1.1)
            if body and bh - (y_cur - by) > 0.60:
                _h2_text(slide, body,
                         left=blk_x + 0.10, top=y_cur + 0.42,
                         width=blk_w - 0.18,
                         height=max(0.28, bh - (y_cur - by) - 0.50),
                         font=font, size_pt=9, color='555555',
                         align='left', line_spacing=1.1)

    return slide


def layout_processflow_v4(prs: Presentation, content: dict, tp: dict):
    """
    Flux de processus — 2 variantes visuelles déterministes.
    v0 : horizontal, boîtes accent pleine zone CT→CB + flèches ▶.
    v1 : vertical, étapes empilées pleine largeur + flèches ↓.
    """
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accents = tp.get('accent_cycle', [
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        accent1,
    ])
    W = tp.get('W', 13.33)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    steps = content.get('steps', [])
    n = min(len(steps), 6)
    if n == 0:
        return slide

    v = _v4_variant(content, 2, tp.get('seed', 0))

    if v == 0:
        # Variante 0 : horizontal, boîtes pleine zone CT→CB
        x_start = _LY.CL - 0.05
        x_end   = _LY.CR + 0.05
        y_mid   = (_LY.CT + _LY.CB) / 2
        box_h   = _LY.CB - _LY.CT          # pleine zone contenu
        box_y   = _LY.CT
        arrow_w = 0.38
        total_w = x_end - x_start
        box_w   = (total_w - (n - 1) * arrow_w) / n

        for i, step in enumerate(steps[:n]):
            color = accents[i % len(accents)]
            bx    = x_start + i * (box_w + arrow_w)
            title = step.get('title', '') if isinstance(step, dict) else str(step)
            body  = step.get('body', '') if isinstance(step, dict) else ''

            _h2_rounded_rect(slide, left=bx, top=box_y,
                              width=box_w, height=box_h, color=color, radius=0.07)
            _h2_text(slide, str(i + 1),
                     left=bx + 0.10, top=box_y + 0.10,
                     width=0.35, height=0.33,
                     font=font, size_pt=11, color='FFFFFF',
                     bold=True, align='left')
            _h2_text(slide, title,
                     left=bx + 0.12, top=box_y + 0.48,
                     width=box_w - 0.22, height=0.55,
                     font=font, size_pt=12, color='FFFFFF',
                     bold=True, align='left', line_spacing=1.1)
            if body:
                _h2_text(slide, body,
                         left=bx + 0.12, top=box_y + 1.10,
                         width=box_w - 0.22, height=box_h - 1.22,
                         font=font, size_pt=10, color='FFFFFF',
                         bold=False, align='left', line_spacing=1.15)
            if i < n - 1:
                ax = bx + box_w + 0.03
                _h2_text(slide, '▶',
                         left=ax, top=y_mid - 0.20,
                         width=arrow_w - 0.06, height=0.40,
                         font=font, size_pt=16, color=accents[(i + 1) % len(accents)],
                         bold=True, align='center')
    else:
        # Variante 1 : vertical, étapes pleine largeur
        arrow_h = 0.30
        total_h = _LY.CB - _LY.CT
        box_h   = (total_h - (n - 1) * arrow_h) / n

        for i, step in enumerate(steps[:n]):
            color = accents[i % len(accents)]
            by    = _LY.CT + i * (box_h + arrow_h)
            title = step.get('title', '') if isinstance(step, dict) else str(step)
            body  = step.get('body', '') if isinstance(step, dict) else ''

            _h2_rounded_rect(slide, left=_LY.CL, top=by,
                              width=_LY.CW, height=box_h, color=color, radius=0.06)
            # Numéro
            _h2_text(slide, str(i + 1),
                     left=_LY.CL + 0.12, top=by + (box_h - 0.50) / 2,
                     width=0.50, height=0.50,
                     font=font, size_pt=18, color='FFFFFF',
                     bold=True, align='center')
            # Titre
            txt_y = by + (box_h - (0.42 if not body else 0.82)) / 2
            _h2_text(slide, title,
                     left=_LY.CL + 0.78, top=txt_y,
                     width=_LY.CW - 0.90, height=0.42,
                     font=font, size_pt=13, color='FFFFFF',
                     bold=True, align='left', line_spacing=1.1)
            if body:
                _h2_text(slide, body,
                         left=_LY.CL + 0.78, top=txt_y + 0.44,
                         width=_LY.CW - 0.90, height=0.38,
                         font=font, size_pt=10, color='FFFFFF',
                         bold=False, align='left', line_spacing=1.1)
            # Flèche ↓
            if i < n - 1:
                ay = by + box_h + 0.01
                _h2_text(slide, '▼',
                         left=_LY.CL + _LY.CW / 2 - 0.20, top=ay,
                         width=0.40, height=arrow_h - 0.02,
                         font=font, size_pt=12, color=accents[(i + 1) % len(accents)],
                         bold=True, align='center')

    return slide


def layout_kpi_grid_v4(prs: Presentation, content: dict, tp: dict):
    """
    Grille de KPIs — 2 variantes visuelles déterministes.
    v0 : valeur centrée + trait accent2 + label (style classique).
    v1 : chaque KPI dans une carte arrondie fond léger, valeur accent en grand.
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accent2 = theme.get('accent2', 'ED0000')
    accents = tp.get('accent_cycle', [
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        accent1,
    ])
    W = tp.get('W', 13.33)
    H = tp.get('H', 7.50)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    kpis = content.get('kpis', [])
    n = min(len(kpis), 6)
    if n == 0:
        return slide

    n_cols  = min(n, 3)
    n_rows  = (n + n_cols - 1) // n_cols
    x_start = _LY.CL
    x_end   = W - 0.45
    y_start = _LY.CT
    y_end   = _LY.CB
    cell_w  = (x_end - x_start) / n_cols
    cell_h  = (y_end - y_start) / n_rows
    v       = _v4_variant(content, 5, tp.get('seed', 0))

    _light_bg = ['F0F0F0', 'F4F4F4', 'EEEEEE']

    # v3 : KPI héros central (premier) + petits en dessous
    if v == 3 and n >= 2:
        hero = kpis[0]
        rest = kpis[1:n]
        h_val   = str(hero.get('value', '')) if isinstance(hero, dict) else str(hero)
        h_lbl   = hero.get('label', '') if isinstance(hero, dict) else ''
        h_sub   = hero.get('sublabel', '') if isinstance(hero, dict) else ''
        h_color = accents[0]
        # Zone héros
        _h2_rounded_rect(slide, left=_LY.CL + 1.5, top=_LY.CT,
                          width=_LY.CW - 3.0, height=2.20,
                          color='F0F0F0', radius=_LY.R_SM)
        _h2_rect(slide, left=_LY.CL + 1.5, top=_LY.CT,
                 width=_LY.CW - 3.0, height=0.06, color=h_color)
        _h2_text(slide, h_val,
                 left=_LY.CL + 1.5, top=_LY.CT + 0.12,
                 width=_LY.CW - 3.0, height=1.0,
                 font=font, size_pt=52, color=h_color, bold=True, align='center')
        if h_lbl:
            _h2_text(slide, h_lbl,
                     left=_LY.CL + 1.5, top=_LY.CT + 1.18,
                     width=_LY.CW - 3.0, height=0.38,
                     font=font, size_pt=_LY.T_LABEL, color=dk1,
                     bold=True, align='center')
        if h_sub:
            _h2_text(slide, h_sub,
                     left=_LY.CL + 1.5, top=_LY.CT + 1.55,
                     width=_LY.CW - 3.0, height=0.30,
                     font=font, size_pt=_LY.T_SMALL, color='888888',
                     bold=False, align='center')
        # KPIs secondaires en ligne
        nr = len(rest)
        sw = _LY.CW / max(nr, 1)
        sy = _LY.CT + 2.35
        sh = _LY.CB - sy
        for j, kpi in enumerate(rest[:5]):
            cx2  = _LY.CL + j * sw
            val2 = str(kpi.get('value', '')) if isinstance(kpi, dict) else str(kpi)
            lbl2 = kpi.get('label', '') if isinstance(kpi, dict) else ''
            sub2 = kpi.get('sublabel', '') if isinstance(kpi, dict) else ''
            c2   = accents[(j + 1) % len(accents)]
            _h2_rounded_rect(slide, left=cx2 + _LY.GAP_SM, top=sy,
                              width=sw - _LY.GAP_SM * 2, height=sh,
                              color=_light_bg[j % len(_light_bg)], radius=_LY.R_SM)
            _h2_rect(slide, left=cx2 + _LY.GAP_SM, top=sy,
                     width=sw - _LY.GAP_SM * 2, height=0.05, color=c2)
            _h2_text(slide, val2,
                     left=cx2 + _LY.GAP_SM, top=sy + 0.10,
                     width=sw - _LY.GAP_SM * 2, height=sh * 0.52,
                     font=font, size_pt=_LY.T_KPI - 6, color=c2,
                     bold=True, align='center')
            if lbl2:
                _h2_text(slide, lbl2,
                         left=cx2 + _LY.GAP_SM, top=sy + sh * 0.56,
                         width=sw - _LY.GAP_SM * 2, height=0.30,
                         font=font, size_pt=_LY.T_SMALL, color=dk1,
                         bold=True, align='center')
            if sub2:
                _h2_text(slide, sub2,
                         left=cx2 + _LY.GAP_SM, top=sy + sh * 0.78,
                         width=sw - _LY.GAP_SM * 2, height=0.22,
                         font=font, size_pt=8, color='888888',
                         bold=False, align='center')
        return slide

    # v4 : barres de progression horizontales sous chaque valeur
    if v == 4:
        for i, kpi in enumerate(kpis[:n]):
            col_i = i % n_cols
            row_i = i // n_cols
            cx  = x_start + col_i * cell_w
            cy  = y_start + row_i * cell_h
            val      = str(kpi.get('value', '')) if isinstance(kpi, dict) else str(kpi)
            label    = kpi.get('label', '') if isinstance(kpi, dict) else ''
            sublabel = kpi.get('sublabel', '') if isinstance(kpi, dict) else ''
            pct_raw  = kpi.get('percent', kpi.get('pct', 0)) if isinstance(kpi, dict) else 0
            try:
                pct = min(1.0, max(0.0, float(str(pct_raw).replace('%', '')) / 100))
            except Exception:
                pct = 0.5  # default 50% si non fourni
            color = accents[i % len(accents)]
            mid_y = cy + cell_h / 2
            _h2_text(slide, val,
                     left=cx + _LY.PAD, top=mid_y - 0.72,
                     width=cell_w - _LY.PAD * 2, height=0.72,
                     font=font, size_pt=_LY.T_KPI, color=color,
                     bold=True, align='center')
            _h2_text(slide, label,
                     left=cx + _LY.PAD, top=mid_y + 0.04,
                     width=cell_w - _LY.PAD * 2, height=0.36,
                     font=font, size_pt=_LY.T_LABEL, color=dk1,
                     bold=False, align='center')
            # Barre de progression
            bar_x = cx + _LY.PAD
            bar_w = cell_w - _LY.PAD * 2
            bar_y = mid_y + 0.46
            bar_h = 0.14
            _h2_rounded_rect(slide, left=bar_x, top=bar_y,
                              width=bar_w, height=bar_h, color='E0E0E0', radius=0.07)
            if pct > 0:
                _h2_rounded_rect(slide, left=bar_x, top=bar_y,
                                  width=max(0.1, bar_w * pct), height=bar_h,
                                  color=color, radius=0.07)
            if sublabel:
                _h2_text(slide, sublabel,
                         left=cx + _LY.PAD, top=mid_y + 0.66,
                         width=cell_w - _LY.PAD * 2, height=0.26,
                         font=font, size_pt=_LY.T_SMALL, color='888888',
                         bold=False, align='center')
            if col_i < n_cols - 1:
                _h2_rect(slide, left=cx + cell_w - 0.005, top=cy + 0.1,
                         width=0.01, height=cell_h - 0.2, color='DDDDDD')
        return slide

    # v0, v1, v2 : boucle commune
    for i, kpi in enumerate(kpis[:n]):
        col_i = i % n_cols
        row_i = i // n_cols
        cx  = x_start + col_i * cell_w
        cy  = y_start + row_i * cell_h
        val      = str(kpi.get('value', '')) if isinstance(kpi, dict) else str(kpi)
        label    = kpi.get('label', '') if isinstance(kpi, dict) else ''
        sublabel = kpi.get('sublabel', '') if isinstance(kpi, dict) else ''
        color    = accents[i % len(accents)]
        mid_y    = cy + cell_h / 2

        if v == 0:
            _h2_text(slide, val,
                     left=cx + _LY.PAD, top=mid_y - 0.78,
                     width=cell_w - _LY.PAD * 2, height=0.80,
                     font=font, size_pt=_LY.T_KPI, color=color,
                     bold=True, align='center')
            sep_w = min(cell_w * 0.45, 1.8)
            _h2_rect(slide, left=cx + (cell_w - sep_w) / 2, top=mid_y + 0.06,
                     width=sep_w, height=_LY.BAR_H, color=accent2)
            _h2_text(slide, label,
                     left=cx + _LY.PAD, top=mid_y + 0.14,
                     width=cell_w - _LY.PAD * 2, height=0.38,
                     font=font, size_pt=_LY.T_LABEL, color=dk1,
                     bold=False, align='center')
            if sublabel:
                _h2_text(slide, sublabel,
                         left=cx + _LY.PAD, top=mid_y + 0.52,
                         width=cell_w - _LY.PAD * 2, height=0.3,
                         font=font, size_pt=_LY.T_SMALL, color='888888',
                         bold=False, align='center')
            if col_i < n_cols - 1:
                _h2_rect(slide, left=cx + cell_w - 0.005, top=cy + 0.1,
                         width=0.01, height=cell_h - 0.2, color='DDDDDD')
        elif v == 1:
            pad = _LY.GAP_SM
            cw  = cell_w - pad * 2
            ch  = cell_h - pad * 2
            _h2_rounded_rect(slide, left=cx + pad, top=cy + pad,
                              width=cw, height=ch,
                              color=_light_bg[i % len(_light_bg)], radius=_LY.R_SM)
            _h2_rect(slide, left=cx + pad, top=cy + pad,
                     width=cw, height=0.06, color=color)
            _h2_text(slide, val,
                     left=cx + pad, top=cy + pad + 0.12,
                     width=cw, height=ch * 0.52,
                     font=font, size_pt=_LY.T_KPI, color=color,
                     bold=True, align='center')
            sep_w = min(cw * 0.5, 1.6)
            _h2_rect(slide, left=cx + pad + (cw - sep_w) / 2,
                     top=cy + pad + ch * 0.55,
                     width=sep_w, height=_LY.BAR_H, color=color)
            _h2_text(slide, label,
                     left=cx + pad, top=cy + pad + ch * 0.60,
                     width=cw, height=0.38,
                     font=font, size_pt=_LY.T_LABEL, color=dk1,
                     bold=False, align='center')
            if sublabel:
                _h2_text(slide, sublabel,
                         left=cx + pad, top=cy + pad + ch * 0.77,
                         width=cw, height=0.28,
                         font=font, size_pt=_LY.T_SMALL, color='888888',
                         bold=False, align='center')
        else:  # v == 2
            val_w = cell_w * 0.38
            lbl_x = cx + val_w + _LY.PAD
            lbl_w = cell_w - val_w - _LY.PAD * 2
            _h2_rect(slide, left=cx, top=cy, width=cell_w, height=cell_h - 0.05,
                     color=_light_bg[i % len(_light_bg)])
            _h2_rect(slide, left=cx, top=cy, width=0.05, height=cell_h - 0.05, color=color)
            _h2_text(slide, val,
                     left=cx + 0.12, top=mid_y - 0.42,
                     width=val_w - 0.15, height=0.84,
                     font=font, size_pt=_LY.T_KPI - 4, color=color,
                     bold=True, align='center')
            _h2_text(slide, label,
                     left=lbl_x, top=mid_y - 0.22,
                     width=lbl_w, height=0.44,
                     font=font, size_pt=_LY.T_LABEL, color=dk1,
                     bold=True, align='left')
            if sublabel:
                _h2_text(slide, sublabel,
                         left=lbl_x, top=mid_y + 0.22,
                         width=lbl_w, height=0.28,
                         font=font, size_pt=_LY.T_SMALL, color='888888',
                         bold=False, align='left')

    return slide


def layout_funnel_v4(prs: Presentation, content: dict, tp: dict):
    """
    Entonnoir (funnel) — 2 variantes visuelles déterministes.
    v0 : rectangles décroissants centrés, accent de plus en plus foncé.
    v1 : barres horizontales proportionnelles, label gauche + valeur dans la barre.
    """
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accents = tp.get('accent_cycle', [
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        accent1,
    ])
    W = tp.get('W', 13.33)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    steps = content.get('steps', content.get('items', []))
    n = min(len(steps), 5)
    if n == 0:
        return slide

    v = _v4_variant(content, 2, tp.get('seed', 0))

    # Couleur : shades d'accent1
    try:
        r0 = int(accent1[0:2], 16)
        g0 = int(accent1[2:4], 16)
        b0 = int(accent1[4:6], 16)
    except Exception:
        r0, g0, b0 = 0, 156, 234

    def _step_data(step):
        if isinstance(step, dict):
            return step.get('label', step.get('title', '')), str(step.get('value', ''))
        return str(step), ''

    if v == 0:
        # Variante 0 : entonnoir centré décroissant
        total_h = _LY.CB - _LY.CT
        level_h = total_h / n - _LY.GAP_XS
        max_w   = W - 1.2
        min_w   = max_w * 0.35

        for i, step in enumerate(steps[:n]):
            label, value = _step_data(step)
            ratio = 1.0 - i * (1.0 - min_w / max_w) / max(n - 1, 1)
            lw    = max_w * ratio
            lx    = (W - lw) / 2
            ly    = _LY.CT + i * (level_h + _LY.GAP_XS)
            dark  = 1.0 - i * 0.18 / max(n - 1, 1)
            color = f'{int(r0*dark):02X}{int(g0*dark):02X}{int(b0*dark):02X}'
            _h2_rect(slide, left=lx, top=ly, width=lw, height=level_h, color=color)
            txt = f'{label}  {value}' if value else label
            _h2_text(slide, txt,
                     left=lx + 0.2, top=ly + (level_h - 0.38) / 2,
                     width=lw - 0.4, height=0.38,
                     font=font, size_pt=13, color='FFFFFF',
                     bold=True, align='center')
    else:
        # Variante 1 : barres horizontales proportionnelles
        total_h  = _LY.CB - _LY.CT
        level_h  = total_h / n - _LY.GAP_XS
        lbl_w    = 2.80
        bar_x    = _LY.CL + lbl_w + 0.20
        max_bar  = _LY.CR - bar_x - 0.10

        for i, step in enumerate(steps[:n]):
            label, value = _step_data(step)
            ly    = _LY.CT + i * (level_h + _LY.GAP_XS)
            # Ratio décroissant (ou use valeur si numeric)
            try:
                ratio = max(0.15, float(value) / 100) if value else 1.0 - i * 0.18 / max(n - 1, 1)
            except ValueError:
                ratio = 1.0 - i * 0.18 / max(n - 1, 1)
            bw    = max_bar * min(ratio, 1.0)
            color = accents[i % len(accents)]

            # Fond gris pleine largeur (barre de fond)
            _h2_rect(slide, left=bar_x, top=ly + _LY.GAP_XS,
                     width=max_bar, height=level_h - _LY.GAP_XS * 2, color='EEEEEE')
            # Barre proportionnelle
            _h2_rect(slide, left=bar_x, top=ly + _LY.GAP_XS,
                     width=bw, height=level_h - _LY.GAP_XS * 2, color=color)
            # Label gauche
            _h2_text(slide, label,
                     left=_LY.CL, top=ly + (level_h - 0.40) / 2,
                     width=lbl_w, height=0.40,
                     font=font, size_pt=12, color=dk1,
                     bold=True, align='right')
            # Valeur dans la barre
            if value:
                _h2_text(slide, value,
                         left=bar_x + 0.10, top=ly + (level_h - 0.38) / 2,
                         width=bw - 0.15, height=0.38,
                         font=font, size_pt=12, color='FFFFFF',
                         bold=True, align='right')

    return slide


# ── V4 Graphiques natifs PowerPoint ─────────────────────────────────────────

def _chart_series_colors(theme: dict) -> list:
    """Retourne la liste des couleurs accent du thème pour les séries."""
    return [
        theme.get('accent1', '009CEA'),
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        theme.get('accent2', 'ED0000'),
        theme.get('accent5', '7030A0'),
        theme.get('accent6', '0070C0'),
    ]


# Zone analyse sous graphique : CT→CHART_BOT puis ANALYSIS_TOP→CB
_CHART_H       = 3.50   # hauteur graphique (top = CT ≈ 1.60)
_ANALYSIS_TOP  = _LY.CT + _CHART_H + 0.15   # ≈ 5.20"
_ANALYSIS_H    = _LY.CB - _ANALYSIS_TOP      # ≈ 1.75"


def _add_chart_analysis(slide, content: dict, tp: dict) -> None:
    """
    Bloc d'analyse / interprétation sous un graphique.
    Cherche 'analysis' > 'interpretation' > 'insight' > 'body'.
    Fond E8EEF5, rayure accent1 gauche, texte dk1.
    N'affiche rien si aucun texte n'est trouvé.
    """
    analysis = (
        content.get('analysis') or
        content.get('interpretation') or
        content.get('insight') or
        content.get('body', '')
    )
    if not analysis:
        return
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    accent = theme.get('accent1', '009CEA')
    _h2_rounded_rect(slide,
                     left=_LY.CL, top=_ANALYSIS_TOP,
                     width=_LY.CW, height=_ANALYSIS_H,
                     color='F0F0F0', radius=_LY.R_SM)
    _h2_rect(slide,
             left=_LY.CL, top=_ANALYSIS_TOP,
             width=0.06, height=_ANALYSIS_H, color=accent)
    _h2_text(slide, analysis,
             left=_LY.CL + 0.18, top=_ANALYSIS_TOP + 0.12,
             width=_LY.CW - 0.24, height=_ANALYSIS_H - 0.20,
             font=font, size_pt=_LY.T_SMALL, color=dk1,
             bold=False, align='left', line_spacing=1.2)


def h2_bar_chart(slide, left: float, top: float, width: float, height: float,
                 categories: list, series: list, font: str, theme: dict):
    """
    Graphique COLUMN_CLUSTERED natif PowerPoint.
    series = [{name, values:[n]}, ...]
    Couleurs des séries = accent1/accent3/accent4.
    Légende en bas si >1 série. Data labels dk1 9pt.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in (categories or [])]
    for s in (series or []):
        vals = tuple(float(v) if v is not None else 0.0 for v in s.get('values', []))
        chart_data.add_series(s.get('name', ''), vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data,
    )
    chart = chart_shape.chart
    colors = _chart_series_colors(theme)

    for i, ser in enumerate(chart.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = _h2_parse_hex(colors[i % len(colors)])

    plot = chart.plots[0]
    plot.has_data_labels = True
    try:
        plot.data_labels.font.size = Pt(9)
        plot.data_labels.font.color.rgb = _h2_parse_hex(theme.get('dk1', '374649'))
    except Exception:
        pass

    if len(series) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False

    return chart_shape


def layout_barchart_v4(prs: Presentation, content: dict, tp: dict):
    """Slide graphique barres groupées (COLUMN_CLUSTERED) + bloc analyse."""
    slide = _blank_v4(prs, tp)
    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    categories = content.get('categories', [])
    series     = content.get('series', [])
    if not categories or not series:
        return slide

    h2_bar_chart(slide,
                 left=0.5, top=_LY.CT, width=12.0, height=_CHART_H,
                 categories=categories, series=series,
                 font=tp.get('font', 'Calibri'), theme=tp.get('theme', {}))
    _add_chart_analysis(slide, content, tp)
    return slide


def h2_line_chart(slide, left: float, top: float, width: float, height: float,
                  categories: list, series: list, font: str, theme: dict):
    """
    Graphique LINE_MARKERS natif PowerPoint.
    Lignes 2.5 pt, couleurs thème, légende en bas si >1 série.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in (categories or [])]
    for s in (series or []):
        vals = tuple(float(v) if v is not None else 0.0 for v in s.get('values', []))
        chart_data.add_series(s.get('name', ''), vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data,
    )
    chart = chart_shape.chart
    colors = _chart_series_colors(theme)

    for i, ser in enumerate(chart.series):
        c = _h2_parse_hex(colors[i % len(colors)])
        try:
            ser.format.line.color.rgb = c
            ser.format.line.width = Pt(2.5)
        except Exception:
            pass
        try:
            ser.marker.format.fill.solid()
            ser.marker.format.fill.fore_color.rgb = c
            ser.marker.size = 7
        except Exception:
            pass

    plot = chart.plots[0]
    plot.has_data_labels = False

    if len(series) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False

    return chart_shape


def layout_linechart_v4(prs: Presentation, content: dict, tp: dict):
    """Slide graphique en lignes (LINE_MARKERS) + bloc analyse."""
    slide = _blank_v4(prs, tp)
    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    categories = content.get('categories', [])
    series     = content.get('series', [])
    if not categories or not series:
        return slide

    h2_line_chart(slide,
                  left=0.5, top=_LY.CT, width=12.0, height=_CHART_H,
                  categories=categories, series=series,
                  font=tp.get('font', 'Calibri'), theme=tp.get('theme', {}))
    _add_chart_analysis(slide, content, tp)
    return slide


def h2_pie_chart(slide, left: float, top: float, width: float, height: float,
                 slices: list, font: str, theme: dict, doughnut: bool = False):
    """
    Graphique PIE ou DOUGHNUT natif PowerPoint.
    slices = [{label, value}, ...]
    Couleurs par segment = accent_cycle. Labels avec pourcentages.
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    if not slices:
        return None

    chart_data = CategoryChartData()
    chart_data.categories = [str(s.get('label', '')) for s in slices]
    chart_data.add_series('', tuple(float(s.get('value', 0)) for s in slices))

    chart_type = XL_CHART_TYPE.DOUGHNUT if doughnut else XL_CHART_TYPE.PIE

    chart_shape = slide.shapes.add_chart(
        chart_type,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data,
    )
    chart = chart_shape.chart

    # Couleurs des segments
    colors = _chart_series_colors(theme)
    try:
        for i, point in enumerate(chart.series[0].points):
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _h2_parse_hex(colors[i % len(colors)])
    except Exception:
        pass

    # Data labels avec pourcentages
    plot = chart.plots[0]
    plot.has_data_labels = True
    try:
        dl = plot.data_labels
        dl.number_format = '0%'
        dl.font.size = Pt(10)
        dl.font.color.rgb = _h2_parse_hex('FFFFFF')
    except Exception:
        pass

    # Légende
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    return chart_shape


def layout_piechart_v4(prs: Presentation, content: dict, tp: dict):
    """Slide graphique camembert / anneau + bloc analyse."""
    slide = _blank_v4(prs, tp)
    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    slices   = content.get('slices', [])
    doughnut = bool(content.get('doughnut', False))
    if not slices:
        return slide

    h2_pie_chart(slide,
                 left=1.5, top=_LY.CT, width=10.0, height=_CHART_H,
                 slices=slices, font=tp.get('font', 'Calibri'),
                 theme=tp.get('theme', {}), doughnut=doughnut)
    _add_chart_analysis(slide, content, tp)
    return slide


def h2_waterfall_chart(slide, left: float, top: float, width: float, height: float,
                       items: list, font: str, theme: dict):
    """
    Cascade financière (waterfall) via COLUMN_STACKED.
    Série base invisible + positifs en accent3 (vert) + négatifs en accent2 (rouge).
    items = [{label, value}, ...]
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    if not items:
        return None

    accent1 = theme.get('accent1', '009CEA')
    accent2 = theme.get('accent2', 'ED0000')
    accent3 = theme.get('accent3', '40A900')

    labels = [str(it.get('label', '')) for it in items]
    raw    = [float(it.get('value', 0)) for it in items]

    # Calcul base, gains, pertes
    base_vals = []
    gain_vals = []
    loss_vals = []
    running = 0.0

    for v in raw:
        if v >= 0:
            base_vals.append(running)
            gain_vals.append(v)
            loss_vals.append(0.0)
            running += v
        else:
            running += v
            base_vals.append(running)
            gain_vals.append(0.0)
            loss_vals.append(abs(v))

    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series('Base',    tuple(base_vals))
    chart_data.add_series('Hausse',  tuple(gain_vals))
    chart_data.add_series('Baisse',  tuple(loss_vals))

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data,
    )
    chart = chart_shape.chart

    series_list = list(chart.series)
    if len(series_list) >= 3:
        # Série base : invisible (fond blanc)
        series_list[0].format.fill.solid()
        series_list[0].format.fill.fore_color.rgb = _h2_parse_hex('FFFFFF')
        series_list[0].format.line.fill.background()

        # Hausse : accent3 (vert)
        series_list[1].format.fill.solid()
        series_list[1].format.fill.fore_color.rgb = _h2_parse_hex(accent3)

        # Baisse : accent2 (rouge)
        series_list[2].format.fill.solid()
        series_list[2].format.fill.fore_color.rgb = _h2_parse_hex(accent2)

    # Data labels visibles sur les séries Hausse et Baisse
    for i, ser in enumerate(series_list[1:], 1):
        try:
            ser.data_labels.font.size = Pt(9)
            ser.data_labels.font.color.rgb = _h2_parse_hex('FFFFFF')
        except Exception:
            pass

    chart.has_legend = False
    return chart_shape


def layout_waterfall_v4(prs: Presentation, content: dict, tp: dict):
    """Slide cascade financière (waterfall) + bloc analyse."""
    slide = _blank_v4(prs, tp)
    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    items = content.get('items', [])
    if not items:
        return slide

    h2_waterfall_chart(slide,
                       left=0.5, top=_LY.CT, width=12.0, height=_CHART_H,
                       items=items, font=tp.get('font', 'Calibri'),
                       theme=tp.get('theme', {}))
    _add_chart_analysis(slide, content, tp)
    return slide


def layout_radar_v4(prs: Presentation, content: dict, tp: dict):
    """
    Graphique radar (RADAR_MARKERS) natif PowerPoint + bloc analyse.
    content: {title, axes:[str], series:[{name, values:[n]}], analysis, footer}
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    slide = _blank_v4(prs, tp)
    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    axes   = content.get('axes', [])
    series = content.get('series', [])
    if not axes or not series:
        return slide

    chart_data = CategoryChartData()
    chart_data.categories = [str(a) for a in axes]
    for s in series:
        vals = tuple(float(v) if v is not None else 0.0 for v in s.get('values', []))
        chart_data.add_series(s.get('name', ''), vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR_MARKERS,
        Inches(1.5), Inches(_LY.CT), Inches(10.0), Inches(_CHART_H),
        chart_data,
    )
    chart  = chart_shape.chart
    colors = _chart_series_colors(tp.get('theme', {}))

    for i, ser in enumerate(chart.series):
        c = _h2_parse_hex(colors[i % len(colors)])
        try:
            ser.format.line.color.rgb = c
            ser.format.line.width = Pt(2.0)
        except Exception:
            pass

    if len(series) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False

    _add_chart_analysis(slide, content, tp)
    return slide


def layout_pyramid_v4(prs: Presentation, content: dict, tp: dict):
    """
    Pyramide hiérarchique (niveaux du haut vers le bas).
    Triangles/trapèzes de largeur croissante, couleur dégradée.
    content: {title, levels:[{label, body}], footer}
    """
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    accent1 = theme.get('accent1', '009CEA')
    W = tp.get('W', 13.33)
    H = tp.get('H', 7.50)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    levels = content.get('levels', [])
    n = min(len(levels), 5)
    if n == 0:
        return slide

    # Zone contenu
    y_start = _LY.CT
    y_end   = _LY.CB
    level_h = (y_end - y_start) / n - _LY.GAP_XS

    # Largeurs croissantes du haut vers le bas
    min_w = W * 0.2
    max_w = W - 1.2

    try:
        r0 = int(accent1[0:2], 16)
        g0 = int(accent1[2:4], 16)
        b0 = int(accent1[4:6], 16)
    except Exception:
        r0, g0, b0 = 0, 156, 234

    for i, level in enumerate(levels[:n]):
        ratio = min_w / max_w + i * (1.0 - min_w / max_w) / max(n - 1, 1)
        lw    = max_w * ratio
        lx    = (W - lw) / 2
        ly    = y_start + i * (level_h + 0.06)

        # Couleur de plus en plus foncée
        dark  = 0.85 - i * 0.15 / max(n - 1, 1)
        color = f'{int(r0 * dark):02X}{int(g0 * dark):02X}{int(b0 * dark):02X}'

        _h2_rect(slide, left=lx, top=ly, width=lw, height=level_h, color=color)

        label = level.get('label', '') if isinstance(level, dict) else str(level)
        body  = level.get('body', '')  if isinstance(level, dict) else ''
        txt   = f'{label}  —  {body}' if body else label

        _h2_text(slide, txt,
                 left=lx + 0.2, top=ly + (level_h - 0.38) / 2,
                 width=lw - 0.4, height=0.38,
                 font=font, size_pt=12, color='FFFFFF',
                 bold=True, align='center')

    return slide


def layout_cycle_v4(prs: Presentation, content: dict, tp: dict):
    """
    Cycle / roue avec étapes circulaires.
    Cercles accent_cycle disposés en cercle avec flèches entre eux.
    content: {title, steps:[{title, body}], footer}
    """
    import math
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    accents = tp.get('accent_cycle', [
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        theme.get('accent1', '009CEA'),
    ])
    W = tp.get('W', 13.33)
    H = tp.get('H', 7.50)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    steps = content.get('steps', [])
    n = min(len(steps), 6)
    if n == 0:
        return slide

    # Centre du cycle dans la zone contenu
    cx_center = W / 2
    cy_center = (_LY.CT + _LY.CB) / 2
    orbit_r   = min((W - 2.0) / 2, (H - 2.2) / 2) * 0.85
    node_r    = 0.55

    for i, step in enumerate(steps[:n]):
        angle = -math.pi / 2 + i * 2 * math.pi / n
        nx    = cx_center + orbit_r * math.cos(angle)
        ny    = cy_center + orbit_r * math.sin(angle)
        color = accents[i % len(accents)]

        title = step.get('title', '') if isinstance(step, dict) else str(step)
        body  = step.get('body', '')  if isinstance(step, dict) else ''

        # Cercle
        _h2_circle(slide, cx=nx, cy=ny, r=node_r, color=color)

        # Numéro
        _h2_text(slide, str(i + 1),
                 left=nx - node_r, top=ny - node_r,
                 width=node_r * 2, height=node_r * 0.7,
                 font=font, size_pt=11, color='FFFFFF',
                 bold=True, align='center')

        # Label sous le cercle (ou dessus selon position)
        label_y = ny + node_r + 0.05
        if label_y + 0.7 > _LY.CB:
            label_y = ny - node_r - 0.7

        _h2_text(slide, title,
                 left=nx - 1.1, top=label_y,
                 width=2.2, height=0.55,
                 font=font, size_pt=11, color=dk1,
                 bold=True, align='center', line_spacing=1.1)

        # Flèche vers le prochain nœud (arc simplifié)
        if n > 1:
            next_angle = angle + 2 * math.pi / n
            mx = cx_center + orbit_r * math.cos(angle + math.pi / n)
            my = cy_center + orbit_r * math.sin(angle + math.pi / n)
            # Simple point de passage — flèche ▶ orientée
            _h2_text(slide, '▸',
                     left=mx - 0.15, top=my - 0.15,
                     width=0.3, height=0.3,
                     font=font, size_pt=12, color=color,
                     bold=True, align='center')

    return slide


def layout_roadmap_v4(prs: Presentation, content: dict, tp: dict):
    """
    Roadmap avec phases et jalons.
    Bande temporelle accent1 + phases distinctes + jalons listés.
    content: {title, phases:[{label, milestones:[str]}], footer}
    """
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accents = tp.get('accent_cycle', [
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        accent1,
    ])
    W = tp.get('W', 13.33)
    H = tp.get('H', 7.50)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    phases = content.get('phases', [])
    n = min(len(phases), 5)
    if n == 0:
        return slide

    # Bande de frise au sommet de la zone contenu (CT)
    y_band   = _LY.CT
    band_h   = _LY.HEAD_H
    x_start  = _LY.CL
    x_end    = _LY.CR
    band_w   = x_end - x_start
    phase_w  = band_w / n

    # Fond de bande (axe du temps)
    _h2_rect(slide, left=x_start, top=y_band, width=band_w, height=band_h, color='F0F0F0')

    # Hauteur disponible pour les jalons (de y_band+band_h jusqu'à CB)
    milestones_zone_h = _LY.CB - (y_band + band_h)

    for i, phase in enumerate(phases[:n]):
        label      = phase.get('label', '') if isinstance(phase, dict) else str(phase)
        milestones = phase.get('milestones', []) if isinstance(phase, dict) else []
        color      = accents[i % len(accents)]
        px         = x_start + i * phase_w

        # Bloc de phase sur la bande
        _h2_rect(slide, left=px + 0.04, top=y_band, width=phase_w - 0.08, height=band_h, color=color)
        _h2_text(slide, label,
                 left=px + 0.08, top=y_band + 0.06,
                 width=phase_w - 0.14, height=band_h - 0.1,
                 font=font, size_pt=11, color='FFFFFF',
                 bold=True, align='center')

        # Jalons — espacés sur toute la zone milestones
        n_ms = min(len(milestones), 5)
        ms_step = milestones_zone_h / max(n_ms, 1)
        for j, ms in enumerate(milestones[:n_ms]):
            my = y_band + band_h + j * ms_step + 0.08
            _h2_rect(slide, left=px + 0.08, top=my + 0.08, width=0.07, height=0.07, color=color)
            _h2_text(slide, str(ms),
                     left=px + 0.22, top=my,
                     width=phase_w - 0.30, height=ms_step - 0.10,
                     font=font, size_pt=10, color=dk1,
                     bold=False, align='left', line_spacing=1.1)

    return slide


# ── V4 Layouts bonus ─────────────────────────────────────────────────────────

def layout_stackedbar_v4(prs: Presentation, content: dict, tp: dict):
    """
    Barres empilées 100 % (COLUMN_STACKED_100).
    Idéal pour parts de marché, répartition, compositions.
    content: {title, categories:[str], series:[{name, values:[n]}], footer}
    """
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    slide = _blank_v4(prs, tp)
    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    categories = content.get('categories', [])
    series     = content.get('series', [])
    if not categories or not series:
        return slide

    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in categories]
    for s in series:
        vals = tuple(float(v) if v is not None else 0.0 for v in s.get('values', []))
        chart_data.add_series(s.get('name', ''), vals)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED_100,
        Inches(0.5), Inches(_LY.CT), Inches(12.0), Inches(_CHART_H),
        chart_data,
    )
    chart  = chart_shape.chart
    colors = _chart_series_colors(tp.get('theme', {}))

    for i, ser in enumerate(chart.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = _h2_parse_hex(colors[i % len(colors)])

    plot = chart.plots[0]
    plot.has_data_labels = True
    try:
        plot.data_labels.number_format = '0%'
        plot.data_labels.font.size = Pt(9)
        plot.data_labels.font.color.rgb = _h2_parse_hex('FFFFFF')
    except Exception:
        pass

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    _add_chart_analysis(slide, content, tp)
    return slide


def layout_beforeafter_v4(prs: Presentation, content: dict, tp: dict):
    """
    Avant / Après — 2 variantes visuelles déterministes.
    v0 : deux colonnes header coloré + fond palette + flèche centrale.
    v1 : colonne gauche EEEEEE + colonne droite E8EEF5, items avec icône ✗/✓.
    content: {title, before:{title,items:[str]}, after:{title,items:[str]}, footer}
    """
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accent3 = theme.get('accent3', '40A900')
    W = tp.get('W', 13.33)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    before = content.get('before', {})
    after  = content.get('after', {})
    v      = _v4_variant(content, 2, tp.get('seed', 0))

    y_top   = _LY.CT
    y_bot   = _LY.CB
    col_h   = y_bot - y_top
    arrow_w = 0.80
    col_w   = (W - _LY.CL - (W - _LY.CR) - arrow_w) / 2
    x_bef   = _LY.CL
    x_arr   = x_bef + col_w
    x_aft   = x_arr + arrow_w
    head_h  = _LY.HEAD_H
    item_h  = _LY.ITEM_H

    before_items = before.get('items', []) if isinstance(before, dict) else []
    after_items  = after.get('items',  []) if isinstance(after,  dict) else []
    n_rows = min(max(len(before_items), len(after_items)),
                 int((col_h - head_h - 0.10) / item_h))

    if v == 0:
        # Variante 0 : header coloré + fond palette + flèches ↔ par ligne
        for cx, col, hdr_color, bg_color, lbl in [
            (x_bef, before, '888888', 'EEEEEE', 'AVANT'),
            (x_aft,  after,  accent1,  'F0F0F0', 'APRÈS'),
        ]:
            col_title = col.get('title', lbl) if isinstance(col, dict) else lbl
            items     = col.get('items', []) if isinstance(col, dict) else []
            _h2_rect(slide, left=cx, top=y_top, width=col_w, height=head_h, color=hdr_color)
            _h2_text(slide, col_title,
                     left=cx + _LY.PAD, top=y_top + 0.07,
                     width=col_w - _LY.PAD * 2, height=head_h - 0.1,
                     font=font, size_pt=_LY.T_HEADER, color='FFFFFF',
                     bold=True, align='left')
            _h2_rect(slide, left=cx, top=y_top + head_h,
                     width=col_w, height=col_h - head_h, color=bg_color)
            for j, item in enumerate(items[:n_rows]):
                iy = y_top + head_h + 0.10 + j * item_h
                _h2_rect(slide, left=cx, top=iy, width=0.06,
                         height=item_h - 0.06, color=hdr_color)
                _h2_text(slide, str(item),
                         left=cx + 0.15, top=iy + 0.06,
                         width=col_w - 0.22, height=item_h - 0.10,
                         font=font, size_pt=11, color=dk1,
                         bold=False, align='left', line_spacing=1.1)
        # Flèche ↔ centrée verticalement sur chaque ligne d'item
        for j in range(n_rows):
            iy = y_top + head_h + 0.10 + j * item_h
            _h2_text(slide, '↔',
                     left=x_arr + 0.05, top=iy + (item_h - 0.40) / 2,
                     width=arrow_w - 0.10, height=0.40,
                     font=font, size_pt=20, color=accent1,
                     bold=True, align='center')
    else:
        # Variante 1 : colonnes fond palette, icônes ✗ (avant) et ✓ (après)
        for cx, col, bg_color, icon, icon_color, lbl in [
            (x_bef, before, 'EEEEEE', '✗', '999999', 'AVANT'),
            (x_aft,  after,  'F0F0F0',  '✓', accent3,  'APRÈS'),
        ]:
            col_title = col.get('title', lbl) if isinstance(col, dict) else lbl
            items     = col.get('items', []) if isinstance(col, dict) else []
            title_color = '888888' if lbl == 'AVANT' else accent1
            _h2_text(slide, col_title,
                     left=cx + _LY.PAD, top=y_top + 0.06,
                     width=col_w - _LY.PAD * 2, height=head_h - 0.1,
                     font=font, size_pt=_LY.T_HEADER, color=title_color,
                     bold=True, align='left')
            _h2_rect(slide, left=cx + _LY.PAD, top=y_top + head_h - 0.04,
                     width=col_w - _LY.PAD * 2, height=0.04, color=title_color)
            _h2_rounded_rect(slide, left=cx, top=y_top + head_h,
                              width=col_w, height=col_h - head_h,
                              color=bg_color, radius=_LY.R_SM)
            for j, item in enumerate(items[:n_rows]):
                iy = y_top + head_h + 0.10 + j * item_h
                _h2_text(slide, icon,
                         left=cx + 0.10, top=iy + 0.04,
                         width=0.28, height=item_h - 0.08,
                         font=font, size_pt=13, color=icon_color,
                         bold=True, align='center')
                _h2_text(slide, str(item),
                         left=cx + 0.42, top=iy + 0.06,
                         width=col_w - 0.48, height=item_h - 0.10,
                         font=font, size_pt=11, color=dk1,
                         bold=False, align='left', line_spacing=1.1)
        # Flèche ↔ centrée verticalement sur chaque ligne d'item
        for j in range(n_rows):
            iy = y_top + head_h + 0.10 + j * item_h
            _h2_text(slide, '↔',
                     left=x_arr + 0.05, top=iy + (item_h - 0.40) / 2,
                     width=arrow_w - 0.10, height=0.40,
                     font=font, size_pt=20, color=accent1,
                     bold=True, align='center')

    return slide


def layout_entity_v4(prs: Presentation, content: dict, tp: dict):
    """
    Comparaison d'entités (pays, acteurs, marques) — 3 variantes déterministes.
    v0 : cartes verticales — barre accent top + grand icon centré + nom + badge + bullets.
    v1 : rangées horizontales — entité label gauche (fond coloré) + bullets droite.
    v2 : tableau comparatif — entités en colonnes, critères en lignes (alternance fond).
    entity schema: {icon?, name, badge?, label?, items:[str], stat_value?, stat_label?}
    content: {title, section_label?, subtitle?, entities:[...], footer}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accents = tp.get('accent_cycle', [
        theme.get('accent1', '009CEA'),
        theme.get('accent2', 'ED0000'),
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
    ])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    entities = content.get('entities', content.get('items', []))
    if not entities:
        return slide

    n = min(len(entities), 4)
    v = _v4_variant(content, 3, tp.get('seed', 0))

    def _ent(e):
        if not isinstance(e, dict):
            return '', str(e), '', str(e), [], '', ''
        return (e.get('icon', ''), e.get('name', e.get('title', '')),
                e.get('badge', e.get('label', '')),
                e.get('label', e.get('badge', '')),
                e.get('items', []), e.get('stat_value', ''), e.get('stat_label', ''))

    if v == 0:
        # Cartes verticales — icon centré + nom + badge pill + bullets
        card_w = (_LY.CW - _LY.GAP_LG * (n - 1)) / n
        card_h = _LY.CB - _LY.CT
        for i in range(n):
            cx    = _LY.CL + i * (card_w + _LY.GAP_LG)
            color = accents[i % len(accents)]
            icon, name, badge, label, items, sv, sl = _ent(entities[i])
            _h2_rounded_rect(slide, left=cx, top=_LY.CT,
                              width=card_w, height=card_h, color='F8F8F8', radius=_LY.R_SM)
            _h2_rect(slide, left=cx, top=_LY.CT, width=card_w, height=0.07, color=color)
            y_cur = _LY.CT + 0.14
            if icon:
                _h2_text(slide, icon, left=cx, top=y_cur,
                         width=card_w, height=0.55,
                         font=font, size_pt=26, color=color, bold=False, align='center')
                y_cur += 0.58
            _h2_text(slide, name, left=cx + _LY.PAD, top=y_cur,
                     width=card_w - _LY.PAD * 2, height=0.36,
                     font=font, size_pt=13, color=dk1, bold=True, align='center')
            y_cur += 0.38
            if badge:
                _h2_rounded_rect(slide, left=cx + (card_w - 1.40) / 2, top=y_cur,
                                  width=1.40, height=0.24, color=color, radius=0.06)
                _h2_text(slide, badge, left=cx + (card_w - 1.40) / 2, top=y_cur + 0.02,
                         width=1.40, height=0.22,
                         font=font, size_pt=8, color='FFFFFF', bold=True, align='center')
                y_cur += 0.30
            _h2_rect(slide, left=cx + _LY.PAD, top=y_cur,
                     width=card_w - _LY.PAD * 2, height=0.03, color='E0E0E0')
            y_cur += 0.10
            body_h = card_h - (y_cur - _LY.CT) - (0.60 if sv else 0.08)
            for j, it in enumerate(items[:6]):
                iy = y_cur + j * (body_h / max(len(items[:6]), 1))
                if iy + 0.26 > _LY.CT + card_h - (0.60 if sv else 0.08):
                    break
                _h2_rect(slide, left=cx + _LY.PAD, top=iy + 0.08,
                         width=0.06, height=0.06, color=color)
                _h2_text(slide, str(it),
                         left=cx + _LY.PAD + 0.12, top=iy,
                         width=card_w - _LY.PAD - 0.20, height=0.28,
                         font=font, size_pt=9, color=dk1, bold=False, align='left')
            if sv:
                _h2_rect(slide, left=cx + _LY.PAD, top=_LY.CB - 0.60,
                         width=card_w - _LY.PAD * 2, height=0.025, color='DDDDDD')
                _h2_text(slide, sv, left=cx + _LY.PAD, top=_LY.CB - 0.57,
                         width=card_w - _LY.PAD * 2, height=0.36,
                         font=font, size_pt=20, color=color, bold=True, align='center')
                if sl:
                    _h2_text(slide, sl, left=cx + _LY.PAD, top=_LY.CB - 0.22,
                             width=card_w - _LY.PAD * 2, height=0.18,
                             font=font, size_pt=7, color='888888', bold=True, align='center')
        return slide

    if v == 1:
        # Rangées horizontales — entité label à gauche + bullets à droite
        n1     = min(n, 4)
        gap    = _LY.GAP_SM
        row_h  = (_LY.CB - _LY.CT - gap * (n1 - 1)) / n1
        lbl_w  = 2.60
        bul_w  = _LY.CW - lbl_w - _LY.GAP_MD
        for i in range(n1):
            cy    = _LY.CT + i * (row_h + gap)
            color = accents[i % len(accents)]
            icon, name, badge, label, items, sv, sl = _ent(entities[i])
            # Fond gauche coloré
            _h2_rounded_rect(slide, left=_LY.CL, top=cy,
                              width=lbl_w, height=row_h, color=color, radius=_LY.R_SM)
            y_lbl = cy + (row_h - 0.44) / 2
            if icon:
                _h2_text(slide, icon, left=_LY.CL, top=cy + 0.08,
                         width=lbl_w, height=0.44,
                         font=font, size_pt=22, color='FFFFFF', bold=False, align='center')
                y_lbl = cy + 0.52
            _h2_text(slide, name, left=_LY.CL + 0.10, top=y_lbl,
                     width=lbl_w - 0.20, height=row_h - (y_lbl - cy) - 0.08,
                     font=font, size_pt=12, color='FFFFFF', bold=True, align='center')
            # Fond fond pale + bullets
            bx = _LY.CL + lbl_w + _LY.GAP_MD
            _h2_rounded_rect(slide, left=bx, top=cy,
                              width=bul_w, height=row_h, color='F8F8F8', radius=_LY.R_SM)
            _h2_rect(slide, left=bx, top=cy, width=0.055, height=row_h, color=color)
            if badge:
                _h2_rounded_rect(slide, left=bx + bul_w - 1.30, top=cy + 0.08,
                                  width=1.20, height=0.22, color=color, radius=0.06)
                _h2_text(slide, badge, left=bx + bul_w - 1.28, top=cy + 0.09,
                         width=1.16, height=0.20,
                         font=font, size_pt=7, color='FFFFFF', bold=True, align='center')
            item_h = row_h / max(len(items[:4]), 1)
            for j, it in enumerate(items[:4]):
                iy = cy + j * item_h + 0.04
                _h2_text(slide, f'• {it}',
                         left=bx + 0.18, top=iy,
                         width=bul_w - 0.24, height=item_h - 0.04,
                         font=font, size_pt=10, color=dk1,
                         bold=False, align='left', line_spacing=1.1)
        return slide

    # Variante 2 : tableau comparatif — entités en colonnes, critères en lignes
    n2       = min(n, 4)
    hdr_h    = 0.70
    body_h   = _LY.CB - _LY.CT - hdr_h
    col_w    = _LY.CW / n2
    criteria = []
    for e in entities[:n2]:
        if isinstance(e, dict):
            for it in e.get('items', []):
                if it not in criteria:
                    criteria.append(it)
    # Limit rows
    max_rows = int(body_h / 0.44)
    criteria = criteria[:max_rows]
    n_rows   = len(criteria)
    row_h    = body_h / max(n_rows, 1) if n_rows else body_h

    # Entêtes de colonnes
    for i in range(n2):
        cx    = _LY.CL + i * col_w
        color = accents[i % len(accents)]
        icon, name, badge, label, items, sv, sl = _ent(entities[i])
        _h2_rect(slide, left=cx, top=_LY.CT, width=col_w, height=hdr_h, color=color)
        x_txt = cx + 0.10
        w_txt = col_w - 0.20
        if icon:
            _h2_text(slide, icon, left=cx, top=_LY.CT + 0.06,
                     width=0.50, height=hdr_h - 0.10,
                     font=font, size_pt=18, color='FFFFFF', bold=False, align='center')
            x_txt = cx + 0.52
            w_txt = col_w - 0.60
        _h2_text(slide, name, left=x_txt, top=_LY.CT + 0.16,
                 width=w_txt, height=hdr_h - 0.20,
                 font=font, size_pt=11, color='FFFFFF', bold=True, align='left')

    # Lignes de critères
    entity_items = [
        (e.get('items', []) if isinstance(e, dict) else []) for e in entities[:n2]
    ]
    for j, crit in enumerate(criteria):
        ry  = _LY.CT + hdr_h + j * row_h
        bg  = 'F0F4F8' if j % 2 == 0 else 'FAFAFA'
        _h2_rect(slide, left=_LY.CL, top=ry, width=_LY.CW, height=row_h, color=bg)
        for i in range(n2):
            cx    = _LY.CL + i * col_w
            color = accents[i % len(accents)]
            val   = crit if crit in entity_items[i] else '–'
            tick  = '✓' if crit in entity_items[i] else '–'
            c_txt = color if tick == '✓' else 'BBBBBB'
            _h2_text(slide, tick, left=cx + 0.10, top=ry + (row_h - 0.28) / 2,
                     width=0.26, height=0.28,
                     font=font, size_pt=12, color=c_txt, bold=True, align='center')
            _h2_text(slide, crit, left=cx + 0.38, top=ry + (row_h - 0.28) / 2,
                     width=col_w - 0.44, height=0.28,
                     font=font, size_pt=9, color=dk1, bold=False, align='left')

    return slide


def layout_conclusion_v4(prs: Presentation, content: dict, tp: dict):
    """
    Slide de synthèse/conclusion — grille 2×2 de cartes (gauche) + sidebar sombre (droite).
    sidebar: {title?, quote?, cta?}  ou champs directs sidebar_title/sidebar_quote/sidebar_cta.
    card schema: {icon?, title, body?}  — jusqu'à 4 cartes.
    content: {title, section_label?, subtitle?, cards:[...], sidebar_title?, sidebar_quote?, sidebar_cta?, footer}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accent2 = theme.get('accent2', 'ED0000')
    accents = tp.get('accent_cycle', [
        theme.get('accent1', '009CEA'),
        theme.get('accent4', 'F66A00'),
        theme.get('accent3', '40A900'),
        theme.get('accent2', 'ED0000'),
    ])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    cards = content.get('cards', [])[:4]

    # Sidebar raw data
    sb = content.get('sidebar', {}) if isinstance(content.get('sidebar'), dict) else {}
    sb_title = content.get('sidebar_title', sb.get('title', ''))
    sb_quote = content.get('sidebar_quote', sb.get('quote', ''))
    sb_cta   = content.get('sidebar_cta',   sb.get('cta',   ''))

    # Layout zones
    sidebar_w = _LY.CW * 0.31
    sidebar_x = _LY.CR - sidebar_w
    grid_w    = sidebar_x - _LY.CL - _LY.GAP_MD
    n         = max(len(cards), 1)
    cols      = 2
    rows      = (n + 1) // 2
    gap       = _LY.GAP_SM
    card_w    = (grid_w - gap) / cols
    card_h    = (_LY.CB - _LY.CT - gap * (rows - 1)) / rows

    for i, card in enumerate(cards):
        cx = _LY.CL + (i % cols) * (card_w + gap)
        cy = _LY.CT + (i // cols) * (card_h + gap)
        color = accents[i % len(accents)]
        icon  = card.get('icon', '') if isinstance(card, dict) else ''
        title = card.get('title', '') if isinstance(card, dict) else str(card)
        body  = card.get('body', '') if isinstance(card, dict) else ''

        # Fond carte + bordure gauche colorée
        _h2_rounded_rect(slide, left=cx, top=cy, width=card_w, height=card_h,
                          color='F8F8F8', radius=_LY.R_SM)
        _h2_rect(slide, left=cx, top=cy, width=0.055, height=card_h, color=color)

        y = cy + 0.14
        if icon:
            _h2_text(slide, icon,
                     left=cx + 0.18, top=y, width=0.38, height=0.36,
                     font=font, size_pt=15, color=dk1, bold=False, align='left')
            y += 0.38

        _h2_text(slide, title,
                 left=cx + 0.18, top=y,
                 width=card_w - 0.26, height=0.40,
                 font=font, size_pt=_LY.T_TITLE, color=dk1, bold=True, align='left')
        y += 0.42

        if body:
            _h2_text(slide, body,
                     left=cx + 0.18, top=y,
                     width=card_w - 0.26, height=cy + card_h - y - 0.10,
                     font=font, size_pt=_LY.T_SMALL, color=dk1,
                     bold=False, align='left', line_spacing=1.2)

    # Sidebar sombre
    sb_h = _LY.CB - _LY.CT
    _h2_rounded_rect(slide, left=sidebar_x, top=_LY.CT,
                      width=sidebar_w, height=sb_h, color=dk1, radius=_LY.R_SM)

    sy = _LY.CT + 0.22
    if sb_title:
        _h2_text(slide, sb_title,
                 left=sidebar_x + 0.20, top=sy,
                 width=sidebar_w - 0.30, height=0.44,
                 font=font, size_pt=14, color='FFFFFF', bold=True, align='left')
        sy += 0.46
        _h2_rect(slide, left=sidebar_x + 0.20, top=sy,
                 width=0.55, height=0.05, color=accent2)
        sy += 0.14

    if sb_quote:
        q_h = min(2.2, max(0.5, len(sb_quote) / 40 * 0.38))
        _h2_rounded_rect(slide, left=sidebar_x + 0.14, top=sy,
                          width=sidebar_w - 0.22, height=q_h + 0.30,
                          color=_darken(dk1, 0.72), radius=_LY.R_SM)
        _h2_text(slide, sb_quote,
                 left=sidebar_x + 0.24, top=sy + 0.12,
                 width=sidebar_w - 0.38, height=q_h,
                 font=font, size_pt=9, color='DDDDDD',
                 bold=False, italic=True, align='left', line_spacing=1.3)
        sy += q_h + 0.44

    if sb_cta:
        cta_y = _LY.CB - 1.40
        _h2_rounded_rect(slide, left=sidebar_x + 0.14, top=cta_y,
                          width=sidebar_w - 0.22, height=1.25,
                          color=_darken(dk1, 0.58), radius=_LY.R_SM)
        _h2_text(slide, sb_cta,
                 left=sidebar_x + 0.14, top=cta_y + 0.38,
                 width=sidebar_w - 0.22, height=0.55,
                 font=font, size_pt=12, color='FFFFFF',
                 bold=True, align='center')

    return slide


def layout_highlight_v4(prs: Presentation, content: dict, tp: dict):
    """
    Message encadré fort — 5 variantes visuelles déterministes.
    v0 : encadré accent1 + barre accent2 gauche + body sous.
    v1 : fond E8EEF5 + règles accent haut/bas + highlight dk1 centré.
    v2 : guillemet 64pt accent1 + encadré E8EEF5 léger.
    v3 : split — panneau dk1 gauche 35% + highlight blanc + body droite.
    v4 : stack — zone highlight E8EEF5 haut + points d'appui liste sous.
    content: {title, highlight, body, points?:[str], footer}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accent2 = theme.get('accent2', 'ED0000')
    W       = tp.get('W', 13.33)
    H       = tp.get('H', 7.50)
    v       = _v4_variant(content, 5, tp.get('seed', 0))

    highlight = content.get('highlight', '')
    body      = content.get('body', '')
    points    = content.get('points', [])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    n_lines = max(1, len(highlight) // 55 + 1)
    txt_sz  = max(16, min(28, int(40 / n_lines)))

    if v == 0:
        # Variante 0 : encadré accent1 + barre accent2 gauche
        box_y = 1.9
        box_h = 2.4 if (body or points) else 2.8
        box_x = _LY.CL + 0.35
        box_w = _LY.CW - 0.7
        _h2_rounded_rect(slide, left=box_x, top=box_y,
                          width=box_w, height=box_h, color=accent1, radius=_LY.R_SM)
        _h2_rect(slide, left=box_x, top=box_y, width=0.12, height=box_h, color=accent2)
        txt_h = min(box_h - 0.4, n_lines * (txt_sz * 0.022 + 0.1))
        txt_y = box_y + (box_h - txt_h) / 2
        _h2_text(slide, highlight,
                 left=box_x + 0.28, top=txt_y,
                 width=box_w - 0.45, height=txt_h,
                 font=font, size_pt=txt_sz, color='FFFFFF',
                 bold=True, align='left', line_spacing=1.25)
        if body or points:
            body_txt = body or '\n'.join(f'• {p}' for p in points[:4])
            _h2_text(slide, body_txt,
                     left=box_x, top=box_y + box_h + _LY.GAP_LG,
                     width=box_w, height=H - 0.65 - (box_y + box_h + _LY.GAP_LG),
                     font=font, size_pt=_LY.T_LABEL, color=dk1,
                     bold=False, align='left', line_spacing=1.3)

    elif v == 1:
        # Variante 1 : fond E8EEF5 + règles accent haut/bas + highlight dk1 centré
        box_y = 1.8
        box_h = 2.5 if (body or points) else 3.0
        box_x = _LY.CL + 0.20
        box_w = _LY.CW - 0.40
        _h2_rounded_rect(slide, left=box_x, top=box_y,
                          width=box_w, height=box_h, color='F0F0F0', radius=_LY.R_SM)
        _h2_rect(slide, left=box_x, top=box_y, width=box_w, height=0.06, color=accent1)
        _h2_rect(slide, left=box_x, top=box_y + box_h - 0.06,
                 width=box_w, height=0.06, color=accent2)
        txt_h = min(box_h - 0.4, n_lines * (txt_sz * 0.022 + 0.1))
        txt_y = box_y + (box_h - txt_h) / 2
        _h2_text(slide, highlight,
                 left=box_x + _LY.PAD, top=txt_y,
                 width=box_w - _LY.PAD * 2, height=txt_h,
                 font=font, size_pt=txt_sz, color=dk1,
                 bold=True, align='center', line_spacing=1.25)
        if body or points:
            body_txt = body or '\n'.join(f'• {p}' for p in points[:4])
            _h2_text(slide, body_txt,
                     left=box_x, top=box_y + box_h + _LY.GAP_LG,
                     width=box_w, height=H - 0.65 - (box_y + box_h + _LY.GAP_LG),
                     font=font, size_pt=_LY.T_LABEL, color=dk1,
                     bold=False, align='left', line_spacing=1.3)

    elif v == 2:
        # Variante 2 : guillemet 64pt + encadré E8EEF5 + highlight centré
        box_y = 1.7
        box_h = 2.8 if (body or points) else 3.2
        box_x = _LY.CL + 0.20
        box_w = _LY.CW - 0.40
        _h2_rounded_rect(slide, left=box_x, top=box_y,
                          width=box_w, height=box_h, color='F0F0F0', radius=_LY.R_SM)
        _h2_rect(slide, left=box_x, top=box_y, width=0.06, height=box_h, color=accent1)
        _h2_text(slide, '\u201c',
                 left=box_x + 0.22, top=box_y + 0.02,
                 width=0.8, height=0.72,
                 font=font, size_pt=64, color=accent1,
                 bold=True, align='left')
        txt_h = min(box_h - 0.5, n_lines * (txt_sz * 0.022 + 0.1))
        txt_y = box_y + (box_h - txt_h) / 2
        _h2_text(slide, highlight,
                 left=box_x + _LY.PAD, top=txt_y,
                 width=box_w - _LY.PAD * 2, height=txt_h,
                 font=font, size_pt=txt_sz, color=dk1,
                 bold=True, align='center', line_spacing=1.25)
        if body or points:
            body_txt = body or '\n'.join(f'• {p}' for p in points[:4])
            _h2_text(slide, body_txt,
                     left=box_x, top=box_y + box_h + _LY.GAP_LG,
                     width=box_w, height=H - 0.65 - (box_y + box_h + _LY.GAP_LG),
                     font=font, size_pt=_LY.T_LABEL, color=dk1,
                     bold=False, align='left', line_spacing=1.3)

    elif v == 3:
        # Variante 3 : split — panneau sombre gauche + highlight blanc centré + body droite
        panel_w = _LY.CW * 0.42
        panel_h = _LY.CB - _LY.CT
        _h2_rounded_rect(slide, left=_LY.CL, top=_LY.CT,
                          width=panel_w, height=panel_h, color=dk1, radius=_LY.R_SM)
        _h2_rect(slide, left=_LY.CL, top=_LY.CT, width=0.06, height=panel_h, color=accent1)
        txt_h = min(panel_h - 0.6, n_lines * (txt_sz * 0.022 + 0.12))
        txt_y = _LY.CT + (panel_h - txt_h) / 2
        _h2_text(slide, highlight,
                 left=_LY.CL + 0.18, top=txt_y,
                 width=panel_w - 0.28, height=txt_h,
                 font=font, size_pt=txt_sz, color='FFFFFF',
                 bold=True, align='center', line_spacing=1.25)
        # Séparateur décoratif accent2
        _h2_rect(slide, left=_LY.CL + panel_w * 0.2, top=txt_y + txt_h + 0.18,
                 width=panel_w * 0.6, height=0.05, color=accent2)
        # Body ou points à droite
        x_r = _LY.CL + panel_w + _LY.GAP_LG
        w_r = _LY.CR - x_r
        if body or points:
            body_txt = body or '\n'.join(f'• {p}' for p in points[:5])
            _h2_text(slide, body_txt,
                     left=x_r, top=_LY.CT + 0.15,
                     width=w_r, height=panel_h - 0.25,
                     font=font, size_pt=_LY.T_LABEL, color=dk1,
                     bold=False, align='left', line_spacing=1.35)

    else:
        # Variante 4 : stack — zone highlight haut + points en liste basse
        hl_h = min(2.6, 1.0 + n_lines * 0.45)
        hl_y = _LY.CT + 0.10
        _h2_rounded_rect(slide, left=_LY.CL, top=hl_y,
                          width=_LY.CW, height=hl_h, color='F0F0F0', radius=_LY.R_SM)
        _h2_rect(slide, left=_LY.CL, top=hl_y, width=_LY.CW, height=0.07, color=accent1)
        txt_h = min(hl_h - 0.3, n_lines * (txt_sz * 0.022 + 0.1))
        txt_y = hl_y + (hl_h - txt_h) / 2
        _h2_text(slide, highlight,
                 left=_LY.CL + _LY.PAD, top=txt_y,
                 width=_LY.CW - _LY.PAD * 2, height=txt_h,
                 font=font, size_pt=txt_sz, color=dk1,
                 bold=True, align='center', line_spacing=1.25)
        # Points / body sous l'encadré
        pts = points if points else ([body] if body else [])
        y_pts = hl_y + hl_h + _LY.GAP_MD
        pts_h = _LY.CB - y_pts
        if pts:
            n_pts  = min(len(pts), 4)
            pt_h   = pts_h / max(n_pts, 1)
            accents_l = tp.get('accent_cycle', [accent1, accent2])
            for j, pt in enumerate(pts[:n_pts]):
                color = accents_l[j % len(accents_l)]
                _h2_circle(slide, cx=_LY.CL + 0.18, cy=y_pts + j * pt_h + pt_h / 2,
                            r=0.10, color=color)
                _h2_text(slide, str(pt),
                         left=_LY.CL + 0.40, top=y_pts + j * pt_h + (pt_h - 0.36) / 2,
                         width=_LY.CW - 0.44, height=0.36,
                         font=font, size_pt=_LY.T_LABEL, color=dk1,
                         bold=False, align='left', line_spacing=1.2)

    return slide


def layout_agenda_v4(prs: Presentation, content: dict, tp: dict):
    """
    Sommaire / Agenda numéroté — 3 variantes visuelles déterministes.
    v0 : cercles numérotés accent_cycle + labels + séparateurs.
    v1 : numéros bold fond coloré + fonds de ligne alternées palette + sous-label.
    v2 : grille 2 colonnes, items répartis gauche/droite, numéro accent + label.
    content: {title, items:[{number,label,sublabel}] ou [str], footer}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accents = tp.get('accent_cycle', [
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        theme.get('accent1', '009CEA'),
    ])
    W = tp.get('W', 13.33)

    _add_template_header_and_footer(slide, content.get('title', 'Sommaire'),
                                    content.get('footer', ''), tp, content)

    items = content.get('items', content.get('agenda_items', []))
    n = min(len(items), 8)
    if n == 0:
        return slide

    v      = _v4_variant(content, 3, tp.get('seed', 0))

    def _item_data(item, idx):
        if isinstance(item, dict):
            return item.get('number', idx + 1), item.get('label', ''), item.get('sublabel', '')
        return idx + 1, str(item), ''

    if v == 2:
        # Variante 2 : 2 colonnes
        n_left  = (n + 1) // 2
        n_right = n - n_left
        col_xs  = [_LY.CL, _LY.CL + _LY.COL_W + _LY.COL_GAP]
        for ci, (x_col, n_col) in enumerate([(col_xs[0], n_left), (col_xs[1], n_right)]):
            start_i = ci * n_left
            row_h   = (_LY.CB - _LY.CT) / max(n_left, 1)
            for ri in range(n_col):
                i        = start_i + ri
                number, label, sublabel = _item_data(items[i], i)
                color    = accents[i % len(accents)]
                iy       = _LY.CT + ri * row_h
                # Fond alternée
                bg = 'F0F0F0' if i % 2 == 0 else 'EEEEEE'
                _h2_rect(slide, left=x_col, top=iy,
                         width=_LY.COL_W, height=row_h - _LY.GAP_XS, color=bg)
                _h2_rect(slide, left=x_col, top=iy,
                         width=_LY.BORDER_W, height=row_h - _LY.GAP_XS, color=color)
                # Numéro
                _h2_text(slide, str(number),
                         left=x_col + 0.12, top=iy + (row_h - 0.40) / 2,
                         width=0.40, height=0.40,
                         font=font, size_pt=16, color=color, bold=True, align='center')
                _h2_text(slide, label,
                         left=x_col + 0.62, top=iy + (row_h - 0.40) / 2,
                         width=_LY.COL_W - 0.68, height=0.40,
                         font=font, size_pt=14, color=dk1, bold=False, align='left')
    else:
        step_h = (_LY.CB - _LY.CT) / n
        for i in range(n):
            number, label, sublabel = _item_data(items[i], i)
            color = accents[i % len(accents)]
            iy    = _LY.CT + i * step_h

            if v == 0:
                # Variante 0 : cercle + label + séparateur
                num_w  = 0.60
                x_num  = _LY.CL
                x_text = x_num + num_w + 0.25
                _h2_circle(slide, cx=x_num + num_w / 2, cy=iy + step_h / 2,
                           r=_LY.R_CIRC, color=color)
                _h2_text(slide, str(number),
                         left=x_num, top=iy + step_h / 2 - 0.22,
                         width=num_w, height=0.44,
                         font=font, size_pt=14, color='FFFFFF', bold=True, align='center')
                _h2_text(slide, label,
                         left=x_text, top=iy + (step_h - 0.42) / 2,
                         width=W - x_text - 0.5, height=0.42,
                         font=font, size_pt=16, color=dk1, bold=False, align='left')
                if i < n - 1:
                    _h2_rect(slide, left=x_num, top=iy + step_h - 0.008,
                             width=W - x_num - 0.5, height=0.008, color='EEEEEE')
            else:
                # Variante 1 : fond ligne palette + badge numéro + sublabel
                row_h  = step_h - _LY.GAP_XS
                bg_row = 'F0F0F0' if i % 2 == 0 else 'EEEEEE'
                num_w  = 0.55
                x_text = _LY.CL + num_w + _LY.PAD * 2 + _LY.GAP_SM
                _h2_rect(slide, left=_LY.CL, top=iy, width=_LY.CW, height=row_h, color=bg_row)
                _h2_rect(slide, left=_LY.CL, top=iy, width=_LY.BORDER_W, height=row_h, color=color)
                badge_h = min(0.45, row_h - 0.06)
                _h2_rounded_rect(slide, left=_LY.CL + _LY.GAP_SM,
                                  top=iy + (row_h - badge_h) / 2,
                                  width=num_w, height=badge_h,
                                  color=color, radius=_LY.R_SM)
                _h2_text(slide, str(number),
                         left=_LY.CL + _LY.GAP_SM,
                         top=iy + (row_h - badge_h) / 2,
                         width=num_w, height=badge_h,
                         font=font, size_pt=16, color='FFFFFF', bold=True, align='center')
                lbl_h = 0.38 if sublabel else min(0.44, row_h - 0.08)
                _h2_text(slide, label,
                         left=x_text, top=iy + (row_h - lbl_h) / 2 - (0.18 if sublabel else 0),
                         width=W - x_text - 0.5, height=lbl_h,
                         font=font, size_pt=15, color=dk1, bold=False, align='left')
                if sublabel:
                    _h2_text(slide, sublabel,
                             left=x_text, top=iy + (row_h - lbl_h) / 2 + lbl_h - 0.12,
                             width=W - x_text - 0.5, height=0.32,
                             font=font, size_pt=_LY.T_SMALL, color='777777',
                             bold=False, align='left')

    return slide


# ── V4 Matrices & Diagrammes ─────────────────────────────────────────────────

def layout_matrix_v4(prs: Presentation, content: dict, tp: dict):
    """
    Matrice 2×2. Chaque quadrant : rounded_rect #F5F5F5 + bordure haute colorée.
    Axes labels en dehors de la grille. Items bullet dans chaque quadrant.
    content: {title, quadrants:[{label,body}]×4, axes:{x,y}, footer}
    """
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accents = tp.get('accent_cycle', [
        theme.get('accent3', '40A900'),
        theme.get('accent4', 'F66A00'),
        accent1,
        theme.get('accent2', 'ED0000'),
    ])
    W = tp.get('W', 13.33)
    H = tp.get('H', 7.50)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    quadrants = content.get('quadrants', [])
    axes      = content.get('axes', {})

    # Grille : 2×2 dans la zone contenu
    y_top    = _LY.CT
    y_bot    = _LY.CB
    gap      = _LY.GAP_SM
    ax_lbl_w = 1.10   # largeur réservée pour label axe Y
    ax_lbl_h = 0.38   # hauteur réservée pour label axe X

    grid_x = _LY.CL + ax_lbl_w
    grid_w = _LY.CR - grid_x
    grid_y = y_top
    grid_h = y_bot - ax_lbl_h - y_top

    cell_w = (grid_w - gap) / 2
    cell_h = (grid_h - gap) / 2

    # Axe Y — label centré sur la hauteur de la grille, aligné à droite
    y_axis = axes.get('y', '')
    if y_axis:
        _h2_text(slide, y_axis,
                 left=_LY.CL, top=grid_y + grid_h / 2 - 0.25,
                 width=ax_lbl_w - 0.10, height=0.50,
                 font=font, size_pt=10, color=dk1,
                 bold=False, align='right')

    # Axe X — label centré sous la grille
    x_axis = axes.get('x', '')
    if x_axis:
        _h2_text(slide, x_axis,
                 left=grid_x + grid_w / 4, top=y_bot - ax_lbl_h,
                 width=grid_w / 2, height=ax_lbl_h,
                 font=font, size_pt=10, color=dk1,
                 bold=False, align='center')

    # 4 quadrants (ordre : TL, TR, BL, BR)
    positions = [
        (grid_x,               grid_y),
        (grid_x + cell_w + gap, grid_y),
        (grid_x,               grid_y + cell_h + gap),
        (grid_x + cell_w + gap, grid_y + cell_h + gap),
    ]
    for i, pos in enumerate(positions):
        if i >= len(quadrants):
            break
        q   = quadrants[i] if isinstance(quadrants[i], dict) else {'label': str(quadrants[i])}
        cx, cy = pos
        color  = accents[i % len(accents)]
        label  = q.get('label', '')
        body   = q.get('body', '')
        items  = q.get('items', [])

        # Fond palette template (pas de F5F5F5 hors palette)
        _h2_rounded_rect(slide, left=cx, top=cy,
                          width=cell_w, height=cell_h,
                          color='EEEEEE', radius=0.04)
        _h2_rect(slide, left=cx, top=cy, width=cell_w, height=0.07, color=color)

        _h2_text(slide, label,
                 left=cx + _LY.PAD, top=cy + 0.13,
                 width=cell_w - _LY.PAD * 2, height=0.38,
                 font=font, size_pt=_LY.T_LABEL, color=dk1,
                 bold=True, align='left')

        # Contenu (body ou items) — dk1 pour cohérence palette
        txt = body
        if not txt and items:
            txt = '\n'.join(f'• {it}' for it in items)
        if txt:
            _h2_text(slide, txt,
                     left=cx + _LY.PAD, top=cy + 0.57,
                     width=cell_w - _LY.PAD * 2, height=cell_h - 0.65,
                     font=font, size_pt=_LY.T_SMALL, color=dk1,
                     bold=False, align='left', line_spacing=1.2)

    return slide


def layout_swot_v4(prs: Presentation, content: dict, tp: dict):
    """
    Analyse SWOT : 4 quadrants colorés (vert/rouge/bleu/orange).
    Header couleur + items bullet dans chaque quadrant.
    content: {title, strengths:[str], weaknesses:[str],
              opportunities:[str], threats:[str], footer}
    """
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    W = tp.get('W', 13.33)
    H = tp.get('H', 7.50)

    _add_template_header_and_footer(slide, content.get('title', 'Analyse SWOT'),
                                    content.get('footer', ''), tp, content)

    # Fonds : palette template uniquement
    quadrant_defs = [
        ('Forces',        'strengths',    '40A900', 'F4F4F4'),
        ('Faiblesses',    'weaknesses',   'ED0000', 'EEEEEE'),
        ('Opportunités',  'opportunities','009CEA', 'F0F0F0'),
        ('Menaces',       'threats',      'F66A00', 'EEEEEE'),
    ]

    y_top  = _LY.CT
    y_bot  = _LY.CB
    gap    = _LY.GAP_SM
    x_left = _LY.CL
    x_right = _LY.CR
    grid_w = x_right - x_left
    grid_h = y_bot - y_top
    cell_w = (grid_w - gap) / 2
    cell_h = (grid_h - gap) / 2
    head_h = _LY.HEAD_H

    positions = [
        (x_left,               y_top),
        (x_left + cell_w + gap, y_top),
        (x_left,               y_top + cell_h + gap),
        (x_left + cell_w + gap, y_top + cell_h + gap),
    ]

    for (label, key, hdr_color, bg_color), (cx, cy) in zip(quadrant_defs, positions):
        items = content.get(key, [])

        # Fond du quadrant (palette template)
        _h2_rounded_rect(slide, left=cx, top=cy,
                          width=cell_w, height=cell_h,
                          color=bg_color, radius=0.04)
        # Header coloré
        _h2_rect(slide, left=cx, top=cy, width=cell_w, height=head_h, color=hdr_color)
        _h2_text(slide, label,
                 left=cx + _LY.PAD, top=cy + _LY.GAP_XS,
                 width=cell_w - _LY.PAD * 2, height=head_h - 0.1,
                 font=font, size_pt=_LY.T_HEADER, color='FFFFFF',
                 bold=True, align='left')

        # Items bullet — texte dk1 (jamais '333333')
        txt = '\n'.join(f'• {it}' for it in items[:8])
        if txt:
            _h2_text(slide, txt,
                     left=cx + _LY.PAD, top=cy + head_h + _LY.GAP_SM,
                     width=cell_w - _LY.PAD * 2, height=cell_h - head_h - _LY.GAP_MD,
                     font=font, size_pt=_LY.T_SMALL, color=dk1,
                     bold=False, align='left', line_spacing=1.25)

    return slide


def layout_proscons_v4(prs: Presentation, content: dict, tp: dict):
    """
    Deux colonnes Pour / Contre.
    Gauche : accent3 (vert) — POUR. Droite : accent2 (rouge) — CONTRE.
    content: {title, pros:[str], cons:[str], footer}
    """
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    accent2 = theme.get('accent2', 'ED0000')
    accent3 = theme.get('accent3', '40A900')
    W = tp.get('W', 13.33)
    H = tp.get('H', 7.50)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    pros = content.get('pros', [])
    cons = content.get('cons', [])

    y_top   = _LY.CT
    y_bot   = _LY.CB
    gap     = _LY.GAP_LG
    x_left  = _LY.CL
    col_w   = (_LY.CW - gap) / 2
    head_h  = _LY.HEAD_H
    item_h  = _LY.ITEM_H

    for col_idx, (items, color, bg, label) in enumerate([
        (pros, accent3, 'EAFAF1', '✓  POUR'),
        (cons, accent2, 'FDEDEC', '✗  CONTRE'),
    ]):
        cx = x_left + col_idx * (col_w + gap)

        # Header
        _h2_rect(slide, left=cx, top=y_top, width=col_w, height=head_h, color=color)
        _h2_text(slide, label,
                 left=cx + _LY.PAD, top=y_top + 0.08,
                 width=col_w - _LY.PAD * 2, height=head_h - 0.1,
                 font=font, size_pt=_LY.T_TITLE, color='FFFFFF',
                 bold=True, align='left')

        # Items avec fonds alternés
        available_h = y_bot - (y_top + head_h) - 0.05
        n = min(len(items), int(available_h / item_h))
        for j, item in enumerate(items[:n]):
            iy     = y_top + head_h + j * item_h
            bg_row = bg if j % 2 == 0 else 'FFFFFF'
            _h2_rect(slide, left=cx, top=iy, width=col_w, height=item_h - 0.04, color=bg_row)
            _h2_rect(slide, left=cx, top=iy, width=0.07, height=item_h - 0.04, color=color)
            _h2_text(slide, str(item),
                     left=cx + 0.16, top=iy + 0.05,
                     width=col_w - 0.22, height=item_h - 0.1,
                     font=font, size_pt=11, color=dk1,
                     bold=False, align='left', line_spacing=1.1)

    return slide


def layout_table_v4(prs: Presentation, content: dict, tp: dict):
    """
    Tableau natif PowerPoint.
    Header : fond accent1 + texte blanc. Lignes alternées #F5F5F5 / #FFFFFF.
    content: {title, headers:[str], rows:[[str]], footer}
    """
    slide  = _blank_v4(prs, tp)
    font   = tp.get('font', 'Calibri')
    theme  = tp.get('theme', {})
    dk1    = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    W = tp.get('W', 13.33)
    H = tp.get('H', 7.50)

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    headers = content.get('headers', [])
    rows    = content.get('rows', [])
    if not headers and not rows:
        return slide

    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    n_rows = len(rows) + 1  # +1 pour l'en-tête
    if n_cols == 0 or n_rows == 0:
        return slide

    # Dimensions
    t_left = 0.55
    t_top  = 1.65
    t_w    = W - t_left - 0.45
    t_h    = min(H - 0.65 - t_top, n_rows * 0.52 + 0.1)
    row_h  = t_h / n_rows

    from pptx.util import Emu as _Emu
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(t_left), Inches(t_top), Inches(t_w), Inches(t_h),
    )
    tbl = table_shape.table

    # Largeurs de colonnes égales
    col_w_emu = int(Inches(t_w) / n_cols)
    for col in tbl.columns:
        col.width = col_w_emu

    def _set_cell(r, c, text, bg_hex, fg_hex, bold=False, size_pt=11):
        cell = tbl.cell(r, c)
        cell.text = str(text) if text is not None else ''
        # Fond
        try:
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = _h2_parse_hex(bg_hex)
        except Exception:
            pass
        # Texte
        try:
            para = cell.text_frame.paragraphs[0]
            run  = para.runs[0] if para.runs else para.add_run()
            run.font.name  = font
            run.font.size  = Pt(size_pt)
            run.font.bold  = bold
            run.font.color.rgb = _h2_parse_hex(fg_hex)
        except Exception:
            pass

    # Ligne d'en-tête
    for c, hdr in enumerate(headers[:n_cols]):
        _set_cell(0, c, hdr, accent1, 'FFFFFF', bold=True, size_pt=12)

    # Lignes de données
    for r, row in enumerate(rows[:n_rows - 1]):
        bg = 'F5F5F5' if r % 2 == 0 else 'FFFFFF'
        for c in range(n_cols):
            val = row[c] if c < len(row) else ''
            _set_cell(r + 1, c, val, bg, dk1, bold=False, size_pt=11)

    return slide


def layout_team_grid_v4(prs: Presentation, content: dict, tp: dict):
    """
    Grille d'équipe — avatars circulaires, nom, rôle, département optionnel.
    v0 : rangée de 3-4 colonnes, cercle XL centré + nom + rôle + stat bas.
    v1 : grille 2×2 cartes arrondies, cercle SM gauche + nom + rôle.
    v2 : 3 colonnes plein hauteur, cercle MD + barre accent + nom + corps.
    content: {title, section_label?, members:[{name,role,department?,icon?,stat_value?,stat_label?,body?}]}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accents = tp.get('accent_cycle', [
        theme.get('accent1', '009CEA'),
        theme.get('accent2', 'ED0000'),
    ])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    members = content.get('members', content.get('items', []))
    n = min(len(members), 6)
    if n == 0:
        return slide

    v = _v4_variant(content, 3, tp.get('seed', 0))

    def _mf(m):
        if not isinstance(m, dict):
            return '👤', str(m), '', '', '', '', ''
        return (m.get('icon', '👤'),
                m.get('name', m.get('title', '')),
                m.get('role', m.get('subtitle', '')),
                m.get('department', m.get('dept', '')),
                m.get('stat_value', ''), m.get('stat_label', ''),
                m.get('body', ''))

    if v == 0:
        # Rangée de colonnes avec grand cercle centré
        n_cols = min(n, 4)
        col_w  = _LY.CW / n_cols
        zone_h = _LY.CB - _LY.CT
        for i, m in enumerate(members[:n_cols]):
            icon, name, role, dept, sv, sl, body = _mf(m)
            color  = accents[i % len(accents)]
            card_x = _LY.CL + i * col_w
            _h2_rounded_rect(slide, left=card_x + _LY.GAP_SM, top=_LY.CT,
                             width=col_w - _LY.GAP_SM * 2, height=zone_h,
                             color='F8F8F8', radius=_LY.R_SM)
            _h2_rect(slide, left=card_x + _LY.GAP_SM, top=_LY.CT,
                     width=col_w - _LY.GAP_SM * 2, height=0.05, color=color)
            cx = card_x + col_w / 2
            r  = min(0.72, col_w * 0.30)
            cy = _LY.CT + r + 0.28
            _h2_circle(slide, cx, cy, r, color)
            _h2_text(slide, icon, cx - r, cy - r * 0.55, r * 2, r * 1.1,
                     font, int(r * 32), 'FFFFFF', bold=True, align='center')
            y = cy + r + 0.20
            _h2_text(slide, name, card_x + _LY.PAD, y,
                     col_w - _LY.PAD * 2, 0.36,
                     font, _LY.T_LABEL, dk1, bold=True, align='center')
            y += 0.38
            if role:
                _h2_text(slide, role, card_x + _LY.PAD, y,
                         col_w - _LY.PAD * 2, 0.26,
                         font, _LY.T_SMALL, '777777', bold=False, align='center')
                y += 0.28
            if dept:
                _h2_text(slide, dept, card_x + _LY.PAD, y,
                         col_w - _LY.PAD * 2, 0.24,
                         font, 9, color, bold=True, align='center')
                y += 0.26
            if body:
                _h2_text(slide, body, card_x + _LY.PAD, y,
                         col_w - _LY.PAD * 2, _LY.CB - y - (0.68 if sv else 0.12),
                         font, _LY.T_SMALL, '666666', bold=False, align='center',
                         line_spacing=1.2)
            if sv:
                _h2_rect(slide, left=card_x + _LY.PAD, top=_LY.CB - 0.66,
                         width=col_w - _LY.PAD * 2, height=0.025, color='DDDDDD')
                _h2_text(slide, sv, card_x + _LY.PAD, _LY.CB - 0.62,
                         col_w - _LY.PAD * 2, 0.38,
                         font, 20, color, bold=True, align='center')
            if sl:
                _h2_text(slide, sl, card_x + _LY.PAD, _LY.CB - 0.24,
                         col_w - _LY.PAD * 2, 0.20,
                         font, 8, '999999', bold=False, align='center')

    elif v == 1:
        # Grille 2×2 ou 2×3 — cercle SM à gauche, texte à droite
        n_cols = 2 if n <= 4 else 3
        n_rows = (n + n_cols - 1) // n_cols
        card_w = (_LY.CW - _LY.GAP_LG * (n_cols - 1)) / n_cols
        card_h = (_LY.CB - _LY.CT - _LY.GAP_SM * (n_rows - 1)) / n_rows
        for i, m in enumerate(members[:n]):
            col_i  = i % n_cols
            row_i  = i // n_cols
            icon, name, role, dept, sv, sl, body = _mf(m)
            color  = accents[i % len(accents)]
            cx     = _LY.CL + col_i * (card_w + _LY.GAP_LG)
            cy     = _LY.CT + row_i * (card_h + _LY.GAP_SM)
            _h2_rounded_rect(slide, left=cx, top=cy, width=card_w, height=card_h,
                             color='F8F8F8', radius=_LY.R_SM)
            _h2_rect(slide, left=cx, top=cy, width=card_w, height=0.05, color=color)
            r  = 0.32
            ax = cx + _LY.PAD + r
            ay = cy + card_h / 2
            _h2_circle(slide, ax, ay, r, color)
            _h2_text(slide, icon, ax - r, ay - r * 0.55, r * 2, r * 1.1,
                     font, 14, 'FFFFFF', bold=True, align='center')
            tx = cx + _LY.PAD + r * 2 + 0.16
            tw = card_w - (tx - cx) - _LY.PAD
            ty = cy + 0.14
            _h2_text(slide, name, tx, ty, tw, 0.34,
                     font, _LY.T_TITLE, dk1, bold=True, align='left')
            ty += 0.36
            if role:
                _h2_text(slide, role, tx, ty, tw, 0.24,
                         font, _LY.T_SMALL, '777777', bold=False, align='left')
                ty += 0.26
            if dept:
                _h2_text(slide, dept, tx, ty, tw, 0.22,
                         font, 9, color, bold=True, align='left')
            if sv:
                _h2_text(slide, sv, cx + card_w - 1.40, cy + 0.12, 1.30, 0.36,
                         font, 18, color, bold=True, align='right')
            if sl:
                _h2_text(slide, sl, cx + card_w - 1.40, cy + 0.50, 1.30, 0.20,
                         font, 8, '999999', bold=False, align='right')

    else:
        # v2 : 3 colonnes plein hauteur — cercle MD + barre accent
        n_cols = min(n, 3)
        col_w  = (_LY.CW - _LY.GAP_LG * (n_cols - 1)) / n_cols
        zone_h = _LY.CB - _LY.CT
        for i, m in enumerate(members[:n_cols]):
            icon, name, role, dept, sv, sl, body = _mf(m)
            color  = accents[i % len(accents)]
            card_x = _LY.CL + i * (col_w + _LY.GAP_LG)
            _h2_rounded_rect(slide, left=card_x, top=_LY.CT, width=col_w, height=zone_h,
                             color=_cbg(tp, i), radius=_LY.R_SM)
            cx = card_x + col_w / 2
            r  = min(0.82, col_w * 0.35)
            cy = _LY.CT + r + 0.24
            _h2_circle(slide, cx, cy, r, color)
            _h2_text(slide, icon, cx - r, cy - r * 0.55, r * 2, r * 1.1,
                     font, int(r * 34), 'FFFFFF', bold=True, align='center')
            _h2_rect(slide, left=card_x + _LY.PAD, top=cy + r + 0.18,
                     width=col_w - _LY.PAD * 2, height=0.03, color=color)
            y = cy + r + 0.26
            _h2_text(slide, name, card_x + _LY.PAD, y,
                     col_w - _LY.PAD * 2, 0.38,
                     font, 15, dk1, bold=True, align='center')
            y += 0.40
            if role:
                _h2_text(slide, role, card_x + _LY.PAD, y,
                         col_w - _LY.PAD * 2, 0.28,
                         font, _LY.T_SMALL, '777777', bold=False, align='center')
                y += 0.30
            if dept:
                _h2_text(slide, dept, card_x + _LY.PAD, y,
                         col_w - _LY.PAD * 2, 0.26,
                         font, 9, color, bold=True, align='center')
                y += 0.28
            if body:
                _h2_text(slide, body, card_x + _LY.PAD, y,
                         col_w - _LY.PAD * 2, _LY.CB - y - (0.68 if sv else 0.12),
                         font, _LY.T_SMALL, '555555', bold=False, align='center',
                         line_spacing=1.25)
            if sv:
                _h2_rect(slide, left=card_x + _LY.PAD, top=_LY.CB - 0.66,
                         width=col_w - _LY.PAD * 2, height=0.025, color='DDDDDD')
                _h2_text(slide, sv, card_x + _LY.PAD, _LY.CB - 0.62,
                         col_w - _LY.PAD * 2, 0.38,
                         font, 22, color, bold=True, align='center')
            if sl:
                _h2_text(slide, sl, card_x + _LY.PAD, _LY.CB - 0.24,
                         col_w - _LY.PAD * 2, 0.20,
                         font, 8, '999999', bold=False, align='center')

    return slide


def layout_stat_banner_v4(prs: Presentation, content: dict, tp: dict):
    """
    Bandeau de 3-4 grandes statistiques — style premium éditorial.
    v0 : cartes arrondies + cercle accent derrière la valeur + icône optionnelle.
    v1 : colonnes pleine hauteur alternées fond accent / fond neutre.
    v2 : valeurs géantes centrées + séparateurs verticaux + icône en cercle.
    content: {title, section_label?, stats:[{value,label,sublabel?,icon?}]}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accents = tp.get('accent_cycle', [
        theme.get('accent1', '009CEA'),
        theme.get('accent2', 'ED0000'),
    ])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    stats = content.get('stats', content.get('kpis', []))
    n = min(len(stats), 4)
    if n == 0:
        return slide

    v      = _v4_variant(content, 3, tp.get('seed', 0))
    cell_w = _LY.CW / n
    zone_h = _LY.CB - _LY.CT

    def _sf(s):
        if not isinstance(s, dict):
            return str(s), '', '', ''
        return (str(s.get('value', s.get('val', ''))),
                s.get('label', ''),
                s.get('sublabel', s.get('sub', '')),
                s.get('icon', ''))

    if v == 0:
        # Carte arrondie + cercle décoratif en fond + valeur + label
        for i, st in enumerate(stats[:n]):
            val, label, sub, icon = _sf(st)
            color  = accents[i % len(accents)]
            card_x = _LY.CL + i * cell_w + _LY.GAP_SM
            card_w = cell_w - _LY.GAP_SM * 2
            card_h = zone_h - 0.10
            _h2_rounded_rect(slide, left=card_x, top=_LY.CT + 0.05,
                             width=card_w, height=card_h, color='F8F8F8', radius=_LY.R_SM)
            # Grand cercle décoratif en arrière-plan (couleur accent si icône, sinon neutre)
            r = min(card_w * 0.36, 0.85)
            cx_card = card_x + card_w / 2
            _h2_circle(slide, cx_card, _LY.CT + 0.05 + r + 0.18, r, color)
            ty = _LY.CT + 0.18
            if icon:
                _h2_text(slide, icon, card_x, ty, card_w, 0.44,
                         font, 24, 'FFFFFF', bold=True, align='center')
                ty += 0.44
            _h2_text(slide, val, card_x, ty, card_w, 0.80,
                     font, 38, 'FFFFFF', bold=True, align='center')
            base_y = _LY.CT + 0.05 + r * 2 + 0.44
            _h2_rect(slide, left=card_x + card_w * 0.20, top=base_y,
                     width=card_w * 0.60, height=0.025, color='DDDDDD')
            base_y += 0.06
            _h2_text(slide, label, card_x + _LY.PAD, base_y,
                     card_w - _LY.PAD * 2, 0.36,
                     font, _LY.T_LABEL, dk1, bold=True, align='center')
            if sub:
                _h2_text(slide, sub, card_x + _LY.PAD, base_y + 0.38,
                         card_w - _LY.PAD * 2, 0.28,
                         font, _LY.T_SMALL, '888888', bold=False, align='center')

    elif v == 1:
        # Colonnes pleine hauteur alternées
        for i, st in enumerate(stats[:n]):
            val, label, sub, icon = _sf(st)
            color  = accents[i % len(accents)]
            is_dark = (i % 2 == 0)
            bg     = color if is_dark else _cbg(tp, i)
            fg     = 'FFFFFF' if is_dark else dk1
            val_c  = 'FFFFFF' if is_dark else color
            _h2_rect(slide, left=_LY.CL + i * cell_w, top=_LY.CT,
                     width=cell_w, height=zone_h, color=bg)
            cy_mid = _LY.CT + zone_h / 2
            y = cy_mid - (0.88 if icon else 0.60)
            if icon:
                _h2_text(slide, icon, _LY.CL + i * cell_w, y, cell_w, 0.46,
                         font, 26, fg, bold=True, align='center')
                y += 0.50
            _h2_text(slide, val, _LY.CL + i * cell_w, y, cell_w, 0.78,
                     font, 44, val_c, bold=True, align='center')
            y += 0.80
            _h2_text(slide, label, _LY.CL + i * cell_w, y, cell_w, 0.34,
                     font, _LY.T_LABEL, fg, bold=True, align='center')
            if sub:
                _h2_text(slide, sub, _LY.CL + i * cell_w, y + 0.36, cell_w, 0.26,
                         font, _LY.T_SMALL, fg, bold=False, align='center')

    else:
        # v2 : valeurs géantes + séparateurs verticaux + icône en cercle
        for i, st in enumerate(stats[:n]):
            val, label, sub, icon = _sf(st)
            color = accents[i % len(accents)]
            cx    = _LY.CL + (i + 0.5) * cell_w
            y     = _LY.CT + zone_h / 2 - (1.0 if icon else 0.70)
            if icon:
                _h2_circle(slide, cx, y + 0.28, 0.30, color)
                _h2_text(slide, icon, cx - 0.30, y + 0.28 - 0.15, 0.60, 0.30,
                         font, 16, 'FFFFFF', bold=True, align='center')
                y += 0.72
            _h2_text(slide, val, _LY.CL + i * cell_w, y, cell_w, 0.82,
                     font, 50, color, bold=True, align='center')
            y += 0.84
            _h2_rect(slide, left=cx - 0.55, top=y, width=1.10, height=0.025, color=color)
            y += 0.06
            _h2_text(slide, label, _LY.CL + i * cell_w, y, cell_w, 0.32,
                     font, _LY.T_LABEL, dk1, bold=True, align='center')
            if sub:
                _h2_text(slide, sub, _LY.CL + i * cell_w, y + 0.34, cell_w, 0.26,
                         font, _LY.T_SMALL, '888888', bold=False, align='center')
            # Séparateur vertical (sauf dernier)
            if i < n - 1:
                _h2_rect(slide, left=_LY.CL + (i + 1) * cell_w - 0.01,
                         top=_LY.CT + 0.40, width=0.02, height=zone_h - 0.80,
                         color='DDDDDD')

    return slide


def layout_icon_row_v4(prs: Presentation, content: dict, tp: dict):
    """
    Rangée de 3-4 features avec icônes en cercles — pattern horizontal connecté.
    v0 : cercles accent plein + ligne connectrice + badge numéro + titre + corps.
    v1 : cercles dans cartes arrondies (fond neutre) + titre + items bullets.
    v2 : cercles contour + connecteurs flèche colorée + titre + corps.
    content: {title, section_label?, items:[{icon,title,body?,items?:[str]}]}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accents = tp.get('accent_cycle', [
        theme.get('accent1', '009CEA'),
        theme.get('accent2', 'ED0000'),
    ])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    items = content.get('items', content.get('columns', []))
    n     = min(len(items), 4)
    if n == 0:
        return slide

    v      = _v4_variant(content, 3, tp.get('seed', 0))
    col_w  = _LY.CW / n
    zone_h = _LY.CB - _LY.CT

    def _if(it):
        if not isinstance(it, dict):
            return '●', str(it), '', []
        body_text = it.get('body', '')
        bullets   = it.get('items', it.get('bullets', []))
        return (it.get('icon', '●'), it.get('title', ''), body_text, bullets)

    if v == 0:
        # Cercles plein + ligne connectrice + badge numéro + corps
        r = 0.52
        icon_cy = _LY.CT + 0.90
        for i, it in enumerate(items[:n]):
            icon, title, body, bullets = _if(it)
            color  = accents[i % len(accents)]
            cx     = _LY.CL + (i + 0.5) * col_w
            # Ligne connectrice (sauf dernier)
            if i < n - 1:
                _h2_rect(slide, left=cx + r, top=icon_cy - 0.015,
                         width=col_w - r, height=0.030, color='DDDDDD')
            # Cercle
            _h2_circle(slide, cx, icon_cy, r, color)
            _h2_text(slide, icon, cx - r, icon_cy - r * 0.55, r * 2, r * 1.1,
                     font, 26, 'FFFFFF', bold=True, align='center')
            # Badge numéro
            br = 0.18
            _h2_circle(slide, cx + r * 0.68, icon_cy + r * 0.68, br, 'FFFFFF')
            _h2_text(slide, str(i + 1), cx + r * 0.68 - br, icon_cy + r * 0.68 - br * 0.55,
                     br * 2, br * 1.1, font, 9, color, bold=True, align='center')
            # Titre
            y = icon_cy + r + 0.24
            tx = _LY.CL + i * col_w + 0.10
            tw = col_w - 0.20
            _h2_text(slide, title, tx, y, tw, 0.42,
                     font, _LY.T_TITLE, dk1, bold=True, align='center')
            y += 0.46
            body_src = body or ('\n'.join(f'• {x}' for x in bullets[:3]))
            if body_src:
                _h2_text(slide, body_src, tx, y, tw, _LY.CB - y - 0.10,
                         font, _LY.T_SMALL, '666666', bold=False, align='center',
                         line_spacing=1.2)

    elif v == 1:
        # Cercles dans cartes arrondies + titre + bullets
        for i, it in enumerate(items[:n]):
            icon, title, body, bullets = _if(it)
            color  = accents[i % len(accents)]
            cx_card = _LY.CL + i * col_w + _LY.GAP_SM
            cw      = col_w - _LY.GAP_SM * 2
            _h2_rounded_rect(slide, left=cx_card, top=_LY.CT,
                             width=cw, height=zone_h, color=_cbg(tp, i), radius=_LY.R_SM)
            _h2_rect(slide, left=cx_card, top=_LY.CT, width=cw, height=0.05, color=color)
            cx = cx_card + cw / 2
            r  = min(0.50, cw * 0.22)
            cy = _LY.CT + r + 0.22
            _h2_circle(slide, cx, cy, r, color)
            _h2_text(slide, icon, cx - r, cy - r * 0.55, r * 2, r * 1.1,
                     font, int(r * 44), 'FFFFFF', bold=True, align='center')
            y = cy + r + 0.22
            _h2_text(slide, title, cx_card + _LY.PAD, y, cw - _LY.PAD * 2, 0.40,
                     font, _LY.T_TITLE, dk1, bold=True, align='center')
            y += 0.44
            body_src = body or ('\n'.join(f'• {x}' for x in bullets[:3]))
            if body_src:
                _h2_text(slide, body_src, cx_card + _LY.PAD, y,
                         cw - _LY.PAD * 2, _LY.CB - y - 0.14,
                         font, _LY.T_SMALL, '555555', bold=False, align='left',
                         line_spacing=1.2)

    else:
        # v2 : cercles contour + flèche colorée connectrice + titre + corps
        r = 0.46
        icon_cy = _LY.CT + 0.88
        for i, it in enumerate(items[:n]):
            icon, title, body, bullets = _if(it)
            color  = accents[i % len(accents)]
            cx     = _LY.CL + (i + 0.5) * col_w
            # Flèche connectrice (ligne + triangulaire)
            if i < n - 1:
                line_x  = cx + r + 0.08
                line_w  = col_w - r * 2 - 0.16
                _h2_rect(slide, left=line_x, top=icon_cy - 0.012,
                         width=line_w * 0.85, height=0.024, color=color)
            # Halo léger derrière le cercle
            _h2_circle(slide, cx, icon_cy, r + 0.08, _cbg(tp, i))
            # Cercle plein
            _h2_circle(slide, cx, icon_cy, r, color)
            _h2_text(slide, icon, cx - r, icon_cy - r * 0.55, r * 2, r * 1.1,
                     font, 22, 'FFFFFF', bold=True, align='center')
            y = icon_cy + r + 0.24
            tx = _LY.CL + i * col_w + 0.08
            tw = col_w - 0.16
            _h2_text(slide, title, tx, y, tw, 0.40,
                     font, _LY.T_TITLE, dk1, bold=True, align='center')
            y += 0.44
            body_src = body or ('\n'.join(f'• {x}' for x in bullets[:3]))
            if body_src:
                _h2_text(slide, body_src, tx, y, tw, _LY.CB - y - 0.08,
                         font, _LY.T_SMALL, '555555', bold=False, align='center',
                         line_spacing=1.2)

    return slide


def layout_section_break_v4(prs: Presentation, content: dict, tp: dict):
    """
    Slide de rupture entre sections — impact visuel fort, typographie XL.
    v0 : moitié gauche sombre (dk1) + grande forme cercle à cheval + titre à droite.
    v1 : bande verticale accent gauche + numéro oversize + titre + cercle décoratif.
    v2 : fond plein dk1 + titre centré blanc + formes géométriques accent.
    content: {title, subtitle?, number?, footer}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accents = tp.get('accent_cycle', [accent1])
    W = tp.get('W', 13.33)
    H = tp.get('H', 7.50)

    title    = content.get('title', '')
    subtitle = content.get('subtitle', '')
    number   = str(content.get('number', ''))
    footer   = content.get('footer', '')
    v = _v4_variant(content, 3, tp.get('seed', 0))

    if v == 0:
        # Demi-fond dk1 gauche + cercle à cheval sur la séparation + titre à droite
        split_x = W * 0.44
        _h2_rect(slide, left=0, top=0, width=split_x, height=H, color=dk1)
        # Grand cercle creux à cheval sur la ligne
        _h2_circle(slide, split_x, H / 2, 1.70, accent1)
        _h2_circle(slide, split_x, H / 2, 1.26, dk1)
        if number:
            _h2_text(slide, number, 0.20, H / 2 - 0.70, split_x - 0.40, 0.80,
                     font, 60, 'FFFFFF', bold=True, align='center')
        tx = split_x + 1.90
        tw = W - tx - 0.50
        ty = H / 2 - (0.60 if not subtitle else 0.80)
        _h2_text(slide, title, tx, ty, tw, 0.82,
                 font, 30, dk1, bold=True, align='left')
        if subtitle:
            _h2_rect(slide, left=tx, top=ty + 0.86, width=1.60, height=0.03, color=accent1)
            _h2_text(slide, subtitle, tx, ty + 0.94, tw, 0.38,
                     font, 13, '777777', bold=False, align='left')

    elif v == 1:
        # Bande verticale accent gauche + numéro + cercle décoratif fond
        band_w = 1.70
        _h2_rect(slide, left=0, top=0, width=band_w, height=H, color=accent1)
        num_txt = number if number else '◆'
        _h2_text(slide, num_txt, 0, H / 2 - 0.60, band_w, 0.80,
                 font, 40, 'FFFFFF', bold=True, align='center')
        # Cercle décoratif léger en fond côté droit
        _h2_circle(slide, W - 1.20, H / 2, 2.30, _cbg(tp, 1))
        cx = band_w + 0.90
        cw = W - cx - 0.55
        _h2_text(slide, title, cx, H / 2 - 0.52, cw, 0.72,
                 font, 32, dk1, bold=True, align='left')
        if subtitle:
            _h2_rect(slide, left=cx, top=H / 2 + 0.24, width=1.80, height=0.03, color=accent1)
            _h2_text(slide, subtitle, cx, H / 2 + 0.32, cw, 0.36,
                     font, 13, '666666', bold=False, align='left')

    else:
        # v2 : fond plein dk1 + titre centré blanc + formes géométriques
        _h2_rect(slide, left=0, top=0, width=W, height=H, color=dk1)
        # Cercle décoratif haut-gauche (semi-transparent via couleur atténuée)
        _h2_circle(slide, 0.40, H * 0.18, 1.50, _darken(accent1, 0.55))
        # Cercle décoratif bas-droit
        _h2_circle(slide, W - 0.40, H * 0.82, 1.90, _darken(dk1, 0.78))
        # Ligne accent centrée
        _h2_rect(slide, left=W * 0.30, top=H * 0.36, width=W * 0.40, height=0.04, color=accent1)
        if number:
            _h2_text(slide, f'{number}.', 0, H * 0.28, W, 0.46,
                     font, 16, accent1, bold=True, align='center')
        _h2_text(slide, title, W * 0.08, H * 0.40, W * 0.84, 0.90,
                 font, 36, 'FFFFFF', bold=True, align='center')
        if subtitle:
            _h2_text(slide, subtitle, W * 0.15, H * 0.64, W * 0.70, 0.38,
                     font, 13, 'CCCCCC', bold=False, align='center')

    if footer:
        fg = '999999' if v < 2 else 'AAAAAA'
        _h2_text(slide, footer, 0.60, H - 0.38, W - 1.20, 0.28,
                 font, 9, fg, bold=False, align='center')

    return slide


def layout_photo_text_v4(prs: Presentation, content: dict, tp: dict):
    """
    Zone photo (placeholder visuel) + contenu structuré — layouts inspirés images 1/2/4.
    v0 : photo gauche (large) + 3 items avec icône et trait à droite.
    v1 : 2 photos empilées gauche + liste d'items à droite avec stats.
    v2 : photo droite + contenu gauche (inverse) — 3 items avec barre accent.
    La zone photo est rendue en fond neutre + barre accent + icône centrée.
    content: {title, section_label?, items:[{icon?,title,body?,stat_value?}], footer}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accent1 = theme.get('accent1', '009CEA')
    accents = tp.get('accent_cycle', [accent1])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    items  = content.get('items', [])
    zone_h = _LY.CB - _LY.CT
    v      = _v4_variant(content, 3, tp.get('seed', 0))

    def _photo_zone(x, y, w, h, color):
        """Placeholder photo : fond neutre + barre accent + icône."""
        _h2_rounded_rect(slide, left=x, top=y, width=w, height=h,
                         color=_cbg(tp, 1), radius=_LY.R_SM)
        _h2_rect(slide, left=x, top=y, width=w, height=0.04, color=color)
        _h2_rect(slide, left=x, top=y + h - 0.40, width=w, height=0.40, color=color)
        _h2_text(slide, '🖼', x, y + h / 2 - 0.28, w, 0.56,
                 font, 28, 'BBBBBB', bold=False, align='center')

    def _item_row(it, ix, iy, iw, ih, color):
        if not isinstance(it, dict):
            it = {'title': str(it)}
        icon  = it.get('icon', '')
        title = it.get('title', '')
        body  = it.get('body', '')
        sv    = it.get('stat_value', '')
        _h2_rounded_rect(slide, left=ix, top=iy + 0.04,
                         width=iw, height=ih - 0.08, color=_cbg(tp, 0), radius=_LY.R_SM)
        _h2_rect(slide, left=ix, top=iy + 0.04, width=0.04, height=ih - 0.08, color=color)
        tx = ix + 0.16
        if icon:
            _h2_icon_circle(slide, tx + 0.24, iy + ih / 2, 0.22, icon, font, color, 'FFFFFF', 11)
            tx += 0.60
        tw_stat = 1.20 if sv else 0
        _h2_text(slide, title, tx, iy + 0.12, iw - (tx - ix) - tw_stat - 0.10, 0.30,
                 font, _LY.T_TITLE, dk1, bold=True, align='left')
        if body:
            _h2_text(slide, body, tx, iy + 0.44, iw - (tx - ix) - tw_stat - 0.10, ih - 0.52,
                     font, _LY.T_SMALL, '666666', bold=False, align='left', line_spacing=1.15)
        if sv:
            _h2_text(slide, sv, ix + iw - tw_stat - 0.06, iy + (ih - 0.40) / 2,
                     tw_stat, 0.40, font, 18, color, bold=True, align='right')

    if v == 0:
        # Photo gauche (4.8") + items à droite
        pw = 4.80
        tx = _LY.CL + pw + _LY.GAP_LG
        tw = _LY.CW - pw - _LY.GAP_LG
        _photo_zone(_LY.CL, _LY.CT, pw, zone_h, accent1)
        n = min(len(items), 3)
        if n > 0:
            ih = zone_h / n
            for i, it in enumerate(items[:n]):
                _item_row(it, tx, _LY.CT + i * ih, tw, ih, accents[i % len(accents)])

    elif v == 1:
        # 2 photos empilées gauche + items droite
        pw = 5.10
        tx = _LY.CL + pw + _LY.GAP_LG
        tw = _LY.CW - pw - _LY.GAP_LG
        ph1 = zone_h * 0.56
        ph2 = zone_h - ph1 - _LY.GAP_SM
        _photo_zone(_LY.CL, _LY.CT, pw, ph1, accent1)
        _photo_zone(_LY.CL, _LY.CT + ph1 + _LY.GAP_SM, pw, ph2,
                    accents[1 % len(accents)])
        n = min(len(items), 4)
        if n > 0:
            ih = zone_h / n
            for i, it in enumerate(items[:n]):
                _item_row(it, tx, _LY.CT + i * ih, tw, ih, accents[i % len(accents)])

    else:
        # v2 : photo droite + items gauche
        pw = 4.60
        tw = _LY.CW - pw - _LY.GAP_LG
        px = _LY.CL + tw + _LY.GAP_LG
        _photo_zone(px, _LY.CT, pw, zone_h, accent1)
        n = min(len(items), 3)
        if n > 0:
            ih = zone_h / n
            for i, it in enumerate(items[:n]):
                _item_row(it, _LY.CL, _LY.CT + i * ih, tw, ih, accents[i % len(accents)])

    return slide


def layout_numbered_features_v4(prs: Presentation, content: dict, tp: dict):
    """
    Features avec grands chiffres décoratifs — style éditorial premium.
    v0 : colonnes (3-4), grand numéro en fond + barre accent + titre + corps.
    v1 : lignes horizontales (3-4), numéro XL à gauche + titre + corps + stat.
    v2 : grille 2×2, cercle numéroté + barre accent + titre + corps.
    content: {title, section_label?, items:[{title,body,number?,stat_value?,stat_label?}]}
    """
    slide   = _blank_v4(prs, tp)
    font    = tp.get('font', 'Calibri')
    theme   = tp.get('theme', {})
    dk1     = theme.get('dk1', '374649')
    accents = tp.get('accent_cycle', [
        theme.get('accent1', '009CEA'),
        theme.get('accent2', 'ED0000'),
    ])

    _add_template_header_and_footer(slide, content.get('title', ''),
                                    content.get('footer', ''), tp, content)

    items  = content.get('items', content.get('features', []))
    n      = min(len(items), 4)
    if n == 0:
        return slide

    v      = _v4_variant(content, 3, tp.get('seed', 0))
    zone_h = _LY.CB - _LY.CT

    def _ff(it, i):
        if not isinstance(it, dict):
            return str(i + 1), str(it), '', '', ''
        return (str(it.get('number', i + 1)),
                it.get('title', ''), it.get('body', ''),
                it.get('stat_value', ''), it.get('stat_label', ''))

    if v == 0:
        # Colonnes — grand numéro en fond (couleur très atténuée) + titre + corps
        n_cols = n
        col_w  = (_LY.CW - _LY.GAP_LG * (n_cols - 1)) / n_cols
        for i, it in enumerate(items[:n]):
            num, title, body, sv, sl = _ff(it, i)
            color  = accents[i % len(accents)]
            cx     = _LY.CL + i * (col_w + _LY.GAP_LG)
            # Fond carte
            _h2_rounded_rect(slide, left=cx, top=_LY.CT, width=col_w, height=zone_h,
                             color='F8F8F8', radius=_LY.R_SM)
            # Grand numéro décoratif en fond (couleur très légère)
            _h2_text(slide, num, cx - 0.06, _LY.CT + 0.06, col_w + 0.12, 0.96,
                     font, 72, _cbg(tp, i + 1), bold=True, align='right')
            # Barre accent
            _h2_rect(slide, left=cx + _LY.PAD, top=_LY.CT + 0.20,
                     width=0.80, height=0.04, color=color)
            # Titre
            _h2_text(slide, title, cx + _LY.PAD, _LY.CT + 0.30,
                     col_w - _LY.PAD * 2, 0.52,
                     font, _LY.T_HEADER, dk1, bold=True, align='left')
            # Séparateur
            _h2_rect(slide, left=cx + _LY.PAD, top=_LY.CT + 0.88,
                     width=col_w - _LY.PAD * 2, height=0.02, color='DDDDDD')
            # Corps
            if body:
                body_h = zone_h - 0.96 - (0.70 if sv else 0.14)
                _h2_text(slide, body, cx + _LY.PAD, _LY.CT + 0.96,
                         col_w - _LY.PAD * 2, body_h,
                         font, _LY.T_SMALL, '555555', bold=False, align='left',
                         line_spacing=1.3)
            if sv:
                _h2_text(slide, sv, cx + _LY.PAD, _LY.CB - 0.62,
                         col_w - _LY.PAD * 2, 0.38,
                         font, 22, color, bold=True, align='left')
            if sl:
                _h2_text(slide, sl, cx + _LY.PAD, _LY.CB - 0.24,
                         col_w - _LY.PAD * 2, 0.20,
                         font, 8, '999999', bold=True, align='left')

    elif v == 1:
        # Lignes — numéro XL à gauche + titre + corps
        n_rows  = min(n, 4)
        row_h   = zone_h / n_rows
        num_w   = 1.30
        for i, it in enumerate(items[:n_rows]):
            num, title, body, sv, sl = _ff(it, i)
            color = accents[i % len(accents)]
            ry    = _LY.CT + i * row_h
            # Barre accent verticale
            _h2_rect(slide, left=_LY.CL, top=ry + 0.08,
                     width=0.04, height=row_h - 0.16, color=color)
            # Grand numéro
            _h2_text(slide, num, _LY.CL + 0.08, ry + 0.04, num_w - 0.08, row_h - 0.08,
                     font, 38, color, bold=True, align='center')
            # Séparateur horizontal (sauf premier)
            if i > 0:
                _h2_rect(slide, left=_LY.CL + num_w, top=ry,
                         width=_LY.CW - num_w, height=0.015, color='EEEEEE')
            tx  = _LY.CL + num_w + 0.30
            tw  = _LY.CW - num_w - 0.40 - (1.30 if sv else 0)
            th  = 0.34 if body else row_h - 0.18
            _h2_text(slide, title, tx, ry + (row_h - th) * 0.25, tw, th,
                     font, _LY.T_HEADER, dk1, bold=True, align='left')
            if body:
                _h2_text(slide, body, tx, ry + 0.44, tw, row_h - 0.50,
                         font, _LY.T_SMALL, '666666', bold=False, align='left',
                         line_spacing=1.2)
            if sv:
                _h2_text(slide, sv, _LY.CR - 1.20, ry + (row_h - 0.46) / 2,
                         1.10, 0.46, font, 22, color, bold=True, align='right')

    else:
        # v2 : grille 2×2 — cercle numéroté + titre + corps
        n_cols = 2
        n_rows = (n + 1) // 2
        card_w = (_LY.CW - _LY.GAP_LG) / 2
        card_h = (zone_h - _LY.GAP_SM * (n_rows - 1)) / n_rows
        for i, it in enumerate(items[:n]):
            col_i  = i % n_cols
            row_i  = i // n_cols
            num, title, body, sv, sl = _ff(it, i)
            color  = accents[i % len(accents)]
            cx     = _LY.CL + col_i * (card_w + _LY.GAP_LG)
            cy     = _LY.CT + row_i * (card_h + _LY.GAP_SM)
            _h2_rounded_rect(slide, left=cx, top=cy, width=card_w, height=card_h,
                             color='F8F8F8', radius=_LY.R_SM)
            # Cercle numéroté
            r  = 0.38
            ax = cx + _LY.PAD + r
            ay = cy + card_h / 2
            _h2_circle(slide, ax, ay, r, color)
            _h2_text(slide, num, ax - r, ay - r * 0.55, r * 2, r * 1.1,
                     font, 20, 'FFFFFF', bold=True, align='center')
            # Barre + titre + corps
            tx = cx + _LY.PAD + r * 2 + 0.20
            tw = card_w - _LY.PAD - r * 2 - 0.30
            _h2_rect(slide, left=tx, top=cy + 0.18, width=0.60, height=0.03, color=color)
            _h2_text(slide, title, tx, cy + 0.26, tw, 0.40,
                     font, _LY.T_HEADER, dk1, bold=True, align='left')
            if body:
                _h2_text(slide, body, tx, cy + 0.70, tw, card_h - 0.78,
                         font, _LY.T_SMALL, '666666', bold=False, align='left',
                         line_spacing=1.2)
            if sv:
                _h2_text(slide, sv, tx + tw - 1.30, cy + 0.28, 1.30, 0.36,
                         font, 18, color, bold=True, align='right')

    return slide


# Types servis par les vrais layouts du template (placeholders natifs)
_V4_NATIVE_TYPES = frozenset({
    # Anciens noms (compat V3)
    'cover_dark', 'cover_split',
    'section',
    'full_text', 'list_numbered', 'list_cards', 'image_split',
    'two_col',
    'kpi_grid', 'kpi_row',
    'timeline_h',
    'quote_dark',
    'stat_hero',
    'closing_dark', 'closing_split',
    # Nouveaux noms V4 (planner mis à jour)
    'cover', 'closing',
    'timeline', 'quote',
    'agenda', 'highlight_box',
    'pros_cons', 'before_after',
})


def _create_slide_v4_native(prs: Presentation,
                             layout_name: str,
                             content: dict,
                             layout_map: dict,
                             tp: dict = None):
    """
    Crée une slide en utilisant le vrai layout du template.
    Remplit les placeholders existants en préservant leur style XML.
    Garantit showMasterSp='1' (logo visible).
    Accepte anciens noms (cover_dark, timeline_h…) et nouveaux noms V4 (cover, timeline…).
    """
    import lxml.etree as _etree

    # Normalisation nouveaux noms → sémantique interne
    _NAME_ALIAS = {
        'cover':   ('cover', None),
        'closing': ('closing', None),
        'timeline': ('timeline_h', None),
        'quote':    ('quote_dark', None),
        'agenda':   ('full_text', None),
        'highlight_box': ('full_text', None),
        'pros_cons':     ('two_col', None),
        'before_after':  ('two_col', None),
    }

    # ── Choisir le layout template selon le type sémantique ──────────────────
    if layout_name in ('cover', 'cover_dark', 'cover_split'):
        idx = layout_map['cover']
    elif layout_name in ('section', 'quote_dark', 'quote'):
        idx = layout_map['section']
    elif layout_name in ('closing', 'closing_dark', 'closing_split'):
        idx = layout_map.get('closing', layout_map['cover'])
    elif layout_name in ('two_col', 'pros_cons', 'before_after'):
        idx = layout_map['two_col']
    else:
        # full_text, list_*, image_split, kpi_*, timeline_h, stat_hero
        idx = layout_map['content']

    layout = prs.slide_layouts[idx]
    slide  = prs.slides.add_slide(layout)

    # Logo garanti
    cSld = slide._element.find(qn('p:cSld'))
    if cSld is not None:
        cSld.set('showMasterSp', '1')

    # ── Formater le texte de body selon le type de slide ─────────────────────
    def _body_text() -> str:
        # Types simples
        if content.get('body'):
            return content['body']
        if content.get('subtitle'):
            return content['subtitle']
        if content.get('paragraphs'):
            return '\n'.join(str(p) for p in content['paragraphs'])
        if content.get('points'):
            return '\n'.join(f'• {p}' for p in content['points'])
        if content.get('items'):
            parts = []
            for item in content['items']:
                if isinstance(item, dict):
                    t = item.get('title', '')
                    b = item.get('body', '')
                    parts.append(f'• {t}' + (f'\n  {b}' if b else ''))
                else:
                    parts.append(f'• {item}')
            return '\n'.join(parts)
        if content.get('cards'):
            parts = []
            for c in content['cards']:
                if isinstance(c, dict):
                    t = c.get('title', '')
                    b = c.get('body', '')
                    parts.append(f'▸ {t}' + (f'\n  {b}' if b else ''))
                else:
                    parts.append(f'▸ {c}')
            return '\n'.join(parts)

        # KPI grid / row → valeur + label alignés
        if content.get('kpis'):
            parts = []
            for kpi in content['kpis']:
                val = kpi.get('value', '')
                lbl = kpi.get('label', '')
                sub = kpi.get('sublabel', '')
                line = f'{val}  —  {lbl}'
                if sub:
                    line += f'  ({sub})'
                parts.append(line)
            return '\n'.join(parts)

        # Timeline → jalons chronologiques
        if content.get('steps'):
            parts = []
            for step in content['steps']:
                date  = step.get('date', '')
                title = step.get('title', '')
                body  = step.get('body', '')
                line  = f'{date}  →  {title}'
                if body:
                    line += f'\n     {body}'
                parts.append(line)
            return '\n'.join(parts)

        # Quote → citation + auteur
        if content.get('quote'):
            text = f'« {content["quote"]} »'
            if content.get('author'):
                text += f'\n\n— {content["author"]}'
            return text

        # Stat hero → valeur + label + contexte
        if content.get('value'):
            parts = [str(content['value'])]
            if content.get('label'):
                parts.append(content['label'])
            if content.get('context'):
                parts.append(content['context'])
            return '\n'.join(parts)

        # pros_cons → pour / contre
        if content.get('pros') or content.get('cons'):
            parts = []
            if content.get('pros'):
                parts.append('✓ POUR')
                for p in content['pros']:
                    parts.append(f'  + {p}')
            if content.get('cons'):
                parts.append('✗ CONTRE')
                for c in content['cons']:
                    parts.append(f'  − {c}')
            return '\n'.join(parts)

        # before_after → avant / après
        if content.get('before') or content.get('after'):
            parts = []
            b = content.get('before', {})
            a = content.get('after', {})
            if isinstance(b, dict):
                parts.append(f"AVANT : {b.get('title','')}")
                for i in b.get('items', []):
                    parts.append(f'  • {i}')
            if isinstance(a, dict):
                parts.append(f"\nAPRÈS : {a.get('title','')}")
                for i in a.get('items', []):
                    parts.append(f'  • {i}')
            return '\n'.join(parts)

        # agenda → sommaire numéroté
        if content.get('agenda_items') or (content.get('items') and layout_name == 'agenda'):
            items = content.get('agenda_items') or content.get('items', [])
            parts = []
            for i, item in enumerate(items, 1):
                if isinstance(item, dict):
                    parts.append(f'{item.get("number", i)}.  {item.get("label", "")}')
                else:
                    parts.append(f'{i}.  {item}')
            return '\n'.join(parts)

        # highlight_box → highlight (grande police) + body
        if content.get('highlight'):
            parts = [content['highlight']]
            if content.get('body'):
                parts.append('')
                parts.append(content['body'])
            return '\n'.join(parts)

        return ''

    def _title_text() -> str:
        """Titre adapté selon le type : stat_hero affiche la valeur dans le titre."""
        if layout_name == 'stat_hero' and content.get('value'):
            return str(content.get('value', ''))
        return content.get('title', '')

    def _subtitle_text() -> str:
        if layout_name == 'stat_hero':
            return content.get('label', '') or content.get('subtitle', '')
        if layout_name in ('quote_dark', 'quote'):
            return ''  # quote va dans body
        return content.get('subtitle', '')

    # ── Remplir les placeholders ──────────────────────────────────────────────
    body_phs = []
    for ph in slide.placeholders:
        idx_ph = ph.placeholder_format.idx
        if idx_ph == 0:
            _fill_placeholder_preserving_style(ph, _title_text())
        elif idx_ph == 2:
            _fill_placeholder_preserving_style(ph, _subtitle_text())
        elif idx_ph == 4:
            if content.get('footer'):
                _fill_placeholder_preserving_style(ph, content['footer'])
        elif idx_ph in (3, 5):
            pass  # date, numéro → laisser
        else:
            body_phs.append(ph)

    # ── Distribuer les body placeholders ─────────────────────────────────────
    def _col_text(col):
        parts = []
        if col.get('title'):
            parts.append(col['title'])
        for i in col.get('items', []):
            parts.append(f'• {i}')
        return '\n'.join(parts)

    if body_phs:
        if layout_name == 'two_col' and content.get('col_a') and content.get('col_b'):
            _fill_placeholder_preserving_style(body_phs[0], _col_text(content['col_a']))
            if len(body_phs) > 1:
                _fill_placeholder_preserving_style(body_phs[1], _col_text(content['col_b']))
            else:
                combined = _col_text(content['col_a']) + '\n\n' + _col_text(content['col_b'])
                _fill_placeholder_preserving_style(body_phs[0], combined)
        else:
            _fill_placeholder_preserving_style(body_phs[0], _body_text())
    else:
        log.warning(f'[V4] layout="{layout_name}" : aucun placeholder body '
                    f'(ph: {[ph.placeholder_format.idx for ph in slide.placeholders]})')

    return slide


# ── Document ingestion ──────────────────────────────────────────────────────

def extract_document_content(file_bytes: bytes, filename: str) -> str:
    """
    Extrait le texte brut d'un document uploadé.
    Formats : PDF (PyMuPDF), DOCX (python-docx), TXT/MD, PPTX.
    Tronqué à ~12 000 caractères pour rester dans les limites de tokens.
    """
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else 'txt'
    text = ''

    try:
        if ext in ('txt', 'md', 'csv'):
            text = file_bytes.decode('utf-8', errors='ignore')

        elif ext == 'pdf':
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(stream=file_bytes, filetype='pdf')
                pages = []
                for page in doc:
                    pages.append(page.get_text())
                text = '\n'.join(pages)
            except ImportError:
                log.warning('[doc] PyMuPDF non disponible — PDF ignoré')
                text = ''

        elif ext == 'docx':
            try:
                from docx import Document as _DocxDoc
                doc  = _DocxDoc(io.BytesIO(file_bytes))
                text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
            except ImportError:
                log.warning('[doc] python-docx non disponible — DOCX ignoré')
                text = ''

        elif ext == 'pptx':
            src_prs = Presentation(io.BytesIO(file_bytes))
            parts   = []
            for slide in src_prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                parts.append(t)
            text = '\n'.join(parts)

    except Exception as e:
        log.warning(f'[doc] Extraction échouée ({filename}): {e}')
        text = ''

    # Tronquer intelligemment (ne pas couper un mot)
    MAX_CHARS = 15_000
    if len(text) > MAX_CHARS:
        cut = text[:MAX_CHARS].rfind('\n')
        text = text[:cut if cut > MAX_CHARS * 0.8 else MAX_CHARS]

    log.info(f'[doc] {filename} → {len(text)} caractères extraits')
    return text.strip()


# ── Planner V4 (async, avec contexte document) ─────────────────────────────

_V4_PLANNER_SYSTEM = """
Tu es un directeur artistique senior, spécialiste des présentations exécutives haut de gamme.
Tu génères des plans de présentation PowerPoint en JSON strict.

PHILOSOPHIE FONDAMENTALE — CHAQUE SLIDE EST UNE AFFICHE :
Chaque slide doit être visuellement complète, dense, travaillée et esthétique.
Pense à chaque slide comme une double page de magazine de qualité : riche en information
visuelle, hiérarchisée, jamais vide, jamais un simple liste de texte.

RÈGLE D'OR : ZÉRO SLIDE PAUVRE
Une slide avec 3 bullets de texte EST UN ÉCHEC DE DESIGN.
Transforme toujours le texte en éléments visuels structurés :
- 3 idées → 3 cartes (list_cards) avec icône + titre + stat
- 4 étapes → process_flow ou list_numbered (timeline/steps)
- Données clés → kpi_grid ou stat_hero
- Comparaison → entity ou two_col ou pros_cons
- Argument central → highlight_box ou quote

RICHESSE VISUELLE OBLIGATOIRE :
- list_cards : TOUJOURS icon (emoji) + title + body (1 phrase) + stat_value/stat_label
- col3 : TOUJOURS icon + label (badge) + title + 3-4 items courts + stat_value/stat_label
- entity : TOUJOURS icon/flag + name + badge + 3-4 items + stat_value/stat_label
- kpi_grid : TOUJOURS value + label + sublabel (source/contexte) pour chaque KPI
- two_col : TOUJOURS title de colonne + subtitle + 3-4 items courts et percutants
- highlight_box : TOUJOURS highlight (accroche forte) + body (1-2 phrases) + points (2-3 bullets)
- Chaque slide DOIT avoir section_label (MAJUSCULES) et subtitle (accroche en 1 ligne)

HIÉRARCHIE DES LAYOUTS (du plus riche au moins riche) :
TIER 1 (préférer) : list_cards, col3, entity, kpi_grid, infographic, stat_hero, conclusion,
                    team_grid, stat_banner, icon_row, numbered_features, photo_text
TIER 2 (utiliser) : two_col, highlight_box, quote, timeline, process_flow, matrix_2x2, swot,
                    section_break
TIER 3 (éviter) : list_numbered, before_after, pros_cons, funnel, pyramid, cycle, roadmap
TIER 4 (dernier recours, max 1 par présentation) : full_text

Les graphiques (bar_chart, etc.) uniquement avec données chiffrées réelles. Max 1-2.
""".strip()

_V4_PLANNER_USER = """
SUJET : {prompt}
NOMBRE DE SLIDES : {nb_slides}
COULEUR PRINCIPALE : #{primary}  |  ACCENT : #{accent}  |  POLICE : {font}

CATALOGUE DES LAYOUTS :

─── SLIDES STRUCTURELLES ───
cover          — Couverture (title, subtitle≤12mots)
section        — Séparateur de section (title, subtitle)
closing        — Clôture / Merci (title, subtitle≤10mots)
agenda         — Sommaire (title, section_label, items:[{{number,label}}] — max 6 items)

─── LAYOUTS VISUELS RICHES [TIER 1 — PRÉFÉRER] ───
list_cards     — 2-4 cartes visuelles COMPLÈTES
                 Champs OBLIGATOIRES : icon(emoji) + title(≤6mots) + body(≤12mots) + stat_value + stat_label
                 Champs optionnels : label(badge≤3mots), subtitle(≤8mots)
                 Exemple : {{"icon":"🌍","label":"EXPORT","title":"Marchés émergents","body":"Croissance 18% en Asie du Sud-Est","stat_value":"18%","stat_label":"CROISSANCE"}}

col3           — 3 colonnes enrichies COMPLÈTES
                 Champs OBLIGATOIRES : icon(emoji) + label(≤3mots) + title(≤6mots) + items:[3-4 items≤8mots] + stat_value + stat_label
                 Exemple : {{"icon":"⚡","label":"ÉNERGIE","title":"Transition verte","subtitle":"2025-2030","items":["Mix 60% renouvelable","Coût -32%","3 pays pilotes"],"stat_value":"60%","stat_label":"RENOUVELABLE"}}

kpi_grid       — Grille de 3-6 KPIs avec chiffres forts
                 Champs OBLIGATOIRES par KPI : value(chiffre+unité) + label(≤5mots) + sublabel(source≤8mots)
                 Optionnel : percent(0-100 pour barre visuelle)

entity         — Comparaison de 2-4 entités (pays, acteurs, marques)
                 Champs OBLIGATOIRES : icon(drapeau/emoji) + name + badge(≤3mots) + items:[3-4 items≤8mots] + stat_value + stat_label

infographic    — Grande stat + décomposition visuelle
                 Champs : value + label + context(≤20mots) + bars:[{{label,percent}}]×3-5

stat_hero      — 1-3 statistiques héros centrées
                 Si 1 stat : value + label + context(≤20mots) + points:[2-3 bullets≤8mots]
                 Si plusieurs : values:[{{value,label,context}}]×2-3

conclusion     — Synthèse finale grille 2×2 + sidebar
                 cards:[{{icon,title,body≤10mots}}]×4 + sidebar_title + sidebar_quote(≤20mots) + sidebar_cta

─── LAYOUTS ANALYSE [TIER 2] ───
two_col        — 2 colonnes comparatives
                 col_a + col_b : chacun title(≤5mots) + subtitle(≤8mots) + items:[3-4 items≤8mots chacun]

highlight_box  — Encadré fort avec message central
                 highlight(accroche≤15mots) + body(≤25mots) + points:[2-3 bullets≤8mots]

quote          — Citation mise en avant
                 quote(≤20mots) + author + source? + category?

timeline       — Frise chronologique
                 steps:[{{date,title≤5mots,body≤10mots}}]×3-6

process_flow   — Flux de processus
                 steps:[{{title≤5mots,body≤10mots}}]×3-5

matrix_2x2     — Matrice 2×2 (Impact/Effort, Urgence/Importance...)
                 quadrants:[{{label,body≤15mots}}]×4 + axes:{{x,y}}

swot           — Analyse SWOT
                 strengths/weaknesses/opportunities/threats:[≤4 items≤8mots chacun]

─── LAYOUTS VISUELS ENRICHIS [TIER 1 — NOUVEAUX] ───
team_grid      — Équipe / personnes avec avatars circulaires
                 Champs OBLIGATOIRES : members:[{icon(emoji),name,role}] — 3 ou 4 membres
                 Optionnel : department, stat_value, stat_label, body
                 Exemple membre : {{"icon":"👩‍💼","name":"Marie Dupont","role":"DG","stat_value":"12 ans","stat_label":"EXP."}}

stat_banner    — Bandeau de 3-4 grandes statistiques (chiffres premium)
                 Champs OBLIGATOIRES : stats:[{value,label}] — 3 ou 4 stats
                 Optionnel : sublabel(source), icon(emoji)
                 Exemple : {{"value":"1 254+","label":"Clients actifs","sublabel":"Depuis 2018","icon":"🌍"}}

icon_row       — 3-4 features en cercles d'icônes connectés horizontalement
                 Champs OBLIGATOIRES : items:[{icon(emoji),title,body}] — 3 ou 4 items
                 Optionnel : items:[str] pour bullets dans chaque feature

numbered_features — Features avec grands numéros décoratifs (style éditorial)
                 Champs OBLIGATOIRES : items:[{title,body}] — 3 ou 4 items
                 Optionnel : number(override), stat_value, stat_label

photo_text     — Zone photo (visuelle) + contenu structuré (3 items)
                 Champs OBLIGATOIRES : items:[{title,body}] — 3 items
                 Optionnel : icon(emoji dans item), stat_value(dans item)
                 Usage : présenter une offre, un projet, un profil avec visuel

section_break  — Slide de rupture dramatique entre sections majeures
                 Champs OBLIGATOIRES : title (max 6 mots) — IMPACT XL
                 Optionnel : subtitle, number(numéro de section)
                 Usage : UNIQUEMENT pour séparer des sections importantes

─── LAYOUTS STRUCTURE [TIER 3 — UTILISER AVEC PARCIMONIE] ───
list_numbered  — Liste 3-4 items UNIQUEMENT si séquence logique stricte
                 items:[{{title≤5mots,body≤10mots}}]×3-4 MAX

pros_cons      — Pour/Contre (pros/cons:[≤4 items≤8mots])
before_after   — Avant/Après (before/after:{{title,items:[≤4 items≤8mots]}})
funnel         — Entonnoir (steps:[{{label,value}}]×3-5)
pyramid        — Pyramide (levels:[{{label,body≤10mots}}])
cycle          — Cycle (steps:[{{title≤5mots,body≤8mots}}]×3-5)
roadmap        — Roadmap (phases:[{{label,milestones:[≤3 jalons]}}]×2-4)

─── DONNÉES [UTILISER SEULEMENT SI DONNÉES RÉELLES] ───
bar_chart/line_chart/pie_chart/stacked_bar/waterfall/radar
  → OBLIGATOIRE : analysis(2 phrases d'interprétation)

─── LAYOUT TEXTE [TIER 4 — DERNIER RECOURS, 1 MAX] ───
full_text      — UNIQUEMENT pour intro/conclusion narrative
                 body≤2 paragraphes de ≤25 mots chacun. JAMAIS de bullets.

─── RÈGLES ABSOLUES ───
1. La première slide est toujours "cover", la dernière "closing".
2. Chaque slide : section_label (MAJUSCULES ≤4 mots) + subtitle (accroche ≤10 mots).
3. "style": alterner 0/1/2/3/4 entre les slides pour maximiser la diversité visuelle.
4. "presentation_seed": entier aléatoire 1-999999 à la racine.
5. footer_text = baseline société (≤8 mots).
6. Répondre UNIQUEMENT avec le JSON demandé, sans commentaire ni markdown.
7. Graphiques : UNIQUEMENT si données chiffrées réelles. Max 1-2 par présentation.
8. JAMAIS de slide avec seulement du texte non structuré. JAMAIS plus de 4 bullets.
9. Chaque slide TIER 1 doit utiliser TOUS ses champs obligatoires — aucun ne peut être omis.
10. Variété obligatoire : pas 2 fois le même layout dans une présentation de ≤10 slides.
11. section_break : utiliser max 1 fois, uniquement si ≥10 slides avec sections distinctes.
12. Nouveaux layouts TIER 1 (team_grid, stat_banner, icon_row, numbered_features, photo_text) :
    utiliser au moins 1 sur 8+ slides — ils apportent une richesse visuelle que list_cards ne peut pas.

FORMAT DE RÉPONSE :
{{
  "presentation_title": "...",
  "footer_text": "...",
  "presentation_seed": 42731,
  "slides": [
    {{
      "layout": "<nom_du_layout>",
      "content": {{
        "title": "...",
        "section_label": "RUBRIQUE",
        "subtitle": "Accroche en une ligne",
        "style": 0,
        ... (tous les champs du layout)
      }}
    }}
  ]
}}
""".strip()

_V4_DOC_INJECT = """

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
DOCUMENT SOURCE (contenu à synthétiser) :
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
{document_content}
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
INSTRUCTION : Synthétise ce document en présentation. Extrais les données
clés, chiffres, arguments. Transforme en contenu visuel — ne recopie pas
mot à mot. Adapte chaque slide au type de layout choisi.
"""


async def plan_presentation_v4(
    prompt: str,
    nb_slides: int,
    tp: dict,
    document_content: str = '',
) -> dict:
    """
    Planifie la présentation V4 via AsyncAnthropic (non-bloquant).
    tp : template profile (theme, font, layout_map…) issu d'analyze_template_v4.
    Accepte aussi un dict palette simple pour compatibilité descendante.
    """
    theme   = tp.get('theme', tp)     # compat : tp peut être un palette dict legacy
    primary = theme.get('accent1') or tp.get('primary', '1A3A6B')
    accent  = theme.get('accent2') or tp.get('accent',  'F0A500')
    font    = tp.get('font') or tp.get('font', 'Calibri')

    user_prompt = _V4_PLANNER_USER.format(
        prompt    = prompt,
        nb_slides = nb_slides,
        primary   = primary,
        accent    = accent,
        font      = font,
    )

    if document_content:
        user_prompt = user_prompt + _V4_DOC_INJECT.format(
            document_content=document_content
        )

    max_tokens = max(5000, nb_slides * 500)

    async with anthropic.AsyncAnthropic(api_key=ANTHROPIC_API_KEY) as async_client:
        for attempt in range(3):
            msg = await async_client.messages.create(
                model      = CLAUDE_MODEL,
                max_tokens = max_tokens,
                system     = _V4_PLANNER_SYSTEM,
                messages   = [{'role': 'user', 'content': user_prompt}],
            )
            try:
                plan = _parse_json_robust(msg.content[0].text.strip(), context='plan_v4')
                log.info(f'[V4] Plan: {len(plan.get("slides",[]))} slides, '
                         f'title="{plan.get("presentation_title","")[:60]}"')
                return plan
            except (ValueError, KeyError) as e:
                log.warning(f'plan_presentation_v4 attempt {attempt+1}/3: {e}')
                if attempt == 2:
                    raise
    raise RuntimeError('plan_presentation_v4 : 3 tentatives échouées')


async def run_pipeline_v4(
    pptx_bytes: bytes,
    prompt: str,
    nb_slides: int,
    document_content: str = '',
    plan: dict = None,
) -> tuple:
    """
    Pipeline V4 — template-native generation.

    Phase 1 : analyze_template_v4 → tp (theme, layout_map, logo_zone, font…)
    Phase 2 : planification narrative async — sauté si plan déjà fourni
    Phase 3 : création slides — native (layout_*_v4) ou V3 fallback
    Phase 4 : suppression slides originales + export

    plan : plan pré-calculé optionnel (évite le double-appel depuis /generate-stream)
    """
    if not ANTHROPIC_API_KEY:
        raise ValueError('Clé API Claude manquante.')

    nb_slides = max(2, min(nb_slides, 30))

    prs        = Presentation(io.BytesIO(pptx_bytes))
    n_original = len(prs.slides)

    # ── Phase 1 ──────────────────────────────────────────────
    log.info(f'[V4] Phase 1 : analyze_template_v4 ({nb_slides} slides demandées)…')
    tp         = analyze_template_v4(prs)
    layout_map = tp['layout_map']

    # Palette V3-compat pour fallback layouts.py
    palette = {
        'primary': tp['theme'].get('accent1', '009CEA'),
        'accent':  tp['theme'].get('accent2', 'ED0000'),
        'dark':    tp['theme'].get('dk1', '374649'),
        'bg':      tp['theme'].get('lt1', 'FFFFFF'),
        'font':    tp['font'],
    }
    brand = extract_brand(prs)

    # ── Phase 2 ──────────────────────────────────────────────
    if plan is None:
        log.info('[V4] Phase 2 : planification narrative…')
        plan = await plan_presentation_v4(prompt, nb_slides, tp, document_content)
    else:
        log.info('[V4] Phase 2 : plan pré-calculé fourni — skip')
    slides_plan = plan.get('slides', [])

    # Seed de présentation — injecté dans tp pour les fonctions layout
    _title_hash = abs(hash(str(plan.get('presentation_title', '')))) % 1000000
    seed = int(plan.get('presentation_seed', _title_hash)) % 1000000
    tp['seed'] = seed
    log.info(f'[V4] Seed de présentation : {seed}')

    # Compléter si Claude en a généré moins que demandé
    fallback_layouts = ['list_cards', 'col3', 'two_col', 'highlight_box']
    while len(slides_plan) < nb_slides:
        fb = fallback_layouts[len(slides_plan) % len(fallback_layouts)]
        slides_plan.append({
            'layout':  fb,
            'content': {
                'title':      'Développement complémentaire',
                'paragraphs': ['Contenu additionnel à personnaliser.'],
                'footer':     plan.get('footer_text', ''),
            },
        })
    slides_plan = slides_plan[:nb_slides]

    # Garantie closing (V4 ou V3 nom)
    _closing_names = {'closing'}
    if not slides_plan or slides_plan[-1].get('layout') not in _closing_names:
        closing_slide = {
            'layout':  'closing',
            'content': {'title': 'Merci', 'subtitle': plan.get('footer_text', '')},
        }
        if len(slides_plan) >= nb_slides and slides_plan:
            slides_plan[-1] = closing_slide
        else:
            slides_plan.append(closing_slide)

    log.info(f'[V4] {len(slides_plan)} slides : {[s.get("layout") for s in slides_plan]}')

    # ── Phase 3 ──────────────────────────────────────────────
    log.info('[V4] Phase 3 : création des slides…')

    # Alias V3 → fonction layouts.py pour types sans implémentation V4 native
    # (pyramid, cycle, roadmap → V3 si disponible, sinon fulltext fallback)
    _V3_ALIAS = {
        'pyramid': 'pyramid',
        'cycle':   'cycle',
        'roadmap': 'roadmap',
    }

    success = 0
    for sp in slides_plan:
        layout_name = sp.get('layout', 'full_text')
        content     = sp.get('content', {})

        if not content.get('footer') and plan.get('footer_text'):
            content['footer'] = plan['footer_text']

        try:
            # ── Routing V4 — layout functions dédiées ────────────────────────
            if layout_name in ('cover', 'cover_dark', 'cover_split'):
                layout_cover_v4(prs, content, tp)
            elif layout_name == 'section':
                layout_section_v4(prs, content, tp)
            elif layout_name in ('closing', 'closing_dark', 'closing_split'):
                layout_closing_v4(prs, content, tp)
            elif layout_name in ('full_text', 'image_split'):
                layout_fulltext_v4(prs, content, tp)
            elif layout_name in ('kpi_grid', 'kpi_row', 'kpi_native'):
                layout_kpi_grid_v4(prs, content, tp)
            elif layout_name in ('quote', 'quote_dark'):
                layout_quote_v4(prs, content, tp)
            elif layout_name in ('list_numbered',):
                layout_list_numbered_v4(prs, content, tp)
            elif layout_name in ('list_cards',):
                layout_list_cards_v4(prs, content, tp)
            elif layout_name in ('two_col',):
                layout_twocol_v4(prs, content, tp)
            elif layout_name in ('col3', 'three_col'):
                layout_col3_v4(prs, content, tp)
            elif layout_name in ('conclusion',):
                layout_conclusion_v4(prs, content, tp)
            elif layout_name in ('entity', 'entity_compare', 'comparison'):
                layout_entity_v4(prs, content, tp)
            elif layout_name in ('stat_hero',):
                layout_stathero_v4(prs, content, tp)
            elif layout_name in ('infographic', 'infograph'):
                layout_infographic_v4(prs, content, tp)
            elif layout_name in ('timeline', 'timeline_h'):
                layout_timeline_v4(prs, content, tp)
            elif layout_name in ('process_flow',):
                layout_processflow_v4(prs, content, tp)
            elif layout_name in ('funnel',):
                layout_funnel_v4(prs, content, tp)
            elif layout_name in ('bar_chart',):
                layout_barchart_v4(prs, content, tp)
            elif layout_name in ('line_chart',):
                layout_linechart_v4(prs, content, tp)
            elif layout_name in ('pie_chart',):
                layout_piechart_v4(prs, content, tp)
            elif layout_name in ('waterfall',):
                layout_waterfall_v4(prs, content, tp)
            elif layout_name in ('matrix_2x2', 'matrix'):
                layout_matrix_v4(prs, content, tp)
            elif layout_name in ('swot',):
                layout_swot_v4(prs, content, tp)
            elif layout_name in ('pros_cons',):
                layout_proscons_v4(prs, content, tp)
            elif layout_name in ('table',):
                layout_table_v4(prs, content, tp)
            elif layout_name in ('radar',):
                layout_radar_v4(prs, content, tp)
            elif layout_name in ('pyramid',):
                layout_pyramid_v4(prs, content, tp)
            elif layout_name in ('cycle',):
                layout_cycle_v4(prs, content, tp)
            elif layout_name in ('roadmap',):
                layout_roadmap_v4(prs, content, tp)
            elif layout_name in ('stacked_bar',):
                layout_stackedbar_v4(prs, content, tp)
            elif layout_name in ('before_after',):
                layout_beforeafter_v4(prs, content, tp)
            elif layout_name in ('highlight_box',):
                layout_highlight_v4(prs, content, tp)
            elif layout_name in ('agenda',):
                layout_agenda_v4(prs, content, tp)
            elif layout_name in ('team_grid', 'team'):
                layout_team_grid_v4(prs, content, tp)
            elif layout_name in ('stat_banner',):
                layout_stat_banner_v4(prs, content, tp)
            elif layout_name in ('icon_row',):
                layout_icon_row_v4(prs, content, tp)
            elif layout_name in ('section_break',):
                layout_section_break_v4(prs, content, tp)
            elif layout_name in ('photo_text',):
                layout_photo_text_v4(prs, content, tp)
            elif layout_name in ('numbered_features',):
                layout_numbered_features_v4(prs, content, tp)

            # ── Routing V3 fallback (résiduel — ne devrait plus être atteint) ─
            else:
                fn_key = _V3_ALIAS.get(layout_name, layout_name)
                layout_fn = LAYOUT_REGISTRY.get(fn_key) or LAYOUT_REGISTRY.get(layout_name)
                if layout_fn:
                    layout_fn(prs, content, palette)
                else:
                    layout_fulltext_v4(prs, content, tp)

            success += 1
            log.info(f'[V4] ✓ {layout_name}')
        except Exception as e:
            log.error(f'[V4] ✗ {layout_name} : {repr(e)}', exc_info=True)
            try:
                layout_fulltext_v4(prs, {
                    'title':  content.get('title', ''),
                    'body':   content.get('body', ''),
                    'footer': content.get('footer', ''),
                }, tp)
                success += 1
            except Exception as e2:
                log.error(f'[V4] ✗ fallback full_text : {repr(e2)}', exc_info=True)

    log.info(f'[V4] Phase 3 terminée : {success}/{len(slides_plan)} OK')

    # ── Phase 4 ──────────────────────────────────────────────
    slides_added = len(prs.slides) - n_original
    if slides_added == 0:
        raise RuntimeError('[V4] Aucune slide créée')

    xml_slides = prs.slides._sldIdLst
    for sld_id in list(prs.slides._sldIdLst)[:n_original]:
        xml_slides.remove(sld_id)
    log.info(f'[V4] {n_original} slides originales supprimées — '
             f'{len(prs.slides)} slides finales')

    # Marqueur V4 dans les métadonnées PPTX (vérifiable via PowerPoint > Infos)
    try:
        import datetime as _dt
        prs.core_properties.keywords = 'VisualCortex-V4'
        prs.core_properties.comments = (
            f'Generated by Visual Cortex V4 — {len(prs.slides)} slides — '
            f'{_dt.datetime.utcnow().strftime("%Y-%m-%d %H:%M")} UTC'
        )
        prs.core_properties.revision = 4
    except Exception:
        pass

    log.info(f'[V4] ✅ Pipeline terminé : {success}/{len(slides_plan)} slides réussis')

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read(), plan, brand, palette


# ══════════════════════════════════════════════════════════════
# ROUTES NIVEAU 2
# ══════════════════════════════════════════════════════════════

@app.get('/profiles')
def get_profiles():
    """Retourne les profils clients disponibles pour le Niveau 2."""
    return {
        key: {'label': p['label'], 'description': p['description']}
        for key, p in CLIENT_PROFILES.items()
    }


@app.post('/generate-v2')
async def generate_presentation_v2(
    request:       Request,
    template:      UploadFile = File(...),
    prompt:        str        = Form(...),
    nb_slides:     str        = Form(default='complet'),
    profile:       str        = Form(default='institutional'),
    document:      UploadFile = File(default=None),
    authorization: str        = Form(default=None),
):
    """
    Route V4 — template-native generation (V3 comme fallback).
    nb_slides : "Essentiel"→6 | "Complet"→10 | "Approfondi"→16 ou entier.
    profile conservé pour compatibilité API future.
    document : fichier source optionnel (PDF, DOCX, TXT, PPTX).
    """
    if not _is_pro(authorization):
        _quota(_ip(request))

    n          = _resolve_nb_slides(nb_slides)
    pptx_bytes = await template.read()

    doc_content = ''
    if document is not None:
        try:
            doc_bytes = await document.read()
            _doc_mb   = len(doc_bytes) / (1024 * 1024)
            if _doc_mb > 20:
                raise HTTPException(400, f'Document trop volumineux ({_doc_mb:.1f} MB, max 20 MB)')
            doc_content = extract_document_content(doc_bytes, document.filename or 'doc')
        except HTTPException:
            raise
        except Exception as e:
            log.warning(f'[/generate-v2] document extraction failed: {e}')

    try:
        final_bytes, plan, brand, palette = await run_pipeline_v4(
            pptx_bytes, prompt, n, doc_content
        )
        log.info('[/generate-v2] Pipeline V4 OK')
    except Exception as e:
        import asyncio as _asyncio
        log.warning(f'[/generate-v2] V4 échoué ({e}) → fallback V3')
        try:
            final_bytes, plan, brand, palette = await _asyncio.to_thread(
                run_pipeline_v3, pptx_bytes, prompt, n
            )
        except Exception as e2:
            log.error(f'[/generate-v2] V3 aussi échoué : {e2}', exc_info=True)
            raise HTTPException(500, f'Erreur génération : {e2}')

    filename = f'visualcortex-{_safe_name(prompt)}.pptx'
    return StreamingResponse(
        io.BytesIO(final_bytes),
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        headers={
            'Content-Disposition': f'attachment; filename={filename}',
            'Content-Length':      str(len(final_bytes)),
        },
    )


@app.post('/generate-v2-preview')
async def generate_v2_preview(
    request:       Request,
    template:      UploadFile = File(...),
    prompt:        str        = Form(...),
    nb_slides:     str        = Form(default='complet'),
    profile:       str        = Form(default='institutional'),
    authorization: str        = Form(default=None),
):
    """Aperçu plan + palette L2 (sans générer le fichier)."""
    pptx_bytes = await template.read()
    prs        = Presentation(io.BytesIO(pptx_bytes))
    brand      = extract_brand(prs)
    palette    = _h2_extract_palette(brand)
    library    = build_layout_library(prs)
    n          = _resolve_nb_slides(nb_slides)
    sel        = select_template_slides(library, n)

    pro = _is_pro(authorization)
    quota_info = (
        {'plan': 'pro'} if pro
        else {'used': _quota(_ip(request))[0], 'total': FREE_QUOTA_PER_IP, 'plan': 'free'}
    )

    plan = plan_presentation(prompt, n, sel, brand)
    return {
        'success':            True,
        'level':              2,
        'profile':            profile,
        'profile_label':      CLIENT_PROFILES.get(profile, {}).get('label', profile),
        'nb_slides_resolved': n,
        'presentation_title': plan.get('presentation_title', prompt[:60]),
        'narrative_arc':      plan.get('narrative_arc', ''),
        'palette':            palette,
        'slides': [
            {
                'index':           s.get('plan_index'),
                'type':            s.get('slide_type'),
                'narrative_angle': s.get('narrative_angle'),
                'key_message':     s.get('key_message'),
            }
            for s in plan.get('slides', [])
        ],
        'quota': quota_info,
    }


# ══════════════════════════════════════════════════════════════
# ENTRY POINT
# ══════════════════════════════════════════════════════════════

if __name__ == '__main__':
    uvicorn.run(
        'main:app',
        host  = '0.0.0.0',
        port  = int(os.environ.get('PORT', 8000)),
        reload = False,
    )

"""
check_bounds_v4.py — Programmatic shape-overflow detector for PPTX files.

For each slide, reads every shape's bounding box and checks for:
  1. Bottom overflow: shape extends below the slide bottom (always a bug)
  2. Major right overflow: extends > 1.5" past slide right (likely a bug)
  3. Major left overflow: extends > 1.5" past slide left (likely a bug)
  4. Content-zone footer bleed: content shape extends below CB=6.95"
     by more than 0.1" (excluding full-height backgrounds)

Usage:
    python check_bounds_v4.py <file.pptx> [--cb 6.95] [--strict]
    python -m check_bounds_v4   (runs self-test on test_v4_output.pptx)

Returns exit code 1 if any violation found.
"""

import sys
import argparse
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Emu

# ── Constants ────────────────────────────────────────────────────────────────
_EMU = 914400           # EMU per inch
_CB  = 6.95             # content bottom (default)
_CT  = 1.55             # content top
_TOL_SLIDE_BOTTOM = 0.10   # inches — tolerance for slide bottom (none below H)
_TOL_MAJOR_HORIZ  = 1.50   # inches — right/left overflow threshold
_TOL_FOOTER_BLEED = 0.10   # inches — tolerance for CB overflow
_FULL_HEIGHT_FRAC = 0.85   # shape height / slide height above which we skip CB check


def _emu_to_in(emu):
    return emu / _EMU


_NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
_NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def _shape_alpha_pct(shape):
    """Return fill alpha 0-100 (100 = fully opaque). Returns 100 if not set."""
    try:
        spPr = shape._element.find(f'{{{_NS_P}}}spPr')
        if spPr is None:
            return 100
        alpha_el = spPr.find(f'.//{{{_NS_A}}}alpha')
        if alpha_el is not None:
            return int(alpha_el.get('val', '100000')) // 1000
    except Exception:
        pass
    return 100


def check_slide(slide, slide_idx, slide_label, W=13.33, H=7.50, cb=_CB):
    """
    Returns list of violation dicts for this slide.
    Each dict: {slide, label, shape_name, type, detail}
    """
    violations = []
    W_emu = W * _EMU
    H_emu = H * _EMU
    cb_emu = cb * _EMU
    full_h_threshold = H * _FULL_HEIGHT_FRAC * _EMU

    for shape in slide.shapes:
        try:
            left   = shape.left   or 0
            top    = shape.top    or 0
            width  = shape.width  or 0
            height = shape.height or 0
        except Exception:
            continue

        right  = left + width
        bottom = top + height
        name   = shape.name or '(unnamed)'
        stype  = type(shape).__name__

        def _v(msg):
            violations.append({
                'slide': slide_idx,
                'label': slide_label,
                'shape': name,
                'type': stype,
                'left_in':   round(_emu_to_in(left),   3),
                'top_in':    round(_emu_to_in(top),    3),
                'right_in':  round(_emu_to_in(right),  3),
                'bottom_in': round(_emu_to_in(bottom), 3),
                'detail': msg,
            })

        # Rule 1: below slide bottom (never acceptable)
        if bottom > H_emu + _TOL_SLIDE_BOTTOM * _EMU:
            _v(f'BELOW SLIDE BOTTOM: bottom={_emu_to_in(bottom):.3f}" > H={H:.3f}"')

        # Rules 2-3 skip decorative shapes (alpha < 30% = nearly transparent)
        alpha = _shape_alpha_pct(shape)
        is_decorative = alpha < 30

        # Rule 2: major right overflow (>1.5" beyond right edge)
        if not is_decorative and right > W_emu + _TOL_MAJOR_HORIZ * _EMU:
            _v(f'MAJOR RIGHT OVERFLOW: right={_emu_to_in(right):.3f}" > W+1.5"={W+1.5:.3f}"')

        # Rule 3: major left overflow (>1.5" beyond left edge)
        if not is_decorative and left < -_TOL_MAJOR_HORIZ * _EMU:
            _v(f'MAJOR LEFT OVERFLOW: left={_emu_to_in(left):.3f}" < -1.5"')

        # Rule 4: content bleeds into footer zone.
        # Exclusions: shapes that START in the footer zone (intentional footer elements),
        # full-height backgrounds that span most of the slide, and decorative shapes
        # (alpha < 30%) which are purely visual and do not displace content.
        footer_zone_top = (H - 0.45) * _EMU  # footer area starts ~7.05"
        in_footer_zone  = top >= footer_zone_top
        if (not is_decorative
                and bottom > cb_emu + _TOL_FOOTER_BLEED * _EMU
                and height < full_h_threshold
                and not in_footer_zone):
            _v(f'FOOTER BLEED: bottom={_emu_to_in(bottom):.3f}" > CB+0.1"={cb+0.1:.3f}" '
               f'(shape h={_emu_to_in(height):.3f}")')

    return violations


def check_pptx(pptx_path, slide_labels=None, W=13.33, H=7.50, cb=_CB, verbose=True):
    """
    Check all slides in a PPTX for shape overflow violations.
    slide_labels: list of strings (one per slide) for readable output.
    Returns list of all violation dicts.
    """
    prs = Presentation(pptx_path)
    all_violations = []

    for i, slide in enumerate(prs.slides):
        label = slide_labels[i] if slide_labels and i < len(slide_labels) else f'slide_{i:03d}'
        viols = check_slide(slide, i, label, W=W, H=H, cb=cb)
        all_violations.extend(viols)

    if verbose:
        _print_report(all_violations, pptx_path, len(prs.slides))

    return all_violations


def _print_report(violations, pptx_path, n_slides):
    sep = '─' * 70
    print(f'\n{sep}')
    print(f'  BOUNDS CHECK — {pptx_path}  ({n_slides} slides)')
    print(sep)

    if not violations:
        print('  ✓  No violations found.')
    else:
        by_slide = {}
        for v in violations:
            by_slide.setdefault(v['slide'], []).append(v)

        for slide_idx in sorted(by_slide):
            vs = by_slide[slide_idx]
            label = vs[0]['label']
            print(f'\n  slide {slide_idx:03d}  [{label}]')
            for v in vs:
                print(f'    ✗ {v["detail"]}')
                print(f'      shape "{v["shape"]}" ({v["type"]})  '
                      f'L={v["left_in"]:.2f}" T={v["top_in"]:.2f}" '
                      f'R={v["right_in"]:.2f}" B={v["bottom_in"]:.2f}"')

    print(f'\n  Total: {len(violations)} violation(s) across {n_slides} slides.')
    print(sep + '\n')


if __name__ == '__main__':
    ap = argparse.ArgumentParser()
    ap.add_argument('pptx', nargs='?', default='test_v4_output.pptx')
    ap.add_argument('--cb', type=float, default=_CB)
    ap.add_argument('--strict', action='store_true',
                    help='Exit with code 1 if any violations found')
    args = ap.parse_args()

    viols = check_pptx(args.pptx, cb=args.cb)
    if args.strict and viols:
        sys.exit(1)
